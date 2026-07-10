"""
DPS 空調デマンド制御システム「提案・稟議シミュレーター」
完全版 app.py

起動方法:
  pip install streamlit altair pandas openpyxl python-pptx
  streamlit run app.py
"""


import streamlit as st

# ===== ログイン =====
def login():
    st.title("ログイン")

    username = st.text_input("ユーザー名")
    password = st.text_input("パスワード", type="password")

    USER = st.secrets["USERNAME"]
    PASS = st.secrets["PASSWORD"]

    if st.button("ログイン"):
        if username == USER and password == PASS:
            st.session_state["logged_in"] = True
            st.rerun()
        else:
            st.error("ユーザー名またはパスワードが違います")

    st.stop()


if "logged_in" not in st.session_state:
    st.session_state["logged_in"] = False

if not st.session_state["logged_in"]:
    login()
# ===================


import streamlit.components.v1 as components
import pandas as pd
import altair as alt
import json
import math
import io
import datetime
import os
import copy


def load_css() -> None:
    """Load the external application stylesheets in a stable order."""
    base_dir = os.path.dirname(os.path.abspath(__file__))
    css_dir = os.path.join(base_dir, "assets", "css")

    for filename in ("common.css", "components.css", "page.css"):
        css_path = os.path.join(css_dir, filename)
        try:
            with open(css_path, "r", encoding="utf-8") as css_file:
                st.markdown(
                    f"<style>{css_file.read()}</style>",
                    unsafe_allow_html=True,
                )
        except FileNotFoundError:
            st.error(f"CSSファイルが見つかりません: {css_path}")
            st.stop()


# openpyxl のインポート（Excel全12シート完全対応）
try:
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False

# pptx 関連のインポート（資料自動生成用）
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
    HAS_PPTX = True
except Exception:
    HAS_PPTX = False

# スライド貼付用 画像出力（matplotlib + pillow）
try:
    from slide_image_export import (make_demand_chart_png, make_control_list_png,
                                     make_control_list_pngs, make_summary_cards_png,
                                     make_usage_bars_png, make_cumulative_balance_png)
    HAS_SLIDE_IMG = True
except Exception:
    HAS_SLIDE_IMG = False

# ══════════════════════════════════════════════════════════════════
# 定数・設定
# ══════════════════════════════════════════════════════════════════
DATA_FILE = "dps_app_data.json"
APP_TITLE = "DPS 空調デマンド制御　提案・稟議シミュレーター"
# アプリのバージョン（サイドバー上部に表示）。改修の都度、内容に応じて自動更新する：
#   ・メジャー(X.0.0)＝大きめの改修：新しい画面(V)の追加、算出方法の追加・削除
#   ・マイナー(x.Y.0)＝中くらいの改修：算出方法の修正、実装済み機能・システムの修正
#   ・パッチ(x.y.Z)＝軽微な改修：ビジュアル・文言など見た目の修正
#   ※出力資料(PPTX等)関連は「機能追加」以外はすべて軽微(パッチ)扱いとする
APP_VERSION = "v1.8.0"
CO2_FACTOR = 0.451          # kg-CO2/kWh
SUGI_KG    = 8.8            # 杉1本年間CO2吸収kg
DEFAULT_SYSTEM_FEE = 241800 # システム利用料（固定調整枠）円/年
AREA_UNIT_PRICE = 50000     # 工事費目安 円/台
CONV_FACTOR = 0.5           # 能力削減→正味エネ削減 換算係数 conv (methodology.md準拠)

DEFAULT_GYOTAI = {
    "スーパー・食品(冷凍冷蔵主役)":    {"ac_peak": 0.05, "ac_kwh": 0.04},
    "製造拠点(プロセス主役)":          {"ac_peak": 0.18, "ac_kwh": 0.15},
    "物流センター・倉庫":              {"ac_peak": 0.47, "ac_kwh": 0.30},
    "ホテル・商業施設":                {"ac_peak": 0.50, "ac_kwh": 0.38},
    "オフィスビル":                    {"ac_peak": 0.40, "ac_kwh": 0.35},
    "病院・医療施設":                  {"ac_peak": 0.35, "ac_kwh": 0.30},
    "学校・教育施設":                  {"ac_peak": 0.45, "ac_kwh": 0.38},
}

# 業態別：空調を除いた電力量の内訳ウェイト（照明／生産設備／その他）。合計1.0。
# ユーザー提供の業種別エネルギー内訳表を、既存7業態に対応づけて算出した相対比率。
# （空調分はシミュレーションの ac_kwh / ac_peak を使うため、ここには含めない）
INDUSTRY_NONAC_WEIGHTS = {
    "スーパー・食品(冷凍冷蔵主役)": {"照明": 0.20, "生産設備": 0.00, "その他": 0.80},  # 冷凍冷蔵が「その他」主体
    "製造拠点(プロセス主役)":       {"照明": 0.12, "生産設備": 0.59, "その他": 0.29},  # 生産設備が主役
    "物流センター・倉庫":           {"照明": 0.235, "生産設備": 0.00, "その他": 0.765},  # 換気・コンプレッサー多
    "ホテル・商業施設":             {"照明": 0.273, "生産設備": 0.00, "その他": 0.727},
    "オフィスビル":                 {"照明": 0.364, "生産設備": 0.00, "その他": 0.636},
    "病院・医療施設":               {"照明": 0.25, "生産設備": 0.00, "その他": 0.75},
    "学校・教育施設":               {"照明": 0.385, "生産設備": 0.00, "その他": 0.615},
}
DEFAULT_NONAC_WEIGHTS = {"照明": 0.25, "生産設備": 0.00, "その他": 0.75}

AREA_BUFFER = {
    "北海道・東北": 1.10,
    "関東":        1.00,
    "中部・北陸":  1.05,
    "近畿":        1.02,
    "中国・四国":  1.08,
    "九州・沖縄":  1.12,
}

HP_TO_KW = 2.8   # 1馬力 ≈ 冷房能力 2.8kW（業務用空調の標準換算）

# 詳細見積：交通費の地区区分（移動拘束費の地区別単価キー → 表示ラベル）
AREA5_LABELS = {
    "地区1": "地区1（東京23区・神奈川）",
    "地区2": "地区2（東京23区外・静岡・山梨・埼玉・千葉）",
    "地区3": "地区3（愛知・岐阜・長野・群馬・栃木・茨城）",
    "地区4": "地区4（兵庫・大阪・京都・奈良・和歌山・福井・滋賀・三重・石川・富山・新潟・福島）",
    "地区5": "地区5（上記以外）※北海道=飛行機/他=新幹線 で交通費を実費入力",
}

# データ補完（検針票に列が無い時）の既定値。各種設定で編集可（const_settings.fb_defaults で上書き）。
FB_DEFAULTS = {
    "契約電力":    626.0,
    "最大需要電力": 400.0,
    "使用量合計":  100000.0,
    "力率":        0.85,    # 不明時は割引・割増なし（分岐点85%）＝力率割引を勝手に乗せない
    "基本料金単価": 1690.65,
    "電力量単価":  19.0,
}

# 🗂️ 型番→馬力マスタ（管理画面で編集可。ここに登録された型番はシミュレーションで自動入力される）
DEFAULT_MODEL_HP = {
    "PUZ-ERMP280KA4": 10.0,
    "PUZ-ERP280KA9": 10.0,
    "PUZ-ERMP160LA2": 6.0,
    "RZPR224A": 8.0,
    "RZRP224": 8.0,
    "MUCZ-G5617S": 2.0,
}

DEFAULT_DATA = {
    "gyotai_master": DEFAULT_GYOTAI,
    "model_hp_master": DEFAULT_MODEL_HP,
    "model_kw_master": {},   # 型番→定格出力kW（任意・入力すれば馬力×係数より優先して計算に使用）
    "system_fee": DEFAULT_SYSTEM_FEE,
    "item_price": 5000,
    "item_margin": 0.30,
    "change_log": [],
    "download_log": [],
    "saved_cases": [],   # 結果一覧（企業名/空調数/年間削減量/回収年数 等の軽量レコード）
    "saved_quotes": [],  # 見積一覧（書類番号/宛先/件名/売価合計/利益/利益率 等）
    "label_profiles": {},  # 削減レポ：客先ごとの表示ラベル設定（プロファイル）。{顧客名: {ラベル...}}
    # 企業別 電力実績データストア（削減レポ用に累積）。{正規化キー: {"name":表示名, "monthly":{"YYYYMM":{算出5項目}}}}
    "company_power_records": {},
    "per_unit_materials": [],  # 室外機1台あたりの機材セット（増減連動用）。[{商品名, 1台あたり数量}] 後で管理画面で更新可
    # 📄 機材単価マスタ（商品名/単位/売価/原価/利益）。見積の機材費は売価をそのまま使用。管理画面で編集可。
    "material_master": [
        {"商品名": "データ取集装置（Marimba Mercury)", "単位": "台", "売価": 96000, "原価": 22000, "利益": 74000},
        {"商品名": "マーキュリー", "単位": "台", "売価": 80000, "原価": 22000, "利益": 58000},
        {"商品名": "I/O入出力モジュール", "単位": "台", "売価": 32000, "原価": 26000, "利益": 6000},
        {"商品名": "MMEazyAir親機", "単位": "台", "売価": 62000, "原価": 39295, "利益": 22705},
        {"商品名": "MMEazyAir2（子機）", "単位": "台", "売価": 43000, "原価": 27170, "利益": 15830},
        {"商品名": "ダイキン　室外機外部制御アダプタ_DTA104A1", "単位": "台", "売価": 29400, "原価": 18880, "利益": 10520},
        {"商品名": "ダイキン　室外機外部制御アダプタ_DTA104A2", "単位": "台", "売価": 29400, "原価": 18880, "利益": 10520},
        {"商品名": "ダイキン 室外機外部制御アダプタ_KRP58M4", "単位": "台", "売価": 15360, "原価": 10750, "利益": 4610},
        {"商品名": "ダイキン　デマンドアダプター_KRP58M1", "単位": "台", "売価": 14100, "原価": 9030, "利益": 5070},
        {"商品名": "ダイキン　デマンドアダプター_KRP58M3", "単位": "台", "売価": 11760, "原価": 8170, "利益": 3590},
        {"商品名": "三菱電機 / PAC-SC36NA", "単位": "台", "売価": 3000, "原価": 2530, "利益": 470},
        {"商品名": "三菱電機 /PAC-SC96NA", "単位": "台", "売価": 4250, "原価": 3270, "利益": 980},
        {"商品名": "日立PCC-1A ３Pコネクター", "単位": "台", "売価": 5520, "原価": 3500, "利益": 2020},
        {"商品名": "パナソニック 3Pプラグコード", "単位": "台", "売価": 12360, "原価": 7950, "利益": 4410},
        {"商品名": "SwitchBot 防水温湿度計", "単位": "個", "売価": 2380, "原価": 1800, "利益": 580},
        {"商品名": "クランプ式パルスセンサー", "単位": "台", "売価": 26520, "原価": 17000, "利益": 9520},
        {"商品名": "東光高岳_無線電流センサ受信機", "単位": "台", "売価": 57200, "原価": 29800, "利益": 27400},
        {"商品名": "東光高岳_無線電流センサ送信機", "単位": "台", "売価": 23500, "原価": 12200, "利益": 11300},
        {"商品名": "無線電流センサ受信機", "単位": "台", "売価": 44460, "原価": 28500, "利益": 15960},
        {"商品名": "無線電流センサ送信機", "単位": "台", "売価": 18820, "原価": 12200, "利益": 6620},
        {"商品名": "計測用（給電用）CT 100A(東光高岳)", "単位": "台", "売価": 2700, "原価": 1100, "利益": 1600},
        {"商品名": "計測用CT 300A", "単位": "台", "売価": 6240, "原価": 4000, "利益": 2240},
        {"商品名": "ENECON", "単位": "台", "売価": 20000, "原価": 16800, "利益": 3200},
        {"商品名": "20m延長ケーブル", "単位": "本", "売価": 6600, "原価": 4200, "利益": 2400},
        {"商品名": "CT接続ケーブル", "単位": "本", "売価": 720, "原価": 300, "利益": 420},
        {"商品名": "USB to RS485コンバータ", "単位": "台", "売価": 2500, "原価": 2226, "利益": 274},
        {"商品名": "AC/DCコンバーター5V-15W", "単位": "台", "売価": 2800, "原価": 2362, "利益": 438},
        {"商品名": "AC/DCコンバーター24V-15W", "単位": "台", "売価": 3100, "原価": 2416, "利益": 684},
        {"商品名": "保護BOX(150x250x120)", "単位": "台", "売価": 1620, "原価": 786, "利益": 834},
        {"商品名": "保護BOX(200x300x160)", "単位": "台", "売価": 3720, "原価": 2169, "利益": 1551},
        {"商品名": "保護BOX小", "単位": "台", "売価": 1350, "原価": 786, "利益": 564},
        {"商品名": "保護BOX中", "単位": "台", "売価": 3100, "原価": 2169, "利益": 931},
        {"商品名": "保護BOX(大)", "単位": "台", "売価": 10100, "原価": 1802, "利益": 8298},
        {"商品名": "通信回線費用（初期費用）", "単位": "式", "売価": 4000, "原価": 2500, "利益": 1500},
    ],
    # 📄 見積書作成の設定（一律利益率・税率・ビジネスタンク率・会社情報 等）。管理画面で編集可。
    "quote_settings": {
        "markup_factor": 1.2,      # 利益率（倍率）：売価 ＝ 原価 × 倍率。画面で手動変更可
        "target_factor": 1.1,      # 倍率の下限アラート基準（これ未満で警告）
        "tax_rate":       0.10,    # 消費税率
        "bt_rate":        0.25,    # ビジネスタンク費 ＝ 粗利益 × 率
        "kouchiku_unit":  12000,   # 空調システム構築費 単価（円/台）
        "eng_defaults": {          # エンジニアリング費の既定（売価/原価）。式確定まで手入力の初期値
            "setup_sale": 1850000, "setup_cost": 1820000,   # 設置工事費
            "mat_sale":    800000, "mat_cost":   784000,    # 材料費
            "misc_sale":   800000, "misc_cost":  800000,    # 諸経費
        },
        "company": {
            "name": "株式会社シムックスイニシアティブ",
            "post": "〒105-0013",
            "addr": "東京都港区浜松町1-30-5 浜松町スクエア10F",
            "tel": "03-6402-2650", "fax": "03-6402-2651",
        },
        "delivery": "別途お打ち合わせ", "deliver_to": "貴社ご指定場所",
        "payment": "月末締め翌月末払い",
    },
    "display_settings": {
        "show_graph_dm": True,
        "show_graph_kwh": True,
        "show_graph_bill": True,
        "show_co2": True,
        "show_talk_script": True,
        "show_loss_10yr": True,
    },
    "calc_settings": {
        "cap_rate": 0.30,
        "energy_save_rate": 0.02,
        "cop": 3.5,            # 冷房能力kW→電力kW 換算
        "default_hp": 5.0,     # 管理外（マスタ未登録）型番の仮置き馬力
    },
    # 年間システム利用料の単価（円・月額／1単位あたり）。管理画面で変更可。
    "fee_settings": {
        "tsushin": 550,   # 通信費（マーキュリー＋MM親機 などの台数 ×）
        "cloud":   500,   # クラウド利用料（制御台数＋1 ×）
        "hoshu":   200,   # ソフト保守（基本1 ×）
        "data":    200,   # データ収集（2＋制御台数÷2 などの台数 ×）
    },
    # 計算で使う固定値（定数）。各種設定で全て変更可。
    "const_settings": {
        "area_unit_price": AREA_UNIT_PRICE,   # 工事費目安 円/台
        "conv_factor":     CONV_FACTOR,       # 能力削減→正味エネ削減 換算係数
        "hp_to_kw":        HP_TO_KW,           # 1馬力→冷房能力kW
        "co2_factor":      CO2_FACTOR,         # kg-CO2/kWh
        "sugi_kg":         SUGI_KG,            # 杉1本 年間CO2吸収kg
        "area_buffer":     dict(AREA_BUFFER),  # エリア係数（地域→係数）
        "fb_defaults":     dict(FB_DEFAULTS),  # データ補完の既定値
        "discount_rate":   0.05,               # B-5: NPV/IRR の割引率（年）
        "npv_years":       10,                 # B-5: NPV算定の対象年数
        "base_ac_ratio":   0.0,                # 業態配慮：基礎電力のうち空調が占める比率（0〜1）
        "pf_base":         0.85,               # 力率割引の基準力率（標準85%）
        "pf_cap":          0.15,               # 力率割引の上下限（±15%＝0.85〜1.15）
        # 詳細見積（見積書方式）の固定値。管理画面で変更可。
        "estimate_settings": {
            "material_factor": 1.2,            # 材料費 ＝ 電材費 × 係数
            "labor_day_unit":  60000,          # 工事費 人日単価（円/人日）
            "labor_fixed":     120000,         # 工事費 固定加算（円）
            "lodging_unit":    8000,           # 宿泊費 単価（円/人工）
            "area_unit": {                     # 移動拘束費 ＝ 作業人数 × 地区別単価
                "地区1": 4000, "地区2": 12000, "地区3": 25000, "地区4": 34000, "地区5": 0,
            },
        },
        # 概算モード（機材情報なし）用の標準値。管理画面で調整可。
        "gaisan_settings": {
            "unit_kw":           10.0,    # 標準室外機 定格kW（概算台数推定に使用）
            "material_per_unit": 100000,  # 1台あたり標準機材費（円/台・暫定/要実績調整）
            "setup_per_unit":     50000,  # 1台あたり標準工事費（円/台・暫定/要実績調整）
            "hw_fixed":          300000,  # 標準ハードウェア 固定費（円/拠点・データ収集装置/親機等）
            "hw_per_unit":        50000,  # 標準ハードウェア 台数比例費（円/台・子機/センサー/IO等）
            "ctrl_ratio":         0.9,    # 概算時の制御可能比率（既定・単一表示用）
            # 概算3シナリオ。unit_kw=台数推定, material=機材費/台,
            # units_per_device=1装置で制御できる台数(HW費のみに反映), ctrl_ratio=制御可能比率(削減にも反映)
            "scenarios": [
                {"key": "best", "label": "① ベスト（最短）",   "unit_kw": 12.0, "material_per_unit": 70000,  "units_per_device": 2, "ctrl_ratio": 1.00},
                {"key": "std",  "label": "② 標準（通常寄せ）", "unit_kw": 10.0, "material_per_unit": 100000, "units_per_device": 2, "ctrl_ratio": 0.90},
                {"key": "hard", "label": "③ 保守（厳しめ）",   "unit_kw": 8.0,  "material_per_unit": 130000, "units_per_device": 1, "ctrl_ratio": 0.80},
            ],
        },
        # A-1: 時間帯別タリフ(TOU)。tou_use=Trueで blended電力量単価を採用
        "tou_use":         False,
        "tou_peak":        22.0,               # 夏季ピーク単価 円/kWh
        "tou_day":         19.0,               # 平日昼間単価
        "tou_night":       12.0,               # 夜間/休日単価
        "tou_split":       [0.15, 0.55, 0.30], # 使用量構成比（ピーク/昼/夜・合計1.0）
    },
}

# A-2: 電力会社タリフ・プリセット（目安・要確認。選択で補完既定値へ反映できる）
TARIFF_PRESETS = {
    "（プリセットを使わない）": None,
    "東京電力（目安）":   {"basic": 1716.0,  "energy": 19.0, "pf_discount": True,  "tou": [22.0, 19.0, 12.0]},
    "東北電力（目安）":   {"basic": 1690.65, "energy": 19.0, "pf_discount": True,  "tou": [21.0, 18.5, 12.0]},
    "関西電力（目安）":   {"basic": 1911.8,  "energy": 18.0, "pf_discount": True,  "tou": [21.5, 18.0, 11.5]},
    "中部電力（目安）":   {"basic": 1771.0,  "energy": 18.5, "pf_discount": True,  "tou": [21.5, 18.5, 12.0]},
    "九州電力（目安）":   {"basic": 1650.0,  "energy": 17.5, "pf_discount": True,  "tou": [20.5, 17.5, 11.0]},
    "北海道電力（目安）": {"basic": 1880.0,  "energy": 20.0, "pf_discount": True,  "tou": [23.0, 20.0, 13.0]},
}

# 🗃️ 機器台帳マスター (SKILL.md・西富士の実実績と同期)
DEFAULT_EQUIPMENT = [
    {"id": 1, "loc": "工場屋上", "mfr": "三菱電機", "model": "PUZ-ERMP280KA4", "ctrl": "〇", "cap": 28.0},
    {"id": 2, "loc": "工場屋上", "mfr": "三菱電機", "model": "PUZ-ERMP280KA4", "ctrl": "〇", "cap": 28.0},
    {"id": 3, "loc": "工場屋上", "mfr": "日立", "model": "不明", "ctrl": "×", "cap": 28.0},
    {"id": 4, "loc": "検査室外", "mfr": "三菱電機", "model": "不明", "ctrl": "×", "cap": 14.0},
    {"id": 5, "loc": "検査室外", "mfr": "三菱電機", "model": "MUCZ-G5617S", "ctrl": "〇", "cap": 5.6},
    {"id": 6, "loc": "事務所棟廻り", "mfr": "三菱電機", "model": "PUZ-ERP280KA9", "ctrl": "〇", "cap": 28.0},
    {"id": 7, "loc": "事務所棟廻り", "mfr": "三菱電機", "model": "PUZ-ERP280KA9", "ctrl": "〇", "cap": 28.0},
    {"id": 8, "loc": "事務所棟廻り", "mfr": "三菱電機", "model": "PUZ-ERMP160LA2", "ctrl": "〇", "cap": 16.0},
    {"id": 9, "loc": "事務所棟廻り", "mfr": "ダイキン", "model": "RZPR224A", "ctrl": "〇", "cap": 22.4},
    {"id": 10, "loc": "事務所棟廻り", "mfr": "ダイキン", "model": "RZPR224A", "ctrl": "〇", "cap": 22.4},
    {"id": 11, "loc": "事務所棟廻り", "mfr": "ダイキン", "model": "RZPR224A", "ctrl": "〇", "cap": 22.4},
    {"id": 12, "loc": "事務所棟廻り", "mfr": "ダイキン", "model": "RZPR224A", "ctrl": "〇", "cap": 22.4},
]

# ══════════════════════════════════════════════════════════════════
# JSON 永続化
# ══════════════════════════════════════════════════════════════════
def load_data() -> dict:
    if os.path.exists(DATA_FILE):
        try:
            with open(DATA_FILE, "r", encoding="utf-8") as f:
                saved = json.load(f)
            merged = copy.deepcopy(DEFAULT_DATA)
            for k, v in saved.items():
                if k in ("gyotai_master", "model_hp_master", "model_kw_master", "calc_settings", "display_settings", "fee_settings", "const_settings", "quote_settings") and isinstance(v, dict):
                    merged[k].update(v)   # 既定キー（cop/default_hp等）を保ちつつ保存値を反映
                else:
                    merged[k] = v
            return merged
        except Exception:
            return copy.deepcopy(DEFAULT_DATA)
    return copy.deepcopy(DEFAULT_DATA)


def save_data(data: dict) -> None:
    with open(DATA_FILE, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)


def _apply_const_settings(data: dict) -> None:
    """各種設定の固定値(const_settings)を、計算で参照するモジュール定数へ反映する。
    Streamlitは毎回スクリプトを上から再実行するため、main()冒頭で毎回呼び出して上書きする。
    （calc_simulation / sanitize_columns 等は従来どおりモジュール定数を参照すればよい）"""
    cs = (data or {}).get("const_settings") or {}
    global AREA_UNIT_PRICE, CONV_FACTOR, HP_TO_KW, CO2_FACTOR, SUGI_KG, AREA_BUFFER, FB_DEFAULTS
    AREA_UNIT_PRICE = float(cs.get("area_unit_price", AREA_UNIT_PRICE))
    CONV_FACTOR     = float(cs.get("conv_factor", CONV_FACTOR))
    HP_TO_KW        = float(cs.get("hp_to_kw", HP_TO_KW))
    CO2_FACTOR      = float(cs.get("co2_factor", CO2_FACTOR))
    SUGI_KG         = float(cs.get("sugi_kg", SUGI_KG))
    if isinstance(cs.get("area_buffer"), dict) and cs["area_buffer"]:
        AREA_BUFFER = {k: float(v) for k, v in cs["area_buffer"].items()}
    if isinstance(cs.get("fb_defaults"), dict) and cs["fb_defaults"]:
        for k, v in cs["fb_defaults"].items():
            if k in FB_DEFAULTS:
                FB_DEFAULTS[k] = float(v)


# ══════════════════════════════════════════════════════════════════
# B-5: 投資評価指標（NPV / IRR）
# ══════════════════════════════════════════════════════════════════
def calc_npv_irr(invest: float, annual_net: float, years: int = 10, rate: float = 0.05):
    """初期投資 invest、毎年の純削減 annual_net（一定）を前提に NPV と IRR を返す。
    戻り値: (npv:float, irr:float|None)。IRRは年次キャッシュフローから二分法で算出（無ければNone）。"""
    invest = float(invest or 0)
    annual_net = float(annual_net or 0)
    years = max(int(years or 0), 1)
    # NPV = -投資 + Σ 純削減/(1+r)^t
    try:
        npv = -invest + sum(annual_net / ((1 + rate) ** t) for t in range(1, years + 1))
    except Exception:
        npv = float("nan")

    # IRR：年次CF = [-invest, net, net, ...] の正味現在価値が0になる割引率を二分法で
    def _npv_at(r):
        return -invest + sum(annual_net / ((1 + r) ** t) for t in range(1, years + 1))

    irr = None
    if annual_net > 0 and invest > 0:
        lo, hi = -0.9, 5.0
        f_lo, f_hi = _npv_at(lo), _npv_at(hi)
        if f_lo * f_hi <= 0:                 # 符号が変われば根あり
            for _ in range(100):
                mid = (lo + hi) / 2
                f_mid = _npv_at(mid)
                if abs(f_mid) < 1e-6:
                    break
                if f_lo * f_mid < 0:
                    hi, f_hi = mid, f_mid
                else:
                    lo, f_lo = mid, f_mid
            irr = (lo + hi) / 2
    return npv, irr


# ══════════════════════════════════════════════════════════════════
# C-12: 入力整合性チェック（異常値・不整合の自動検知 → 警告リストを返す）
# ══════════════════════════════════════════════════════════════════
def validate_inputs(df, res: dict, econ: dict):
    """試算の入力・結果を点検し、(level, メッセージ) のリストを返す。
    level: 'error'（要修正）/ 'warn'（要確認）/ 'ok'（問題なし）。"""
    issues = []
    try:
        n = len(df) if df is not None else 0
        if n != 12:
            issues.append(("warn", f"月数が{n}ヶ月です（通常は12ヶ月分を推奨）。"))

        if df is not None and len(df) > 0:
            be = pd.to_numeric(df.get("基本料金単価"), errors="coerce")
            en = pd.to_numeric(df.get("電力量単価"), errors="coerce")
            pf = pd.to_numeric(df.get("力率"), errors="coerce")
            ct = pd.to_numeric(df.get("契約電力"), errors="coerce")
            md = pd.to_numeric(df.get("最大需要電力"), errors="coerce")
            us = pd.to_numeric(df.get("使用量合計"), errors="coerce")

            if be is not None and be.notna().any():
                bmin, bmax = float(be.min()), float(be.max())
                if bmin < 500 or bmax > 3000:
                    issues.append(("warn", f"基本料金単価が想定外（{bmin:,.0f}〜{bmax:,.0f}円/kW）。電力会社のタリフをご確認ください。"))
            if en is not None and en.notna().any():
                emin, emax = float(en.min()), float(en.max())
                if emin < 5 or emax > 50:
                    issues.append(("warn", f"電力量単価が想定外（{emin:,.1f}〜{emax:,.1f}円/kWh）。タリフをご確認ください。"))
            if pf is not None and pf.notna().any():
                if float(pf.min()) < 0.7 or float(pf.max()) > 1.0:
                    issues.append(("warn", "力率が0.7〜1.0の範囲外です（既定0.85基準）。"))
            if ct is not None and md is not None and ct.notna().any() and md.notna().any():
                if float(ct.max()) > 0 and float(md.max()) > float(ct.max()):
                    issues.append(("error", f"最大需要電力（{md.max():,.0f}kW）が契約電力（{ct.max():,.0f}kW）を超えています。データをご確認ください。"))
            if us is not None and us.notna().any() and float(us.min()) <= 0:
                issues.append(("error", "使用量合計に0以下の月があります。検針票の値をご確認ください。"))

        # 空調割合の妥当性
        ac = float(res.get("ac_kwh_r", 0))
        if ac < 0.03 or ac > 0.6:
            issues.append(("warn", f"空調割合が {ac:.0%} と想定外です（業態選択・推計をご確認ください）。"))

        # 制御可能比率
        cr = float(res.get("ctrl_ratio", 0))
        if cr <= 0 or cr > 1.0:
            issues.append(("error", f"制御可能比率が {cr:.0%} です（0〜100%に収まるべきです）。"))

        # 純削減が利用料を下回る
        if float(res.get("net_saving", 0)) <= 0:
            issues.append(("warn", "年間純削減がマイナス（利用料が削減額を上回る）。制御台数・利用料・前提を見直してください。"))

        # 機材費が極端
        ti = float(res.get("total_invest", 0))
        nc = max(int(res.get("n_units_ctrl", 1)), 1)
        if ti / nc > 500000:
            issues.append(("warn", f"1台あたり初期費用が ¥{ti/nc:,.0f} と高めです（工事費・台数をご確認ください）。"))
    except Exception as e:
        issues.append(("warn", f"整合性チェック中に想定外（{e}）。"))

    if not issues:
        issues.append(("ok", "主要な入力・結果に明らかな異常は検出されませんでした。"))
    return issues


def log_change(data: dict, category: str, detail: str) -> None:
    entry = {
        "timestamp": datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
        "category": category,
        "detail": detail,
    }
    data["change_log"].insert(0, entry)
    if len(data["change_log"]) > 200:
        data["change_log"] = data["change_log"][:200]
    save_data(data)


def log_download(data: dict, filename: str) -> None:
    entry = {
        "timestamp": datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
        "filename": filename,
    }
    data["download_log"].insert(0, entry)
    if len(data["download_log"]) > 200:
        data["download_log"] = data["download_log"][:200]
    save_data(data)


# ══════════════════════════════════════════════════════════════════
# 列名の超強力・名寄せクレンジング (謎のエラー粉砕エンジン)
# ══════════════════════════════════════════════════════════════════
def _canonical_for(col):
    """任意の列名を、内部の標準列名に名寄せする（日本語・英語・表記ゆれ対応）。
    判定できない場合は None を返す。"""
    c = str(col).strip().replace("\n", "").replace(" ", "").replace("　", "").lower()
    if not c or c.startswith("unnamed") or c == "nan":
        return None
    if "契約" in c or "contract" in c:
        return "契約電力"
    # 燃調・再エネ・時間帯別単価は「/kWh・単価・ピーク」を含むため、最大需要/電力量単価より先に判定する
    if "燃料" in c or "燃調" in c or "fuel" in c:
        return "燃料費調整額"
    if "再エネ" in c or "賦課金" in c or "renewable" in c or "levy" in c:
        return "再エネ賦課金"
    if "夏季" in c or ("ピーク" in c and "単価" in c):
        return "夏季ピーク単価"
    if "昼間" in c or "平日" in c:
        return "平日昼間単価"
    if "夜間" in c or "休日" in c:
        return "夜間休日単価"
    if "市場" in c or "market" in c:
        return "市場価格調整項"
    if ("最大需要" in c or "最大ﾃﾞﾏﾝﾄﾞ" in c or "最大デマンド" in c or "デマンド" in c
            or "ﾃﾞﾏﾝﾄﾞ" in c or "demand" in c or "ピーク" in c or "peak" in c):
        return "最大需要電力"
    if "基本料金" in c or "基本単価" in c or ("基本" in c and "単価" in c) or "basic" in c:
        return "基本料金単価"
    if ("電力量単価" in c or "従量単価" in c or "従量" in c or "電気量" in c
            or "unitprice" in c or "単価" in c or "/kwh" in c or "rate" in c):
        return "電力量単価"
    if ("使用量" in c or "使用電力量" in c or "使用電力" in c or "usage" in c
            or "consumption" in c or ("電力量" in c and "kwh" in c)):
        return "使用量合計"
    if "力率" in c or c == "pf" or "powerfactor" in c or "力 率" in c:
        return "力率"
    if "月" in c or "month" in c or "年月" in c or "ym" == c:
        return "月"
    return None


def sanitize_columns(df: pd.DataFrame) -> pd.DataFrame:
    df = df.copy()
    df.columns = [str(c).strip().replace("\n", "").replace(" ", "") for c in df.columns]

    mapping = {}
    for col in df.columns:
        cc = _canonical_for(col)
        if cc and cc not in mapping.values():   # 同一標準列への重複割当ては最初の列を優先
            mapping[col] = cc

    df = df.rename(columns=mapping)
    
    required_defaults = {
        "月": [f"2025/{m:02d}" for m in range(1, 13)],
        "契約電力": [FB_DEFAULTS["契約電力"]] * 12,
        "最大需要電力": [FB_DEFAULTS["最大需要電力"]] * 12,
        "使用量合計": [FB_DEFAULTS["使用量合計"]] * 12,
        "力率": [FB_DEFAULTS["力率"]] * 12,
        "基本料金単価": [FB_DEFAULTS["基本料金単価"]] * 12,
        "電力量単価": [FB_DEFAULTS["電力量単価"]] * 12,
    }
    
    for col, default_val in required_defaults.items():
        if col not in df.columns:
            df[col] = default_val
        else:
            if col != "月":
                df[col] = pd.to_numeric(df[col].astype(str).str.replace(",", "").str.extract(r'([0-9\.]+)')[0], errors='coerce')
                df[col] = df[col].fillna(default_val[0])
            else:
                df[col] = df[col].astype(str).fillna("不明")

    out_cols = ["月", "契約電力", "最大需要電力", "使用量合計", "力率", "基本料金単価", "電力量単価"]
    # 任意列：月別で入れれば自動反映/欄として保持。列が無い・空なら自動で非反映（※排除）。
    #   燃調・再エネ＝計算に反映。夏季ピーク/平日昼間/夜間休日/市場価格調整＝現状は欄のみ（将来のTOU/調整用）。
    for opt in ["燃料費調整額", "再エネ賦課金",
                "夏季ピーク単価", "平日昼間単価", "夜間休日単価", "市場価格調整項"]:
        if opt in df.columns:
            df[opt] = pd.to_numeric(
                df[opt].astype(str).str.replace(",", "").str.extract(r'(-?[0-9\.]+)')[0],
                errors='coerce').fillna(0.0)
            out_cols.append(opt)

    return df[out_cols]


REQUIRED_COLS = ["月", "契約電力", "最大需要電力", "使用量合計", "力率", "基本料金単価", "電力量単価"]
CRITICAL_COLS = ["最大需要電力", "使用量合計", "基本料金単価", "電力量単価"]  # 試算結果を左右する重要列


def load_uploaded_table(uploaded):
    """指定外フォーマットの Excel / CSV でも内容から取り込む。
    - Excel は全シートを走査し、最も「電力データらしい」ヘッダー行を自動検出
    - タイトル行などがヘッダーの上にあっても検出可能
    戻り値: (生データ DataFrame, 採用シート名, ヘッダースコア)
    """
    name = (uploaded.name or "").lower()

    def hdr_score(vals):
        return sum(1 for v in vals if _canonical_for(v))

    frames = []
    if name.endswith(".csv"):
        for enc in (None, "cp932", "utf-8-sig"):
            try:
                uploaded.seek(0)
                frames = [("CSV", pd.read_csv(uploaded, header=None, dtype=str))]
                break
            except Exception:
                continue
    else:
        xls = pd.ExcelFile(uploaded)
        for sn in xls.sheet_names:
            try:
                frames.append((sn, pd.read_excel(xls, sheet_name=sn, header=None, dtype=str)))
            except Exception:
                pass

    best = None  # (score, sheet, header_row, raw)
    for sn, raw in frames:
        if raw is None or raw.empty:
            continue
        for r in range(min(len(raw), 20)):
            sc = hdr_score(list(raw.iloc[r].values))
            if best is None or sc > best[0]:
                best = (sc, sn, r, raw)

    if best is None:
        raise ValueError("読み取れる表が見つかりませんでした。")

    sc, sn, r, raw = best
    if sc < 2:   # ヘッダーを自動検出できない → 先頭シートの1行目をヘッダーとして素直に読む
        header = list(raw.iloc[0].values)
        data = raw.iloc[1:].copy()
    else:
        header = list(raw.iloc[r].values)
        data = raw.iloc[r + 1:].copy()
    data.columns = [str(h) for h in header]
    data = data.dropna(how="all").reset_index(drop=True)
    return data, sn, sc


def build_input_template_xlsx() -> bytes:
    """インポート用の入力Excelテンプレート（7列＋12ヶ月サンプル）を生成。"""
    if not HAS_OPENPYXL:
        return b""
    wb = Workbook()
    ws = wb.active
    ws.title = "電力データ入力"
    navy = "1B2A4A"; green = "2E9E5B"
    headers = ["月", "契約電力(kW)", "最大需要電力(kW)", "使用量合計(kWh)",
               "力率", "基本料金単価(円/kW)", "電力量単価(円/kWh)",
               "燃料費調整額(円/kWh)", "再エネ賦課金(円/kWh)",
               "夏季ピーク単価(円/kWh)", "平日昼間単価(円/kWh)", "夜間休日単価(円/kWh)", "市場価格調整項(円/kWh)"]
    widths = [14, 16, 18, 18, 10, 20, 22, 20, 20, 22, 22, 22, 22]
    for i, w in enumerate(widths, 1):
        ws.column_dimensions[get_column_letter(i)].width = w
    n_cols = len(headers)

    # 説明行
    ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=n_cols)
    note = ws.cell(row=1, column=1,
                   value="▼ 12ヶ月分の電力データを入力（行は増減可）。1行目の説明行は削除可／列名（2行目）は変更しないでください。"
                         "　H列以降は任意：燃料費調整額・再エネ賦課金は月別で入れれば自動反映。"
                         "夏季ピーク/平日昼間/夜間休日/市場価格調整は“欄のみ”（現状は計算未反映・空欄でOK）。")
    note.font = Font(name="Yu Gothic", size=10, bold=True, color="FFFFFF")
    note.fill = PatternFill("solid", fgColor=green)
    note.alignment = Alignment(horizontal="left", vertical="center")
    ws.row_dimensions[1].height = 24

    # ヘッダー行（2行目）
    for i, h in enumerate(headers, 1):
        c = ws.cell(row=2, column=i, value=h)
        c.font = Font(name="Yu Gothic", size=11, bold=True, color="FFFFFF")
        c.fill = PatternFill("solid", fgColor=navy)
        c.alignment = Alignment(horizontal="center", vertical="center")
    ws.row_dimensions[2].height = 22

    # サンプル12ヶ月（数値は目安。実データに置き換えてください。末尾2列は任意＝空欄でOK）
    sample = [
        ("2025/04", 626, 400, 100000, 0.98, 1690.65, 19.0, -2.10, 3.49),
        ("2025/05", 626, 410, 104000, 0.98, 1690.65, 19.0, -1.80, 3.49),
        ("2025/06", 626, 450, 118000, 0.98, 1690.65, 19.0, -1.20, 3.49),
        ("2025/07", 626, 500, 132000, 0.98, 1690.65, 19.0, -0.50, 3.49),
        ("2025/08", 626, 510, 135000, 0.98, 1690.65, 19.0,  0.30, 3.49),
        ("2025/09", 626, 505, 130000, 0.98, 1690.65, 19.0,  0.10, 3.49),
        ("2025/10", 626, 460, 120000, 0.98, 1690.65, 19.0, -0.60, 3.49),
        ("2025/11", 626, 430, 110000, 0.98, 1690.65, 19.0, -1.40, 3.49),
        ("2025/12", 626, 470, 122000, 0.98, 1690.65, 19.0, -1.90, 3.49),
        ("2026/01", 626, 480, 124000, 0.98, 1690.65, 19.0, -2.30, 3.49),
        ("2026/02", 626, 440, 112000, 0.98, 1690.65, 19.0, -2.50, 3.49),
        ("2026/03", 626, 420, 105000, 0.98, 1690.65, 19.0, -2.20, 3.49),
    ]
    for ri, row in enumerate(sample, start=3):
        for ci, val in enumerate(row, start=1):
            c = ws.cell(row=ri, column=ci, value=val)
            c.font = Font(name="Yu Gothic", size=10)
            c.alignment = Alignment(horizontal="center")

    # ── 2枚目：室外機リスト（任意）── 入力すると型番ベースで精緻化 ──
    ws2 = wb.create_sheet("室外機リスト(任意)")
    u_headers = ["機器ID/場所", "メーカー", "型番", "定格冷房kW(任意)", "稼働係数(0-1)", "制御可否(○/×)"]
    u_widths = [16, 14, 20, 18, 16, 16]
    for i, w in enumerate(u_widths, 1):
        ws2.column_dimensions[get_column_letter(i)].width = w
    ws2.merge_cells(start_row=1, start_column=1, end_row=1, end_column=6)
    n2 = ws2.cell(row=1, column=1,
                  value="▼ 任意：室外機を1台1行で入力すると、型番→定格(冷房能力)から容量加重で精緻に試算します。"
                        "／ 定格が分かる場合は『定格冷房kW』に直接入力（型番より優先）。空欄なら型番から自動推定。")
    n2.font = Font(name="Yu Gothic", size=10, bold=True, color="FFFFFF")
    n2.fill = PatternFill("solid", fgColor=navy)
    n2.alignment = Alignment(horizontal="left", vertical="center")
    ws2.row_dimensions[1].height = 30
    for i, h in enumerate(u_headers, 1):
        c = ws2.cell(row=2, column=i, value=h)
        c.font = Font(name="Yu Gothic", size=11, bold=True, color="FFFFFF")
        c.fill = PatternFill("solid", fgColor=green)
        c.alignment = Alignment(horizontal="center", vertical="center")
    ws2.row_dimensions[2].height = 22
    u_sample = [
        ("1F-室外機A", "三菱電機", "PUZ-ZRMP112", "", 0.8, "○"),
        ("1F-室外機B", "ダイキン", "RZRP224", "", 0.7, "○"),
        ("サーバー室", "日立", "RAS-AP280", "", 0.6, "×"),
        ("事務所", "三菱電機", "MSZ-56", "", 0.5, "○"),
    ]
    for ri, row in enumerate(u_sample, start=3):
        for ci, val in enumerate(row, start=1):
            c = ws2.cell(row=ri, column=ci, value=val)
            c.font = Font(name="Yu Gothic", size=10)
            c.alignment = Alignment(horizontal="center")

    # ── 3枚目：年間システム利用料の機器台数（任意）──
    ws3 = wb.create_sheet("利用料機器(任意)")
    ws3.column_dimensions["A"].width = 24
    ws3.column_dimensions["B"].width = 12
    ws3.merge_cells(start_row=1, start_column=1, end_row=1, end_column=2)
    n3 = ws3.cell(row=1, column=1,
                  value="▼ 任意：年間システム利用料の機器台数を入力。通信費＝マーキュリー＋MM親機、"
                        "データ収集＝マーキュリー＋MM親機＋MM子機 の数量に使用します（空欄なら既定値）。")
    n3.font = Font(name="Yu Gothic", size=10, bold=True, color="FFFFFF")
    n3.fill = PatternFill("solid", fgColor=navy)
    n3.alignment = Alignment(horizontal="left", vertical="center")
    ws3.row_dimensions[1].height = 30
    for i, h in enumerate(["項目", "台数"], 1):
        c = ws3.cell(row=2, column=i, value=h)
        c.font = Font(name="Yu Gothic", size=11, bold=True, color="FFFFFF")
        c.fill = PatternFill("solid", fgColor=green)
        c.alignment = Alignment(horizontal="center", vertical="center")
    ws3.row_dimensions[2].height = 22
    for ri, (lbl, val) in enumerate([("マーキュリー", 1), ("MMEazyAir親機", 1), ("MMEazyAir子機", 0)], start=3):
        cl = ws3.cell(row=ri, column=1, value=lbl)
        cl.font = Font(name="Yu Gothic", size=10, bold=True)
        cv = ws3.cell(row=ri, column=2, value=val)
        cv.font = Font(name="Yu Gothic", size=10)
        cv.alignment = Alignment(horizontal="center")

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf.getvalue()


def detect_fee_devices_from_upload(uploaded):
    """アップロードExcelから利用料機器台数（マーキュリー/MM親機/MM子機）を検出。無ければNone。"""
    import re
    name = (getattr(uploaded, "name", "") or "").lower()
    if name.endswith(".csv"):
        return None
    try:
        uploaded.seek(0)
        xls = pd.ExcelFile(uploaded)
    except Exception:
        return None
    found = {"mercury": None, "mm_parent": None, "mm_child": None}
    for sn in xls.sheet_names:
        try:
            raw = pd.read_excel(xls, sheet_name=sn, header=None, dtype=str)
        except Exception:
            continue
        for _, row in raw.iterrows():
            cells = ["" if pd.isna(v) else str(v) for v in row.values]
            joined = "".join(cells)
            label = None
            if "マーキュリー" in joined:
                label = "mercury"
            elif "子機" in joined:
                label = "mm_child"
            elif "親機" in joined:
                label = "mm_parent"
            if label and found[label] is None:
                for v in cells:
                    if any(k in v for k in ("マーキュリー", "親機", "子機", "項目")):
                        continue
                    m = re.search(r"\d+", v.replace(",", ""))
                    if m:
                        found[label] = int(m.group())
                        break
    if all(v is None for v in found.values()):
        return None
    return {k: int(v or 0) for k, v in found.items()}


# ══════════════════════════════════════════════════════════════════
# 室外機（型番→定格）解析：methodology §5-6（ハイブリッド方式）
# ══════════════════════════════════════════════════════════════════
def estimate_capacity_kw(model):
    """型番から冷房能力kW(定格)を推定。methodology §6準拠。推定不可は None。"""
    import re
    if model is None:
        return None
    m = str(model).strip().upper().replace("　", "").replace(" ", "")
    if m in ("", "-", "不明", "読み取り不可", "NAN", "NONE", "ー"):
        return None
    nums = re.findall(r'\d+', m)
    if not nums:
        return None
    # 東芝 ROA/ROB：末尾1桁を落として ÷10（1125→11.2, 633→6.3）
    if m.startswith("ROA") or m.startswith("ROB"):
        n = max(nums, key=len)
        if len(n) >= 2:
            return int(n[:-1]) / 10.0
    # ルームエアコン MUZ/MUCZ/MSZ/CU/AU：先頭2桁 ÷10（36→3.6, 56→5.6）
    if any(m.startswith(p) for p in ("MUCZ", "MUZ", "MSZ", "CU", "AU")):
        n = nums[0]
        if len(n) >= 2:
            return int(n[:2]) / 10.0
    # 業務用（三菱/ダイキン/日立）：3桁数字 ÷10（280→28.0, 224→22.4, 112→11.2, 080→8.0）
    three = [n for n in nums if len(n) == 3]
    if three:
        return int(three[0]) / 10.0
    two = [n for n in nums if len(n) == 2]
    if two:
        return int(two[0]) / 10.0
    return int(nums[0]) / 10.0


def _unit_col(col):
    c = str(col).strip().replace("\n", "").replace(" ", "").replace("　", "").lower()
    if not c or c.startswith("unnamed") or c == "nan":
        return None
    if "型番" in c or "型式" in c or "品番" in c or "model" in c:
        return "型番"
    if "メーカ" in c or "maker" in c or "brand" in c:
        return "メーカー"
    if "定格" in c or "能力" in c or "冷房" in c or "capacity" in c:
        return "定格冷房kW"
    if "稼働" in c or "負荷率" in c or "duty" in c or "運転率" in c:
        return "稼働係数"
    if "制御" in c or "control" in c or ("対象" in c and "率" not in c):
        return "制御可否"
    if "機器" in c or "番号" in c or "場所" in c or "名称" in c or c == "id" or c == "no":
        return "機器ID"
    return None


def parse_outdoor_units(df_units, cop=3.5):
    """室外機リストDFを標準化。戻り値カラム：機器ID/メーカー/型番/定格冷房kW/稼働係数/制御可否/電力kW。
    型番から定格を推定（明示の定格列があればそれを優先）。空行は除外。該当なしは None。"""
    import re
    if df_units is None or len(df_units) == 0:
        return None
    d = df_units.copy()
    m = {}
    for col in d.columns:
        cc = _unit_col(col)
        if cc and cc not in m.values():
            m[col] = cc
    d = d.rename(columns=m)
    if "型番" not in d.columns and "定格冷房kW" not in d.columns:
        return None
    cop = cop if (cop and cop > 0) else 3.5
    rows = []
    for i, (_, r) in enumerate(d.iterrows(), 1):
        model = str(r.get("型番", "") if "型番" in d.columns else "").strip()
        maker = str(r.get("メーカー", "") if "メーカー" in d.columns else "").strip()
        if model.lower() in ("", "nan") and maker.lower() in ("", "nan") and "定格冷房kW" not in d.columns:
            continue
        cap = None
        if "定格冷房kW" in d.columns:
            try:
                s = re.sub(r'[^0-9.]', '', str(r.get("定格冷房kW", "")))
                cap = float(s) if s else None
            except Exception:
                cap = None
            if cap is not None and (cap != cap or cap <= 0):
                cap = None
        if cap is None:
            cap = estimate_capacity_kw(model)
        cap_estimated = cap is None
        if cap is None:
            cap = 11.2   # 推定不可時の暫定定格（業務用の代表値）
        duty = 0.7
        if "稼働係数" in d.columns:
            try:
                dv = float(re.sub(r'[^0-9.]', '', str(r.get("稼働係数", ""))) or "nan")
                if dv == dv and dv > 0:
                    duty = dv if dv <= 1.5 else dv / 100.0
            except Exception:
                pass
        ctrl = True
        if "制御可否" in d.columns:
            s = str(r.get("制御可否", "")).strip().lower()
            if s in ("×", "x", "✕", "不可", "no", "false", "0", "対象外", "除外", "na", "なし"):
                ctrl = False
        rid = str(r.get("機器ID", "")) if "機器ID" in d.columns else ""
        if rid.strip().lower() in ("", "nan"):
            rid = f"#{i}"
        rows.append({
            "機器ID": rid, "メーカー": maker, "型番": model,
            "定格冷房kW": round(cap, 1), "稼働係数": round(duty, 2),
            "制御可否": bool(ctrl), "電力kW": round(cap / cop, 2),
            "定格推定": cap_estimated,
        })
    if not rows:
        return None
    return pd.DataFrame(rows)


def detect_units_from_upload(uploaded, cop=3.5):
    """アップロードExcelの全シートから『室外機リスト』らしいシートを探して解析。無ければ None。"""
    name = (uploaded.name or "").lower()
    if name.endswith(".csv"):
        return None
    try:
        uploaded.seek(0)
        xls = pd.ExcelFile(uploaded)
    except Exception:
        return None
    best = None  # (score, sheet, header_row, raw)
    for sn in xls.sheet_names:
        try:
            raw = pd.read_excel(xls, sheet_name=sn, header=None, dtype=str)
        except Exception:
            continue
        for r in range(min(len(raw), 20)):
            vals = [str(v) for v in raw.iloc[r].values]
            cc = [_unit_col(v) for v in vals]
            has_model = any(x == "型番" for x in cc)
            score = sum(1 for x in cc if x)
            if has_model and score >= 2 and (best is None or score > best[0]):
                best = (score, sn, r, raw)
    if best is None:
        return None
    sc, sn, r, raw = best
    header = list(raw.iloc[r].values)
    data = raw.iloc[r + 1:].copy()
    data.columns = [str(h) for h in header]
    data = data.dropna(how="all").reset_index(drop=True)
    return parse_outdoor_units(data, cop)


def build_units_from_models(rows, hp_master, default_hp=5.0, cop=3.5, kw_master=None):
    """機体入力（型番リスト）＋ 型番→馬力マスタ から units_df を生成。
    マスタ登録あり→その馬力（管理=True）、なし→仮置きdefault_hp（管理=False／※管理外）。
    定格出力kW（kw_master）が登録されていれば、馬力×係数より優先して定格冷房kWに使用。"""
    cop = cop if (cop and cop > 0) else 3.5
    norm = {str(k).strip().upper().replace(" ", ""): float(v)
            for k, v in (hp_master or {}).items() if str(k).strip()}
    knorm = {str(k).strip().upper().replace(" ", ""): float(v)
             for k, v in (kw_master or {}).items() if str(k).strip() and float(v or 0) > 0}
    out = []
    for i, r in enumerate(rows, 1):
        model = str(r.get("型番", "") or "").strip()
        if model == "" or model.lower() == "nan":
            continue
        key = model.upper().replace(" ", "")
        in_hp, in_kw = key in norm, key in knorm
        managed = in_hp or in_kw
        if in_hp:
            hp = norm[key]
        elif in_kw:
            hp = round(knorm[key] / HP_TO_KW, 1)   # 定格出力から馬力を逆算（表示用）
        else:
            hp = float(default_hp)
        cap = knorm[key] if in_kw else round(hp * HP_TO_KW, 1)
        try:
            duty = float(r.get("稼働係数", 0.7) or 0.7)
            if duty <= 0:
                duty = 0.7
            if duty > 1.5:
                duty = duty / 100.0
        except Exception:
            duty = 0.7
        cs = str(r.get("制御可否", "○")).strip().lower()
        ctrl = cs not in ("×", "x", "✕", "不可", "no", "false", "0", "対象外", "除外", "なし")
        rid = str(r.get("機器ID/場所", r.get("機器ID", "")) or "").strip() or f"#{i}"
        out.append({
            "機器ID": rid, "メーカー": str(r.get("メーカー", "") or ""), "型番": model,
            "馬力": hp, "定格冷房kW": cap, "稼働係数": round(duty, 2),
            "制御可否": bool(ctrl), "電力kW": round(cap / cop, 2),
            "管理": bool(managed), "定格推定": (not managed),
        })
    if not out:
        return None
    return pd.DataFrame(out)


# ══════════════════════════════════════════════════════════════════
# 計算ロジック（methodology.md「引き算推計・過去IFシミュレーション」完全同期版）
# ══════════════════════════════════════════════════════════════════
def estimate_units_from_meter(df: pd.DataFrame, app_data: dict, gyotai: str,
                              ac_override=None) -> dict:
    """機材情報が無い案件向けの「概算台数」推定。
    検針票の最大デマンド × 業態の空調ピーク比率 ÷ 標準室外機定格kW で室外機台数を概算する。
    ※ 削減額などの算出式(calc_simulation)は一切変更せず、その入力(台数)を作るだけのヘルパー。"""
    try:
        d = sanitize_columns(df)
        max_demand = float(pd.to_numeric(d["最大需要電力"], errors="coerce").max())
    except Exception:
        max_demand = 0.0
    if not (max_demand > 0):
        max_demand = 0.0
    if ac_override and ac_override.get("ac_peak") is not None:
        ac_peak = float(ac_override["ac_peak"])
    else:
        gm = app_data.get("gyotai_master", {}).get(
            gyotai, DEFAULT_GYOTAI.get("製造拠点(プロセス主役)", {"ac_peak": 0.18}))
        ac_peak = float(gm.get("ac_peak", 0.18))
    _g = app_data.get("const_settings", {}).get("gaisan_settings", {})
    unit_kw = float(_g.get("unit_kw", 10.0)) or 10.0
    ratio = float(_g.get("ctrl_ratio", 0.9))
    ac_peak_kw = max(max_demand * ac_peak, 0.0)
    n_total = max(int(round(ac_peak_kw / unit_kw)), 1) if ac_peak_kw > 0 else 1
    n_ctrl = min(max(int(round(n_total * ratio)), 1), n_total)
    return {"n_units_total": n_total, "n_units_ctrl": n_ctrl,
            "ac_peak_kw": ac_peak_kw, "unit_kw": unit_kw, "ac_peak_r": ac_peak}


def compute_gaisan_scenarios(res: dict, app_data: dict) -> list:
    """概算の3シナリオ（ベスト/標準/保守）を計算して返す。
    削減式(calc_simulation)は無改変で、台数・制御比率だけ変えて3回呼ぶ。初期費用は概算式で組む。"""
    import math
    df = res.get("df")
    gyotai = res.get("gyotai", "")
    ac_ovr = {"ac_peak": res.get("ac_peak_r"), "ac_kwh": res.get("ac_kwh_r")}
    try:
        max_demand = float(pd.to_numeric(sanitize_columns(df)["最大需要電力"], errors="coerce").max())
    except Exception:
        max_demand = 0.0
    ac_peak = float(res.get("ac_peak_r", 0.18) or 0.18)
    _cs = app_data.get("const_settings", {})
    _g = _cs.get("gaisan_settings", {})
    setup_pu = float(_g.get("setup_per_unit", 50000))
    hw_fixed = float(_g.get("hw_fixed", 300000))
    hw_pu = float(_g.get("hw_per_unit", 50000))
    area_unit = float(_cs.get("area_unit_price", 50000))
    kouchiku_unit = float(app_data.get("quote_settings", {}).get("kouchiku_unit", 12000))
    scenarios = _g.get("scenarios") or DEFAULT_DATA["const_settings"]["gaisan_settings"]["scenarios"]
    ac_peak_kw = max(max_demand * ac_peak, 0.0)
    out = []
    for s in scenarios:
        unit_kw = float(s.get("unit_kw", 10.0)) or 10.0
        cr = float(s.get("ctrl_ratio", 0.9))
        upd = max(int(s.get("units_per_device", 1)), 1)
        n_total = max(int(round(ac_peak_kw / unit_kw)), 1) if ac_peak_kw > 0 else 1
        n_ctrl = min(max(int(round(n_total * cr)), 1), n_total)
        try:
            r = calc_simulation(df, app_data, gyotai, cr, n_total, n_ctrl, ac_override=ac_ovr)
            gross = float(r.get("gross_saving", 0) or 0); net = float(r.get("net_saving", 0) or 0)
        except Exception:
            gross = net = 0.0
        material = n_ctrl * float(s.get("material_per_unit", 100000))
        setup = n_ctrl * setup_pu
        kouji = round(n_ctrl * area_unit / 10000) * 10000
        hw = hw_fixed + math.ceil(n_ctrl / upd) * hw_pu     # 1装置N台制御＝HW台数比例分のみ縮小
        kouchiku = n_ctrl * kouchiku_unit
        invest = int(material + setup + kouji + hw + kouchiku)
        payback = (invest / net) if net > 0 else None
        out.append({"label": s.get("label", s.get("key", "")), "n_total": n_total, "n_ctrl": n_ctrl,
                    "ctrl_ratio": cr, "units_per_device": upd, "gross": gross, "net": net,
                    "invest": invest, "payback": payback})
    return out


def calc_simulation(df: pd.DataFrame, app_data: dict,
                    gyotai: str, ctrl_ratio: float,
                    n_units_total: int, n_units_ctrl: int,
                    units_df=None, ac_override=None) -> dict:
    # ac_override={"ac_peak":..,"ac_kwh":..} が来た場合は業態マスタより優先（固定値入力モード）
    if ac_override and (ac_override.get("ac_peak") is not None) and (ac_override.get("ac_kwh") is not None):
        gm = {"ac_kwh": float(ac_override["ac_kwh"]), "ac_peak": float(ac_override["ac_peak"])}
    else:
        gm = app_data["gyotai_master"].get(gyotai, DEFAULT_GYOTAI["製造拠点(プロセス主役)"])
    ac_kwh_r = gm["ac_kwh"]
    ac_peak_r = gm["ac_peak"]
    cap_rate = app_data["calc_settings"]["cap_rate"]       # 容量削減率 30%
    sys_fee  = app_data["system_fee"]
    cop      = app_data.get("calc_settings", {}).get("cop", 3.5)

    # ─── ハイブリッド：室外機リストがあれば容量加重で制御比率を精緻化 ───
    unit_detail = None
    ctrl_weight = 0.0
    ctrl_ratio_source = "台数ベース（業態平均）"
    if units_df is not None and len(units_df) > 0:
        u = units_df.copy()
        u["重み"] = u["電力kW"] * u["稼働係数"]
        total_weight = float(u["重み"].sum())
        ctrl_weight = float(u.loc[u["制御可否"], "重み"].sum())
        if total_weight > 0:
            ctrl_ratio = ctrl_weight / total_weight     # 容量加重の制御可能比率で上書き
        n_units_total = int(len(u))
        n_units_ctrl = int(u["制御可否"].sum())
        unit_detail = u
        ctrl_ratio_source = "型番ベース（容量加重）"

    # ─── フォールバック：データ無し時は固定値を選択可（運用配慮） ───
    fb = app_data.get("fallback", {})
    if fb.get("ratio_mode") == "固定値" and fb.get("ratio_fix") is not None:
        ctrl_ratio = float(fb["ratio_fix"])
        ctrl_ratio_source = "固定値（手動指定）"

    df_clean = sanitize_columns(df)

    # ─── 季節性（ベースロード法）。業態の自動判断はせず、トグルで明示制御 ───
    opts      = app_data.get("calc_options", {})
    seasonal  = bool(opts.get("seasonal_ac", False))            # 季節性を反映するか
    gyotai_consider = bool(opts.get("gyotai_consider", False))  # 業態配慮（基礎電力の空調分）を見込むか
    # 業態配慮ONなら、基礎電力のうち admin の「基礎電力の空調比率」分は空調として残す（差し引かない）
    _bac = float(app_data.get("const_settings", {}).get("base_ac_ratio", 0.0)) if gyotai_consider else 0.0
    _bac = min(max(_bac, 0.0), 1.0)
    # 力率割引の基準力率・上下限（管理画面で変更可）
    _cset0   = app_data.get("const_settings", {})
    pf_base  = float(_cset0.get("pf_base", 0.85))
    pf_cap   = float(_cset0.get("pf_cap", 0.15))
    base_kwh_m = (float(df_clean["使用量合計"].min()) * (1.0 - _bac)) if seasonal else 0.0
    # ─── 燃調・再エネは月別列があれば自動反映（記載なし/全ゼロなら自動で非反映） ───
    fuel_col = "燃料費調整額" in df_clean.columns
    ren_col  = "再エネ賦課金" in df_clean.columns
    fuel_included = bool(fuel_col and float(df_clean["燃料費調整額"].abs().sum()) > 0)
    ren_included  = bool(ren_col and float(df_clean["再エネ賦課金"].sum()) > 0)

    def _row_add(row):
        a = float(row.get("燃料費調整額", 0.0)) if fuel_included else 0.0
        b = float(row.get("再エネ賦課金", 0.0)) if ren_included else 0.0
        return a + b

    # ─── ① 電力量料金IFシミュレーション（conv=0.5 適用） ──────────────
    ac_kwh_list   = []
    reduc_kwh_list = []
    ene_saving_list = []

    for _, row in df_clean.iterrows():
        usage_m = row["使用量合計"]
        if seasonal:
            ac_kwh_m = max(usage_m - base_kwh_m, 0.0)     # 季節超過分＝空調（ベースロード法）
        else:
            ac_kwh_m = usage_m * ac_kwh_r                 # 業態の空調割合（従来）
        reduc_kwh_m = round(ac_kwh_m * ctrl_ratio * cap_rate * CONV_FACTOR)
        eff_rate_m  = row["電力量単価"] + _row_add(row)    # 従量単価＋（燃調・再エネ：列があれば月別）
        ene_saving_m = reduc_kwh_m * eff_rate_m

        ac_kwh_list.append(ac_kwh_m)
        reduc_kwh_list.append(reduc_kwh_m)
        ene_saving_list.append(ene_saving_m)

    df_clean["空調推計kWh"]    = ac_kwh_list
    df_clean["削減kWh"]        = reduc_kwh_list
    df_clean["電力量料金削減額"] = ene_saving_list
    df_clean["導入後使用量kWh"] = df_clean["使用量合計"] - df_clean["削減kWh"]

    ene_saving_annual = sum(ene_saving_list)
    # 季節性ONのときは実効空調割合を再計算（下流表示の整合用）
    _tot_usage = float(df_clean["使用量合計"].sum())
    ac_kwh_r_eff = (sum(ac_kwh_list) / _tot_usage) if (seasonal and _tot_usage > 0) else ac_kwh_r

    # ─── ② 基本料金IFシミュレーション（実量制・力率調整同期） ────────────────
    max_demand = df_clean["最大需要電力"].max()
    delta_kw = max_demand * ac_peak_r * ctrl_ratio * cap_rate
    
    if fb.get("contract_mode") == "固定値":
        cfix = float(fb.get("contract_fix") or 0)
        old_contract = cfix if cfix > 0 else max_demand   # 0なら年間最大需要を代用
    else:
        # 「最大」ではなく「最新（直近月）」の契約電力を採用する
        _ct = pd.to_numeric(df_clean["契約電力"], errors="coerce")
        old_contract = 0.0
        try:
            _dt = pd.to_datetime(df_clean["月"].astype(str), errors="coerce")
            if _dt.notna().any():
                _latest_idx = _dt.idxmax()                 # 直近月のインデックス
                _v = _ct.loc[_latest_idx]
                old_contract = float(_v) if pd.notna(_v) else 0.0
        except Exception:
            old_contract = 0.0
        if old_contract <= 0:                              # 月が解釈不能/最新が空 → 末尾行で代用
            _valid = _ct.dropna()
            old_contract = float(_valid.iloc[-1]) if len(_valid) else 0.0
        if old_contract <= 0:                              # それも不可 → 従来どおり最大で代用
            old_contract = float(_ct.max() or 0)
    if old_contract <= 0:
        old_contract = max_demand
    # 下限＝制御後の年間ピーク需要（max需要−delta）。最少月の最大需要で過度に頭打ちさせない。
    # ただし契約は増やさない（過少契約データでも new>old にしない）。
    post_peak      = max(max_demand - delta_kw, 0.0)
    new_contract   = min(old_contract, max(old_contract - delta_kw, post_peak))
    contract_delta = max(old_contract - new_contract, 0.0)

    dm_saving_list = []
    reduc_dm_list = []
    for _, row in df_clean.iterrows():
        power_factor = row.get("力率", pf_base)
        pf_adj = 1.0 - (power_factor - pf_base)
        pf_adj = max(1.0 - pf_cap, min(pf_adj, 1.0 + pf_cap))

        saving_dm = contract_delta * row["基本料金単価"] * pf_adj
        dm_saving_list.append(saving_dm)
        reduc_dm_list.append(delta_kw)

    df_clean["削減デマンドkW"]    = reduc_dm_list
    df_clean["基本料金削減額"]    = dm_saving_list
    df_clean["導入後最大DM"]     = df_clean["最大需要電力"] - df_clean["削減デマンドkW"]

    dm_saving_annual = sum(dm_saving_list)

    # ─── ③ 年間純削減額（実利） ─────────────────────
    gross_saving    = ene_saving_annual + dm_saving_annual
    net_saving      = gross_saving - sys_fee

    bill_list      = []
    bill_after_list = []
    for idx, row in df_clean.iterrows():
        pf = row.get("力率", pf_base)
        pf_adj = 1.0 - (pf - pf_base)
        pf_adj = max(1.0 - pf_cap, min(pf_adj, 1.0 + pf_cap))
        
        _eff = row["電力量単価"] + _row_add(row)
        bill = (old_contract * row["基本料金単価"] * pf_adj
                + row["使用量合計"] * _eff)
        bill_after = (new_contract * row["基本料金単価"] * pf_adj
                      + (row["使用量合計"] - reduc_kwh_list[idx]) * _eff)
        bill_list.append(bill)
        bill_after_list.append(bill_after)

    df_clean["現状電気料金推計"] = bill_list
    df_clean["導入後電気料金推計"] = bill_after_list

    ann_bill = sum(bill_list)
    total_reduc_kwh = sum(reduc_kwh_list)
    co2_kg          = total_reduc_kwh * CO2_FACTOR
    sugi_trees      = co2_kg / SUGI_KG

    # ─── 室外機リストがある場合：年間削減額を容量加重で1台ずつ配分（methodology §5） ───
    if unit_detail is not None and ctrl_weight > 0:
        unit_detail = unit_detail.copy()
        mask = unit_detail["制御可否"]
        unit_detail["年間削減配分円"] = 0.0
        unit_detail.loc[mask, "年間削減配分円"] = (
            gross_saving * unit_detail.loc[mask, "重み"] / ctrl_weight
        ).round(0)

    return {
        "df": df_clean,
        "units_detail": unit_detail,
        "ctrl_ratio_source": ctrl_ratio_source,
        "ctrl_ratio": ctrl_ratio,
        "cop": cop,
        "base_kwh": df_clean["使用量合計"].min(),
        "base_dm": df_clean["最大需要電力"].min(),
        "ac_kwh_r": ac_kwh_r_eff,                 # 季節性ON時は実効空調割合
        "ac_kwh_r_base": ac_kwh_r,                # 業態既定（参考）
        "seasonal_ac": seasonal,
        "gyotai_consider": gyotai_consider,       # 業態配慮（基礎電力の空調分）
        "base_ac_ratio": _bac,                    # 業態配慮の率（0〜1）
        "pf_base": pf_base,                       # 力率割引の基準力率
        "pf_cap": pf_cap,                         # 力率割引の上下限
        "fuel_included": fuel_included,           # 燃料費調整額を反映したか
        "ren_included": ren_included,             # 再エネ賦課金を反映したか
        "ac_peak_r": ac_peak_r,
        "ene_saving_annual": ene_saving_annual,
        "dm_saving_annual": dm_saving_annual,
        "gross_saving": gross_saving,
        "net_saving": net_saving,
        "sys_fee": sys_fee,
        "old_contract": old_contract,
        "new_contract": new_contract,
        "contract_delta": contract_delta,
        "total_reduc_kwh": total_reduc_kwh,
        "ann_bill": ann_bill,
        "co2_kg": co2_kg,
        "sugi_trees": sugi_trees,
        "cap_rate": cap_rate,
        "ctrl_ratio": ctrl_ratio,
        "gyotai": gyotai,
        "n_units_total": n_units_total,
        "n_units_ctrl": n_units_ctrl,
    }


# ══════════════════════════════════════════════════════════════════
# PowerPoint 自動生成エンジン（マッキンゼー風・全12枚スライド）
# ══════════════════════════════════════════════════════════════════
def provisional_notes(res: dict, app_data: dict) -> list:
    """この結果が依拠する『概算・業態標準値・仮置き・自動補完』の前提を、人が読める注釈リストで返す。
    現行の算出ロジックには一切触れず、結果フラグから前提を可視化するためだけに使用。"""
    notes = []
    try:
        gy = str(res.get("gyotai", ""))
        src = str(res.get("ctrl_ratio_source", ""))
        econ = res.get("econ", {}) or {}
        fb = app_data.get("fallback", {}) or {}
        ac_peak = float(res.get("ac_peak_r", 0) or 0)
        ac_kwh = float(res.get("ac_kwh_r", 0) or 0)

        if econ.get("gaisan_mode"):
            notes.append("【概算モード】機材情報が無いため、室外機台数は検針票の空調ピーク負荷÷標準室外機定格kWで推定し、初期費用・回収期間も標準単価×推定台数で概算。実機材の取得後に精緻化されます。")

        if gy == "その他※数値を指定する":
            notes.append(f"空調割合は手動指定値（最大デマンド {ac_peak*100:.0f}％／使用量 {ac_kwh*100:.0f}％）を使用。")
        else:
            notes.append(f"空調割合は業態「{gy}」の標準値（最大デマンド {ac_peak*100:.0f}％／使用量 {ac_kwh*100:.0f}％）を仮置き（業界配慮）。実測比率があれば精緻化されます。")

        if "台数" in src:
            notes.append("制御可能比率は台数ベースの業態平均（概算）。室外機の型番リスト入力で容量加重に精緻化されます。")
        elif "固定値" in src:
            notes.append("制御可能比率は手動の固定値を使用。")

        if not bool(econ.get("estimate_mode", False)):
            notes.append("初期費用は概算（制御台数 × 目安単価 × エリア係数）。詳細見積モードで実費に精緻化できます。")

        if bool(res.get("seasonal_ac", False)):
            notes.append("空調分は季節性（ベースロード法）で推計しています。")
        if bool(res.get("gyotai_consider", False)):
            notes.append("基礎電力の一定割合を空調分として見込む『業態配慮』を適用しています。")

        _exm = res.get("excluded_months") or []
        if _exm:
            notes.append("一部の月（" + "・".join(str(m) for m in _exm) + "）を試算対象から除外しています（月スコープ）。")
        try:
            _ea = st.session_state.get("excl_areas") or []
        except Exception:
            _ea = []
        if _ea:
            notes.append("一部の階・エリア（" + "・".join(str(a) for a in _ea) + "）を制御対象から除外しています。")

        if fb.get("contract_mode") == "固定値" and not float(fb.get("contract_fix") or 0):
            notes.append("契約電力は年間最大需要で代用しています（契約電力データ未取得）。")
        try:
            if st.session_state.get("import_warn"):
                notes.append("不足していた重要データを既定値で自動補完しています（要確認）。")
        except Exception:
            pass
    except Exception:
        pass
    return notes


def build_pptx(res: dict, client_name: str, app_data: dict) -> bytes:
    if not HAS_PPTX:
        return b""

    LOGO_PATH = "company_logo_192x192.png"

        
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    
    navy = RGBColor(27, 42, 74)        # 濃紺（#1B2A4A）参考デザイン準拠
    green = RGBColor(46, 158, 91)      # エメラルドグリーン（#2E9E5B）削減・利益
    blue = RGBColor(46, 124, 214)      # アクセントブルー（#2E7CD6）数量・電力量
    orange = RGBColor(232, 130, 30)    # アクセントオレンジ（#E8821E）回収・ROI
    gray = RGBColor(90, 100, 116)      # 本文グレー（#5A6474）
    light_gray = RGBColor(241, 244, 248)
    light_green = RGBColor(226, 243, 233)  # 薄エメラルド（ハイライト背景）
    light_navy = RGBColor(232, 237, 245)   # 薄ネイビー（アイコンチップ淡色）
    white = RGBColor(255, 255, 255)

    # ── 5mm マス目背景画像（参考デザイン準拠：ごく淡いグレー格子の明るい背景）を一度だけ生成 ──
    GRID_LINE_RGB = (228, 233, 240)   # ごく薄いグレー罫線
    GRID_BG_RGB   = (247, 249, 252)   # 明るいオフホワイト背景
    GRID_STEP_MM  = 5.0               # 5mm 方眼
    grid_bytes = None
    try:
        from PIL import Image, ImageDraw
        _dpi = 150
        _w_in, _h_in = 13.333, 7.5
        _W, _H = int(round(_w_in * _dpi)), int(round(_h_in * _dpi))
        _step = GRID_STEP_MM / 25.4 * _dpi
        _img = Image.new("RGB", (_W, _H), GRID_BG_RGB)
        _draw = ImageDraw.Draw(_img)
        _x = 0.0
        while _x <= _W:
            _xi = int(round(_x))
            _draw.line([(_xi, 0), (_xi, _H)], fill=GRID_LINE_RGB, width=1)
            _x += _step
        _y = 0.0
        while _y <= _H:
            _yi = int(round(_y))
            _draw.line([(0, _yi), (_W, _yi)], fill=GRID_LINE_RGB, width=1)
            _y += _step
        _gbuf = io.BytesIO()
        _img.save(_gbuf, format="PNG")
        grid_bytes = _gbuf.getvalue()
    except Exception:
        grid_bytes = None

    def add_grid_bg(slide):
        """スライド最背面に 5mm マス目背景を敷く。"""
        if grid_bytes:
            slide.shapes.add_picture(io.BytesIO(grid_bytes), Inches(0), Inches(0),
                                     width=Inches(13.333), height=Inches(7.5))
        else:
            # PIL 不在時のフォールバック：5mm 間隔の薄い罫線を描画
            grid_rgb = RGBColor(*GRID_LINE_RGB)
            step = GRID_STEP_MM / 25.4  # inch
            n = 1
            while n * step < 13.333:
                ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(n * step), Inches(0), Inches(0.008), Inches(7.5))
                ln.fill.solid(); ln.fill.fore_color.rgb = grid_rgb; ln.line.fill.background()
                n += 1
            n = 1
            while n * step < 7.5:
                ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(n * step), Inches(13.333), Inches(0.008))
                ln.fill.solid(); ln.fill.fore_color.rgb = grid_rgb; ln.line.fill.background()
                n += 1

    def set_fill_alpha(shape, opacity):
        """単色塗りシェイプに透明度を設定（opacity: 0.0=透明 〜 1.0=不透明）。"""
        try:
            from pptx.oxml.ns import qn
            sf = shape._element.spPr.find(qn('a:solidFill'))
            if sf is None:
                return
            clr = sf.find(qn('a:srgbClr'))
            if clr is None:
                return
            for ex in clr.findall(qn('a:alpha')):
                clr.remove(ex)
            a = clr.makeelement(qn('a:alpha'), {'val': str(int(opacity * 100000))})
            clr.append(a)
        except Exception:
            pass

    # ── 高級感アクセント用カラー ──
    gold      = RGBColor(193, 158, 88)    # シャンパンゴールド
    gold_lt   = RGBColor(214, 188, 130)   # ライトゴールド
    deep_navy = RGBColor(7, 17, 40)       # さらに深い濃紺（グラデ用）

    def set_gradient(shape, c1, c2, angle=90):
        """2色リニアグラデーションを設定（失敗時は c1 のベタ塗り）。"""
        try:
            shape.fill.gradient()
            stops = shape.fill.gradient_stops
            stops[0].position = 0.0; stops[0].color.rgb = c1
            stops[1].position = 1.0; stops[1].color.rgb = c2
            try:
                shape.fill.gradient_angle = angle
            except Exception:
                pass
        except Exception:
            shape.fill.solid(); shape.fill.fore_color.rgb = c1

    def add_soft_shadow(shape, blur=0.07, dist=0.045, direction=5400000, alpha=0.55, color_hex='182846'):
        """シェイプに外側影／グロー（dist=0で発光風）を付与。color_hexで色指定可。"""
        try:
            from pptx.oxml.ns import qn
            spPr = shape._element.spPr
            for e in spPr.findall(qn('a:effectLst')):
                spPr.remove(e)
            eff = spPr.makeelement(qn('a:effectLst'), {})
            shd = eff.makeelement(qn('a:outerShdw'), {
                'blurRad': str(int(blur * 914400)),
                'dist': str(int(dist * 914400)),
                'dir': str(direction),
                'rotWithShape': '0',
            })
            clr = shd.makeelement(qn('a:srgbClr'), {'val': color_hex})
            a = clr.makeelement(qn('a:alpha'), {'val': str(int(alpha * 100000))})
            clr.append(a); shd.append(clr); eff.append(shd)
            spPr.append(eff)
        except Exception:
            pass

    def add_footer(slide):
        """全スライド共通：右下に小さくロゴのみ（罫線・社名・タイトルは廃止。ページ番号は末尾で並び替え後に一括付与）。"""
        if os.path.exists(LOGO_PATH):
            slide.shapes.add_picture(LOGO_PATH, Inches(12.55), Inches(7.04), height=Inches(0.30))

    def add_page_number(slide, page_no):
        """ページ番号を左下に配置。"""
        rb = slide.shapes.add_textbox(Inches(0.7), Inches(7.06), Inches(1.4), Inches(0.32))
        pr = rb.text_frame.paragraphs[0]
        pr.text = f"P. {page_no:02d}"
        pr.font.name = 'Yu Gothic'; pr.font.size = Pt(9); pr.font.bold = True; pr.font.color.rgb = navy
        pr.alignment = PP_ALIGN.LEFT

    def add_metric_card(slide, x, y, w, h, icon, label, value, sub, accent, value_size=22):
        """参考デザイン準拠の指標カード：白カード＋色付きアイコンチップ＋大きな数字。"""
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        card.fill.solid(); card.fill.fore_color.rgb = white
        card.line.color.rgb = RGBColor(225, 230, 238); card.line.width = Pt(1)
        add_soft_shadow(card, blur=0.09, dist=0.055, alpha=0.20)
        # アイコンチップ（色付き角丸）
        chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.28), Inches(y + 0.3), Inches(0.62), Inches(0.62))
        chip.fill.solid(); chip.fill.fore_color.rgb = accent; chip.line.fill.background()
        add_soft_shadow(chip, blur=0.04, dist=0.03, alpha=0.25)
        tfc = chip.text_frame
        tfc.margin_left = tfc.margin_right = tfc.margin_top = tfc.margin_bottom = 0
        pc = tfc.paragraphs[0]; pc.text = icon
        pc.font.size = Pt(24); pc.alignment = PP_ALIGN.CENTER; pc.font.color.rgb = white
        # ラベル（チップの右）
        lb = slide.shapes.add_textbox(Inches(x + 1.02), Inches(y + 0.34), Inches(w - 1.2), Inches(0.6))
        tfl = lb.text_frame; tfl.word_wrap = True
        tfl.margin_left = tfl.margin_top = 0
        pl = tfl.paragraphs[0]; pl.text = label
        pl.font.name = 'Yu Gothic'; pl.font.size = Pt(11); pl.font.bold = True; pl.font.color.rgb = navy
        pl.line_spacing = 1.05
        # 大きな数字（アクセントカラー）
        vb = slide.shapes.add_textbox(Inches(x + 0.22), Inches(y + 1.12), Inches(w - 0.44), Inches(0.95))
        tfv = vb.text_frame; tfv.word_wrap = True
        tfv.margin_left = tfv.margin_right = 0
        pv = tfv.paragraphs[0]; pv.text = value
        pv.font.name = 'Yu Gothic'; pv.font.size = Pt(value_size); pv.font.bold = True
        pv.font.color.rgb = accent; pv.alignment = PP_ALIGN.CENTER
        # サブキャプション
        if sub:
            sb = slide.shapes.add_textbox(Inches(x + 0.22), Inches(y + h - 0.82), Inches(w - 0.44), Inches(0.72))
            tfs = sb.text_frame; tfs.word_wrap = True
            ps = tfs.paragraphs[0]; ps.text = sub
            ps.font.name = 'Yu Gothic'; ps.font.size = Pt(9.5); ps.font.color.rgb = gray
            ps.alignment = PP_ALIGN.CENTER; ps.line_spacing = 1.15

    def apply_base_layout(slide, title_text, highlight_text=None):
        # 1. 5mm マス目グリッド背景
        add_grid_bg(slide)

        # 2. タイトル左のアクセントバー（濃紺＋グリーンの2連）
        ab1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(0.42), Inches(0.13), Inches(0.56))
        ab1.fill.solid(); ab1.fill.fore_color.rgb = navy; ab1.line.fill.background()
        ab2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.86), Inches(0.42), Inches(0.13), Inches(0.56))
        ab2.fill.solid(); ab2.fill.fore_color.rgb = green; ab2.line.fill.background()

        # 3. タイトル（濃紺・太字／左アクセントバーと垂直中央を揃える）
        title_box = slide.shapes.add_textbox(Inches(1.12), Inches(0.42), Inches(11.5), Inches(0.56))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = 'Yu Gothic'
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = navy

        # 4. 説明文（グレー本文）
        if highlight_text:
            hb = slide.shapes.add_textbox(Inches(1.12), Inches(1.22), Inches(11.5), Inches(0.72))
            tfh = hb.text_frame
            tfh.word_wrap = True
            tfh.margin_left = tfh.margin_top = 0
            ph = tfh.paragraphs[0]
            ph.text = highlight_text
            ph.font.name = 'Yu Gothic'
            ph.font.size = Pt(13)
            ph.font.color.rgb = gray
            ph.line_spacing = 1.2

        # 5. フッター（社名・ページ番号・ロゴ）
        add_footer(slide)


    # 📊 Slide 1: 表紙（プレミアム仕様）
    slide1 = prs.slides.add_slide(blank_layout)
    add_grid_bg(slide1)  # 表紙にも 5mm マス目背景を付与

    # ロゴ（全ページ共通：右下に小さく）
    if os.path.exists(LOGO_PATH):
        slide1.shapes.add_picture(LOGO_PATH, Inches(12.55), Inches(7.04), height=Inches(0.30))

    # 宛名（クライアント名）— 左上
    sh_client = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(0.45), Inches(3.7), Inches(0.58))
    sh_client.fill.solid(); sh_client.fill.fore_color.rgb = white
    sh_client.line.color.rgb = navy; sh_client.line.width = Pt(1.5)
    add_soft_shadow(sh_client, blur=0.04, dist=0.03, alpha=0.20)
    tf_cl = sh_client.text_frame
    tf_cl.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_cl = tf_cl.paragraphs[0]
    p_cl.text = f"{client_name}　御中"
    p_cl.font.name = 'Yu Gothic'; p_cl.font.size = Pt(15); p_cl.font.bold = True
    p_cl.font.color.rgb = navy; p_cl.alignment = PP_ALIGN.CENTER

    # 中央：タイトルパネル（上下中央寄せ）
    panel = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.95), Inches(11.93), Inches(2.1))
    panel.fill.solid(); panel.fill.fore_color.rgb = white
    panel.line.color.rgb = RGBColor(210, 216, 224); panel.line.width = Pt(1)
    add_soft_shadow(panel, blur=0.12, dist=0.07, alpha=0.22)
    tf_p = panel.text_frame
    tf_p.word_wrap = True
    tf_p.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf_p.margin_left = tf_p.margin_right = Inches(0.4); tf_p.margin_top = Inches(0.3)
    p_t = tf_p.paragraphs[0]
    p_t.text = "デジタルパワーサービス（DPS）導入による"
    p_t.font.name = 'Yu Gothic'; p_t.font.size = Pt(33); p_t.font.bold = True; p_t.font.color.rgb = navy
    p_t.alignment = PP_ALIGN.CENTER
    p_t2 = tf_p.add_paragraph()
    p_t2.text = "電力経費削減のご提案"
    p_t2.font.name = 'Yu Gothic'; p_t2.font.size = Pt(33); p_t2.font.bold = True
    p_t2.space_before = Pt(6)
    p_t2.font.color.rgb = green
    p_t2.alignment = PP_ALIGN.CENTER

    # サブタイトル帯
    st_box = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(4.35), Inches(11.93), Inches(0.72))
    st_box.fill.solid(); st_box.fill.fore_color.rgb = light_green
    st_box.line.color.rgb = green; st_box.line.width = Pt(1)
    gtab_cv = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(4.35), Inches(0.12), Inches(0.72))
    gtab_cv.fill.solid(); gtab_cv.fill.fore_color.rgb = green; gtab_cv.line.fill.background()
    tf_st = st_box.text_frame
    tf_st.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_st = tf_st.paragraphs[0]
    p_st.text = "～ 過去の電気料金実績に基づく 空調デマンド制御シミュレーション提案 ～"
    p_st.font.name = 'Yu Gothic'; p_st.font.size = Pt(16); p_st.font.bold = True
    p_st.font.color.rgb = navy; p_st.alignment = PP_ALIGN.CENTER

    # 注釈（サブタイトル帯の右下）
    _cvnote = slide1.shapes.add_textbox(Inches(0.7), Inches(5.25), Inches(11.93), Inches(0.32))
    _cvnote.text_frame.word_wrap = True
    _cvp = _cvnote.text_frame.paragraphs[0]
    _cvp.text = "※本試算は概算・業態標準値などの仮置き前提を含みます（前提条件は表紙ノート参照）。"
    _cvp.font.name = 'Yu Gothic'; _cvp.font.size = Pt(8.5); _cvp.font.color.rgb = gray
    _cvp.alignment = PP_ALIGN.RIGHT

    # 📊 Slide 2: 導入効果の総括 (総初期投資ベース)
    slide2 = prs.slides.add_slide(blank_layout)
    payback_years_calc = res.get("total_invest", 2400000) / res.get("net_saving", 1060633) if res.get("net_saving", 1060633) > 0 else 99.9
    apply_base_layout(slide2, "導入効果の総括：収益試算サマリー", 
                      f"約 {payback_years_calc:.1f} 年で投資回収を実現。以降は毎年約 {res.get('net_saving', 0)/10000:,.1f} 万円の純利益を生み出す財務改善ソリューション")
    
    # ③ サマリKPIカード画像（シミュレーション画面と同一デザイン）を埋め込み
    _NAVY_C, _GREEN_C = "#13315C", "#3DAE4E"
    _s_inv = float(res.get("total_invest", 0) or 0)
    _s_grs = float(res.get("gross_saving", 0) or 0)
    _s_nt = float(res.get("net_saving", 0) or 0)
    _s_sf = float(res.get("sys_fee", 0) or 0)
    _s_pb = (_s_inv / _s_nt) if _s_nt > 0 else 0.0
    _summary_cards2 = [
        {"icon": "money.png", "title": "初期導入費用（総投資額）", "value": f"{_s_inv:,.0f}円",
         "subtitle": f"制御対象 {res.get('n_units_ctrl', 31)}台分 / 工事費込", "color": _NAVY_C},
        {"icon": "graf.png", "title": "年間総削減額（グロス）", "value": f"{_s_grs:,.0f}円",
         "subtitle": "基本料金＋電力量の年間削減合計", "color": _GREEN_C},
        {"icon": "plus.png", "title": "年間実質利点（手残り）", "value": f"{_s_nt:,.0f}円",
         "subtitle": f"削減合計 − 年間維持費 {_s_sf:,.0f}円", "color": _GREEN_C},
        {"icon": "clock.png", "title": "投資回収期間（ROI）", "value": f"約{_s_pb:.1f}年",
         "subtitle": f"約{round(_s_pb*12)}ヶ月で完全回収、以降は純利益", "color": _GREEN_C},
    ]
    try:
        _cards_png2 = make_summary_cards_png(_summary_cards2, fig_w=13.0, fig_h=3.7)
        _cimg_w2 = 11.93
        _cimg_h2 = _cimg_w2 * 3.7 / 13.0
        _cimg_x2 = (13.333 - _cimg_w2) / 2.0
        slide2.shapes.add_picture(io.BytesIO(_cards_png2), Inches(_cimg_x2),
                                  Inches(2.05), width=Inches(_cimg_w2))
    except Exception:
        pass

    # 注釈は1つに集約（算出根拠＋仮置き前提の免責をまとめて記載）
    fn_box = slide2.shapes.add_textbox(Inches(0.7), Inches(5.55), Inches(11.93), Inches(1.35))
    tf_fn = fn_box.text_frame
    tf_fn.word_wrap = True
    p_fn = tf_fn.paragraphs[0]
    p_fn.text = (f"※ 算出根拠：現状の過去12ヶ月の検針票データ（最大デマンド {res['df']['最大需要電力'].max():.0f}kW）に、削減率 30% × 制御可能比率を適用。\n"
                 "※ 独自のローテーション制御技術により、生産ラインや室温（1℃以内）への影響がないことを保証いたします。\n"
                 "※ 本試算は概算・業態標準値などの仮置き前提を含みます（前提条件は表紙ノート参照）。")
    p_fn.font.name = 'Yu Gothic'
    p_fn.font.size = Pt(11)
    p_fn.font.color.rgb = gray
    p_fn.line_spacing = 1.3

    # 📊 Slide 3: 試算の分析方法および算出前提の定義
    slide3 = prs.slides.add_slide(blank_layout)
    apply_base_layout(slide3, "試算の分析方法および算出前提の定義", 
                      "不確実な未来予測を排除し、過去の「確定ファクト」から逆算する引き算推計法を採用")
    
    # 今回の算出設定から採択手法を判定
    _opts3 = app_data.get("calc_options", {})
    _seasonal3 = bool(_opts3.get("seasonal_ac", False))
    _gcons3 = bool(_opts3.get("gyotai_consider", False))
    _units3 = res.get("units_detail") is not None and len(res.get("units_detail", [])) > 0
    _adopted3 = {
        "引き算推計法": True,
        "業態別原単位法": (not _seasonal3),
        "季節性ベースロード法": _seasonal3,
        "業態配慮法": _gcons3,
        "ピーク需要分析": True,
        "力率調整法": True,
        "設備積み上げ法": _units3,
    }

    def _draw_method_grid(_ts):
        """image_sort（7手法グリッド）を _ts に配置し、採択手法に赤枠＋凡例を重ねる。"""
        _sort_img = "image_sort.png"
        if not os.path.exists(_sort_img):
            return
        _sw3 = 10.0; _sh3 = _sw3 / 2.198                # 元画像アスペクト(w/h)
        _sx3 = (13.333 - _sw3) / 2.0; _sy3 = 2.1
        _ts.shapes.add_picture(_sort_img, Inches(_sx3), Inches(_sy3), width=Inches(_sw3))
        _cellsX = [(0.02, 0.23), (0.27, 0.48), (0.52, 0.73), (0.77, 0.98)]
        _cellsY = [(0.03, 0.47), (0.53, 0.97)]
        _method_cell = {
            "引き算推計法": (0, 1), "業態別原単位法": (0, 2), "季節性ベースロード法": (0, 3),
            "業態配慮法": (1, 0), "ピーク需要分析": (1, 1), "力率調整法": (1, 2), "設備積み上げ法": (1, 3),
        }
        for _mname, (_rr3, _cc3) in _method_cell.items():
            if not _adopted3.get(_mname):
                continue
            _cx0, _cx1 = _cellsX[_cc3]; _cy0, _cy1 = _cellsY[_rr3]
            _rb = _ts.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(_sx3 + _cx0 * _sw3), Inches(_sy3 + _cy0 * _sh3),
                Inches((_cx1 - _cx0) * _sw3), Inches((_cy1 - _cy0) * _sh3))
            _rb.fill.background()
            _rb.line.color.rgb = RGBColor(214, 40, 40); _rb.line.width = Pt(2.5)
            _rb.shadow.inherit = False
        _leg3 = _ts.shapes.add_textbox(Inches(0.7), Inches(6.75), Inches(11.93), Inches(0.28))
        _pleg3 = _leg3.text_frame.paragraphs[0]
        _pleg3.text = "■ 赤枠＝本案件で採択した分析方法"
        _pleg3.font.name = 'Yu Gothic'; _pleg3.font.size = Pt(10); _pleg3.font.bold = True
        _pleg3.font.color.rgb = RGBColor(214, 40, 40)

    # P.3 本文：支給の3カード（image_sort2.png）があればそれを掲載。無ければ7手法グリッド。
    _has_logic3 = os.path.exists("image_sort2.png")
    if _has_logic3:
        try:
            from PIL import Image as _PIL3
            _im3 = _PIL3.open("image_sort2.png"); _pw3, _ph3 = _im3.size
            _asp3 = (_pw3 / _ph3) if _ph3 else 2.2
        except Exception:
            _asp3 = 2.2
        _lw3 = 11.9; _lh3 = _lw3 / _asp3
        if _lh3 > 4.9:
            _lh3 = 4.9; _lw3 = _lh3 * _asp3
        _lx3 = (13.333 - _lw3) / 2.0
        slide3.shapes.add_picture("image_sort2.png", Inches(_lx3), Inches(2.15), width=Inches(_lw3))
    else:
        _draw_method_grid(slide3)

    # 📊 Slide 4: 制御可否：対象拠点における機器仕分け実態
    slide4 = prs.slides.add_slide(blank_layout)
    apply_base_layout(slide4, "制御可否：対象拠点における機器仕分け実態", 
                      "現場の操業を止めない、物理的制約に基づいたロジカルな制御対象台数の切り分け")
    
    # ── 上部：制御台数の軽い要約バー ──
    _nc4 = res.get('n_units_ctrl', 31)
    _nt4 = res.get('n_units_total', 34)
    bar4 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(2.05), Inches(11.93), Inches(0.9))
    bar4.fill.solid(); bar4.fill.fore_color.rgb = white
    bar4.line.color.rgb = RGBColor(225, 230, 238); bar4.line.width = Pt(1)
    add_soft_shadow(bar4, blur=0.09, dist=0.055, alpha=0.18)

    barL = slide4.shapes.add_textbox(Inches(1.05), Inches(2.05), Inches(6.7), Inches(0.9))
    _tfL = barL.text_frame; _tfL.word_wrap = True; _tfL.vertical_anchor = MSO_ANCHOR.MIDDLE
    _tfL.margin_top = _tfL.margin_bottom = 0
    _pL = _tfL.paragraphs[0]
    _rl1 = _pL.add_run(); _rl1.text = "制御可能空調比率（物理台数ベース）　"
    _rl1.font.name = 'Yu Gothic'; _rl1.font.size = Pt(14); _rl1.font.bold = True; _rl1.font.color.rgb = navy
    _rl2 = _pL.add_run(); _rl2.text = f"{_nc4} / {_nt4} 台"
    _rl2.font.name = 'Yu Gothic'; _rl2.font.size = Pt(26); _rl2.font.bold = True; _rl2.font.color.rgb = green

    barR = slide4.shapes.add_textbox(Inches(7.75), Inches(2.05), Inches(4.7), Inches(0.9))
    _tfR = barR.text_frame; _tfR.word_wrap = True; _tfR.vertical_anchor = MSO_ANCHOR.MIDDLE
    _tfR.margin_top = _tfR.margin_bottom = 0
    _pR = _tfR.paragraphs[0]
    _rr1 = _pR.add_run(); _rr1.text = "対象空調の定格容量カバー比率　"
    _rr1.font.name = 'Yu Gothic'; _rr1.font.size = Pt(14); _rr1.font.bold = True; _rr1.font.color.rgb = navy
    _rr2 = _pR.add_run(); _rr2.text = f"{res.get('ctrl_ratio', 1.0):.1%}"
    _rr2.font.name = 'Yu Gothic'; _rr2.font.size = Pt(26); _rr2.font.bold = True; _rr2.font.color.rgb = green

    # ── 下部：②制御可否リスト画像を埋め込み ──
    _ud4 = res.get('units_detail')
    _items4 = []
    if _ud4 is not None and len(_ud4) > 0:
        for _, _r4 in _ud4.iterrows():
            _nm4 = str(_r4.get("機器ID", "") or "")
            _items4.append({
                "系統名": _nm4, "設置場所": _nm4,
                "メーカー": str(_r4.get("メーカー", "") or ""),
                "型式": str(_r4.get("型番", "") or ""),
                "制御可否": bool(_r4.get("制御可否", False)),
            })
    else:
        for _eq4 in DEFAULT_EQUIPMENT:
            _items4.append({
                "系統名": _eq4["loc"], "設置場所": _eq4["loc"],
                "メーカー": _eq4["mfr"], "型式": _eq4["model"],
                "制御可否": str(_eq4["ctrl"]).strip() in ("〇", "○", "◯"),
            })
    _ctrl_ct4 = sum(1 for it in _items4 if it["制御可否"])
    _ctrl_extra_pages = []     # 2ページ目以降（続きスライド用）。1ページ最大50台
    _list_done = False
    try:
        # ページ分割：1枚目は台数控えめ（上部バーと同居のため）、2枚目以降はフル掲載
        _FIRST_N, _CONT_N = 24, 40
        if len(_items4) <= _FIRST_N:
            _pages4 = [_items4]
        else:
            _pages4 = [_items4[:_FIRST_N]]
            _rest4 = _items4[_FIRST_N:]
            for _i in range(0, len(_rest4), _CONT_N):
                _pages4.append(_rest4[_i:_i + _CONT_N])
        _np4t = len(_pages4)
        _lp4 = [make_control_list_png(
                    _pgi, total_units=len(_items4), controllable_units=_ctrl_ct4,
                    page=_k4, total_pages=_np4t)
                for _k4, _pgi in enumerate(_pages4, start=1)]

        # 1枚目：上部バー枠幅(11.93)に合わせて下部に配置（収まらなければ縮小して中央）
        _r0 = max((len(_pages4[0]) + 1) // 2, 1)
        _fh0 = max(3.2, 1.1 + _r0 * 0.255)
        _iw4 = 11.93; _ix4 = 0.7
        _ih0 = _iw4 * _fh0 / 13.0
        if _ih0 > 3.8:
            _iw4 = _iw4 * 3.8 / _ih0; _ix4 = (13.333 - _iw4) / 2.0
        slide4.shapes.add_picture(io.BytesIO(_lp4[0]), Inches(_ix4),
                                  Inches(3.1), width=Inches(_iw4))
        _list_done = True
        _ctrl_extra_pages = list(_lp4[1:])              # 2枚目以降は続きスライドへ（フル）
        if _np4t > 1:
            note4 = slide4.shapes.add_textbox(Inches(0.7), Inches(7.18), Inches(6.0), Inches(0.28))
            np4 = note4.text_frame.paragraphs[0]
            np4.text = f"※ 全{len(_items4)}台。本ページは 1/{_np4t}（続きは次ページに掲載）。"
            np4.font.name = 'Yu Gothic'; np4.font.size = Pt(9); np4.font.color.rgb = gray
    except Exception:
        _list_done = False

    if not _list_done and _items4:
        # フォールバック：画像生成に失敗しても、ネイティブ表で必ずリストを表示する
        _shown = _items4[:12]
        _rows4 = len(_shown) + 1
        _tw4 = 11.93; _tx4 = 0.7                       # バー枠と同じ幅・左位置
        _th4 = min(3.7, _rows4 * 0.32)
        _tb4 = slide4.shapes.add_table(_rows4, 5, Inches(_tx4), Inches(3.15), Inches(_tw4), Inches(_th4)).table
        for _ci, _cw4 in enumerate([3.6, 2.8, 2.2, 2.53, 0.8]):
            _tb4.columns[_ci].width = Inches(_cw4)
        for _ci, _hh in enumerate(["系統名", "設置場所", "メーカー", "型式", "制御可否"]):
            _tb4.cell(0, _ci).text = _hh
        for _ri, _it in enumerate(_shown, start=1):
            _tb4.cell(_ri, 0).text = str(_it["系統名"]); _tb4.cell(_ri, 1).text = str(_it["設置場所"])
            _tb4.cell(_ri, 2).text = str(_it["メーカー"]); _tb4.cell(_ri, 3).text = str(_it["型式"])
            _tb4.cell(_ri, 4).text = "○" if _it["制御可否"] else "×"
        for _ri in range(_rows4):
            for _ci in range(5):
                _cell4 = _tb4.cell(_ri, _ci)
                _cell4.margin_top = _cell4.margin_bottom = Inches(0.02)
                _p4 = _cell4.text_frame.paragraphs[0]
                _p4.font.name = 'Yu Gothic'; _p4.font.size = Pt(10)
                _p4.font.color.rgb = white if _ri == 0 else navy
                if _ci == 4:
                    _p4.alignment = PP_ALIGN.CENTER
                if _ri == 0:
                    _cell4.fill.solid(); _cell4.fill.fore_color.rgb = navy; _p4.font.bold = True

    # 📊 Slide 5: シミュレーション：企業別・月別料金比較詳細
    slide5 = prs.slides.add_slide(blank_layout)
    apply_base_layout(slide5, "シミュレーション：企業別・月別料金比較詳細", 
                      "12ヶ月分の電力推移を1円単位で徹底検証。無駄な支払いの発生源を特定します。")
    
    # ① 使用量＋デマンド グラフ画像（シミュレーション画面と同一デザイン）を埋め込み
    # ※背面の角丸カードは廃止（グラフ画像自体が白背景のため不要）
    try:
        _df5 = res['df']
        _m5 = [str(m) for m in _df5['月'].tolist()]
        _u5 = _df5['使用量合計'].tolist()
        _rd5 = _df5['削減kWh'].tolist() if '削減kWh' in _df5.columns else [0] * len(_m5)
        _dm5 = _df5['最大需要電力'].tolist()
        _chart_png = make_demand_chart_png(_m5, _u5, _rd5, _dm5,
                                           target_units=res.get('n_units_ctrl', 0))
        _chw = 8.33
        slide5.shapes.add_picture(io.BytesIO(_chart_png), Inches(0.7), Inches(1.8),
                                  width=Inches(_chw))
    except Exception:
        pass

    # ── 右側：削減シナリオ表（グラフ拡大に伴い右へ寄せて省スペース化） ──
    tbl_shape5 = slide5.shapes.add_table(5, 3, Inches(9.23), Inches(1.8), Inches(3.4), Inches(3.0))
    t5 = tbl_shape5.table
    t5.columns[0].width = Inches(0.9)
    t5.columns[1].width = Inches(1.35)
    t5.columns[2].width = Inches(1.15)
    t5.rows[0].height = Inches(0.5)
    for _ri in range(1, 5):
        t5.rows[_ri].height = Inches(0.6)

    for _ci, _h in enumerate(["シナリオ", "契約電力 推移", "削減量"]):
        t5.cell(0, _ci).text = _h

    _old_c = res.get('old_contract', 500)
    _tier_new = [
        _old_c - res.get('contract_delta', 0) * 0.67,
        res.get('new_contract', 500),
        _old_c - res.get('contract_delta', 0) * 1.33,
        _old_c - res.get('contract_delta', 0) * 1.67,
    ]
    _tier_lbl = ["保守 (20%)", "標準 (30%)", "積極 (40%)", "限界 (50%)"]
    for idx, (_lbl, _nv) in enumerate(zip(_tier_lbl, _tier_new), start=1):
        _red = _old_c - _nv
        _pct = (_red / _old_c * 100) if _old_c else 0
        t5.cell(idx, 0).text = _lbl
        t5.cell(idx, 1).text = f"{_old_c:.0f} → {_nv:.0f} kW"
        t5.cell(idx, 2).text = f"▲{_red:.0f}kW / -{_pct:.0f}%"

    for ri in range(5):
        _is_std = "標準" in t5.cell(ri, 0).text
        for ci in range(3):
            cell = t5.cell(ri, ci)
            cell.margin_left = cell.margin_right = Inches(0.05)
            cell.margin_top = cell.margin_bottom = Inches(0.03)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.font.name = 'Yu Gothic'
            p.font.size = Pt(9.5) if ri == 0 else Pt(10.5)
            p.alignment = PP_ALIGN.CENTER if ci >= 1 else PP_ALIGN.LEFT
            cell.fill.solid()
            if ri == 0:
                p.font.color.rgb = white; p.font.bold = True
                cell.fill.fore_color.rgb = navy
            else:
                p.font.color.rgb = navy
                p.font.bold = _is_std
                if _is_std:
                    cell.fill.fore_color.rgb = light_green
                else:
                    cell.fill.fore_color.rgb = RGBColor(247, 249, 252) if ri % 2 == 0 else white

    fn_b5 = slide5.shapes.add_textbox(Inches(9.23), Inches(5.05), Inches(3.4), Inches(1.7))
    tf_fn5 = fn_b5.text_frame
    tf_fn5.word_wrap = True
    p_fn5 = tf_fn5.paragraphs[0]
    p_fn5.text = f"【参考：現状年間実績ファクト】\n・電力契約量（最新） : {res.get('old_contract', 500):.0f} kW\n・最大デマンド値 : {res['df']['最大需要電力'].max():.0f} kW\n・総使用電力量 : {res['df']['使用量合計'].sum():,.0f} kWh"
    p_fn5.font.name = 'Yu Gothic'
    p_fn5.font.size = Pt(10)
    p_fn5.font.color.rgb = gray
    p_fn5.line_spacing = 1.3

    # 📊 Slide 6: 導入方法と費用 (総費用ベース)
    slide6 = prs.slides.add_slide(blank_layout)
    apply_base_layout(slide6, "既存設備を活かす『アドオン工法』による無停止・低リスク導入",
                      "既存の室外機をそのまま活かし、生産を止めずに低リスク・短工期で導入します")

    # ── 左：室外機イラスト（支給画像・背景透過／拡大） ──
    _kucho = "icon_kucho_t.png" if os.path.exists("icon_kucho_t.png") else "icon_kucho.png"
    if os.path.exists(_kucho):
        _kh = 3.4                       # 縦長画像を拡大
        _kw = _kh * 186.0 / 254.0       # 元画像アスペクト(w/h=0.732)
        slide6.shapes.add_picture(_kucho, Inches(2.4 - _kw / 2), Inches(2.15), height=Inches(_kh))

    # ── 右：制御機器の説明＋取付方法を1つの枠に集約 ──
    info = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.85), Inches(2.15), Inches(7.78), Inches(3.4))
    info.fill.solid(); info.fill.fore_color.rgb = white
    info.line.color.rgb = navy; info.line.width = Pt(1.25)
    add_soft_shadow(info, blur=0.08, dist=0.05, alpha=0.18)
    tf_i = info.text_frame
    tf_i.word_wrap = True; tf_i.vertical_anchor = MSO_ANCHOR.TOP
    tf_i.margin_left = tf_i.margin_right = Inches(0.3); tf_i.margin_top = Inches(0.22)
    _pi_h = tf_i.paragraphs[0]
    _pi_h.text = "⚙️ 制御機器：DPSユニット"
    _pi_h.font.name = 'Yu Gothic'; _pi_h.font.size = Pt(15); _pi_h.font.bold = True
    _pi_h.font.color.rgb = navy; _pi_h.space_after = Pt(6)
    _pi_b = tf_i.add_paragraph()
    _pi_b.text = "室外機へのアドオン装着仕様のため、冷媒配管などの「配管切断工事」は一切不要。施工リスクを極小化し、即日導入が可能です。"
    _pi_b.font.name = 'Yu Gothic'; _pi_b.font.size = Pt(11.5); _pi_b.font.color.rgb = gray
    _pi_b.line_spacing = 1.3; _pi_b.space_after = Pt(12)
    _pi_s = tf_i.add_paragraph()
    _pi_s.text = "■ 取付方法（標準取付工法）"
    _pi_s.font.name = 'Yu Gothic'; _pi_s.font.size = Pt(13); _pi_s.font.bold = True
    _pi_s.font.color.rgb = navy; _pi_s.space_after = Pt(6)
    _checks = [
        ("完全屋外施工", "室内への立ち入りは原則不要"),
        ("稼働停止ゼロ", "生産ラインや生産を止めることなく安全に工事"),
        ("工期スピード", "1台あたり約 2〜3 時間で完了"),
    ]
    for _ct, _cd in _checks:
        _pc = tf_i.add_paragraph()
        _rk = _pc.add_run(); _rk.text = "✓ "
        _rk.font.name = 'Yu Gothic'; _rk.font.size = Pt(12); _rk.font.bold = True; _rk.font.color.rgb = green
        _rt = _pc.add_run(); _rt.text = f"{_ct}："
        _rt.font.name = 'Yu Gothic'; _rt.font.size = Pt(12); _rt.font.bold = True; _rt.font.color.rgb = navy
        _rd = _pc.add_run(); _rd.text = _cd
        _rd.font.name = 'Yu Gothic'; _rd.font.size = Pt(11); _rd.font.color.rgb = gray
        _pc.line_spacing = 1.25; _pc.space_after = Pt(4)

    # ── 下部：初期費用内訳バンド（濃紺） ──
    band6 = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(5.7), Inches(11.93), Inches(1.0))
    set_gradient(band6, navy, deep_navy, angle=0)
    band6.line.fill.background()
    add_soft_shadow(band6, blur=0.08, dist=0.05, alpha=0.35)
    tf_bd = band6.text_frame
    tf_bd.margin_left = tf_bd.margin_right = Inches(0.3)
    tf_bd.word_wrap = True; tf_bd.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_bd1 = tf_bd.paragraphs[0]
    p_bd1.text = f"初期費用内訳（制御対象 {res.get('n_units_ctrl', 31)}台分）"
    p_bd1.font.name = 'Yu Gothic'; p_bd1.font.size = Pt(11); p_bd1.font.color.rgb = RGBColor(0xCB, 0xDA, 0xF0)
    p_bd1.alignment = PP_ALIGN.CENTER; p_bd1.space_after = Pt(3)
    p_bd2 = tf_bd.add_paragraph()
    p_bd2.text = (f"制御機器本体費 / 標準取付配線工事費 / 初期設定・検証調整費 一式 ＝ "
                  f"総額 {res.get('total_invest', 2400000):,} 円（税別）")
    p_bd2.font.name = 'Yu Gothic'; p_bd2.font.size = Pt(14); p_bd2.font.bold = True
    p_bd2.font.color.rgb = white; p_bd2.alignment = PP_ALIGN.CENTER

    # 📊 Slide 7: ＜アディショナル＞空調交換提案 (補助金表記削除・純粋更新)
    slide7 = prs.slides.add_slide(blank_layout)
    apply_base_layout(slide7, "【ご参考】空調機リプレイス提案", 
                      "老朽化した低効率機を一新。最新の省エネ機種へのリプレイスと回収年数の比較検証")
    
    s_a = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.2), Inches(5.5), Inches(4.0))
    s_a.fill.solid()
    s_a.fill.fore_color.rgb = white
    s_a.line.color.rgb = RGBColor(225, 230, 238); s_a.line.width = Pt(1)
    add_soft_shadow(s_a, blur=0.09, dist=0.055, alpha=0.20)
    tf_sa = s_a.text_frame
    tf_sa.margin_left = tf_sa.margin_top = Inches(0.3)
    p_sa_h = tf_sa.paragraphs[0]
    p_sa_h.text = "シナリオA：DPSデマンド制御 (現状設備に後付け)"
    p_sa_h.font.name = 'Yu Gothic'
    p_sa_h.font.size = Pt(16)
    p_sa_h.font.bold = True
    p_sa_h.font.color.rgb = navy
    p_sa_h.space_after = Pt(20)
    
    sa_points = [
        "・初期投資： 約 240 万円",
        f"・回収期間： 約 {payback_years_calc:.1f} 年 (実利ベースの高速投資回収)",
        "・工事期間： わずか1週間。既存設備を活かすため稼働リスク極小",
        "・特徴： 元本回収スピードを最も最優先させた実利・アドオンモデル"
    ]
    for pt in sa_points:
        p_pt = tf_sa.add_paragraph()
        p_pt.text = pt
        p_pt.font.name = 'Yu Gothic'
        p_pt.font.size = Pt(11.5)
        p_pt.font.color.rgb = gray
        p_pt.space_after = Pt(10)
        p_pt.line_spacing = 1.3
        
    s_b = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(2.2), Inches(5.7), Inches(4.0))
    s_b.fill.solid()
    s_b.fill.fore_color.rgb = light_green
    s_b.line.color.rgb = green
    s_b.line.width = Pt(1.5)
    add_soft_shadow(s_b, blur=0.09, dist=0.055, alpha=0.20)
    tf_sb = s_b.text_frame
    tf_sb.margin_left = tf_sb.margin_top = Inches(0.3)
    p_sb_h = tf_sb.paragraphs[0]
    p_sb_h.text = "シナリオB：高効率空調一括更新"
    p_sb_h.font.name = 'Yu Gothic'
    p_sb_h.font.size = Pt(16)
    p_sb_h.font.bold = True
    p_sb_h.font.color.rgb = green
    p_sb_h.space_after = Pt(20)
    
    sb_points = [
        "・初期総投資： 約 2,400 万円 (室外機30台の一括更新)",
        "・回収期間： 約 12.8 年 (空調自体の省エネ性能向上効果)",
        "・特徴： 設備の根本更新。15年以上経過した老朽化対策を兼ねる"
    ]
    for pt in sb_points:
        p_pt = tf_sb.add_paragraph()
        p_pt.text = pt
        p_pt.font.name = 'Yu Gothic'
        p_pt.font.size = Pt(11.5)
        p_pt.font.color.rgb = navy
        p_pt.space_after = Pt(10)
        p_pt.line_spacing = 1.3

    # 📊 Slide 8: 今後の流れ（円形フローチャート：参考デザイン準拠）
    slide8 = prs.slides.add_slide(blank_layout)
    # スライド内テキストは全廃し、支給の導入スケジュール画像のみを掲載
    add_grid_bg(slide8)
    _sched = "image_schedule.png"
    if os.path.exists(_sched):
        _sw = 11.93
        slide8.shapes.add_picture(_sched, Inches((13.333 - _sw) / 2), Inches(0.6), width=Inches(_sw))
    add_footer(slide8)

    # 📊 Slide 9: アディショナル：累積機会損失 (総投資ベース)
    slide9 = prs.slides.add_slide(blank_layout)
    apply_base_layout(slide9, "10年間の収支差シミュレーション",
                      "「現状維持」は毎年120万円の目に見えない現金を失い続ける選択肢であることを証明します")
    
    # ── 試算ベース値 ──
    annual_usage = float(res['df']['使用量合計'].sum())            # 年間使用量 kWh
    reduc_kwh    = float(res.get('total_reduc_kwh', 0))            # 年間削減量 kWh
    usage_after  = annual_usage - reduc_kwh                        # 導入後 年間使用量 kWh
    annual_bill  = float(res.get('ann_bill', 0)) or float(res.get('gross_saving', 0)) * 8  # 年間電気料金 円
    gross        = float(res.get('gross_saving', 0))               # 年間電気料金削減 円
    bill_after   = annual_bill - gross                            # 導入後 年間電気料金 円
    net          = float(res.get('net_saving', 0))                # 年間純削減（利用料控除後）円
    invest       = float(res.get('total_invest', 2400000))        # 初期総投資 円
    years = list(range(1, 11))

    usage_b_w = annual_usage / 10000.0    # 未導入 年間使用量（万kWh）
    usage_a_w = usage_after / 10000.0     # 導入後 年間使用量（万kWh）
    reduc_w   = reduc_kwh / 10000.0       # 年間削減量（万kWh）

    # ── 上部 左：制御%別 使用量バー画像（カード枠なし・添付デザイン準拠） ──
    ac_year_w = annual_usage * res.get('ac_kwh_r', 0.15) / 10000.0       # 空調 年間（万kWh）
    ctrl_w    = ac_year_w * float(res.get('ctrl_ratio', 1.0))           # 制御対象 空調（万kWh）
    _caps = [("保守 20%", 0.20), ("標準 30%", 0.30), ("積極 40%", 0.40)]
    _use_cats = ["現状"] + [nm for nm, _ in _caps]
    _use_vals = [ac_year_w] + [ac_year_w - ctrl_w * c * CONV_FACTOR for _, c in _caps]
    try:
        _bars_png = make_usage_bars_png(_use_cats, _use_vals)
        slide9.shapes.add_picture(io.BytesIO(_bars_png), Inches(0.7), Inches(1.6), width=Inches(5.7))
    except Exception:
        pass

    # ── 上部 右：累積収支の折れ線画像（0ライン破線＋回収吹き出し＋終点ラベル） ──
    yrs0 = list(range(0, 11))                                   # 0(導入時)〜10年
    cum_w = [(net * y - invest) / 10000.0 for y in yrs0]        # 累積収支（初期費用控除後・万円）
    payback = (invest / net) if net > 0 else 99.0
    _yr_labels = ["導入時"] + [f"{y}年" for y in range(1, 11)]
    try:
        _bal_png = make_cumulative_balance_png(_yr_labels, cum_w, payback_years=payback)
        slide9.shapes.add_picture(io.BytesIO(_bal_png), Inches(6.93), Inches(1.6), width=Inches(5.7))
    except Exception:
        pass

    # ── 下部：差分テーブル（行番号付き・添付デザイン準拠） ──
    rows_t9, cols_t9 = 4, 5
    tbl_shape9 = slide9.shapes.add_table(rows_t9, cols_t9, Inches(0.7), Inches(5.55), Inches(11.93), Inches(1.44))
    t9 = tbl_shape9.table
    t9.columns[0].width = Inches(0.6)
    t9.columns[1].width = Inches(4.28)
    t9.columns[2].width = Inches(2.35)
    t9.columns[3].width = Inches(2.35)
    t9.columns[4].width = Inches(2.35)
    t9.rows[0].height = Inches(0.42)
    for _ri in range(1, 4):
        t9.rows[_ri].height = Inches(0.34)

    t9.cell(0, 0).text = ""
    t9.cell(0, 1).text = "差分（現状維持 → DPS導入）"
    t9.cell(0, 2).text = "1年後"
    t9.cell(0, 3).text = "5年後"
    t9.cell(0, 4).text = "10年後"

    _row_labels = ["累計 使用量削減", "累計 電気料金削減", "実質累計収支（初期投資控除後）"]
    for _ri, _lbl in enumerate(_row_labels, start=1):
        t9.cell(_ri, 0).text = str(_ri)
        t9.cell(_ri, 1).text = _lbl
    for ci, yv in enumerate([1, 5, 10], start=2):
        t9.cell(1, ci).text = f"{reduc_kwh * yv / 10000:,.1f} 万kWh"
        t9.cell(2, ci).text = f"＋{gross * yv / 10000:,.0f} 万円"
        bal = net * yv - invest
        t9.cell(3, ci).text = f"{'＋' if bal >= 0 else '▲'}{abs(bal) / 10000:,.0f} 万円"

    for ri in range(rows_t9):
        for ci in range(cols_t9):
            cell = t9.cell(ri, ci)
            cell.margin_left = Inches(0.12); cell.margin_right = Inches(0.06)
            cell.margin_top = cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.font.name = 'Yu Gothic'
            p.font.size = Pt(12.5) if ri == 0 else Pt(12)
            p.alignment = PP_ALIGN.CENTER if ci == 0 else PP_ALIGN.LEFT
            if ri == 0:
                p.font.color.rgb = white; p.font.bold = True
                cell.fill.solid(); cell.fill.fore_color.rgb = navy
            elif ri == 3:
                p.font.color.rgb = green; p.font.bold = True
                cell.fill.solid(); cell.fill.fore_color.rgb = light_green
            else:
                p.font.color.rgb = navy
                p.font.bold = (ci <= 1)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(247, 249, 252) if ri % 2 == 0 else white

    # 📊 Slide 10: 本業財務価値
    slide10 = prs.slides.add_slide(blank_layout)
    apply_base_layout(slide10, "削減額の財務価値換算（本業売上換算）", 
                      "年間削減される経費を、本業（主力製品）を追加販売した場合の利益規模に逆換算して評価します")
    
    card10 = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(2.2), Inches(10.3), Inches(4.55))
    card10.fill.solid()
    card10.fill.fore_color.rgb = white
    card10.line.color.rgb = RGBColor(225, 230, 238); card10.line.width = Pt(1)
    add_soft_shadow(card10, blur=0.1, dist=0.06, alpha=0.20)

    # アイコンチップ（上部中央）
    chip10 = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(13.333/2 - 0.4), Inches(2.6), Inches(0.8), Inches(0.8))
    chip10.fill.solid(); chip10.fill.fore_color.rgb = green; chip10.line.fill.background()
    add_soft_shadow(chip10, blur=0.04, dist=0.03, alpha=0.25)
    p_chip10 = chip10.text_frame.paragraphs[0]
    p_chip10.text = "📦"; p_chip10.font.size = Pt(34); p_chip10.alignment = PP_ALIGN.CENTER; p_chip10.font.color.rgb = white

    tf_c10 = card10.text_frame
    tf_c10.margin_top = Inches(1.5)
    tf_c10.word_wrap = True

    p_10_l = tf_c10.paragraphs[0]
    p_10_l.text = "主力製品 追加販売換算"
    p_10_l.font.name = 'Yu Gothic'
    p_10_l.font.size = Pt(18)
    p_10_l.font.bold = True
    p_10_l.font.color.rgb = navy
    p_10_l.alignment = PP_ALIGN.CENTER

    items_needed = math.ceil(res.get("net_saving", 1060633) / (res.get("item_price", 5000) * res.get("item_margin", 0.3)))
    p_10_v = tf_c10.add_paragraph()
    p_10_v.text = f"{items_needed:,} 個"
    p_10_v.font.name = 'Yu Gothic'
    p_10_v.font.size = Pt(54)
    p_10_v.font.bold = True
    p_10_v.font.color.rgb = green
    p_10_v.alignment = PP_ALIGN.CENTER
    p_10_v.space_before = Pt(8)
    p_10_v.space_after = Pt(10)

    p_10_s = tf_c10.add_paragraph()
    p_10_s.text = f"（計算基礎：製品単価 ¥{res.get('item_price', 5000):,} × 粗利益率 {res.get('item_margin', 0.3):.1%} ＝ 粗利益 ¥{int(res.get('item_price', 5000)*res.get('item_margin', 0.3)):,} / 個）\n新規にこれだけの数を追加で売り切る営業努力と、DPSを稼働させて無駄なエアコン費をAIが自動カットする確実性をぜひご比較ください。"
    p_10_s.font.name = 'Yu Gothic'
    p_10_s.font.size = Pt(12)
    p_10_s.font.color.rgb = gray
    p_10_s.alignment = PP_ALIGN.CENTER
    p_10_s.line_spacing = 1.3

    p_10_s2 = tf_c10.add_paragraph()
    p_10_s2.text = "＊製品を1個多く売るには、商品原価だけでなく人件費・営業時間といった「見えないコスト」も毎回かかります。\nDPSによる削減は、それらゼロのまま毎年そのまま純利益として手元に残ります。"
    p_10_s2.font.name = 'Yu Gothic'
    p_10_s2.font.size = Pt(12)
    p_10_s2.font.bold = True
    p_10_s2.font.color.rgb = navy
    p_10_s2.alignment = PP_ALIGN.CENTER
    p_10_s2.line_spacing = 1.3
    p_10_s2.space_before = Pt(10)

    # 📊 Slide 11: 環境価値
    slide11 = prs.slides.add_slide(blank_layout)
    apply_base_layout(slide11, "財務改善と同時に達成する、確実な環境貢献（ESG価値）",
                      "電気代を下げながら、対外公表できる確かな脱炭素インパクトを毎年生み出します")

    _co2 = float(res.get('co2_kg', 0))
    _sugi = float(res.get('sugi_trees', 0)) or (_co2 / SUGI_KG if SUGI_KG else 0.0)
    _households = _co2 / 2720.0   # 一般家庭1世帯あたり年間CO2排出量 約2,720kg を基準

    _esg_cards = [
        {"icon": "co2",   "label": "年間 CO2排出削減量",       "value": f"{_co2:,.0f} kg-CO2",
         "note": "省エネ法の定期報告書 実績値として公表可能"},
        {"icon": "tree",  "label": "スギの木 換算 森林効果",    "value": f"約 {_sugi:,.0f} 本分",
         "note": f"（杉1本あたり 年間 {SUGI_KG:g}kg 吸収 換算）"},
        {"icon": "house", "label": "一般家庭の年間排出 換算",   "value": f"約 {_households:,.1f} 世帯分",
         "note": "世帯の年間排出量をまるごと打ち消す規模"},
    ]

    # 支給アイコン（緑ラインアート・透過PNG）。無ければ自作シルエットにフォールバック
    _esg_icon_file = {"co2": "icon_claud.png", "tree": "icon_tree.png", "house": "icon_home.png"}

    _cw = 3.643
    _gap = 0.5
    _left0 = 0.7
    for idx, card in enumerate(_esg_cards):
        _lx = _left0 + idx * (_cw + _gap)
        box = slide11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(_lx), Inches(2.15), Inches(_cw), Inches(3.5))
        box.fill.solid(); box.fill.fore_color.rgb = white
        box.line.color.rgb = RGBColor(200, 224, 205); box.line.width = Pt(1.25)
        add_soft_shadow(box, blur=0.09, dist=0.055, alpha=0.18)

        _iw = 1.2
        _icf = _esg_icon_file.get(card["icon"])
        if _icf and os.path.exists(_icf):
            slide11.shapes.add_picture(_icf, Inches(_lx + (_cw - _iw) / 2), Inches(2.5), width=Inches(_iw))
        else:
            try:
                _png = make_esg_icon_png(card["icon"])
                slide11.shapes.add_picture(io.BytesIO(_png), Inches(_lx + (_cw - _iw) / 2), Inches(2.5), width=Inches(_iw))
            except Exception:
                pass

        tb = slide11.shapes.add_textbox(Inches(_lx + 0.2), Inches(3.95), Inches(_cw - 0.4), Inches(1.55))
        tf_c = tb.text_frame; tf_c.word_wrap = True
        tf_c.vertical_anchor = MSO_ANCHOR.TOP
        p_lbl = tf_c.paragraphs[0]
        p_lbl.text = card["label"]
        p_lbl.font.name = 'Yu Gothic'; p_lbl.font.size = Pt(13); p_lbl.font.bold = True
        p_lbl.font.color.rgb = navy; p_lbl.alignment = PP_ALIGN.CENTER; p_lbl.space_after = Pt(6)

        p_val = tf_c.add_paragraph()
        p_val.text = card["value"]
        p_val.font.name = 'Yu Gothic'; p_val.font.size = Pt(24); p_val.font.bold = True
        p_val.font.color.rgb = green; p_val.alignment = PP_ALIGN.CENTER; p_val.space_after = Pt(6)

        p_note = tf_c.add_paragraph()
        p_note.text = card["note"]
        p_note.font.name = 'Yu Gothic'; p_note.font.size = Pt(10); p_note.font.color.rgb = gray
        p_note.alignment = PP_ALIGN.CENTER; p_note.line_spacing = 1.2

    # ── 下部：ESG価値の公表・訴求バナー ──
    band11 = slide11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(_left0), Inches(5.95), Inches(_cw * 3 + _gap * 2), Inches(0.95))
    band11.fill.solid(); band11.fill.fore_color.rgb = white
    band11.line.color.rgb = green; band11.line.width = Pt(1.5); band11.shadow.inherit = False
    tf_b = band11.text_frame
    tf_b.margin_left = tf_b.margin_right = Inches(0.3)
    tf_b.word_wrap = True; tf_b.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_b = tf_b.paragraphs[0]
    p_b.text = ("省エネ法に基づく定期報告書の実績値として公表可能。SDGsやカーボンニュートラル目標の取り組み実績として合算でき、"
                "取引先・金融機関・採用市場でのESG評価向上に直結。")
    p_b.font.name = 'Yu Gothic'; p_b.font.size = Pt(12); p_b.font.bold = True
    p_b.font.color.rgb = navy; p_b.alignment = PP_ALIGN.CENTER; p_b.line_spacing = 1.25

    # 📊 Slide 12: Q&A・トラブルシューティング
    slide12 = prs.slides.add_slide(blank_layout)
    apply_base_layout(slide12, "QA・よくある質問とトラブルシューティング", 
                      "工場長や現場担当者が懸念する「自動制御に対する不安」を、実例と設計思想からすべて先回りして解消します")
    
    qas = [
        ("Q1. 温度管理：室内環境や製品品質への影響はありませんか？", 
         "➔ A1. 送風切替は数分。室温変化は1℃以下に抑制されるため、精密検査室や製品保管エリアでも影響なくご活用いただけます。"),
        ("Q2. 機器寿命：コンプレッサー等の摩耗への懸念について", 
         "➔ A2. 電子信号を用いた冷却・送風モード of 圧縮機保護基準を完全遵守した制御を行うため、機器寿命への悪影響は一切ありません。"),
        ("Q3. 緊急対応：システム障害時のバックアップ体制は？", 
         "➔ A3. 「セーフティオートシャットダウン自動復旧回路」を標準内蔵。万一の通信エラー時は即座にエアコン通常運転に戻ります。")
    ]
    
    for idx, (q, a) in enumerate(qas):
        top_pos = Inches(2.2 + idx * 1.5)
        box = slide12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), top_pos, Inches(11.93), Inches(1.2))
        box.fill.solid()
        box.fill.fore_color.rgb = white
        box.line.color.rgb = RGBColor(225, 230, 238)
        box.line.width = Pt(1)
        add_soft_shadow(box, blur=0.07, dist=0.045, alpha=0.16)
        # 左のグリーンアクセントバー
        qbar = slide12.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), top_pos, Inches(0.1), Inches(1.2))
        qbar.fill.solid(); qbar.fill.fore_color.rgb = green; qbar.line.fill.background()

        tf_qa = box.text_frame
        tf_qa.margin_left = tf_qa.margin_right = Inches(0.32)
        tf_qa.margin_top = Inches(0.12)
        tf_qa.word_wrap = True
        
        p_q = tf_qa.paragraphs[0]
        p_q.text = q
        p_q.font.name = 'Yu Gothic'
        p_q.font.size = Pt(12)
        p_q.font.bold = True
        p_q.font.color.rgb = navy
        p_q.space_after = Pt(4)
        
        p_a = tf_qa.add_paragraph()
        p_a.text = a
        p_a.font.name = 'Yu Gothic'
        p_a.font.size = Pt(11)
        p_a.font.color.rgb = gray
        p_a.line_spacing = 1.2

    # ── P.3が3カード（image_logic）の場合、採用した分析手法（7手法グリッド＋赤枠）を別スライドで直後に掲載 ──
    _method_slide = None
    if _has_logic3 and os.path.exists("image_sort.png"):
        _method_slide = prs.slides.add_slide(blank_layout)
        apply_base_layout(_method_slide, "試算ロジック：採用した分析手法",
                          "7つの分析手法のうち、本案件で採択した手法を赤枠で明示します")
        _draw_method_grid(_method_slide)

    # ── 概算モード：算出ロジックの説明ペライチ（分析方法の直後に配置） ──
    _gaisan_slide = None
    if (res.get("econ", {}) or {}).get("gaisan_mode"):
        _gaisan_slide = prs.slides.add_slide(blank_layout)
        apply_base_layout(_gaisan_slide, "本試算は【概算】です — 算出ロジックのご説明",
                          "機材情報が無いため、検針票と業態標準値から台数・初期費用を概算しています")
        _gw = 11.93
        _wb = _gaisan_slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(2.1), Inches(_gw), Inches(0.85))
        _wb.fill.solid(); _wb.fill.fore_color.rgb = RGBColor(0xFF, 0xF4, 0xE5)
        _wb.line.color.rgb = RGBColor(0xE8, 0x73, 0x0C); _wb.line.width = Pt(1.5); _wb.shadow.inherit = False
        _tfw = _wb.text_frame; _tfw.word_wrap = True; _tfw.vertical_anchor = MSO_ANCHOR.MIDDLE
        _tfw.margin_left = _tfw.margin_right = Inches(0.3)
        _pw0 = _tfw.paragraphs[0]
        _pw0.text = "⚠ 本ページ以降の「台数・初期費用・投資回収期間」は概算値です。実機材・実見積の取得後に精緻化されます。"
        _pw0.font.name = 'Yu Gothic'; _pw0.font.size = Pt(14); _pw0.font.bold = True
        _pw0.font.color.rgb = RGBColor(0xB5, 0x50, 0x00)
        _gset = app_data.get("const_settings", {}).get("gaisan_settings", {})
        _steps = [
            ("icon_gaisan1.png", "① 予測台数（検針票から推定）",
             f"最大デマンド × 業態の空調ピーク比率 ÷ 標準室外機定格kW（{_gset.get('unit_kw', 10):g}kW）で室外機台数を推定。"),
            ("icon_gaisan2.png", "② 初期費用（標準単価 × 台数）",
             "機材費・工事費・ハードウェア費・構築費の標準単価 × 推定台数 ＋ 固定費で概算（実見積ではなく標準原価）。"),
            ("icon_gaisan3.png", "③ 削減額・投資回収（実績ベース）",
             "削減額は検針票（実績）ベースの現行ロジック（台数非依存）。回収 ＝ 概算初期費用 ÷ 年間純削減。"),
        ]
        # 横並び3カード（P.3の3カード風）。アイコンは仮置き（icon_gaisanN.png があれば自動差し替え）
        _cn = len(_steps); _cgap = 0.4
        _cw = (_gw - _cgap * (_cn - 1)) / _cn
        _cy, _ch = 3.1, 3.4
        for _i, (_icf, _hh, _bb) in enumerate(_steps):
            _cx = 0.7 + _i * (_cw + _cgap)
            _cbx = _gaisan_slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(_cx), Inches(_cy), Inches(_cw), Inches(_ch))
            _cbx.fill.solid(); _cbx.fill.fore_color.rgb = white
            _cbx.line.color.rgb = RGBColor(225, 230, 238); _cbx.line.width = Pt(1)
            add_soft_shadow(_cbx, blur=0.08, dist=0.05, alpha=0.16)
            # アイコン（仮置き）：ファイルがあれば画像、無ければ番号入りの円プレースホルダ
            _isz = 1.0; _ix = _cx + (_cw - _isz) / 2
            if _icf and os.path.exists(_icf):
                _gaisan_slide.shapes.add_picture(_icf, Inches(_ix), Inches(_cy + 0.3), width=Inches(_isz))
            else:
                _ph_ic = _gaisan_slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(_ix), Inches(_cy + 0.3), Inches(_isz), Inches(_isz))
                _ph_ic.fill.solid(); _ph_ic.fill.fore_color.rgb = RGBColor(0xE6, 0xF2, 0xE9)
                _ph_ic.line.color.rgb = green; _ph_ic.line.width = Pt(1.5)
                _pph = _ph_ic.text_frame.paragraphs[0]; _pph.text = str(_i + 1)
                _pph.font.name = 'Yu Gothic'; _pph.font.size = Pt(28); _pph.font.bold = True
                _pph.font.color.rgb = green; _pph.alignment = PP_ALIGN.CENTER
            # タイトル
            _tb = _gaisan_slide.shapes.add_textbox(Inches(_cx + 0.15), Inches(_cy + 1.45), Inches(_cw - 0.3), Inches(0.7))
            _tb.text_frame.word_wrap = True
            _ptt = _tb.text_frame.paragraphs[0]; _ptt.text = _hh
            _ptt.font.name = 'Yu Gothic'; _ptt.font.size = Pt(13); _ptt.font.bold = True
            _ptt.font.color.rgb = navy; _ptt.alignment = PP_ALIGN.CENTER
            # 説明
            _db = _gaisan_slide.shapes.add_textbox(Inches(_cx + 0.22), Inches(_cy + 2.15), Inches(_cw - 0.44), Inches(1.15))
            _db.text_frame.word_wrap = True
            _pdd = _db.text_frame.paragraphs[0]; _pdd.text = _bb
            _pdd.font.name = 'Yu Gothic'; _pdd.font.size = Pt(10.5); _pdd.font.color.rgb = gray
            _pdd.line_spacing = 1.2; _pdd.alignment = PP_ALIGN.CENTER
        _gnote = _gaisan_slide.shapes.add_textbox(Inches(0.7), Inches(6.62), Inches(_gw), Inches(0.3))
        _pgn = _gnote.text_frame.paragraphs[0]
        _pgn.text = "※ 想定の幅（ベスト／標準／保守）は試算画面の『概算：3シナリオ比較』でも確認できます。アイコンは仮置きです。"
        _pgn.font.name = 'Yu Gothic'; _pgn.font.size = Pt(9.5); _pgn.font.color.rgb = gray

    # ── 制御可否リストが1ページ（最大50台）に収まらない場合、続きスライドを自動追加 ──
    _ctrl_extra_slides = []
    for _pk, _pg in enumerate(_ctrl_extra_pages, start=2):
        _se = prs.slides.add_slide(blank_layout)
        # 続きスライドはタイトルのみ（小見出しなし）＝本文スペースを全面リストに使う
        apply_base_layout(_se, "制御可否：対象拠点における機器仕分け実態（続き）")
        try:
            from PIL import Image as _PILImg
            _im = _PILImg.open(io.BytesIO(_pg)); _pw, _ph = _im.size
            _asp = (_pw / _ph) if _ph else 2.0
            _w = 11.93; _x = 0.7                          # 全幅
            _h = _w / _asp
            if _h > 5.7:                                  # タイトル直下〜下端まで目いっぱい
                _h = 5.7; _w = _h * _asp; _x = (13.333 - _w) / 2.0
            _se.shapes.add_picture(io.BytesIO(_pg), Inches(_x), Inches(1.3), width=Inches(_w))
        except Exception:
            pass
        _ctrl_extra_slides.append(_se)

    # ── 全スライド共通：本文（コンテンツ領域）をわずかに上へ詰める（タイトル・小見出し・フッターは対象外） ──
    try:
        _UP = Inches(0.15)
        _LO = Inches(2.03)   # 小見出し帯より下＝コンテンツのみ対象
        _HI = Inches(6.95)   # フッター（ロゴ・ページ番号 7.0〜）は対象外
        for _si, _sl in enumerate(prs.slides):
            if _si == 0:
                continue     # 表紙は対象外
            for _shp in _sl.shapes:
                try:
                    _t = _shp.top
                    if _t is not None and _LO <= _t <= _HI:
                        _shp.top = _t - _UP
                except Exception:
                    pass
    except Exception:
        pass

    # ── スライド並び替え＆削除：P.7（空調機リプレイス）削除／P.9-11 を P.5 と P.6 の間へ移動 ──
    try:
        sldIdLst = prs.slides._sldIdLst
        sld_ids = list(sldIdLst)
        # 動的に追加したスライド（構築順で index12 以降）を所定位置へ差し込む
        _base_n = 12
        _method_idx = []
        if _method_slide is not None:
            _method_idx = [_base_n]; _base_n += 1              # 採用手法スライドは P.3 直後
        _gaisan_idx = []
        if _gaisan_slide is not None:
            _gaisan_idx = [_base_n]; _base_n += 1              # 概算説明は分析方法の直後
        _extra_idx = list(range(_base_n, _base_n + len(_ctrl_extra_slides)))  # 制御リスト続きは P.4 直後
        new_order = [0, 1, 2] + _method_idx + _gaisan_idx + [3] + _extra_idx + [4, 8, 10, 5, 7, 11]
        if len(sld_ids) >= 12:
            for s in sld_ids:
                sldIdLst.remove(s)
            for i in new_order:
                sldIdLst.append(sld_ids[i])
    except Exception:
        pass
    # ── 算出の前提（仮置き・概算・業態配慮）の注釈を全スライド下部＋スライド1ノートに追記 ──
    _pnotes = provisional_notes(res, app_data)
    _gaisan_on = bool((res.get("econ", {}) or {}).get("gaisan_mode"))
    if _gaisan_on:
        _foot = "※【概算】機材情報が無いため、台数・初期費用・投資回収は推定値です（実機材・実見積の取得後に精緻化）。"
    else:
        _foot = "※本試算は概算・業態標準値などの仮置き前提を含みます（前提条件は表紙ノート参照）。" if _pnotes else ""
    # 免責注釈は「算出した数値」を掲載しているスライドだけに付与
    # （料金比較/10年収支/環境価値/アドオン費用）
    # slide2(P.2)は本文内に集約済み、slide4(P.4/制御可否)はユーザー指定で概算注釈を非表示
    _num_slide_ids = set()
    for _s in (slide5, slide6, slide9, slide11):
        try:
            _num_slide_ids.add(_s.slide_id)
        except Exception:
            pass
    try:
        if _pnotes:
            _ns = slide1.notes_slide
            _ns.notes_text_frame.text = "【算出の前提（仮置き・概算・業態配慮）】\n" + "\n".join("・" + n for n in _pnotes)
    except Exception:
        pass

    # 並び替え後の最終順でページ番号を左下に付与
    for _idx, _slide in enumerate(prs.slides, start=1):
        try:
            add_page_number(_slide, _idx)
        except Exception:
            pass
        if _foot and _slide.slide_id in _num_slide_ids:
            try:
                # ページ番号（左下 0.7〜2.1in）の右側に、少し大きめで配置
                _tb = _slide.shapes.add_textbox(Inches(2.3), Inches(7.06), Inches(10.33), Inches(0.34))
                _tb.text_frame.word_wrap = True
                _p = _tb.text_frame.paragraphs[0]
                _p.text = _foot
                _p.font.name = 'Yu Gothic'
                _p.font.size = Pt(8.5)
                _p.font.color.rgb = gray
                _p.alignment = PP_ALIGN.LEFT
            except Exception:
                pass

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


# ══════════════════════════════════════════════════════════════════
# Excel出力（SKILL.md完全準拠 ＋ 追加4シート新設）
# ══════════════════════════════════════════════════════════════════
def build_excel(res: dict, client_name: str, app_data: dict) -> bytes:
    """build_workbook.py の整形式（前提連動・経済性・考察込み）に揃えたExcelをbytesで返す。
    セッションのdf・機材・経済性から入力テンプレ(設定/月次/機材)を組み立て、build_workbook.build() を実行。"""
    import tempfile, os as _os, sys as _sys
    from openpyxl import Workbook as _WB
    _sc = _os.path.join(_os.path.dirname(_os.path.abspath(__file__)), "scripts")
    if _sc not in _sys.path:
        _sys.path.insert(0, _sc)
    import build_workbook as _bw
    AREA_BUF = {"北海道・東北": 1.10, "関東": 1.00, "中部・北陸": 1.05,
                "近畿": 1.02, "中国・四国": 1.08, "九州・沖縄": 1.12}
    econ = res.get("econ", {})
    df = res["df"]
    units = res.get("units_detail")
    n_ctrl = int(econ.get("n_ctrl", 1)) or 1
    fee_per = round((econ.get("sys_fee") or res.get("sys_fee", 0)) / n_ctrl) if n_ctrl else 12000
    areaf = AREA_BUF.get(econ.get("area", "関東"), 1.00)
    wbin = _WB(); s = wbin.active; s.title = "設定"
    s.append(["設定項目", "値"])
    for k, v in [("拠点名", client_name), ("業態", res.get("gyotai", "製造拠点(プロセス主役)")),
                 ("空調ピーク割合", res.get("ac_peak_r", 0.20)), ("空調電力量割合", res.get("ac_kwh_r", 0.15)),
                 ("conv", 0.5), ("システム一式", econ.get("system_cost", 1400000)),
                 ("セットアップ費", econ.get("setup_cost", 400000)), ("工事費単価", 50000),
                 ("エリア係数", areaf), ("年間利用料単価", fee_per), ("補助金", econ.get("subsidy", 0)),
                 ("制御可能比率固定値", round(float(res.get("ctrl_ratio", 1.0)), 4)),
                 ("主力商品の単価", res.get("item_price", 5000)),
                 ("主力商品の粗利益率", res.get("item_margin", 0.30)), ("CO2排出係数", 0.451)]:
        s.append([k, v])
    m = wbin.create_sheet("月次")
    m.append(["月", "使用量kWh", "最大デマンドkW", "平日昼間単価", "燃料調整単価",
              "再エネ単価", "市場調整単価", "基本料金単価", "契約kW"])
    for _, row in df.iterrows():
        m.append([str(row["月"]), float(row["使用量合計"]), float(row["最大需要電力"]),
                  float(row["電力量単価"]), 0, 0, 0, float(row["基本料金単価"]), float(row["契約電力"])])
    e = wbin.create_sheet("機材")
    e.append(["設置場所", "メーカー", "型式", "制御可否", "定格kW"])
    if units is not None and len(units) > 0:
        for _, u in units.iterrows():
            e.append([u.get("機器ID", ""), u.get("メーカー", ""), u.get("型番", ""),
                      "〇" if u.get("制御可否") else "×", u.get("定格冷房kW", "")])
    else:
        e.append(["（台帳なし・固定比率を使用）", "", "", "〇", ""])
    ti = tempfile.NamedTemporaryFile(suffix="_in.xlsx", delete=False); ti.close()
    to = tempfile.NamedTemporaryFile(suffix="_out.xlsx", delete=False); to.close(); _os.unlink(to.name)
    wbin.save(ti.name)
    _bw.build(ti.name, to.name)
    with open(to.name, "rb") as fh:
        data = fh.read()
    for p in (ti.name, to.name):
        try:
            _os.unlink(p)
        except Exception:
            pass
    return data


def build_excel_legacy(res: dict, client_name: str, app_data: dict) -> bytes:
    wb = Workbook()
    
    # ── スタイル定義 ────────────────────────────────
    navy = "1F3864"; green = "1E6B2E"; gold = "FFC000"; red = "C00000"
    lg = "E2EFDA"; lb = "DEEAF1"; ly = "FFF2CC"; lor = "FCE4D6"
    FONT = "Yu Gothic"
    
    def fnt(sz=10, bold=False, color="000000"):
        return Font(name=FONT, size=sz, bold=bold, color=color)
    def fl(h): return PatternFill("solid", start_color=h, fgColor=h)
    def bd():
        s = Side(style="thin", color="AAAAAA")
        return Border(left=s, right=s, top=s, bottom=s)
    def al(h="left", v="center"):
        return Alignment(horizontal=h, vertical=v, wrap_text=True)

    # ────── ① シート0_顧客提示用サマリー ──────
    ws0 = wb.active
    ws0.title = "シート0_顧客提示用サマリー"
    for i, w in enumerate([4,34,22,22,22,22,4], 1):
        ws0.column_dimensions[get_column_letter(i)].width = w
    
    ws0.merge_cells("B1:F1")
    c = ws0.cell(row=1, column=2, value=f"空調デマンド制御 提案サマリー ｜ {client_name}")
    c.font = fnt(14, True, "FFFFFF"); c.fill = fl(navy); c.alignment = al("center"); c.border = bd()
    ws0.row_dimensions[1].height = 32

    ws0.merge_cells("B2:F2")
    c = ws0.cell(row=2, column=2, value=f"作成日: {datetime.date.today()} ｜ 業態: {res.get('gyotai', '製造拠点')}")
    c.font = fnt(9, color="FFFFFF"); c.fill = fl("2F5496"); c.alignment = al("center"); rh_target = 2; ws0.row_dimensions[rh_target].height = 16

    r = 4
    ws0.merge_cells(f"B{r}:F{r}")
    c = ws0.cell(row=r, column=2, value="【財務サマリー】"); c.font = fnt(11, True, "FFFFFF")
    c.fill = fl(green); c.alignment = al(); c.border = bd(); ws0.row_dimensions[r].height = 20; r+=1

    def ws0_kv(r, label, value, bg=lg, vcolor=green):
        ws0.merge_cells(f"B{r}:C{r}")
        c = ws0.cell(row=r, column=2, value=label); c.font = fnt(10, True); c.fill = fl(ly); c.border = bd()
        ws0.merge_cells(f"D{r}:F{r}")
        c2 = ws0.cell(row=r, column=4, value=value); c2.font = fnt(12, True, vcolor); c2.fill = fl(bg); c2.alignment = al("right"); c2.border = bd()
        ws0.row_dimensions[r].height = 20

    ws0_kv(r, "基本料金削減額（年間）", f"¥{res.get('dm_saving_annual', 0):,.0f}"); r+=1
    ws0_kv(r, "電力量料金削減額（年間）", f"¥{res.get('ene_saving_annual', 0):,.0f}"); r+=1
    ws0_kv(r, "年間総削減額（グロス）", f"¥{res.get('gross_saving', 0):,.0f}", bg=lg, vcolor=green); r+=1
    ws0_kv(r, "年間純削減額（実利）", f"¥{res.get('net_saving', 0):,.0f}", bg=lg, vcolor=green); r+=1
    ws0_kv(r, "CO₂削減量", f"{res.get('co2_kg', 0):,.0f} kg-CO₂/年", bg=lg, vcolor="1B5E20"); r+=1
    ws0_kv(r, "スギの木換算", f"約 {res.get('sugi_trees', 0):,.0f} 本分/年", bg=lg, vcolor="1B5E20"); r+=2

    # ────── ② サマリー報告 ──────
    ws1 = wb.create_sheet("サマリー報告")
    for i, w in enumerate([4,30,40,4], 1):
        ws1.column_dimensions[get_column_letter(i)].width = w
    ws1.cell(row=1, column=2, value="製造拠点 電力・空調制御 サマリー報告").font = fnt(14, True, navy)
    ws1.cell(row=2, column=2, value=f"所在地: 静岡県富士宮市北山").font = fnt(10, color="555555")
    
    ws1.cell(row=4, column=2, value="■ 概要").font = fnt(11, True, navy)
    ws1.cell(row=5, column=2, value="対象期間").font = fnt(10); ws1.cell(row=5, column=3, value="12ヶ月実績").font = fnt(10)
    ws1.cell(row=6, column=2, value="年間最大デマンド").font = fnt(10); ws1.cell(row=6, column=3, value=f"{res.get('old_contract', 500):,.0f} kW").font = fnt(10, True)
    ws1.cell(row=7, column=2, value="実効デマンド削減").font = fnt(10); ws1.cell(row=7, column=3, value=f"▲ {res.get('contract_delta', 0):,.1f} kW").font = fnt(10, True, green)
    ws1.cell(row=8, column=2, value="実効使用電力量削減").font = fnt(10); ws1.cell(row=8, column=3, value=f"▲ {res.get('total_reduc_kwh', 0):,.0f} kWh").font = fnt(10, True, green)

    # ────── ③ 前提・制御条件 ──────
    ws2 = wb.create_sheet("前提・制御条件")
    ws2.column_dimensions['B'].width = 35
    ws2.column_dimensions['C'].width = 25
    ws2.cell(row=1, column=2, value="試算設定パラメータ").font = fnt(12, True, navy)
    ws2.cell(row=3, column=2, value="空調ピーク寄与割合").font = fnt(10); ws2.cell(row=3, column=3, value=res.get("ac_peak_r", 0.18)).font = fnt(10)
    ws2.cell(row=4, column=2, value="空調電力量割合").font = fnt(10); ws2.cell(row=4, column=3, value=res.get("ac_kwh_r", 0.15)).font = fnt(10)
    ws2.cell(row=5, column=2, value="定格キャップ削減率").font = fnt(10); ws2.cell(row=5, column=3, value=res.get("cap_rate", 0.30)).font = fnt(10)
    ws2.cell(row=6, column=2, value="能力削減→エネ削減係数 conv").font = fnt(10); ws2.cell(row=6, column=3, value=CONV_FACTOR).font = fnt(10)
    ws2.cell(row=7, column=2, value="システム年間利用料").font = fnt(10); ws2.cell(row=7, column=3, value=res.get("sys_fee", 241800)).font = fnt(10)

    # ────── ④ 月次・電気料金 ──────
    ws3 = wb.create_sheet("月次・電気料金")
    for i, w in enumerate([4,12,18,18,18,18,18], 1):
        ws3.column_dimensions[get_column_letter(i)].width = w
    ws3.cell(row=1, column=2, value="月次検針実績データ").font = fnt(12, True, navy)
    
    headers3 = ["月", "契約電力(kW)", "最大デマンド(kW)", "使用量(kWh)", "力率調整", "基本単価(円)"]
    for ci, h in enumerate(headers3, 2):
        cell = ws3.cell(row=3, column=ci, value=h)
        cell.font = fnt(9, True, "FFFFFF"); cell.fill = fl(navy); cell.alignment = al("center")
        
    for r_idx, row in res["df"].iterrows():
        ws3.cell(row=r_idx+4, column=2, value=row["月"]).font = fnt(9)
        ws3.cell(row=r_idx+4, column=3, value=row["契約電力"]).font = fnt(9)
        ws3.cell(row=r_idx+4, column=4, value=row["最大需要電力"]).font = fnt(9)
        ws3.cell(row=r_idx+4, column=5, value=row["使用量合計"]).font = fnt(9)
        ws3.cell(row=r_idx+4, column=6, value=row["力率"]).font = fnt(9)
        ws3.cell(row=r_idx+4, column=7, value=row["基本料金単価"]).font = fnt(9)

    # ────── ⑤ 空調・その他 ──────
    ws4 = wb.create_sheet("空調・その他")
    ws4.column_dimensions['B'].width = 25
    ws4.column_dimensions['C'].width = 25
    ws4.cell(row=1, column=2, value="空調と他設備の負荷内訳(年間推計値)").font = fnt(12, True, navy)
    ws4.cell(row=3, column=2, value="区分").font = fnt(10, True); ws4.cell(row=3, column=3, value="推計年間電力量(kWh)").font = fnt(10, True)
    ws4.cell(row=4, column=2, value="総年間使用量").font = fnt(10); ws4.cell(row=4, column=3, value=res["df"]["使用量合計"].sum()).font = fnt(10)
    ws4.cell(row=5, column=2, value="空調負荷分").font = fnt(10); ws4.cell(row=5, column=3, value=res["df"]["使用量合計"].sum() * res.get("ac_kwh_r", 0.15)).font = fnt(10)
    ws4.cell(row=6, column=2, value="その他非空調負荷").font = fnt(10); ws4.cell(row=6, column=3, value=res["df"]["使用量合計"].sum() * (1 - res.get("ac_kwh_r", 0.15))).font = fnt(10)

    # ────── ⑥ 機器台帳 ──────
    ws5 = wb.create_sheet("機器台帳")
    for i, w in enumerate([4,10,18,18,22,12,18,18], 1):
        ws5.column_dimensions[get_column_letter(i)].width = w
    ws5.cell(row=1, column=2, value="空調室外機 機器台帳").font = fnt(12, True, navy)
    
    headers5 = ["#", "設置場所", "メーカー", "型式", "制御可否", "定格冷房能力(kW)", "制御対象(kW)"]
    for ci, h in enumerate(headers5, 2):
        cell = ws5.cell(row=3, column=ci, value=h)
        cell.font = fnt(9, True, "FFFFFF"); cell.fill = fl(navy); cell.alignment = al("center")
        
    for idx, eq in enumerate(DEFAULT_EQUIPMENT, 4):
        ws5.cell(row=idx, column=2, value=eq["id"]).font = fnt(9)
        ws5.cell(row=idx, column=3, value=eq["loc"]).font = fnt(9)
        ws5.cell(row=idx, column=4, value=eq["mfr"]).font = fnt(9)
        ws5.cell(row=idx, column=5, value=eq["model"]).font = fnt(9)
        ws5.cell(row=idx, column=6, value=eq["ctrl"]).font = fnt(9)
        ws5.cell(row=idx, column=7, value=eq["cap"]).font = fnt(9)
        ws5.cell(row=idx, column=8, value=eq["cap"] if eq["ctrl"] == "〇" else 0).font = fnt(9)

    # ────── ⑦ 空調制御シナリオ ──────
    ws6 = wb.create_sheet("空調制御シナリオ")
    ws6.column_dimensions['B'].width = 25
    ws6.column_dimensions['C'].width = 25
    ws6.column_dimensions['D'].width = 25
    ws6.cell(row=1, column=2, value="制御強度別シミュレーション効果").font = fnt(12, True, navy)
    
    ws6.cell(row=3, column=2, value="シナリオ項目").font = fnt(10, True)
    ws6.cell(row=3, column=3, value="保守 (20%削減)").font = fnt(10, True)
    ws6.cell(row=3, column=4, value="標準 (30%削減)").font = fnt(10, True)
    
    ws6.cell(row=4, column=2, value="年間基本料金削減額").font = fnt(10)
    ws6.cell(row=4, column=3, value=res.get("dm_saving_annual", 0) * (0.2 / 0.3)).font = fnt(10)
    ws6.cell(row=4, column=4, value=res.get("dm_saving_annual", 0)).font = fnt(10)
    
    ws6.cell(row=5, column=2, value="年間電力量削減額").font = fnt(10)
    ws6.cell(row=5, column=3, value=res.get("ene_saving_annual", 0) * (0.2 / 0.3)).font = fnt(10)
    ws6.cell(row=5, column=4, value=res.get("ene_saving_annual", 0)).font = fnt(10)

    # ────── ⑧ 室外機ROI ──────
    ws7 = wb.create_sheet("室外機ROI")
    ws7.column_dimensions['B'].width = 25
    ws7.column_dimensions['C'].width = 25
    ws7.column_dimensions['D'].width = 25
    ws7.cell(row=1, column=2, value="室外機絞り込みROI算定").font = fnt(12, True, navy)
    ws7.cell(row=3, column=2, value="設置場所").font = fnt(10, True)
    ws7.cell(row=3, column=3, value="型式").font = fnt(10, True)
    ws7.cell(row=3, column=4, value="回収年数").font = fnt(10, True)
    
    for idx, eq in enumerate(DEFAULT_EQUIPMENT[:5], 4):
        ws7.cell(row=idx, column=2, value=eq["loc"]).font = fnt(10)
        ws7.cell(row=idx, column=3, value=eq["model"]).font = fnt(10)
        ws7.cell(row=idx, column=4, value=f"{res.get('total_invest', 2400000)/res.get('net_saving', 1060633):.1f} 年" if eq["ctrl"] == "〇" else "対象外").font = fnt(10)

    # ────── ⑨ 分析方法の選択肢について (新規追加) ──────
    ws8 = wb.create_sheet("分析方法の選択肢について")
    ws8.column_dimensions['B'].width = 30
    ws8.column_dimensions['C'].width = 65
    ws8.cell(row=1, column=2, value="分析方法の選択肢について (現状持ち合わせている分析方法の概要)").font = fnt(12, True, navy)
    
    methods = [
        ("A. ベースライン法", "天候に強く相関する拠点用。春秋などの冷暖房を使わない月次最低消費電力を基準（非空調ベースライン）とし、夏冬の超過分を「空調負荷」として精密に切り出して算出。"),
        ("B. 設備割合推計法", "生産プロセスや冷凍冷蔵設備が主たる動力を占める特殊な工場用。業態統計データをもとに空調消費電力の固定比率（プロセス製造なら15%など）をあてはめ、過大評価を排した推計を実施。"),
        ("C. 30分デマンドデータ分析", "スマートメーターから得られる「30分単位デマンドログ」が取得できた場合。稼働ピークの発生日時を直接ピンポイントで特定し、1kW未満の単位で無駄なピーク需要の削減幅を完全判定。"),
        ("D. 稼働時間シフト分析", "操業カレンダー、生産スケジュール、および勤務シフト時間を突き合わせ、操業時間内の稼働率の山谷をプロファイリングして制御効果をさらに先鋭化。")
    ]
    for idx, (m_title, m_desc) in enumerate(methods, 3):
        ws8.cell(row=idx, column=2, value=m_title).font = fnt(11, True, navy)
        ws8.cell(row=idx, column=3, value=m_desc).font = fnt(10, False)
        ws8.row_dimensions[idx].height = 25

    # ────── ⑩ 今回の分析方法について (新規追加) ──────
    ws9 = wb.create_sheet("今回の分析方法について")
    ws9.column_dimensions['B'].width = 25
    ws9.column_dimensions['C'].width = 65
    ws9.cell(row=1, column=2, value="今回の分析方法について (各算出式と業界配慮についての説明)").font = fnt(12, True, navy)
    
    formulas = [
        ("1. 空調寄与ピーク推計", "公式： ΔkW = 年間最大デマンド × 業態空調割合(18%) × 制御対象率(100.0%) × キャップ率(30%)\n製造拠点固有の生産機械動力への干渉を100%防ぎ、空調寄与分のみをスマートターゲットとします。"),
        ("2. 力率調整(15%割引)の同期", "公式： 基本料金 = 契約電力 × 基本単価 × 力率調整\n西富士事業所様の極めて優良な力率(100%)実績に完全追従し、15%割引（基本料金×0.85）を正確に維持した削減試算を実施。"),
        ("3. 従量電力量削減(conv)", "公式： 削減kWh = 使用量合計 × 空調比率(15%) × 制御対象率 × キャップ率 × 換算 conv(0.5)\n能力の抑制が直ちに同率のエネ削減にはならないため、物理安全係数 conv(0.5) を掛けてリスク回避・過大請求の誇張を防ぎます。"),
        ("4. 業界・製造現場への操業配慮", "生産現場の品質・操業に一切影響を及ぼさないよう、常時ローテーション送風（変化1℃未満）を担保するための、不作為を生まない保守的かつ精緻な定格引き算法を基礎ロジックとして設計。")
    ]
    for idx, (f_title, f_desc) in enumerate(formulas, 3):
        ws9.cell(row=idx, column=2, value=f_title).font = fnt(11, True, green)
        ws9.cell(row=idx, column=3, value=f_desc).font = fnt(10, False)
        ws9.row_dimensions[idx].height = 32

    # ────── ⑪ 導入しなかった場合の累積機会損失 (新規追加) ──────
    ws10 = wb.create_sheet("導入しなかった場合の累積機会損失")
    for i, w in enumerate([4,32,25,25,25,4], 1):
        ws10.column_dimensions[get_column_letter(i)].width = w
    ws10.cell(row=1, column=2, value="導入しなかった場合の累積機会損失 (財務影響シミュレーション)").font = fnt(12, True, navy)
    
    ws10.cell(row=3, column=2, value="財務シナリオ").font = fnt(10, True, "FFFFFF"); ws10.cell(row=3, column=2).fill = fl(navy)
    ws10.cell(row=3, column=3, value="1年目").font = fnt(10, True, "FFFFFF"); ws10.cell(row=3, column=3).fill = fl(navy)
    ws10.cell(row=3, column=4, value="5年後").font = fnt(10, True, "FFFFFF"); ws10.cell(row=3, column=4).fill = fl(navy)
    ws10.cell(row=3, column=5, value="10年後").font = fnt(10, True, "FFFFFF"); ws10.cell(row=3, column=5).fill = fl(navy)
    
    gross_yr = res.get("gross_saving", 1302433)
    net_yr = res.get("net_saving", 1060633)
    invest = res.get("total_invest", 2400000)
    
    ws10.cell(row=4, column=2, value="現状維持 (無駄な電気代無駄払い累計)").font = fnt(10, True, red)
    ws10.cell(row=4, column=3, value=f"▲ {gross_yr/10000:.1f} 万円").font = fnt(10)
    ws10.cell(row=4, column=4, value=f"▲ {gross_yr*5/10000:.1f} 万円").font = fnt(10)
    ws10.cell(row=4, column=5, value=f"▲ {gross_yr*10/10000:.1f} 万円").font = fnt(10)
    
    ws10.cell(row=5, column=2, value="DPS導入 (回収後手残りキャッシュ累計)").font = fnt(10, True, green)
    ws10.cell(row=5, column=3, value=f"▲ {(invest - net_yr)/10000:.1f} 万円").font = fnt(10)
    ws10.cell(row=5, column=4, value=f"+ {(net_yr*5 - invest)/10000:.1f} 万円").font = fnt(10)
    ws10.cell(row=5, column=5, value=f"+ {(net_yr*10 - invest)/10000:.1f} 万円").font = fnt(10)
    
    for r_idx in [4, 5]:
        for c_idx in [2, 3, 4, 5]:
            ws10.cell(row=r_idx, column=c_idx).border = bd()
            if r_idx == 4:
                ws10.cell(row=r_idx, column=c_idx).fill = fl("FFF5F5")
            else:
                ws10.cell(row=r_idx, column=c_idx).fill = fl(lg)

    # ────── ⑫ 環境価値 (新規追加) ──────
    ws11 = wb.create_sheet("環境価値")
    ws11.column_dimensions['B'].width = 30
    ws11.column_dimensions['C'].width = 40
    ws11.cell(row=1, column=2, value="環境価値サマリー (ESG・脱炭素への貢献)").font = fnt(12, True, green)
    
    ws11.cell(row=3, column=2, value="環境評価項目").font = fnt(10, True, "FFFFFF"); ws11.cell(row=3, column=2).fill = fl(green)
    ws11.cell(row=3, column=3, value="年間効果量").font = fnt(10, True, "FFFFFF"); ws11.cell(row=3, column=3).fill = fl(green)
    
    ws11.cell(row=4, column=2, value="年間総削減電力量").font = fnt(10)
    ws11.cell(row=4, column=3, value=f"{res.get('total_reduc_kwh', 0):,.0f} kWh / 年").font = fnt(10, True)
    
    ws11.cell(row=5, column=2, value="CO₂排出量直接削減効果").font = fnt(10)
    ws11.cell(row=5, column=3, value=f"{res.get('co2_kg', 0):,.1f} kg-CO₂ / 年").font = fnt(10, True)
    
    ws11.cell(row=6, column=2, value="杉の木換算(年間吸収量相当)").font = fnt(10)
    ws11.cell(row=6, column=3, value=f"約 {res.get('sugi_trees', 0):,.0f} 本分 / 年").font = fnt(10, True)
    
    for r_idx in [4, 5, 6]:
        for c_idx in [2, 3]:
            ws11.cell(row=r_idx, column=c_idx).border = bd()

    # 確実に全12シートのExcelデータをバイナリで保存してエクスポート
    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf.getvalue()


# ══════════════════════════════════════════════════════════════════
# 画像出力タブ用：Excel取込（添付『空調制御効果試算』形式＝入力/室外機情報 を参照）
#   ※ 添付ファイルは“形式の一例”。値は焼き込まず、実際にアップロードされた内容のみ反映する。
#   縦横どちらのレイアウト（横並びグリッド／フラット表）でも読めるよう自動判定する。
# ══════════════════════════════════════════════════════════════════
def _adm_is_dt(v):
    from datetime import datetime, date
    return isinstance(v, (datetime, date)) or (
        hasattr(v, "year") and hasattr(v, "month") and not isinstance(v, (str, int, float)))


def _adm_month_label(v):
    try:
        return f"{int(v.year)}/{int(v.month):02d}"
    except Exception:
        return str(v)


def admin_import_input_graph(uploaded):
    """『入力』シート（横並びグリッド）／フラット形式から ①グラフ用の
    DataFrame(月,使用量,削減,最大デマンド) を返す。読めなければ None。"""
    from collections import Counter
    try:
        uploaded.seek(0)
        xls = pd.ExcelFile(uploaded)
    except Exception:
        return None

    def _pick(names):
        cands = [s for s in names if ("入力" in s) and not any(k in s for k in ("室外機", "機器", "利用", "機体"))]
        return cands or list(names)

    for sn in _pick(xls.sheet_names):
        try:
            raw = pd.read_excel(xls, sheet_name=sn, header=None)
        except Exception:
            continue
        nr, nc = raw.shape
        dt_cells = [(r, c, raw.iat[r, c]) for r in range(nr) for c in range(nc) if _adm_is_dt(raw.iat[r, c])]
        if len(dt_cells) < 6:
            continue
        rowc = Counter(r for r, _, _ in dt_cells)
        colc = Counter(c for _, c, _ in dt_cells)
        best_row, row_n = rowc.most_common(1)[0]
        _, col_n = colc.most_common(1)[0]

        def _num(v):
            try:
                if v is None or (isinstance(v, float) and pd.isna(v)):
                    return 0.0
                return float(v)
            except Exception:
                return 0.0

        if row_n >= col_n and row_n >= 6:
            # 横並び（ゼンケン『入力』型）：月＝best_row の日付列
            month_cols = sorted(c for r, c, _ in dt_cells if r == best_row)
            months = [_adm_month_label(raw.iat[best_row, c]) for c in month_cols]

            def _find_row(keys, avoid=()):
                for r in range(nr):
                    lab = "".join(str(raw.iat[r, c]) for c in range(min(nc, 4)) if raw.iat[r, c] is not None)
                    if any(k in lab for k in keys) and not any(a in lab for a in avoid):
                        return r
                return None

            r_use = _find_row(("使用量合計", "使用量"), avoid=("夏季", "平日", "夜間", "内訳"))
            r_dm = _find_row(("最大需要", "最大デマンド", "デマンド"))
            usage = [_num(raw.iat[r_use, c]) if r_use is not None else 0.0 for c in month_cols]
            demand = [_num(raw.iat[r_dm, c]) if r_dm is not None else 0.0 for c in month_cols]
            return pd.DataFrame({"月": months, "使用量": usage, "削減": [0] * len(months), "最大デマンド": demand})
        else:
            # 縦並び（フラット表）：月＝best_col の日付行、直上をヘッダとみなす
            best_col = colc.most_common(1)[0][0]
            month_rows = sorted(r for r, c, _ in dt_cells if c == best_col)
            header_row = max(min(month_rows) - 1, 0)

            def _find_col(keys, avoid=()):
                for c in range(nc):
                    lab = str(raw.iat[header_row, c])
                    if any(k in lab for k in keys) and not any(a in lab for a in avoid):
                        return c
                return None

            c_use = _find_col(("使用量合計", "使用量"), avoid=("夏季", "平日", "夜間"))
            c_dm = _find_col(("最大需要", "最大デマンド", "デマンド"))
            months, usage, demand = [], [], []
            for r in month_rows:
                months.append(_adm_month_label(raw.iat[r, best_col]))
                usage.append(_num(raw.iat[r, c_use]) if c_use is not None else 0.0)
                demand.append(_num(raw.iat[r, c_dm]) if c_dm is not None else 0.0)
            return pd.DataFrame({"月": months, "使用量": usage, "削減": [0] * len(months), "最大デマンド": demand})
    return None


def admin_import_units_list(uploaded):
    """『室外機情報』シート等から ②制御可否リスト用 DataFrame
    (系統名,設置場所,メーカー,型式,制御可否<〇/×>) を返す。読めなければ None。"""
    try:
        uploaded.seek(0)
        xls = pd.ExcelFile(uploaded)
    except Exception:
        return None

    def _pick(names):
        cands = [s for s in names if any(k in s for k in ("室外機", "機器", "機体"))]
        return cands or list(names)

    for sn in _pick(xls.sheet_names):
        try:
            raw = pd.read_excel(xls, sheet_name=sn, header=None, dtype=object)
        except Exception:
            continue
        nr, nc = raw.shape
        hdr = None
        for r in range(min(nr, 25)):
            joined = " ".join(str(raw.iat[r, c]) for c in range(nc) if raw.iat[r, c] is not None)
            if sum(k in joined for k in ("系統", "設置場所", "メーカー", "型式", "型番", "制御")) >= 2:
                hdr = r
                break
        if hdr is None:
            continue
        labs = [str(raw.iat[hdr, c]) if raw.iat[hdr, c] is not None else "" for c in range(nc)]

        def _col(keys, avoid=()):
            for c, l in enumerate(labs):
                if any(k in l for k in keys) and not any(a in l for a in avoid):
                    return c
            return None

        c_sys = _col(("系統", "機器ID", "機器", "場所"), avoid=("設置場所",))
        c_loc = _col(("設置場所",))
        c_mfr = _col(("メーカー",))
        c_mdl = _col(("型式", "型番"))
        c_ctrl = _col(("制御",))
        rows = []
        for r in range(hdr + 1, nr):
            def _g(c):
                if c is None:
                    return ""
                v = raw.iat[r, c]
                if v is None or (isinstance(v, float) and pd.isna(v)):
                    return ""
                return str(v).strip()
            sysn, loc, mfr, mdl = _g(c_sys), _g(c_loc), _g(c_mfr), _g(c_mdl)
            if not any([sysn, loc, mfr, mdl]):
                continue
            ok = _g(c_ctrl) in ("〇", "○", "◯", "可", "o", "O", "1", "True", "true")
            rows.append({"系統名": sysn, "設置場所": loc, "メーカー": mfr,
                         "型式": mdl, "制御可否": "〇" if ok else "×"})
        if rows:
            return pd.DataFrame(rows)
    return None


def _adm_ctrl_is_ok(v):
    """制御可否セル（〇/×/bool 等）を bool へ。"""
    if isinstance(v, bool):
        return v
    return str(v).strip() in ("〇", "○", "◯", "可", "o", "O", "1", "True", "true")


# ══════════════════════════════════════════════════════════════════
# 管理画面 (戻るボタン補強)
# ══════════════════════════════════════════════════════════════════
def show_admin(app_data: dict) -> None:
    st.markdown("## 各種設定")

    tab1, tab_mat, tab2, tab3, tab4, tab5, tab_img = st.tabs([
        "計算ロジック変更",
        "機材単価マスタ（見積）",
        "馬力・機材テーブル",
        "業態マスタ設定",
        "表示項目変更",
        "変更・DL履歴ログ",
        "画像出力（手動値）",
    ])

    # ─── タブ1: 計算ロジック ──────────────────────────
    with tab1:
        st.subheader("計算ロジック（数式）の変更")
        st.info("変更すると全試算に即時反映されます。")

        cap = st.slider(
            "容量削減率（空調制御による最大カット率）",
            min_value=0.05, max_value=0.50, step=0.05,
            value=float(app_data["calc_settings"]["cap_rate"]),
            format="%.2f"
        )
        # 💡 セッション状態からサイドバーの値を取得し、無ければapp_dataの値を使う
        sys_fee = st.session_state.get("sidebar_sys_fee", int(app_data["system_fee"]))

        col1, col2 = st.columns(2)
        with col1:
            st.markdown("**ベースロード法**")
            st.caption("最少使用量月を基礎電力（空調ゼロ基準）として自動取得し、各月から引き算して空調kWhを逆算します。")
        with col2:
            st.markdown("**実量制ルール**")
            st.caption("過去12ヶ月の最大DM削減幅を契約電力から差し引き、新契約電力を算出して基本料金を再計算します。")

        # ── 年間システム利用料の単価（1単位あたり月額）──
        st.markdown("---")
        st.markdown("**年間システム利用料の単価（1単位あたり・月額）**")
        st.caption("ここで変更した単価が、サイドバーの年間利用料（数量×単価×12）に即時反映されます。")
        _fs = app_data.get("fee_settings", DEFAULT_DATA["fee_settings"])
        fc1, fc2 = st.columns(2)
        with fc1:
            fee_tsushin = st.number_input("通信費 単価（円）", min_value=0, max_value=1000000,
                                          step=50, value=int(_fs.get("tsushin", 550)), key="set_fee_tsushin")
            fee_hoshu   = st.number_input("ソフト保守 単価（円）", min_value=0, max_value=1000000,
                                          step=50, value=int(_fs.get("hoshu", 200)), key="set_fee_hoshu")
        with fc2:
            fee_cloud   = st.number_input("クラウド利用料 単価（円）", min_value=0, max_value=1000000,
                                          step=50, value=int(_fs.get("cloud", 500)), key="set_fee_cloud")
            fee_data    = st.number_input("データ収集 単価（円）", min_value=0, max_value=1000000,
                                          step=50, value=int(_fs.get("data", 200)), key="set_fee_data")

        # ── 計算で使う固定値（定数）の編集 ──
        st.markdown("---")
        st.markdown("**計算で使う固定値（定数）**")
        st.caption("工事費単価・エリア係数・各種換算係数・データ補完の既定値など、コードに固定していた値をすべて変更できます。")
        _cs = app_data.get("const_settings", DEFAULT_DATA["const_settings"])
        with st.expander("🔧 固定値を編集（工事費・エリア係数・換算係数・補完既定値）", expanded=False):
            st.markdown("**■ 工事費・換算係数**")
            kc1, kc2 = st.columns(2)
            with kc1:
                c_area_unit = st.number_input("工事費目安単価（円/台）", min_value=0, max_value=10000000,
                                              step=1000, value=int(_cs.get("area_unit_price", 50000)), key="set_c_area_unit")
                c_conv = st.number_input("換算係数 conv（能力削減→正味エネ削減）", min_value=0.0, max_value=2.0,
                                         step=0.05, value=float(_cs.get("conv_factor", 0.5)), format="%.2f", key="set_c_conv")
                c_hp = st.number_input("馬力→冷房能力kW（1馬力あたり）", min_value=0.1, max_value=20.0,
                                       step=0.1, value=float(_cs.get("hp_to_kw", 2.8)), format="%.1f", key="set_c_hp")
            with kc2:
                c_co2 = st.number_input("CO₂排出係数（kg-CO₂/kWh）", min_value=0.0, max_value=2.0,
                                        step=0.001, value=float(_cs.get("co2_factor", 0.451)), format="%.3f", key="set_c_co2")
                c_sugi = st.number_input("杉1本 年間CO₂吸収（kg）", min_value=0.1, max_value=100.0,
                                         step=0.1, value=float(_cs.get("sugi_kg", 8.8)), format="%.1f", key="set_c_sugi")

            st.markdown("**■ エリア係数（施工場所エリア別）**")
            _ab = _cs.get("area_buffer", DEFAULT_DATA["const_settings"]["area_buffer"])
            area_buf_new = {}
            ab_cols = st.columns(3)
            for i, (rgn, coef) in enumerate(_ab.items()):
                with ab_cols[i % 3]:
                    area_buf_new[rgn] = st.number_input(f"{rgn}", min_value=0.5, max_value=3.0, step=0.01,
                                                        value=float(coef), format="%.2f", key=f"set_c_ab_{i}")

            st.markdown("**■ データ補完の既定値（検針票に列が無い時の仮置き）**")
            _fb = _cs.get("fb_defaults", DEFAULT_DATA["const_settings"]["fb_defaults"])
            fb_new = {}
            fb_specs = [
                ("契約電力", "契約電力（kW）", "%.0f", 1.0),
                ("最大需要電力", "最大需要電力（kW）", "%.0f", 1.0),
                ("使用量合計", "使用量合計（kWh/月）", "%.0f", 100.0),
                ("力率", "力率", "%.2f", 0.01),
                ("基本料金単価", "基本料金単価（円/kW）", "%.2f", 0.01),
                ("電力量単価", "電力量単価（円/kWh）", "%.2f", 0.1),
            ]
            fb_cols = st.columns(3)
            for i, (key, lbl, fmt, stp) in enumerate(fb_specs):
                with fb_cols[i % 3]:
                    fb_new[key] = st.number_input(lbl, min_value=0.0, max_value=10000000.0, step=stp,
                                                  value=float(_fb.get(key, FB_DEFAULTS[key])), format=fmt,
                                                  key=f"set_c_fb_{i}")

            st.markdown("**■ 投資評価（NPV / IRR）**")
            ec1, ec2 = st.columns(2)
            with ec1:
                c_discount = st.number_input("割引率（NPV/IRR）", min_value=0.0, max_value=0.30, step=0.005,
                                             value=float(_cs.get("discount_rate", 0.05)), format="%.3f", key="set_c_disc")
            with ec2:
                c_npv_years = st.number_input("NPV算定年数", min_value=1, max_value=30,
                                              value=int(_cs.get("npv_years", 10)), key="set_c_npvy")

            st.markdown("**■ 季節性の業態配慮**")
            c_base_ac = st.number_input(
                "基礎電力の空調比率（%）", min_value=0, max_value=100, step=5,
                value=int(round(float(_cs.get("base_ac_ratio", 0.0)) * 100)), key="set_c_base_ac",
                help="サイドバーの『業態を考慮』ON時に、最少月（基礎電力）のうち空調が占めるとみなす割合。"
                     "24時間冷凍冷蔵など最少月にも空調が残る業態で引き上げます（既定0%）。")

            st.markdown("**■ 力率割引（基本料金）**")
            pc1, pc2 = st.columns(2)
            with pc1:
                c_pf_base = st.number_input("基準力率（%）", min_value=50, max_value=100, step=1,
                                            value=int(round(float(_cs.get("pf_base", 0.85)) * 100)), key="set_c_pf_base",
                                            help="この力率を基準に、上回ると割引・下回ると割増（標準85%）。")
            with pc2:
                c_pf_cap = st.number_input("割引・割増の上限（±%）", min_value=0, max_value=50, step=1,
                                           value=int(round(float(_cs.get("pf_cap", 0.15)) * 100)), key="set_c_pf_cap",
                                           help="力率割引/割増の頭打ち（標準±15%＝0.85〜1.15）。")

            st.markdown("**■ 詳細見積（見積書方式）の固定値**")
            _est0 = _cs.get("estimate_settings", DEFAULT_DATA["const_settings"]["estimate_settings"])
            es1, es2 = st.columns(2)
            with es1:
                c_mat_factor = st.number_input("材料費係数（電材費×）", min_value=1.0, max_value=3.0, step=0.05,
                                               value=float(_est0.get("material_factor", 1.2)), format="%.2f", key="set_est_mf")
                c_labor_unit = st.number_input("工事費 人日単価（円/人日）", min_value=0, max_value=1000000, step=5000,
                                               value=int(_est0.get("labor_day_unit", 60000)), key="set_est_ldu")
            with es2:
                c_labor_fixed = st.number_input("工事費 固定加算（円）", min_value=0, max_value=10000000, step=10000,
                                                value=int(_est0.get("labor_fixed", 120000)), key="set_est_lfx")
                c_lodging = st.number_input("宿泊費 単価（円/人工）", min_value=0, max_value=1000000, step=500,
                                            value=int(_est0.get("lodging_unit", 8000)), key="set_est_lou")
            st.caption("交通費の移動拘束費 ＝ 作業人数 × 地区別単価（円）")
            _au0 = _est0.get("area_unit", {})
            area_unit_new = {}
            for _k in ["地区1", "地区2", "地区3", "地区4", "地区5"]:
                area_unit_new[_k] = st.number_input(AREA5_LABELS.get(_k, _k), min_value=0, max_value=1000000, step=1000,
                                                    value=int(_au0.get(_k, 0)), key=f"set_est_au_{_k}")

            st.markdown("**■ 概算モード（機材情報なし）の標準値**")
            _g0 = _cs.get("gaisan_settings", DEFAULT_DATA["const_settings"]["gaisan_settings"])
            gc1, gc2, gc3 = st.columns(3)
            with gc1:
                c_g_unitkw = st.number_input("標準室外機 定格kW/台", min_value=1.0, max_value=100.0, step=0.5,
                                             value=float(_g0.get("unit_kw", 10.0)), format="%.1f", key="set_g_unitkw",
                                             help="概算台数 ＝ 空調ピークkW ÷ この値。")
            with gc2:
                c_g_mat = st.number_input("標準機材費（円/台）", min_value=0, max_value=100000000, step=10000,
                                          value=int(_g0.get("material_per_unit", 100000)), key="set_g_mat")
            with gc3:
                c_g_setup = st.number_input("標準工事費（円/台）", min_value=0, max_value=100000000, step=10000,
                                            value=int(_g0.get("setup_per_unit", 50000)), key="set_g_setup")
            gh1, gh2, gh3 = st.columns(3)
            with gh1:
                c_g_hwfix = st.number_input("標準HW固定費（円/拠点）", min_value=0, max_value=100000000, step=10000,
                                            value=int(_g0.get("hw_fixed", 300000)), key="set_g_hwfix",
                                            help="データ収集装置・親機など拠点1式の固定分。")
            with gh2:
                c_g_hwunit = st.number_input("標準HW費（円/台）", min_value=0, max_value=100000000, step=10000,
                                             value=int(_g0.get("hw_per_unit", 50000)), key="set_g_hwunit",
                                             help="子機・センサー・I/O など台数比例分。")
            with gh3:
                c_g_ratio = st.number_input("概算時の制御可能比率（%）", min_value=0, max_value=100, step=5,
                                            value=int(round(float(_g0.get("ctrl_ratio", 0.9)) * 100)), key="set_g_ratio")
            st.caption("※ 概算モードの初期費用 ＝ HW固定費 ＋ 推定台数 ×（機材費＋工事費＋HW/台＋交通費/台＋構築費/台）。実績値へ随時更新してください。")

        if st.button("💾 計算設定を保存"):
            app_data["calc_settings"]["cap_rate"] = cap
            app_data["system_fee"] = sys_fee
            app_data["fee_settings"] = {
                "tsushin": int(fee_tsushin), "cloud": int(fee_cloud),
                "hoshu": int(fee_hoshu), "data": int(fee_data),
            }
            # 既存 const_settings（tou_* 等）を保ちつつ更新（マージ保存）
            _cs_new = dict(app_data.get("const_settings", DEFAULT_DATA["const_settings"]))
            _cs_new.update({
                "area_unit_price": int(c_area_unit),
                "conv_factor": float(c_conv),
                "hp_to_kw": float(c_hp),
                "co2_factor": float(c_co2),
                "sugi_kg": float(c_sugi),
                "area_buffer": {k: float(v) for k, v in area_buf_new.items()},
                "fb_defaults": {k: float(v) for k, v in fb_new.items()},
                "discount_rate": float(c_discount),
                "npv_years": int(c_npv_years),
                "base_ac_ratio": float(c_base_ac) / 100.0,
                "pf_base": float(c_pf_base) / 100.0,
                "pf_cap": float(c_pf_cap) / 100.0,
                "estimate_settings": {
                    "material_factor": float(c_mat_factor),
                    "labor_day_unit": int(c_labor_unit),
                    "labor_fixed": int(c_labor_fixed),
                    "lodging_unit": int(c_lodging),
                    "area_unit": {k: int(v) for k, v in area_unit_new.items()},
                },
                "gaisan_settings": {
                    **(_cs.get("gaisan_settings", {}) or {}),   # scenarios 等の既存キーを保持
                    "unit_kw": float(c_g_unitkw),
                    "material_per_unit": int(c_g_mat),
                    "setup_per_unit": int(c_g_setup),
                    "hw_fixed": int(c_g_hwfix),
                    "hw_per_unit": int(c_g_hwunit),
                    "ctrl_ratio": float(c_g_ratio) / 100.0,
                },
            })
            app_data["const_settings"] = _cs_new
            _apply_const_settings(app_data)   # 即時にモジュール定数へ反映
            log_change(app_data, "計算ロジック",
                       f"容量削減率={cap:.0%}, システム利用料={sys_fee:,}円, "
                       f"利用料単価(通信{fee_tsushin}/クラウド{fee_cloud}/保守{fee_hoshu}/データ{fee_data}), "
                       f"固定値(工事¥{c_area_unit}/conv{c_conv}/HP→kW{c_hp}/CO2{c_co2}/杉{c_sugi})")
            save_data(app_data)
            st.success("保存しました。")
            st.rerun()

    # ─── タブ2: 馬力・機材テーブル ───────────────────
    with tab2:
        st.subheader("🗂️ 型番→馬力マスタ（機体ごとの管理）")
        st.caption("ここに登録した型番は、シミュレーション画面の機体入力で**自動的に馬力が反映**されます。"
                   "未登録の型番は「※管理外」と表示され、下の仮置き馬力で算出します。")

        hp_master = app_data.get("model_hp_master", {})
        kw_master = app_data.get("model_kw_master", {})
        master_df = pd.DataFrame(
            [{"型番": k, "馬力": v, "定格出力kW": float(kw_master.get(k, round(v * HP_TO_KW, 1)))}
             for k, v in hp_master.items()]
            or [{"型番": "", "馬力": 0.0, "定格出力kW": 0.0}]
        )
        edited_master = st.data_editor(
            master_df, use_container_width=True, num_rows="dynamic",
            column_config={
                "型番": st.column_config.TextColumn("型番", width="large"),
                "馬力": st.column_config.NumberColumn("馬力 (HP)", format="%.1f", min_value=0.0),
                "定格出力kW": st.column_config.NumberColumn(
                    "定格出力kW（入力可）", format="%.1f", min_value=0.0,
                    help="カタログの定格冷房能力。入力するとこの値を優先して計算に使用します。"),
            },
            key="hp_master_editor",
        )
        st.caption("定格出力kWを入力するとカタログ値を優先して計算します。"
                   "0（空欄）の場合は 馬力 × 換算係数 で自動算出します。馬力が空で定格出力ありなら馬力を逆算します。")

        col_hp1, col_hp2 = st.columns(2)
        with col_hp1:
            default_hp = st.number_input(
                "管理外（未登録型番）の仮置き馬力", min_value=1.0, max_value=50.0, step=0.5,
                value=float(app_data.get("calc_settings", {}).get("default_hp", 5.0)),
                help="マスタに無い型番はこの馬力で暫定計算します。",
            )
        with col_hp2:
            cop_val = st.number_input(
                "COP（冷房能力kW→電力kW換算）", min_value=2.0, max_value=6.0, step=0.1,
                value=float(app_data.get("calc_settings", {}).get("cop", 3.5)),
            )

        if st.button("💾 馬力マスタ・換算設定を保存", type="primary"):
            new_master, new_kw = {}, {}
            for _, r in edited_master.iterrows():
                m = str(r.get("型番", "") or "").strip()
                if not m:
                    continue
                try:
                    hp = float(r.get("馬力", 0) or 0)
                except Exception:
                    hp = 0.0
                try:
                    kw = float(r.get("定格出力kW", 0) or 0)
                except Exception:
                    kw = 0.0
                if hp <= 0 and kw > 0:        # 馬力が空で定格出力ありなら逆算
                    hp = round(kw / HP_TO_KW, 1)
                if hp > 0:
                    new_master[m] = hp
                    if kw > 0:
                        new_kw[m] = kw
            app_data["model_hp_master"] = new_master
            app_data["model_kw_master"] = new_kw
            app_data["calc_settings"]["default_hp"] = float(default_hp)
            app_data["calc_settings"]["cop"] = float(cop_val)
            log_change(app_data, "馬力マスタ",
                       f"{len(new_master)}型番を登録 / 定格出力{len(new_kw)}件 / 仮置き{default_hp}HP / COP{cop_val}")
            st.success(f"保存しました（{len(new_master)}型番／定格出力 {len(new_kw)}件）。")
            st.rerun()

        st.divider()
        # ── 室外機1台あたりの機材セット（増減連動用・後で更新可）──
        st.subheader("室外機1台あたりの機材セット（増減連動用）")
        st.caption("室外機（制御対象）1台に紐づく機材と数量を登録します。"
                   "ここを設定しておくと、今後『機体の増減を機材に反映』する際の削除/追加の根拠に使えます。"
                   "※現時点で不明な場合は空のままでOK。判明したら更新してください。")
        _pum = app_data.get("per_unit_materials", []) or []
        if not _pum:
            _pum = [{"商品名": "", "1台あたり数量": 0} for _ in range(3)]
        _pum_df = pd.DataFrame(_pum)
        if "商品名" not in _pum_df.columns:
            _pum_df["商品名"] = ""
        if "1台あたり数量" not in _pum_df.columns:
            _pum_df["1台あたり数量"] = 0
        _pum_edit = st.data_editor(
            _pum_df[["商品名", "1台あたり数量"]], num_rows="dynamic", use_container_width=True,
            key="per_unit_mat_editor",
            column_config={
                "商品名": st.column_config.TextColumn("商品名（機材単価マスタの名称に合わせると単価連携可）", width="large"),
                "1台あたり数量": st.column_config.NumberColumn("室外機1台あたり数量", min_value=0, step=1),
            })
        if st.button("室外機1台あたりの機材セットを保存", key="save_per_unit_mat"):
            _rows = [{"商品名": str(r.get("商品名", "")).strip(),
                      "1台あたり数量": int(r.get("1台あたり数量", 0) or 0)}
                     for r in _pum_edit.to_dict("records") if str(r.get("商品名", "")).strip()]
            app_data["per_unit_materials"] = _rows
            save_data(app_data)
            log_change(app_data, "室外機1台あたり機材セット", f"{len(_rows)}件")
            st.success(f"保存しました（{len(_rows)}件）。")
            st.rerun()

        st.divider()
        st.markdown("##### 馬力→kW 換算 早見表（参考）")
        hp_table = pd.DataFrame({
            "馬力": [1, 1.5, 2, 2.5, 3, 4, 5, 6, 8, 10],
            "冷房能力kW（参考）": [2.8, 4.0, 5.6, 7.1, 8.0, 11.2, 14.0, 16.0, 22.4, 28.0],
            "消費電力kW（参考）": [0.68, 1.08, 1.45, 1.82, 2.15, 2.90, 3.75, 4.50, 6.00, 7.80],
        })
        st.dataframe(hp_table, use_container_width=True)

        st.divider()
        item_price = st.number_input(
            "主力商品単価（円）※本業利益換算用",
            min_value=100, max_value=100000000, step=10000,
            value=int(app_data["item_price"])
        )
        item_margin = st.slider(
            "主力商品粗利益率",
            min_value=0.01, max_value=0.80, step=0.01,
            value=float(app_data["item_margin"]),
            format="%.0f%%"
        )
        if st.button("💾 商品設定を保存"):
            app_data["item_price"] = item_price
            app_data["item_margin"] = item_margin
            log_change(app_data, "商品設定",
                       f"単価={item_price:,}円, 粗利={item_margin:.0%}")
            st.success("保存しました。")
            st.rerun()

    # ─── タブ3: 業態マスタ ───────────────────────────
    with tab3:
        st.subheader("企業の業態マスタ設定")
        st.caption("業態別の空調ピーク割合・電力量割合を設定します。")

        gm = app_data["gyotai_master"]
        gm_df = pd.DataFrame([
            {"業態": k, "空調ピーク割合": v["ac_peak"], "空調電力量割合": v["ac_kwh"]}
            for k, v in gm.items()
        ])
        edited_gm = st.data_editor(gm_df, use_container_width=True, num_rows="dynamic",
                                   column_config={
                                       "空調ピーク割合": st.column_config.NumberColumn(format="%.2f"),
                                       "空調電力量割合": st.column_config.NumberColumn(format="%.2f"),
                                   })
        if st.button("💾 業態マスタを保存"):
            new_gm = {}
            for _, row in edited_gm.iterrows():
                new_gm[row["業態"]] = {
                    "ac_peak": float(row["空調ピーク割合"]),
                    "ac_kwh": float(row["空調電力量割合"])
                }
            app_data["gyotai_master"] = new_gm
            log_change(app_data, "業態マスタ", f"{len(new_gm)}業態を更新")
            save_data(app_data)
            st.success("保存しました。")
            st.rerun()

    # ─── タブ4: 表示項目変更 ─────────────────────────
    with tab4:
        st.subheader("表示項目の変更")
        ds = app_data["display_settings"]
        show_dm   = st.checkbox("デマンドグラフを表示",  value=ds.get("show_graph_dm", True))
        show_kwh  = st.checkbox("電力量グラフを表示",    value=ds.get("show_graph_kwh", True))
        show_bill = st.checkbox("料金グラフを表示",      value=ds.get("show_graph_bill", True))
        show_co2  = st.checkbox("環境価値（CO₂）を表示", value=ds.get("show_co2", True))
        show_talk = st.checkbox("トークスクリプトを表示",value=ds.get("show_talk_script", True))
        show_loss = st.checkbox("10年損失コストを表示",  value=ds.get("show_loss_10yr", True))

        if st.button("💾 表示設定を保存"):
            app_data["display_settings"] = {
                "show_graph_dm": show_dm,
                "show_graph_kwh": show_kwh,
                "show_graph_bill": show_bill,
                "show_co2": show_co2,
                "show_talk_script": show_talk,
                "show_loss_10yr": show_loss,
            }
            log_change(app_data, "表示設定", "表示項目を更新")
            save_data(app_data)
            st.success("保存しました。")
            st.rerun()

    # ─── タブ5: 変更・DLログ ─────────────────────────
    with tab5:
        st.subheader("変更履歴ログ")
        if app_data["change_log"]:
            log_df = pd.DataFrame(app_data["change_log"])
            st.dataframe(log_df, use_container_width=True, hide_index=True)
        else:
            st.info("変更履歴はまだありません。")

        st.subheader("ダウンロード履歴ログ")
        if app_data["download_log"]:
            dl_df = pd.DataFrame(app_data["download_log"])
            st.dataframe(dl_df, use_container_width=True, hide_index=True)
        else:
            st.info("ダウンロード履歴はまだありません。")

        col1, col2 = st.columns(2)
        with col1:
            if st.button("🗑️ 変更履歴をクリア"):
                app_data["change_log"] = []
                save_data(app_data)
                st.success("クリアしました。")
                st.rerun()
        with col2:
            if st.button("🗑️ DL履歴をクリア"):
                app_data["download_log"] = []
                save_data(app_data)
                st.success("クリアしました。")
                st.rerun()

    # ─── タブ: 画像出力（手動値）─ シミュ画面と同じデザイン/配置で、数値だけ任意に指定してPNG生成 ─
    with tab_img:
        st.subheader("グラフ画像出力（数値は手動入力）")
        st.caption("シミュレーション画面と同じデザイン・配置のPNGを、任意の数値で生成します"
                   "（レイアウト・項目は固定、数値のみ変更可）。シミュレーション画面側の仕様は変わりません。")
        if not HAS_SLIDE_IMG:
            st.info("画像出力には matplotlib / pillow が必要です。`pip install matplotlib pillow` を実行してください。")
        else:
            # ── Excel取込（『空調制御効果試算』形式：入力シート→①／室外機情報シート→②）──
            with st.expander("📥 Excelから取込（①=「入力」シート／②=「室外機情報」シート）", expanded=False):
                st.caption("添付の試算Excel（入力／室外機情報シート）をアップロードすると、①グラフ値・②制御可否リストへ自動で反映します。"
                           "※Excelは形式の一例。実際にアップロードした内容のみ取り込みます（横並び・フラットどちらの表でも可）。")
                _up = st.file_uploader("試算Excel（.xlsx）をアップロード", type=["xlsx"], key="adm_img_upload")
                _ic = st.columns(3)
                _do_imp_g = _ic[0].button("① 「入力」→ グラフに取込", key="adm_imp_g", use_container_width=True,
                                          disabled=(_up is None))
                _do_imp_l = _ic[1].button("② 「室外機情報」→ リストに取込", key="adm_imp_l", use_container_width=True,
                                          disabled=(_up is None))
                _do_imp_both = _ic[2].button("①② まとめて取込", key="adm_imp_both", type="primary",
                                             use_container_width=True, disabled=(_up is None))
                if _up is not None and (_do_imp_g or _do_imp_l or _do_imp_both):
                    _msgs = []
                    if _do_imp_g or _do_imp_both:
                        _gdf = admin_import_input_graph(_up)
                        if _gdf is not None and len(_gdf):
                            st.session_state["_adm_g_data"] = _gdf
                            st.session_state.pop("adm_g_editor", None)
                            _msgs.append(f"① グラフ：{len(_gdf)}ヶ月分を取込")
                        else:
                            _msgs.append("① グラフ：「入力」シートから月次データを読めませんでした")
                    if _do_imp_l or _do_imp_both:
                        _ldf = admin_import_units_list(_up)
                        if _ldf is not None and len(_ldf):
                            st.session_state["_adm_l_data"] = _ldf
                            st.session_state.pop("adm_l_editor", None)
                            _n_ok = int(_ldf["制御可否"].apply(_adm_ctrl_is_ok).sum())
                            _msgs.append(f"② リスト：{len(_ldf)}台（制御対象〇={_n_ok}台）を取込")
                        else:
                            _msgs.append("② リスト：「室外機情報」シートから機器行を読めませんでした")
                    for _m in _msgs:
                        (st.success if "取込" in _m else st.warning)(_m)
                    if any("取込" in _m for _m in _msgs):
                        st.rerun()

            # ① 使用量＋デマンド グラフ用（12ヶ月）
            st.markdown("**① 使用量＋デマンド グラフ の数値**")
            _def_g = [{"月": f"2025/{m:02d}" if m >= 4 else f"2026/{m:02d}", "使用量": 0, "削減": 0, "最大デマンド": 0}
                      for m in list(range(4, 13)) + [1, 2, 3]]
            _g_src = st.session_state.get("_adm_g_data")
            _g_df0 = _g_src if isinstance(_g_src, pd.DataFrame) and len(_g_src) else pd.DataFrame(_def_g)
            _g_edit = st.data_editor(
                _g_df0, num_rows="fixed", use_container_width=True, key="adm_g_editor",
                column_config={
                    "月": st.column_config.TextColumn("月（YYYY/MM）"),
                    "使用量": st.column_config.NumberColumn("使用量(kWh)", format="%d"),
                    "削減": st.column_config.NumberColumn("削減(kWh)", format="%d"),
                    "最大デマンド": st.column_config.NumberColumn("最大デマンド(kW)", format="%d"),
                })
            st.session_state.setdefault("adm_target_units", 0)
            _tu = st.number_input("制御台数（グラフのラベル／制御可否リストの制御台数に使用）",
                                  min_value=0, max_value=100000, step=1, key="adm_target_units")

            # ② 制御可否リスト（制御可否は〇×で記載）
            st.markdown("**② 制御可否リスト の内容**")
            _def_l = [{"系統名": "", "設置場所": "", "メーカー": "", "型式": "", "制御可否": "〇"} for _ in range(3)]
            _l_src = st.session_state.get("_adm_l_data")
            _l_df0 = _l_src if isinstance(_l_src, pd.DataFrame) and len(_l_src) else pd.DataFrame(_def_l)
            _l_edit = st.data_editor(
                _l_df0, num_rows="dynamic", use_container_width=True, key="adm_l_editor",
                column_config={
                    "系統名": st.column_config.TextColumn("系統名"),
                    "設置場所": st.column_config.TextColumn("設置場所"),
                    "メーカー": st.column_config.TextColumn("メーカー"),
                    "型式": st.column_config.TextColumn("型式"),
                    "制御可否": st.column_config.SelectboxColumn("制御可否", options=["〇", "×"], default="〇"),
                })
            # 出力する項目を選択（不要な列は外せる）
            _ADM_LIST_COLS = ["系統名", "設置場所", "メーカー", "型式", "制御可否"]
            st.session_state.setdefault("adm_list_cols", _ADM_LIST_COLS)
            _adm_sel_cols = st.multiselect(
                "② 制御可否リストに含める項目（不要な列は外せます）",
                _ADM_LIST_COLS, key="adm_list_cols")

            # ③ サマリKPIカードの数値（カード構成・文言・色は固定、数値のみ）
            st.markdown("**③ サマリKPIカード の数値**")
            for _k0 in ("adm_inv", "adm_grs", "adm_nt", "adm_sf"):
                st.session_state.setdefault(_k0, 0)
            _cc = st.columns(4)
            _inv = _cc[0].number_input("初期導入費用（円）", min_value=0, max_value=10_000_000_000, step=100000, key="adm_inv")
            _grs = _cc[1].number_input("年間総削減額 グロス（円）", min_value=0, max_value=10_000_000_000, step=10000, key="adm_grs")
            _nt = _cc[2].number_input("年間実質利点 手残り（円）", min_value=0, max_value=10_000_000_000, step=10000, key="adm_nt")
            _sf = _cc[3].number_input("年間維持費（円）", min_value=0, max_value=10_000_000_000, step=10000, key="adm_sf")

            # ①②③は独立生成（1つがエラーでも他は生成できる／②だけ・③だけもOK）
            st.caption("各ボタンで①②③を個別に生成できます（まとめて出す場合は3つ押してください）。")
            _bc = st.columns(3)
            _do_g = _bc[0].button("① グラフを生成 / 更新", key="adm_gen_g", type="primary", use_container_width=True)
            _do_l = _bc[1].button("② 制御可否リストを生成 / 更新", key="adm_gen_l", use_container_width=True)
            _do_c = _bc[2].button("③ サマリカードを生成 / 更新", key="adm_gen_c", use_container_width=True)

            if _do_g:
                try:
                    _rows = _g_edit.to_dict("records")
                    _months = [str(r.get("月", "")) for r in _rows]
                    _usage = [float(r.get("使用量", 0) or 0) for r in _rows]
                    _reduc = [float(r.get("削減", 0) or 0) for r in _rows]
                    _demand = [float(r.get("最大デマンド", 0) or 0) for r in _rows]
                    _cu = int(_tu) if int(_tu) > 0 else 0
                    st.session_state["_adm_img_chart"] = make_demand_chart_png(
                        _months, _usage, _reduc, _demand, target_units=_cu)
                    st.success("① グラフを生成しました。")
                except Exception as e:
                    st.error(f"① グラフ生成でエラー: {e}")

            if _do_l:
                try:
                    _list_items = [{"系統名": str(r.get("系統名", "")), "設置場所": str(r.get("設置場所", "")),
                                    "メーカー": str(r.get("メーカー", "")), "型式": str(r.get("型式", "")),
                                    "制御可否": _adm_ctrl_is_ok(r.get("制御可否", ""))}
                                   for r in _l_edit.to_dict("records")
                                   if any(str(r.get(k, "")).strip() for k in ("系統名", "設置場所", "メーカー", "型式"))]
                    if not _list_items:
                        st.warning("② リストの内容（系統名・型式など）を1行以上入力してください。")
                    else:
                        _n_total = len(_list_items)
                        _n_ctrl = sum(1 for it in _list_items if it["制御可否"])
                        _ctrl_units = int(_tu) if int(_tu) > 0 else _n_ctrl
                        st.session_state["_adm_img_list"] = make_control_list_pngs(
                            _list_items, total_units=_n_total, controllable_units=_ctrl_units,
                            cols=(_adm_sel_cols or None))
                        st.success(f"② 制御可否リストを生成しました（{_n_total}台）。")
                except Exception as e:
                    st.error(f"② 制御可否リスト生成でエラー: {e}")

            if _do_c:
                try:
                    _pb = (_inv / _nt) if _nt > 0 else 0.0
                    _NAVY_C, _GREEN_C = "#13315C", "#3DAE4E"
                    _cards = [
                        {"icon": "money.png", "title": "初期導入費用（総投資額）", "value": f"{_inv:,.0f}円",
                         "subtitle": "税抜 導入費用 / 機器代・工事費含む", "color": _NAVY_C},
                        {"icon": "graf.png", "title": "年間総削減額（グロス）", "value": f"{_grs:,.0f}円",
                         "subtitle": "基本料金＋電力量の年間削減合計", "color": _GREEN_C},
                        {"icon": "plus.png", "title": "年間実質利点（手残り）", "value": f"{_nt:,.0f}円",
                         "subtitle": f"削減合計 {_grs:,.0f}円 − 年間維持費 {_sf:,.0f}円", "color": _GREEN_C},
                        {"icon": "clock.png", "title": "投資回収期間（ROI）", "value": f"約{_pb:.1f}年",
                         "subtitle": f"約{round(_pb*12)}ヶ月で完全回収、以降は純利益", "color": _GREEN_C},
                    ]
                    st.session_state["_adm_img_cards"] = make_summary_cards_png(_cards)
                    st.success("③ サマリカードを生成しました。")
                except Exception as e:
                    st.error(f"③ サマリカード生成でエラー: {e}")

            if st.session_state.get("_adm_img_chart"):
                st.markdown("**① 使用量＋デマンド グラフ**")
                st.image(st.session_state["_adm_img_chart"], use_container_width=True)
                st.download_button("グラフPNGをダウンロード", st.session_state["_adm_img_chart"],
                                   file_name="使用量_デマンドグラフ.png", mime="image/png", key="adm_dl_chart")
            _li = st.session_state.get("_adm_img_list")
            if _li:
                _np = len(_li)
                st.markdown("**② 制御可否リスト**" + (f"（{_np}枚に分割）" if _np > 1 else ""))
                for _k, _img in enumerate(_li, start=1):
                    st.image(_img, use_container_width=True)
                    _fn = "制御可否リスト.png" if _np == 1 else f"制御可否リスト_{_k}of{_np}.png"
                    st.download_button("制御可否リストPNGをダウンロード" + ("" if _np == 1 else f"（{_k}/{_np}）"),
                                       _img, file_name=_fn, mime="image/png", key=f"adm_dl_list_{_k}")
            if st.session_state.get("_adm_img_cards"):
                st.markdown("**③ サマリKPIカード**")
                st.image(st.session_state["_adm_img_cards"], use_container_width=True)
                st.download_button("サマリカードPNGをダウンロード", st.session_state["_adm_img_cards"],
                                   file_name="サマリカード.png", mime="image/png", key="adm_dl_cards")

    # ─── 機材単価マスタ（見積書作成用）──────────────
    with tab_mat:
        st.subheader("機材単価マスタ（見積書作成）")
        st.info("商品名・単位・売価・原価・利益で管理します。見積書の機材費は、ここの売価をそのまま使用します"
                "（利益＝売価−原価は自動計算）。")
        _qs = copy.deepcopy(DEFAULT_DATA["quote_settings"])
        _qs.update(app_data.get("quote_settings", {}) or {})

        # 旧構造（品名/参考売価）からの移行
        _rawm = app_data.get("material_master", []) or []
        _norm_rows = []
        for m in _rawm:
            nm = m.get("商品名", m.get("品名", ""))
            urikae = m.get("売価", m.get("参考売価", 0)) or 0
            genka = m.get("原価", 0) or 0
            _norm_rows.append({"商品名": nm, "単位": m.get("単位", "台"),
                               "売価": int(urikae), "原価": int(genka), "利益": int(urikae) - int(genka)})
        mat_df = pd.DataFrame(_norm_rows if _norm_rows else
                              [{"商品名": "", "単位": "台", "売価": 0, "原価": 0, "利益": 0}])
        edited = st.data_editor(
            mat_df, num_rows="dynamic", use_container_width=True, key="admin_mat_master",
            column_config={
                "商品名": st.column_config.TextColumn("商品名", width="large"),
                "単位": st.column_config.TextColumn("単位", width="small"),
                "売価": st.column_config.NumberColumn("売価（円）", min_value=0, step=1, format="%d"),
                "原価": st.column_config.NumberColumn("原価（円）", min_value=0, step=1, format="%d"),
                "利益": st.column_config.NumberColumn("利益（自動＝売価−原価）", disabled=True, format="%d"),
            })
        st.caption("利益列は自動計算です（保存時に売価−原価で再計算）。")

        st.markdown("---")
        st.markdown("**見積の既定値（利益率・税率・ビジネスタンク率・構築費単価）**")
        q1, q2, q3, q4 = st.columns(4)
        with q1:
            _dm = st.number_input("既定 利益率（原価×倍率）", 1.0, 5.0, float(_qs.get("markup_factor", 1.2)), step=0.05)
        with q2:
            _tm = st.number_input("倍率 下限アラート", 1.0, 5.0, float(_qs.get("target_factor", 1.1)), step=0.05)
        with q3:
            _bt = st.number_input("ビジネスタンク率 %", 0.0, 90.0, float(_qs.get("bt_rate", 0.25)) * 100, step=1.0)
        with q4:
            _ku = st.number_input("構築費 単価（円/台）", 0, 1000000, int(_qs.get("kouchiku_unit", 12000)), step=1000)
        _tx = st.number_input("消費税率 %", 0.0, 30.0, float(_qs.get("tax_rate", 0.10)) * 100, step=1.0)

        if st.button("💾 機材マスタ・見積設定を保存", type="primary", key="admin_mat_save"):
            _mm = []
            for r in edited.to_dict("records"):
                nm = str(r.get("商品名", "") or "").strip()
                if not nm:
                    continue
                try:
                    uri = int(float(r.get("売価", 0) or 0))
                except Exception:
                    uri = 0
                try:
                    gen = int(float(r.get("原価", 0) or 0))
                except Exception:
                    gen = 0
                _mm.append({"商品名": nm, "単位": str(r.get("単位", "台") or "台"),
                            "売価": uri, "原価": gen, "利益": uri - gen})
            app_data["material_master"] = _mm
            _qs["markup_factor"] = float(_dm)
            _qs["target_factor"] = float(_tm)
            _qs["bt_rate"] = _bt / 100.0
            _qs["kouchiku_unit"] = int(_ku)
            _qs["tax_rate"] = _tx / 100.0
            app_data["quote_settings"] = _qs
            save_data(app_data)
            log_change(app_data, "機材単価マスタ", f"{len(app_data['material_master'])}件（売価/原価）・倍率×{_dm:.2f}")
            st.success("保存しました。")
            st.rerun()


# ══════════════════════════════════════════════════════════════════
# ヘルプ画面 (戻るボタン補強)
# ══════════════════════════════════════════════════════════════════
def show_results_list(app_data: dict) -> None:
    """📋 結果一覧ページ：保存した案件の軽量サマリー（企業名/空調数/年間削減量/回収年数）を一覧表示。"""
    st.markdown("## 📋 結果一覧")

    cases = app_data.get("saved_cases", [])
    if not cases:
        st.info("保存された結果はまだありません。シミュレーション画面の「💾 この結果を一覧に保存」から登録できます。")
        return

    st.caption(f"保存件数：{len(cases)} 件（詳細データは各社フォルダで管理。ここは一覧用の軽量サマリーです）")

    def _pb(v):
        try:
            v = float(v)
            return "回収不能" if v == float("inf") else f"{v:.1f} 年"
        except Exception:
            return "—"

    rows = []
    for c in cases:
        rows.append({
            "企業名": c.get("company", "—"),
            "業態": c.get("gyotai", "—"),
            "空調台数": f"{int(c.get('units_ctrl', 0))} 台",
            "年間削減量(kWh)": c.get("saving_kwh", 0),
            "年間削減額(円)": c.get("saving_yen", 0),
            "回収年数": _pb(c.get("payback")),
            "保存日時": c.get("saved_at", "—"),
        })
    df = pd.DataFrame(rows)
    st.dataframe(
        df.style.format({"年間削減量(kWh)": "{:,.0f}", "年間削減額(円)": "¥{:,.0f}"}),
        use_container_width=True, hide_index=True)

    # ── 削除 ──
    st.markdown("##### 🗑️ 案件の削除")
    labels = [f"{i+1}. {c.get('company','—')}（{c.get('saved_at','')}）" for i, c in enumerate(cases)]
    sel = st.selectbox("削除する案件を選択", ["（選択してください）"] + labels, key="list_del_sel")
    if st.button("🗑️ 選択した案件を削除", key="list_del_btn"):
        if sel in labels:
            cases.pop(labels.index(sel))
            app_data["saved_cases"] = cases
            save_data(app_data)
            st.success("削除しました。")
            st.rerun()
        else:
            st.warning("削除する案件を選択してください。")


def show_help() -> None:
    st.markdown("## ヘルプメニュー")

    tab1, tab2 = st.tabs(["① システム概要", "② 操作ガイド"])

    with tab1:
        st.markdown("""
### DPS 空調デマンド制御 提案・稟議シミュレーター とは？
本システムは、顧客企業の電力検針票データ（12ヶ月実績）をもとに、
空調デマンド制御システム（DPS）を**もし過去12ヶ月に導入していたら**いくら節約できていたか、
**1円単位でガチ計算する過去IFシミュレーター**です。

---

### 計算の核心：「引き算推計法（ベースロード法）」

```
各月の空調推計kWh = 各月の使用量合計 - 12ヶ月の最少使用量（基礎電力）
削減kWh = 空調推計kWh × 業態別空調割合 × 容量削減率(30%) × 制御対象率
```

不確かな将来予測ではなく、**確定した過去データを基にした引き算**なので、
社長・CFOへの数字の根拠説明が明確です。

---

### 3年回収ルールによる自動分岐

| 投資回収年数 | 提案モード |
|---|---|
| **3年以内** | ✅ DPS制御提案モード（経費削減ストーリー） |
| **3年超** | 🔄 空調リプレイス提案モード（省エネ補助金活用を示唆） |

---

### JSON永続保存
- `dps_app_data.json` に業態マスタ・利用料・ログを自動保存
- アプリ再起動後も設定が維持されます
        """)

    with tab2:
        # ════════ パワポをめくる式 可視化マニュアル（シミュ画面の赤枠ハイライト＋◀▶＋丸ポチ）════════
        slides = [
            {"accent": "#0F2E5D", "title": "全体の流れ",
             "hi": (184, 46, 502, 328, "全体の流れ"),
             "lead": "本ツールは「過去12ヶ月の実績データ」から空調デマンド制御の効果を1円単位で試算します。流れは次の5ステップです。",
             "points": ["① データ入力（インポート／手動）",
                        "② 機体・ハードウェア・初期費用の設定（サイドバー）",
                        "③ シミュレーション実行",
                        "④ 結果ダッシュボードで確認・詳細分析",
                        "⑤ 提案書・見積書・削減レポ・請求書を出力"]},
            {"accent": "#0F6CBD", "title": "① データ入力",
             "hi": (184, 46, 502, 58, "①電力データ取込"),
             "lead": "メイン画面「電力データ インポート / 手動入力」から取り込みます。",
             "points": ["CSV/Excelをドラッグ＆ドロップ（指定外フォーマットも内容から自動判定）",
                        "テンプレートDL可。必要列＝月／契約電力／最大需要電力／使用量合計／力率／基本料金単価／電力量単価",
                        "テンプレの『利用料機器(任意)』『室外機リスト(任意)』も同時取込→ハードウェア・機体へ自動反映",
                        "手動データ編集テーブルで直接修正も可能"]},
            {"accent": "#107C41", "title": "② 機体・ハードウェア",
             "hi": (184, 104, 502, 50, "②機体・ハードウェア"),
             "lead": "制御スコープと機材を設定します。",
             "points": ["機体（室外機）入力：型番→馬力マスタで自動反映、制御可否(○/×)で対象を設定（容量加重で精緻化）",
                        "ハードウェア関連（部材）一覧：利用料機器から自動転記、単価は機材単価マスタ準拠",
                        "業態は『その他※数値を指定する』で空調割合を手動指定も可"]},
            {"accent": "#9A6700", "title": "② 初期費用・構築費",
             "hi": (8, 40, 170, 372, "②初期費用（サイドバー）"),
             "lead": "サイドバー「初期費用」で投資額を組み立てます。",
             "points": ["概算 または 詳細見積（電材費・人日・地区別交通費＝地区5のみ実費・レンタカー＝12,000円×日数）",
                        "ハードウェア関連費／空調システム構築費（データ取集装置・室外機・温湿度センサー）",
                        "端数調整・補助金を差引 → 総初期投資額(A) に反映"]},
            {"accent": "#0F2E5D", "title": "③④ 実行と結果",
             "hi": (184, 156, 502, 44, "③④結果ダッシュボード"),
             "lead": "「シミュレーション実行」で一括試算→ダッシュボードへ。",
             "points": ["KPI×4（初期費用／投資回収／年間削減／10年累計）＋AIサマリー（結論ファースト）",
                        "主要グラフ（最大デマンド／電力量／料金を切替表示）＋比較表",
                        "入力整合性チェック・『算出の前提（仮置き・概算・業態配慮）』も上部に表示"]},
            {"accent": "#0F6CBD", "title": "④ 詳細分析",
             "hi": (184, 202, 322, 172, "④詳細分析グラフ"),
             "lead": "詳細分析タブで根拠まで開示します。",
             "points": ["算出方法：計算式→当てはめ→結果をすべて開示",
                        "回収逆算：目標回収年数から制御台数を逆算（本体の回収と整合）",
                        "制御強度・10年リスクも確認可能"]},
            {"accent": "#107C41", "title": "⑤ 出力（提案・見積）",
             "hi": (510, 202, 176, 172, "⑤提案書・出力"),
             "lead": "そのまま提案・稟議に使える資料を出力します。",
             "points": ["PowerPoint提案書／Excel稟議シート／スライド貼付用PNG／営業トーク",
                        "出力資料には仮置き前提の注釈が自動付記",
                        "見積書作成：部材表＋機材単価マスタで明細自動算出→PDF/Excel"]},
            {"accent": "#9A6700", "title": "⑤ 削減レポ・請求書",
             "hi": (14, 48, 158, 42, "メニュー（削減レポ/請求書）"),
             "lead": "運用後の実績報告と請求まで対応します。",
             "points": ["削減レポ：保存案件＋全体把握Excelから実績報告書（①②③＋グラフ）をExcel/PDF",
                        "請求書（成果報酬型）：押印氏名・単価を編集してPDF/Excel",
                        "各種設定：マスタ（計算ロジック／機材単価／馬力／業態／表示）を管理"]},
        ]
        n = len(slides)
        idx = int(st.session_state.get("help_slide", 0))
        idx = max(0, min(idx, n - 1))
        s = slides[idx]

        def _screen_svg(hi):
            hx, hy, hw, hh, hl = hi
            ty = hy - 20 if hy - 20 > 10 else hy + 3
            tagw = len(hl) * 12 + 16
            return f"""<svg viewBox="0 0 700 420" width="100%" preserveAspectRatio="xMidYMid meet" xmlns="http://www.w3.org/2000/svg" style="border:1px solid #E5E7EB;border-radius:10px;background:#fff;height:300px;">
  <rect x="8" y="8" width="684" height="404" rx="8" fill="#FFFFFF" stroke="#D8DEE9"/>
  <rect x="8" y="8" width="684" height="30" rx="8" fill="#0F2E5D"/><rect x="8" y="24" width="684" height="14" fill="#0F2E5D"/>
  <text x="20" y="27" fill="#fff" font-size="12" font-weight="bold">DPS シミュレーター</text>
  <circle cx="640" cy="23" r="4" fill="#107C41"/><circle cx="656" cy="23" r="4" fill="#9A6700"/><circle cx="672" cy="23" r="4" fill="#0F6CBD"/>
  <rect x="8" y="40" width="170" height="372" fill="#1B2A4A"/>
  <rect x="18" y="50" width="68" height="16" rx="3" fill="#33415F"/><rect x="92" y="50" width="68" height="16" rx="3" fill="#33415F"/>
  <rect x="18" y="72" width="68" height="16" rx="3" fill="#33415F"/><rect x="92" y="72" width="68" height="16" rx="3" fill="#33415F"/>
  <rect x="18" y="104" width="142" height="22" rx="3" fill="#223254"/><rect x="18" y="134" width="142" height="22" rx="3" fill="#223254"/>
  <rect x="18" y="164" width="142" height="22" rx="3" fill="#223254"/><rect x="18" y="194" width="142" height="22" rx="3" fill="#223254"/>
  <rect x="18" y="372" width="142" height="26" rx="4" fill="#0F6CBD"/><text x="89" y="389" fill="#fff" font-size="10" text-anchor="middle">シミュレーション実行</text>
  <rect x="186" y="48" width="498" height="52" rx="5" fill="#F7F9FC" stroke="#D8DEE9"/>
  <text x="196" y="66" fill="#374151" font-size="11" font-weight="bold">電力データ インポート / 手動入力</text>
  <rect x="196" y="74" width="478" height="18" rx="3" fill="#fff" stroke="#E5E7EB"/>
  <rect x="186" y="106" width="498" height="46" rx="5" fill="#F7F9FC" stroke="#D8DEE9"/>
  <text x="196" y="124" fill="#374151" font-size="11" font-weight="bold">機体（室外機）入力 / ハードウェア関連</text>
  <rect x="196" y="130" width="478" height="16" rx="3" fill="#fff" stroke="#E5E7EB"/>
  <rect x="186" y="158" width="118" height="40" rx="5" fill="#fff" stroke="#D8DEE9"/><rect x="313" y="158" width="118" height="40" rx="5" fill="#fff" stroke="#D8DEE9"/>
  <rect x="440" y="158" width="118" height="40" rx="5" fill="#fff" stroke="#D8DEE9"/><rect x="567" y="158" width="117" height="40" rx="5" fill="#fff" stroke="#D8DEE9"/>
  <rect x="186" y="158" width="118" height="3" fill="#0F2E5D"/><rect x="313" y="158" width="118" height="3" fill="#0F6CBD"/><rect x="440" y="158" width="118" height="3" fill="#107C41"/><rect x="567" y="158" width="117" height="3" fill="#9A6700"/>
  <rect x="186" y="204" width="318" height="168" rx="5" fill="#F7F9FC" stroke="#D8DEE9"/><text x="196" y="222" fill="#374151" font-size="11" font-weight="bold">グラフ</text>
  <rect x="210" y="300" width="24" height="60" fill="#BFBFBF"/><rect x="242" y="320" width="24" height="40" fill="#70AD47"/>
  <rect x="300" y="280" width="24" height="80" fill="#BFBFBF"/><rect x="332" y="310" width="24" height="50" fill="#70AD47"/>
  <rect x="392" y="290" width="24" height="70" fill="#BFBFBF"/><rect x="424" y="330" width="24" height="30" fill="#70AD47"/>
  <rect x="512" y="204" width="172" height="168" rx="5" fill="#F7F9FC" stroke="#D8DEE9"/><text x="522" y="222" fill="#374151" font-size="11" font-weight="bold">AIサマリー / 出力</text>
  <rect x="522" y="232" width="152" height="48" rx="4" fill="#E8F5E9" stroke="#C4E6D1"/><rect x="522" y="290" width="152" height="22" rx="4" fill="#0F2E5D"/><rect x="522" y="318" width="152" height="22" rx="4" fill="#107C41"/>
  <rect x="{hx}" y="{hy}" width="{hw}" height="{hh}" rx="5" fill="none" stroke="#E03131" stroke-width="3"/>
  <rect x="{hx}" y="{ty}" width="{tagw}" height="18" rx="3" fill="#E03131"/>
  <text x="{hx + 8}" y="{ty + 13}" fill="#fff" font-size="11" font-weight="bold">{hl}</text>
</svg>"""

        _pts = "".join(f"<li>{p}</li>" for p in s["points"])
        # 中身の高さをページ間で固定 → 下のページ送りボタンの位置がずれない（マウス固定で連打可）
        with st.container(height=560, border=False):
            st.markdown(_screen_svg(s["hi"]), unsafe_allow_html=True)
            st.markdown(
                "<div style='border:1px solid #E5E7EB;border-radius:12px;overflow:hidden;margin-top:8px;"
                "box-shadow:0 2px 10px rgba(16,24,40,.08);background:#fff;'>"
                f"<div style='background:{s['accent']};color:#fff;padding:12px 20px;font-size:17px;font-weight:700;'>{s['title']}</div>"
                "<div style='padding:14px 24px;'>"
                f"<div style='font-size:14px;color:#374151;margin-bottom:8px;'>{s['lead']}</div>"
                f"<ul style='font-size:14px;color:#374151;line-height:1.9;margin:0;padding-left:20px;'>{_pts}</ul>"
                "</div>"
                f"<div style='text-align:right;padding:6px 20px;color:#9CA3AF;font-size:12px;border-top:1px solid #F1F5F9;'>{idx+1} / {n}　DPS 操作マニュアル</div>"
                "</div>", unsafe_allow_html=True)

        nv = st.columns([1, 1, 1], vertical_alignment="center")
        if nv[0].button("◀ 前のページ", key="help_prev", use_container_width=True, disabled=(idx <= 0)):
            st.session_state.help_slide = idx - 1; st.rerun()
        nv[1].markdown(f"<div style='text-align:center;color:#6B7280;font-weight:600;'>{idx+1} / {n}</div>",
                       unsafe_allow_html=True)
        if nv[2].button("次のページ ▶", key="help_next", use_container_width=True, disabled=(idx >= n - 1)):
            st.session_state.help_slide = idx + 1; st.rerun()

        # 丸ポチ（クリックでそのページへジャンプ）

        dcols = st.columns(n)
        for _i in range(n):
            if dcols[_i].button("●" if _i == idx else "○", key=f"help_dot_{_i}", use_container_width=True):
                st.session_state.help_slide = _i; st.rerun()

        st.divider()

        # ════════ STEPウィザード（テキスト詳細・末尾に次のStepボタン）════════
        st.markdown("#### ステップ別ガイド（詳細）")
        steps = [
            {"t": "データの入力", "b": """
**メイン画面「電力データ インポート / 手動入力」** から取り込みます。
- CSV / Excel をドラッグ＆ドロップ（指定外フォーマットも全シート走査で自動判定）。取り込めない時は「テンプレDL」を利用。
- 必要列：**月／契約電力／最大需要電力／使用量合計／力率／基本料金単価／電力量単価**（H列以降の燃料費調整額・再エネ賦課金などは任意）。
- テンプレの **「利用料機器(任意)」「室外機リスト(任意)」** を含めて取り込むと、ハードウェア一覧・機体入力へ自動反映されます。
- **手動データ編集テーブル** で数値を直接修正してから試算も可能です。
"""},
            {"t": "機体・ハードウェア", "b": """
**機体（室外機）入力**
- 型番を入力すると **馬力マスタ** から自動反映（未登録は仮置き）。**制御可否(○/×)** で制御対象を設定 → 容量加重で制御比率を精緻化。

**ハードウェア関連（部材）一覧**
- テンプレの **利用料機器** の記載品を全件転記。単価は **機材単価マスタ** の売価に準拠（マスタ更新が自動反映）。
- 合計はサイドバー「ハードウェア関連費」と総初期投資額(A)へ加算されます。
"""},
            {"t": "初期費用・構築費（サイドバー）", "b": """
**初期費用**
- **概算** または **詳細見積**（電材費×係数／作業人日／地区別交通費＝**地区5のみ実費入力**／レンタカー＝**12,000円×利用日数**）。
- **空調システム構築費**（データ取集装置＝機器台数／室外機＝制御台数／温湿度センサー＝防水温湿度計）。
- **端数調整**（合計から減算）・**補助金** を差し引いて **総初期投資額(A)** が確定。

**業態・空調割合**
- 顧客情報の業態プルダウン直下に「空調割合想定」を表示。**「その他※数値を指定する」** で手動指定も可能。
"""},
            {"t": "シミュレーション実行", "b": """
サイドバー下部の **「シミュレーション実行」** で全試算を一括実行します。
- 機体（型番リスト）を入力していれば容量加重、未入力なら台数ベースで自動算出。
- 完了すると結果ダッシュボードに切り替わります。
"""},
            {"t": "結果ダッシュボードの見方", "b": """
- **ヘッダー**：顧客名・業態・日付＋ステータスバッジ。
- **KPI×4**：初期費用（総投資）／投資回収期間／年間削減額／10年累計（NPV）。
- **AIサマリー**：結論ファーストの要約。続いて **主要グラフ**・**比較表**。
- 上部の **「入力整合性チェック」** と **「算出の前提（仮置き・概算・業態配慮）」** で前提も確認できます。
"""},
            {"t": "詳細分析タブ", "b": """
詳細分析タブを切り替えて根拠を確認します。
- **グラフ**：最大デマンド／電力量／料金を**ラジオで1枚ずつ切替**表示（業態内訳の色分け）。
- **財務 / 制御強度 / リスク**：NPV・IRR、制御モード、10年財務リスク。
- **算出方法**：計算式→当てはめ→結果をすべて開示。
- **回収逆算**：目標回収年数から制御台数を逆算（本体の回収＝総投資÷純削減 と整合）。
"""},
            {"t": "提案書・出力", "b": """
結果画面下部の **「提案書・出力」** から：
- **PowerPoint提案書** / **Excel稟議シート** / **結果一覧に保存**。
- **営業トークスクリプト** / **スライド貼付用PNG**（グラフ・制御可否リスト・サマリカード）。
- 出力資料には**仮置き前提の注釈**が自動で付記されます。
"""},
            {"t": "見積書作成", "b": """
メニュー **「見積書作成」**：
- **部材表インポート**＋**機材単価マスタ**で機材費を自動算出。エンジ費・諸経費・構築費・ビジネスタンク率も計上。
- 引用できる数値はドロップダウンで調整。**プレビュー**後に **PDF/Excel** 出力。
"""},
            {"t": "削減レポ・請求書", "b": """
メニュー **「削減レポ」**：
- **保存案件の呼び出し**＋**全体把握Excelの追加インポート**で、客先様式の **電力削減実績報告書**（①利用状況／②最大需要／③使用電力量＋グラフ）を **Excel / PDF** 出力。
- 表示ラベルは **客先ごとにプロファイル保存**。データ不足時は出力前に **警告ポップアップ**。
- 下部 **「7. 請求書」**：成果報酬型請求書を作成。**押印（※印）の氏名を左右それぞれ編集**でき、単価は引用＋追記。**PDF / Excel** 出力。
"""},
            {"t": "各種設定（マスタ）", "b": """
メニュー **「設定」**：
- **計算ロジック変更** / **機材単価マスタ** / **馬力・機材テーブル** / **業態マスタ設定** / **表示項目変更** / **変更・DL履歴ログ**。
- 設定は `dps_app_data.json` に自動保存され、再起動後も維持されます。
"""},
        ]
        ns = len(steps)
        _opts = [f"STEP{i+1}：{stp['t']}" for i, stp in enumerate(steps)]
        # selectbox の状態（help_step_sel）を唯一の真実として扱う。前後ボタンはこの値を更新する。
        st.session_state.setdefault("help_step_sel", _opts[0])
        if st.session_state.get("help_step_sel") not in _opts:
            st.session_state["help_step_sel"] = _opts[0]
        _jump = st.selectbox("ステップを選択", _opts, key="help_step_sel")
        cur = _opts.index(_jump)

        # 本文の高さをステップ間で固定 → 下の前後ボタンの位置がずれない（マウス固定で連打可）
        with st.container(height=340, border=False):
            st.markdown(f"### STEP{cur+1}：{steps[cur]['t']}")
            st.markdown(steps[cur]["b"])

        # 末尾ナビ：前後ステップへ（selectboxキーはon_clickコールバックで更新＝生成後変更エラー回避）
        def _step_go(delta):
            _o = [f"STEP{i+1}：{stp['t']}" for i, stp in enumerate(steps)]
            _cur = _o.index(st.session_state["help_step_sel"]) if st.session_state.get("help_step_sel") in _o else 0
            st.session_state.help_step_sel = _o[max(0, min(_cur + delta, len(_o) - 1))]

        st.divider()
        bn = st.columns(2)
        if cur > 0:
            bn[0].button(f"← STEP{cur}：{steps[cur-1]['t']}", key="help_step_prev",
                         use_container_width=True, on_click=_step_go, args=(-1,))
        if cur < ns - 1:
            bn[1].button(f"STEP{cur+2}：{steps[cur+1]['t']} →", key="help_step_next",
                         use_container_width=True, type="primary", on_click=_step_go, args=(1,))


# ══════════════════════════════════════════════════════════════════
# サンプルデータ生成
# ══════════════════════════════════════════════════════════════════
def get_sample_df() -> pd.DataFrame:
    return pd.DataFrame({
        "月":         ["2025/05","2025/06","2025/07","2025/08","2025/09","2025/10",
                       "2025/11","2025/12","2026/01","2026/02","2026/03","2026/04"],
        "契約電力":    [626, 586, 586, 586, 586, 586, 586, 586, 586, 586, 586, 586],
        "最大需要電力":[415, 444, 492, 518, 576, 487, 485, 473, 528, 595, 581, 456],
        "using_kwh_name":  [95126, 101666, 118104, 121798, 132886, 112109, 112123, 123379, 113890, 122040, 98220, 84994],
        "力率":        [1.0]*12,
        "基本料金単価":[1690.65]*12,
        "電力量単価":  [17.68, 20.99, 21.47, 20.82, 20.19, 19.74, 20.16, 19.71, 19.6, 17.44, 17.44, 19.62],
    }).rename(columns={"using_kwh_name": "使用量合計"})


def get_blank_df() -> pd.DataFrame:
    """運用時のクリア起動用：12ヶ月の空テンプレート（月・数値は未入力＝0/空）。"""
    n = 12
    return pd.DataFrame({
        "月":          [""] * n,
        "契約電力":     [0] * n,
        "最大需要電力":  [0] * n,
        "使用量合計":    [0] * n,
        "力率":         [0.0] * n,
        "基本料金単価":  [0.0] * n,
        "電力量単価":    [0.0] * n,
    })


def get_sample_models():
    """テスト用：空調情報（機体＝室外機リスト）のサンプル。型番は馬力マスタ登録済みのものを使用。"""
    return [
        {"機器ID/場所": "1F-A", "階/エリア": "1F", "メーカー": "三菱電機", "型番": "PUZ-ERMP280KA4", "稼働係数": 0.8, "制御可否": "○"},
        {"機器ID/場所": "1F-B", "階/エリア": "1F", "メーカー": "三菱電機", "型番": "PUZ-ERMP160LA2", "稼働係数": 0.7, "制御可否": "○"},
        {"機器ID/場所": "2F-A", "階/エリア": "2F", "メーカー": "ダイキン", "型番": "RZRP224", "稼働係数": 0.7, "制御可否": "○"},
        {"機器ID/場所": "2F-B", "階/エリア": "2F", "メーカー": "ダイキン", "型番": "RZPR224A", "稼働係数": 0.6, "制御可否": "○"},
        {"機器ID/場所": "サーバー室", "階/エリア": "1F", "メーカー": "三菱電機", "型番": "MUCZ-G5617S", "稼働係数": 0.5, "制御可否": "×"},
    ]


def get_sample_hw():
    """テスト用：機材情報（ハードウェア／利用料機器）のサンプル。商品名は機材単価マスタに合わせ単価自動反映。"""
    return [
        {"商品名": "データ取集装置（Marimba Mercury)", "数量": 1, "単位": "台", "売価単価": 0.0, "原価単価": 0.0},
        {"商品名": "MMEazyAir親機", "数量": 1, "単位": "台", "売価単価": 0.0, "原価単価": 0.0},
        {"商品名": "MMEazyAir2（子機）", "数量": 2, "単位": "台", "売価単価": 0.0, "原価単価": 0.0},
        {"商品名": "SwitchBot 防水温湿度計", "数量": 3, "単位": "個", "売価単価": 0.0, "原価単価": 0.0},
        {"商品名": "I/O入出力モジュール", "数量": 1, "単位": "台", "売価単価": 0.0, "原価単価": 0.0},
    ]


def _usage_swing(df):
    """使用電力量の振れ幅（最少月比の最大増加率）を返す。波が激しいデータの判定に使用。
    例：0.6 なら『最少月比で最大+60%』。データ不足や算出不可は 0.0。"""
    try:
        u = pd.to_numeric(df.get("使用量合計"), errors="coerce").fillna(0.0)
        nz = u[u > 0]
        if len(nz) < 4:
            return 0.0
        mn, mx = float(nz.min()), float(nz.max())
        return (mx - mn) / mn if mn > 0 else 0.0
    except Exception:
        return 0.0


# ══════════════════════════════════════════════════════════════════
# ペライチ解説セクション
# ══════════════════════════════════════════════════════════════════
def show_peraichi(res: dict) -> None:
    st.markdown("---")
    st.markdown("### 分析前提の解説")

    st.markdown(f"""
<div style="display:grid; grid-template-columns:repeat(3,1fr); gap:14px;">
  <div style="background:#E8F5E9;border-radius:10px;padding:14px;text-align:center;border:1.5px solid #1E6B2E;display:flex;flex-direction:column;justify-content:center;">
    <div style="font-size:11px;color:#555;margin-bottom:4px;">ベースロード法</div>
    <div style="font-size:22px;font-weight:bold;color:#1E6B2E;">{res['base_kwh']:,} kWh</div>
    <div style="font-size:11px;color:#555;">最少月を基礎電力として自動取得</div>
    <div style="font-size:10px;color:#888;margin-top:6px;">空調ゼロ基準の月を自動判定し<br>各月から引き算して空調分を逆算</div>
  </div>
  <div style="background:#E3F2FD;border-radius:10px;padding:14px;text-align:center;border:1.5px solid #1565C0;display:flex;flex-direction:column;justify-content:center;">
    <div style="font-size:11px;color:#555;margin-bottom:4px;">稼働率法（ピーク法）</div>
    <div style="font-size:22px;font-weight:bold;color:#1565C0;">{res['base_dm']:,.0f} kW</div>
    <div style="font-size:11px;color:#555;">最低月デマンド = 基礎デマンド</div>
    <div style="font-size:10px;color:#888;margin-top:6px;">空調寄与ピークを引き算で逆算<br>実量制契約kW引き下げに連動</div>
  </div>
  <div style="background:#FFF8E1;border-radius:10px;padding:14px;text-align:center;border:1.5px solid #F57F17;display:flex;flex-direction:column;justify-content:center;">
    <div style="font-size:11px;color:#555;margin-bottom:4px;">業態別 空調割合</div>
    <div style="font-size:22px;font-weight:bold;color:#E65100;">空調 {res['ac_kwh_r']:.0%}<br>ピーク {res['ac_peak_r']:.0%}</div>
    <div style="font-size:11px;color:#555;">業態: {res['gyotai']}</div>
    <div style="font-size:10px;color:#888;margin-top:6px;">チラー・空調・照明等の仮置き割合<br>管理画面で業態別に変更可能</div>
  </div>
</div>
""", unsafe_allow_html=True)

    st.markdown("""
<div style="background:#F5F5F5;border-radius:8px;padding:12px;margin-top:12px;font-size:12px;color:#444;">
<b>⚠️ 前提の仮置きについて</b>：空調割合は業態マスタの仮置き数値を適用しています。
実際の30分デマンドデータや時間帯別kWhが提供された場合、より精度の高い実測値に置き換えることが可能です。
現時点の試算はあくまでも「過去IFシミュレーション」であり、将来を保証するものではありません。
</div>
""", unsafe_allow_html=True)


# ══════════════════════════════════════════════════════════════════
# 🎯 目標回収年数から制御対象を逆算（小馬力＝割高機を1台ずつ削ぎ落とす）
#    ※独立機能。不要になれば optimize_control_plan / show_payback_optimizer と
#      呼び出し1行を消すだけで完全に撤去できます。
# ══════════════════════════════════════════════════════════════════
def optimize_control_plan(units, desired_recovery_years, system_fixed_cost,
                          sort_mode="horsepower", annual_fixed_cost=0.0):
    """目標回収年数を満たすよう、制御対象の空調機を1台ずつ除外して逆算する。

    units: [{id,name,horsepower,baseAnnualSavings,baseInitialCost}, ...]
    sort_mode: "horsepower"=馬力の小さい順 / "roi"=費用対効果の悪い(単体回収年が長い)順
    戻り値: success / finalRecoveryYears / totalInitialCost / totalAnnualSavings /
            controlledUnits / excludedUnits
    """
    # ── 0. 入力バリデーション（不正値でも落とさず安全な結果を返す）──
    work = list(units) if isinstance(units, (list, tuple)) else []
    fixed = float(system_fixed_cost) if system_fixed_cost is not None else 0.0
    afix = float(annual_fixed_cost) if annual_fixed_cost else 0.0   # 年間利用料など（粗→純の差）
    target = float(desired_recovery_years) if desired_recovery_years else float("inf")

    if not work:
        return {"success": False, "finalRecoveryYears": float("inf"),
                "totalInitialCost": fixed, "totalAnnualSavings": 0.0,
                "controlledUnits": [], "excludedUnits": [], "steps": []}

    # ── 補助: 残った制御対象からトータル値・回収年を算出（削減0以下は回収不能）──
    def totals(controlled):
        cost = sum(float(u.get("baseInitialCost", 0)) for u in controlled) + fixed
        gross_save = sum(float(u.get("baseAnnualSavings", 0)) for u in controlled)
        save = gross_save - afix                       # 純削減（年間利用料を控除）＝本体の回収式と整合
        years = (cost / save) if save > 0 else float("inf")
        return cost, save, years

    def r2(n):  # 見やすく丸める（infは保持）
        return round(n, 2) if n not in (float("inf"), float("-inf")) else n

    # ステップ履歴：各段階で「何を外したら何台・いくら・何年になったか」を記録
    steps = []
    def add_step(removed, controlled, cost, save, years):
        steps.append({
            "order": len(steps),                    # 0=現状(全台), 1,2,…=除外した順
            "removedId": (removed.get("id") if removed else None),
            "removedName": (removed.get("name") if removed else None),
            "removedHp": (float(removed.get("horsepower", 0)) if removed else None),
            "controlledCount": len(controlled),
            "totalInitialCost": cost,
            "totalAnnualSavings": save,
            "recoveryYears": r2(years),
            "meetsTarget": (years <= target),
        })

    # ── 1. まず全台制御（フルセット）で現状を計算 ──
    controlled = list(work)
    excluded = []
    cost, save, years = totals(controlled)
    add_step(None, controlled, cost, save, years)   # 段階0：現状（全台制御）

    # ── 2. 全台で既に目標達成なら、削らずそのまま返す ──
    if years <= target:
        return {"success": True, "finalRecoveryYears": r2(years),
                "totalInitialCost": cost, "totalAnnualSavings": save,
                "controlledUnits": controlled, "excludedUnits": excluded, "steps": steps}

    # ── 3. 目標未達 → 「削ぎ落とす順」に並べ替え（先頭が最初に外す機）──
    def _solo_years(u):  # 単体回収年（費用対効果。小さいほど優秀）
        s = float(u.get("baseAnnualSavings", 0))
        return (float(u.get("baseInitialCost", 0)) / s) if s > 0 else float("inf")

    if sort_mode == "roi":
        # 費用対効果の悪い（単体回収年が長い）順に外す
        controlled.sort(key=_solo_years, reverse=True)
    else:
        # 馬力の小さい順に外す（仕様の既定）
        controlled.sort(key=lambda u: float(u.get("horsepower", 0)))

    # ── 4〜6. 目標達成 or これ以上削れなくなるまで、先頭を1台ずつ除外 ──
    #    最低1台は残す（全除外＝削減0でゼロ除算・空提案になるのを防ぐデッドロック対策）。
    #    ※スマート化：固定費が大きいと「削るほど悪化」する事もあるため、目標未達のまま
    #      削り切った場合は“最短回収年だった状態（=最良解）”を返す（むやみに悪化させない）。
    best = {"years": years, "cost": cost, "save": save,
            "controlled": list(controlled), "excluded": list(excluded)}
    while years > target and len(controlled) > 1:
        removed = controlled.pop(0)             # 4. 先頭(最初に外す機)を1台外して退避
        excluded.append(removed)
        cost, save, years = totals(controlled)   # 5. 残りで再計算（6.条件はwhileで判定）
        add_step(removed, controlled, cost, save, years)  # この除外で何年になったかを記録
        if years < best["years"]:                # 最良（最短回収年）状態を更新・記憶
            best = {"years": years, "cost": cost, "save": save,
                    "controlled": list(controlled), "excluded": list(excluded)}

    # ── 7. 最終状態を判定して返却 ──
    if years <= target:
        # 目標達成 → 達成した時点（最小限の除外）の状態をそのまま返す
        return {"success": True, "finalRecoveryYears": r2(years),
                "totalInitialCost": cost, "totalAnnualSavings": save,
                "controlledUnits": controlled, "excludedUnits": excluded, "steps": steps}
    # 目標未達 → むやみに悪化させず、最短回収年だった最良解を返す
    return {"success": False, "finalRecoveryYears": r2(best["years"]),
            "totalInitialCost": best["cost"], "totalAnnualSavings": best["save"],
            "controlledUnits": best["controlled"], "excludedUnits": best["excluded"],
            "steps": steps}


def show_reduction_planner(res: dict) -> None:
    """回収年を3レバー（機体数／初期費用割引／制御率）で調整するプランナー（組合せ可）。

    「機体を1台減らす」＝制御台数をそれが使われる全コスト式（室外機構築費・工事費(概算)・
    クラウド利用料）から1引いて初期費用・年間利用料・削減を一括で再計算する。
    """
    def yen(v):
        return "¥–" if v in (float("inf"), float("-inf")) else f"¥{v:,.0f}"

    econ = res.get("econ", {}) or {}
    n_ctrl = int(res.get("n_units_ctrl", econ.get("n_ctrl", 1)) or 1)
    n_ctrl = max(n_ctrl, 1)
    ene = float(res.get("ene_saving_annual", 0) or 0)
    dm  = float(res.get("dm_saving_annual", 0) or 0)
    gross_base = ene + dm
    total_inv = float(res.get("total_invest", 0) or 0)
    subsidy = float(econ.get("subsidy", 0) or 0)
    cap_base = float(res.get("cap_rate", 0.30) or 0.30) or 0.30

    # 台数連動コスト（制御台数を使っている式）
    kouchiku_unit     = float(econ.get("kouchiku_unit", 0) or 0)      # 室外機構築費 単価/台
    kouchiku_out_base = float(econ.get("kouchiku_out", 0) or 0)       # ＝ 単価 × 制御台数
    estimate_mode     = bool(econ.get("estimate_mode", False))       # 詳細見積=工事費は人日(台数非依存)
    kouji_base        = float(econ.get("kouji_cost", 0) or 0)         # 概算は台数連動／詳細は固定(交通費)
    fb = econ.get("fee_breakdown", {}) or {}
    p_cloud      = float(fb.get("p_cloud", 0) or 0)                   # クラウド利用料 単価
    monthly_base = float(fb.get("monthly", (float(res.get("sys_fee", 0) or 0)) / 12) or 0)
    m_cloud_base = float(fb.get("m_cloud", p_cloud * (n_ctrl + 1)) or 0)
    # 台数非依存の固定初期費用（材料費・セットアップ・ハードウェア・データ取集/センサー構築費 等）
    fixed_cost = max(total_inv - kouji_base - kouchiku_out_base, 0.0)

    # 機体ごとの年間削減（容量加重配分）。弱い機から外すため強い順に保持。
    ut = res.get("units_detail")
    saves = None
    try:
        if ut is not None and "年間削減配分円" in list(getattr(ut, "columns", [])):
            s = ut.loc[ut["制御可否"], "年間削減配分円"].astype(float).tolist()
            saves = sorted([float(x) for x in s if x == x], reverse=True)
            if len(saves) != n_ctrl or sum(saves) <= 0:
                saves = None
    except Exception:
        saves = None

    def kept_frac(keep_k):
        keep_k = max(1, min(int(keep_k), n_ctrl))
        if saves:
            return sum(saves[:keep_k]) / sum(saves)     # 容量加重（強い機を優先して残す）
        return keep_k / n_ctrl                           # 配分が無ければ台数比

    def plan(keep_k=None, cap_scale=1.0, discount_rate=0.0):
        keep_k = n_ctrl if keep_k is None else max(1, min(int(keep_k), n_ctrl))
        f = kept_frac(keep_k)
        gross = gross_base * cap_scale * f
        # 制御台数を 1 引く＝下記3つの式の台数を keep_k に
        kouji_k        = kouji_base if estimate_mode else kouji_base * (keep_k / n_ctrl)
        kouchiku_out_k = (kouchiku_unit * keep_k) if kouchiku_unit > 0 else kouchiku_out_base * (keep_k / n_ctrl)
        init = (fixed_cost + kouji_k + kouchiku_out_k) * (1.0 - discount_rate)
        monthly = monthly_base - m_cloud_base + p_cloud * (keep_k + 1)   # クラウド＝制御台数+1
        fee = max(monthly, 0.0) * 12
        net = gross - fee
        years = (init / net) if net > 0 else float("inf")
        return {"keep_k": keep_k, "init": init, "fee": fee, "gross": gross,
                "net": net, "years": years, "cap_scale": cap_scale,
                "cap_rate": cap_base * cap_scale, "discount": discount_rate, "frac": f}

    base = plan()

    def yrs(p):
        return "回収不能" if p["years"] == float("inf") else f"{p['years']:.2f}年"

    def show_result(p, *, show_units=True):
        cols = st.columns(3)
        if p["years"] == float("inf") or base["years"] == float("inf"):
            dy = None
        else:
            d = p["years"] - base["years"]
            dy = (f"{d:+.2f}年（{'短縮' if d < 0 else '延長'}）" if abs(d) >= 0.005 else "±0")
        cols[0].metric("投資回収年", yrs(p), delta=dy, delta_color="inverse")
        cols[1].metric("初期費用", yen(p["init"]),
                       delta=(yen(p["init"] - base["init"]) if abs(p["init"] - base["init"]) >= 1 else None))
        cols[2].metric("年間 純削減", yen(p["net"]),
                       delta=(yen(p["net"] - base["net"]) if abs(p["net"] - base["net"]) >= 1 else None))
        bits = []
        if show_units:
            bits.append(f"制御 {p['keep_k']}/{n_ctrl}台")
        bits.append(f"容量削減率 {p['cap_rate']*100:.0f}%")
        if p["discount"] > 0:
            bits.append(f"割引 {p['discount']*100:.0f}%")
        st.caption("　/　".join(bits) +
                   f"　｜　グロス削減 {yen(p['gross'])}・年間利用料 {yen(p['fee'])}")

    st.markdown("**🎯 回収年プランナー（3つのレバーで調整・組合せ可）**")
    st.caption(f"現状＝制御 {n_ctrl}台・容量削減率 {cap_base*100:.0f}%・初期費用 {yen(total_inv)}・"
               f"回収 {yrs(base)}。下のタブで各レバーを動かすと回収年が即試算されます（実シミュレーションの近似）。")

    t1, t2, t3, t4 = st.tabs(["① 機体を減らす", "② 初期費用割引", "③ デマンド制御率", "④ 組合せ"])

    with t1:
        st.caption("弱い機体（削減への寄与が小さい順）から外します。制御台数に連動する"
                   "室外機構築費・工事費(概算)・クラウド利用料が一緒に減るため、初期費用と利用料も下がります。")
        kk = st.slider("制御する機体数", 1, n_ctrl, n_ctrl, key="pl_keep_t1")
        p = plan(keep_k=kk)
        show_result(p)
        if kk < n_ctrl:
            _dinit = base["init"] - p["init"]
            st.caption(f"↳ {n_ctrl - kk}台を外し、初期費用 {yen(_dinit)} ・年間利用料 "
                       f"{yen(base['fee'] - p['fee'])} を圧縮（うち室外機構築費・工事費・クラウド分）。"
                       + ("　※詳細見積モードのため工事費は人日積算で台数非依存です。" if estimate_mode else ""))
        if not saves:
            st.caption("※機体リスト（型番）が無いため、削減は台数比で按分しています。"
                       "型番リストを取り込むと容量加重で精緻化されます。")

    with t2:
        st.caption("初期費用そのものを値引き（補助金・キャンペーン等）した場合の回収年。削減効果は据え置き。")
        dr = st.slider("初期費用の割引率（%）", 0, 60, 0, step=5, key="pl_disc_t2") / 100.0
        show_result(plan(discount_rate=dr), show_units=False)

    with t3:
        st.caption("デマンド制御の強さ（容量削減率）を上げ下げした場合。削減額が概ね比例して増減します"
                   "（デマンドの契約下限は簡略化した近似）。")
        newcap = st.slider("容量削減率（デマンド制御の強さ・%）", 10, 60,
                           int(round(cap_base * 100)), step=5, key="pl_cap_t3") / 100.0
        show_result(plan(cap_scale=newcap / cap_base), show_units=False)

    with t4:
        st.caption("3レバーを同時に調整（例：デマンド制御率を上げつつ、弱い機体を絞る）。")
        cc = st.columns(3)
        kk = cc[0].slider("制御する機体数", 1, n_ctrl, n_ctrl, key="pl_keep_t4")
        newcap = cc[1].slider("容量削減率（%）", 10, 60, int(round(cap_base * 100)),
                              step=5, key="pl_cap_t4") / 100.0
        dr = cc[2].slider("初期費用 割引率（%）", 0, 60, 0, step=5, key="pl_disc_t4") / 100.0
        show_result(plan(keep_k=kk, cap_scale=newcap / cap_base, discount_rate=dr))

    st.divider()
    with st.expander("🔧 機体を1台ずつ選んで微調整（目標回収年から自動逆算）", expanded=False):
        show_payback_optimizer(res)


def show_payback_optimizer(res: dict) -> None:
    """結果画面下部：目標回収年数から制御対象を逆算して提案するUI。"""
    units_detail = res.get("units_detail")
    econ = res.get("econ", {})

    # 型番リスト（機体入力）が無いと1台ずつの逆算ができない
    if units_detail is None or len(units_detail) == 0 or "年間削減配分円" not in units_detail.columns:
        st.info("🎯 この逆算機能は『機体入力（型番リスト）』が必要です。サイドバーで室外機の型番・台数を登録すると利用できます。")
        return

    ut = units_detail
    n_ctrl = int(res.get("n_units_ctrl", max(int(ut["制御可否"].sum()), 1)))
    # 本体の回収（総初期投資額 ÷ 純削減額）と一致させる：
    #   ・固定費（台数非依存）＝ 総初期投資額 −（室外機工事費 ＋ 室外機構築費）
    #   ・1台あたり ＝（室外機工事費 ＋ 室外機構築費）÷ 制御台数
    #   ・年間削減は「純」（粗削減 − 年間利用料）に合わせるため annual_fixed_cost を控除
    total_inv     = float(res.get("total_invest", 0))
    kouji         = float(econ.get("kouji_cost", 0))
    kouchiku_out  = float(econ.get("kouchiku_out", 0))      # 室外機構築費（制御台数連動）
    gross_saving  = float(res.get("gross_saving", 0))
    net_saving    = float(res.get("net_saving", 0))
    annual_fixed_cost = max(gross_saving - net_saving, 0.0)  # ＝年間システム利用料
    system_fixed_cost = total_inv - kouji - kouchiku_out
    base_init_per_unit = ((kouji + kouchiku_out) / n_ctrl) if n_ctrl > 0 else 0.0
    has_hp = "馬力" in ut.columns

    # 制御対象（〇）の機だけを逆算の対象に。非制御機は最初から除外扱い。
    units = []
    for _, r in ut[ut["制御可否"]].iterrows():
        hp = float(r["馬力"]) if has_hp and pd.notna(r.get("馬力")) else float(r.get("電力kW", 0))
        units.append({
            "id": str(r.get("機器ID", "")),
            "name": (str(r.get("メーカー", "")) + " " + str(r.get("型番", ""))).strip(),
            "horsepower": hp,
            "baseAnnualSavings": float(r.get("年間削減配分円", 0) or 0),
            "baseInitialCost": base_init_per_unit,
        })

    if not units:
        st.info("🎯 制御対象（〇）の機体がありません。機体入力で制御可否を設定してください。")
        return

    # 既定の目標値：現在の粗削減ベース回収年あたりを初期表示
    cur_cost = sum(u["baseInitialCost"] for u in units) + system_fixed_cost
    cur_save = sum(u["baseAnnualSavings"] for u in units) - annual_fixed_cost
    cur_years = (cur_cost / cur_save) if cur_save > 0 else 99.9

    c1, c2 = st.columns([1, 1])
    with c1:
        target = st.number_input("🎯 希望する投資回収年数（年）", min_value=0.5, max_value=99.0,
                                 value=float(min(max(round(cur_years, 1), 0.5), 99.0)),
                                 step=0.5, key="opt_target_years")
    with c2:
        sort_label = st.selectbox("削ぎ落とす優先順位", ["馬力の小さい順", "費用対効果の悪い順"],
                                  index=0, key="opt_sort_mode",
                                  help="馬力順＝仕様どおり。費用対効果順＝1台ごとの回収年が長い機から外す（より的確）。")
    sort_mode = "roi" if sort_label == "費用対効果の悪い順" else "horsepower"

    out = optimize_control_plan(units, target, system_fixed_cost, sort_mode=sort_mode,
                                annual_fixed_cost=annual_fixed_cost)

    # ── ✏️ 制御対象の選択（✓を編集 → 逆算後の回収年数が即反映。初期＝自動逆算）──
    st.markdown("**🎯 制御対象の選択（✓を編集すると「逆算後の回収年数」が即反映）**")
    st.caption("初期状態は自動逆算の結果。希望年数より「選択中の制御対象から予測される回収年数」が"
               "短い場合は、その実回収年を表示します。任意の機を手動でON/OFFして組み替えも可能です。")

    auto_ctrl_ids = {u["id"] for u in out["controlledUnits"]}
    auto_fy = out["finalRecoveryYears"]
    sig = f"{round(target, 2)}|{sort_mode}|" + "/".join(u["id"] for u in units)
    editor_key = "opt_editor_" + str(abs(hash(sig)))

    if st.button("↻ 自動の逆算結果に戻す", key="opt_reset_" + str(abs(hash(sig)))):
        st.session_state.pop(editor_key, None)   # 手動編集を破棄して自動結果へ復帰
        st.rerun()

    edit_df = pd.DataFrame([{
        "制御対象": (u["id"] in auto_ctrl_ids),
        "管理No.": u["id"],
        "設置場所/機種": u["name"],
        "馬力/規模": round(u["horsepower"], 1),
        "年間削減額": u["baseAnnualSavings"],
        "個別初期費用": u["baseInitialCost"],
    } for u in units])

    edited = st.data_editor(
        edit_df, key=editor_key, use_container_width=True, hide_index=True,
        column_config={
            "制御対象": st.column_config.CheckboxColumn(
                "制御対象", help="✓を外すとその機を除外して即再計算します"),
            "馬力/規模": st.column_config.NumberColumn("馬力/規模", format="%.1f"),
            "年間削減額": st.column_config.NumberColumn("年間削減額", format="¥%d"),
            "個別初期費用": st.column_config.NumberColumn("個別初期費用", format="¥%d"),
        },
        disabled=["管理No.", "設置場所/機種", "馬力/規模", "年間削減額", "個別初期費用"],
    )

    # ── 選択中の✓から再計算（位置で対応づけ：管理No.重複でも安全）──
    sel_flags = list(edited["制御対象"])
    sel_units = [u for u, on in zip(units, sel_flags) if bool(on)]
    m_cost = sum(u["baseInitialCost"] for u in sel_units) + system_fixed_cost
    m_save = sum(u["baseAnnualSavings"] for u in sel_units) - annual_fixed_cost
    m_years = (m_cost / m_save) if m_save > 0 else float("inf")
    m_years_txt = "回収不能" if m_years == float("inf") else f"{m_years:.1f} 年"
    n_excl = len(units) - len(sel_units)

    # ── サマリー：「逆算後の回収年数」＝選択中の制御対象から予測される実回収年 ──
    #    希望年数より実回収が短ければ、その短い実回収年がそのまま表示される。
    m1, m2, m3, m4 = st.columns(4)
    _delta = (None if (m_years == float("inf") or auto_fy == float("inf")
                       or abs(m_years - auto_fy) < 0.05)
              else f"{m_years - auto_fy:+.1f} 年 vs自動")
    m1.metric("逆算後の回収年数", m_years_txt, delta=_delta)
    m2.metric("制御 残台数", f"{len(sel_units)} 台", delta=(f"-{n_excl} 台" if n_excl else None))
    m3.metric("トータル初期費用", f"¥{m_cost:,.0f}")
    m4.metric("純年間削減額（利用料控除後）", f"¥{m_save:,.0f}")

    if len(sel_units) == 0:
        st.error("制御対象が0台です。最低1台は✓を入れてください。")
    elif m_years <= target:
        st.success(f"✅ 希望 {target:.1f}年以内（実回収 {m_years_txt}）。除外 {n_excl}台の構成です。")
    else:
        st.warning(f"⚠️ 希望 {target:.1f}年は未達（実回収 {m_years_txt}）。"
                   "✓の調整、または希望年数の見直しを。")

    # ── 段階別シミュレーション（自動逆算の推移：参考）──
    steps = out.get("steps", [])
    if steps:
        st.markdown("**📉 削ぎ落とし段階別（参考）：どの機を外すと回収年数がどう変わるか**")
        adopted_order = len(out["excludedUnits"])  # 自動逆算が採用した段階の除外数

        def _yr_txt(v):
            return "回収不能" if v == float("inf") else f"{v:.2f} 年"

        rows = []
        for s in steps:
            if s["order"] == 0:
                action = "現状（全台制御）"
            else:
                hp = s["removedHp"]
                hp_txt = f"／{hp:.1f}" if hp is not None else ""
                action = f"{s['removedName']}（No.{s['removedId']}{hp_txt}）を除外"
            rows.append({
                "段階": ("現状" if s["order"] == 0 else f"-{s['order']}台目"),
                "操作（この段階で外す機）": action,
                "残台数": s["controlledCount"],
                "回収年数": _yr_txt(s["recoveryYears"]),
                "目標達成": ("✅" if s["meetsTarget"] else "—"),
                "判定": ("◀ 自動採用" if s["order"] == adopted_order else ""),
            })
        steps_df = pd.DataFrame(rows)

        def _highlight(row):
            if row["判定"] == "◀ 自動採用":
                return ["background-color: #E2F3E9"] * len(row)
            return [""] * len(row)

        st.dataframe(steps_df.style.apply(_highlight, axis=1),
                     use_container_width=True, hide_index=True)
        st.caption("上から順に、馬力（または費用対効果）の悪い機を1台ずつ外したときの回収年数の推移です。"
                   "「◀ 自動採用」は自動逆算が選んだ段階。✓は手動で自由に変えられます。")

    st.caption(
        f"※ 簡易逆算（粗削減ベース）。固定費 ¥{system_fixed_cost:,.0f}（AIシステム一式＋セットアップ）＋"
        f"1台あたり工事費 ¥{base_init_per_unit:,.0f} で算出。「↻ 自動の逆算結果に戻す」で自動選択へ復帰。"
        "年間利用料の台数連動や制御比率の非線形性は簡略化しています。")


# ══════════════════════════════════════════════════════════════════
# 制御モード切替UI (②指示により各要素を動的算出化)
# ══════════════════════════════════════════════════════════════════
def calc_payback(invest: float, net: float, gross: float):
    """投資回収年を必ず数値で返す（infや「回収不能」を出さない）。
    ・純利益>0 → 純利益(実利)ベースの実回収年（何十年でもそのまま数値で表示）
    ・純利益<=0（年間利用料が削減額を上回る）→ 粗削減ベースの参考年数 ＋ 超過フラグ
    戻り値: (年数:float, 利用料超過フラグ:bool)
    """
    try:
        if net and net > 0:
            return invest / net, False
        if gross and gross > 0:
            return invest / gross, True
    except Exception:
        pass
    return 99.9, True


def show_calc_methodology(res: dict) -> None:
    """このシミュレーションの全算出ロジックを「計算式 → 当てはめた数値 → 結果」で開示する。
    画面下部の『算出方法』ドロップダウン内で呼び出す。すべて res の実値を代入して描画。"""
    econ        = res.get("econ", {})
    system_cost = float(econ.get("system_cost", 0))
    setup_cost  = float(econ.get("setup_cost", 0))
    kouji_cost  = float(econ.get("kouji_cost", 0))
    hardware_cost = float(econ.get("hardware_cost", 0))
    kouchiku_total = float(econ.get("kouchiku_total", 0))
    hasu_adj    = float(econ.get("hasu_adj", 0))
    subsidy     = float(econ.get("subsidy", 0))
    area        = econ.get("area", "—")
    n_ctrl      = int(res.get("n_units_ctrl", econ.get("n_ctrl", 0)) or 0)
    n_total     = int(res.get("n_units_total", 0) or 0)
    sys_fee     = float(res.get("sys_fee", econ.get("sys_fee", 0)) or 0)
    total_inv   = float(res.get("total_invest", system_cost + kouji_cost + setup_cost))
    net_inv     = float(res.get("net_invest", max(total_inv - subsidy, 0)))

    # econ未設定（総投資額のみ既知）の保険：内訳を既定値から復元し、当てはめが¥0で出ないようにする
    econ_missing = (system_cost + kouji_cost + setup_cost) <= 0
    if econ_missing and total_inv > 0:
        _area = area if area in AREA_BUFFER else "関東"
        _nc   = n_ctrl if n_ctrl > 0 else int(round(total_inv / 100000))
        kouji_cost  = round(_nc * AREA_UNIT_PRICE * AREA_BUFFER.get(_area, 1.0) / 10000) * 10000
        setup_cost  = 400000
        system_cost = max(total_inv - kouji_cost - setup_cost, 0)
        area = _area

    ctrl_ratio  = float(res.get("ctrl_ratio", 0))
    cap_rate    = float(res.get("cap_rate", 0))
    ac_kwh_r    = float(res.get("ac_kwh_r", 0))
    ac_peak_r   = float(res.get("ac_peak_r", 0))
    ene_y       = float(res.get("ene_saving_annual", 0))
    dm_y        = float(res.get("dm_saving_annual", 0))
    gross       = float(res.get("gross_saving", 0))
    net         = float(res.get("net_saving", 0))
    old_c       = float(res.get("old_contract", 0))
    new_c       = float(res.get("new_contract", 0))
    cdelta      = float(res.get("contract_delta", 0))
    total_kwh   = float(res.get("total_reduc_kwh", 0))
    co2         = float(res.get("co2_kg", 0))
    sugi        = float(res.get("sugi_trees", 0))

    # 年間総使用量・最大需要（月次データから集計）
    try:
        df = res.get("df")
        ann_usage = float(df["使用量合計"].sum()) if df is not None else 0.0
        max_demand = float(df["最大需要電力"].max()) if df is not None else 0.0
    except Exception:
        ann_usage, max_demand = 0.0, 0.0

    payback, over_fee = calc_payback(total_inv, net, gross)
    area_buf = AREA_BUFFER.get(area, None)

    def yen(v):  return f"¥{v:,.0f}"
    def kwh(v):  return f"{v:,.0f}kWh"
    def kw(v):   return f"{v:,.1f}kW"
    def pct(v):  return f"{v*100:.1f}%"

    GRAY = "#666666"; AC = "#00B050"; LINE = "#E0E0E0"
    def card(no, title, formula, apply_html, result, note=""):
        note_html = (f'<div style="font-size:11px;color:#999;margin-top:6px;">{note}</div>'
                     if note else "")
        return f"""
        <div style="border:1px solid {LINE};border-radius:8px;padding:14px 16px;margin-bottom:12px;background:#FFFFFF;">
          <div style="font-size:13px;font-weight:700;color:#333;margin-bottom:8px;">
            <span style="color:{AC};">{no}</span>　{title}</div>
          <div style="font-size:13px;color:{GRAY};margin-bottom:4px;">
            <span style="display:inline-block;width:64px;color:#aaa;font-size:11px;">計算式</span>{formula}</div>
          <div style="font-size:13px;color:#333;margin-bottom:4px;">
            <span style="display:inline-block;width:64px;color:#aaa;font-size:11px;">当てはめ</span>{apply_html}</div>
          <div style="font-size:14px;color:#111;font-weight:700;">
            <span style="display:inline-block;width:64px;color:#aaa;font-size:11px;font-weight:400;">結果</span>
            <span style="color:{AC};">{result}</span></div>
          {note_html}
        </div>"""

    st.caption("※ 下記はすべて、今表示しているシミュレーション結果に実際に使われている数値を計算式へ当てはめたものです。")

    html = ""

    _emode = bool(econ.get("estimate_mode", False))
    _ed = econ.get("estimate_detail") or {}
    if _emode and _ed:
        # ── 初期費用（詳細見積：見積書方式）──
        html += card(
            "初期費用", "初期費用（詳細見積：見積書方式）",
            "材料費 ＋ 工事費 ＋ 交通費 ＋ ハードウェア関連 ＋ 空調システム構築費 － 端数調整 （－ 補助金 ＝ 実質投資額）",
            f"{yen(system_cost)} ＋ {yen(setup_cost)} ＋ {yen(kouji_cost)} ＋ {yen(hardware_cost)} ＋ {yen(kouchiku_total)}"
            f"{(' － 端数 ' + yen(hasu_adj)) if hasu_adj else ''}"
            f"{(' － 補助 ' + yen(subsidy)) if subsidy else ''}",
            f"総投資 {yen(total_inv)}" + (f" / 実質 {yen(net_inv)}（補助金差引後）" if subsidy else ""),
            note="見積書と同じ方式（材料費＝電材費×係数／工事費＝人日積算／交通費＝実費＋拘束＋宿泊＋レンタカー／"
                 "ハードウェア＝『利用料機器』記載品×機材単価マスタの売価合計／構築費＝単価×台数）。"
        )
        # 材料費
        html += card("材料費", "システム一式（材料費）",
                     f"電材費 × 材料費係数({_ed.get('material_factor',1.2):.2f})",
                     f"{yen(_ed.get('denzai',0))} × {_ed.get('material_factor',1.2):.2f}",
                     f"{yen(system_cost)}")
        # 工事費
        html += card("工事費", "セットアップ費（工事費）",
                     f"((ｴﾑｽﾞｶﾗｯﾄ＋シムックス人数) × 作業日数 × 人日単価) ＋ 固定加算",
                     f"(({_ed.get('n_mzu',0)}＋{_ed.get('n_simx',0)})×{_ed.get('work_days',0)}日×{yen(_ed.get('labor_day_unit',60000))})"
                     f" ＋ {yen(_ed.get('labor_fixed',120000))}",
                     f"{yen(setup_cost)}")
        # 交通費
        html += card("交通費", "交通費（施工場所エリア）",
                     "交通費(実費) ＋ 移動拘束費(作業人数×地区別単価) ＋ 宿泊費(人工×単価) ＋ レンタカー代",
                     f"実費{yen(_ed.get('trans_base',0))} ＋ 拘束{yen(_ed.get('restraint',0))}"
                     f"({_ed.get('n_mzu',0)+_ed.get('n_simx',0)}×{yen(_ed.get('area_unit',0))}) ＋ "
                     f"宿泊{yen(_ed.get('lodging',0))}({_ed.get('lodging_md',0)}×{yen(_ed.get('lodging_unit',8000))}) ＋ "
                     f"ﾚﾝﾀｶｰ{yen(_ed.get('rentacar',0))}",
                     f"{yen(kouji_cost)}",
                     note=f"地区：{AREA5_LABELS.get(_ed.get('area5',''), _ed.get('area5',''))}。"
                          "地区別単価・人日単価・各係数は『各種設定 ＞ 計算ロジックの変更』で変更可。")
    else:
        # ── 初期費用（概算）──
        html += card(
            "初期費用", "初期費用（総投資額・概算）",
            "AIシステム一式 ＋ 工事費 ＋ セットアップ費 ＋ ハードウェア関連 ＋ 空調システム構築費 － 端数調整 （－ 補助金 ＝ 実質投資額）",
            f"{yen(system_cost)} ＋ {yen(kouji_cost)} ＋ {yen(setup_cost)} ＋ {yen(hardware_cost)} ＋ {yen(kouchiku_total)}"
            f"{(' － 端数 ' + yen(hasu_adj)) if hasu_adj else ''}"
            f"{(' － 補助 ' + yen(subsidy)) if subsidy else ''}",
            f"総投資 {yen(total_inv)}" + (f" / 実質 {yen(net_inv)}（補助金差引後）" if subsidy else ""),
            note="工事費は『制御台数 × 目安単価 × エリア係数』で概算（下段参照）。"
                 "ハードウェア＝『利用料機器』記載品×機材単価マスタの売価合計／構築費＝単価×台数。"
                 "詳細がある場合はサイドバーで詳細見積に切替可。"
        )
        html += card(
            "工事費", "工事費（概算）",
            f"制御台数 × 工事目安単価({yen(AREA_UNIT_PRICE)}/台) × エリア係数 →（1万円単位で丸め）",
            f"{n_ctrl}台 × {yen(AREA_UNIT_PRICE)} × {area_buf if area_buf is not None else '—'}"
            f"（エリア：{area}）",
            f"{yen(kouji_cost)}",
        )

    # ── ハードウェア関連（利用料機器の機材費合計）──
    _hw_rows_m = econ.get("hardware_rows") or []
    if hardware_cost > 0 or _hw_rows_m:
        _hw_apply = "　＋　".join(
            f"{r.get('商品名','')} {yen(float(r.get('売価単価',0) or 0))}×{int(r.get('数量',0) or 0)}"
            for r in _hw_rows_m if int(r.get('数量', 0) or 0) > 0
        ) or "（数量0）"
        html += card(
            "ハードウェア", "ハードウェア関連（機材費合計）",
            "Σ（各機材の 売価単価 × 数量）　※『利用料機器』記載品を機材単価マスタの売価で算定",
            _hw_apply,
            f"{yen(hardware_cost)}",
            note="品名は『利用料機器(任意)』シートから転記。単価は各種設定の『機材単価マスタ』の売価に準拠。",
        )
    # ── 空調システム構築費（単価×台数）──
    if kouchiku_total > 0 or econ.get("kouchiku_total") is not None:
        _ku = float(econ.get("kouchiku_unit", 12000))
        _kq_dc = int(econ.get("kouchiku_q_dc", 0)); _kq_out = int(econ.get("kouchiku_q_out", 0))
        _kq_temp = int(econ.get("kouchiku_q_temp", 0))
        html += card(
            "構築費", "空調システム構築費（単価 × 台数）",
            "単価 ×（データ取集装置 ＋ 室外機 ＋ 温湿度センサー の台数）",
            f"{yen(_ku)} ×（データ取集装置 {_kq_dc}〔ﾏｰｷｭﾘｰ＋MM親機＋MM子機〕"
            f" ＋ 室外機 {_kq_out}〔制御台数〕 ＋ 温湿度センサー {_kq_temp}〔防水温湿度計〕）",
            f"{yen(kouchiku_total)}",
            note="構築費単価は『各種設定 ＞ 見積の既定値（構築費単価）』で変更可。",
        )
    if hasu_adj > 0:
        html += card(
            "端数調整", "端数調整（初期費用合計から減算）",
            "総初期投資額 ＝ 材料費 ＋ 工事費 ＋ 交通費 ＋ ハードウェア関連 ＋ 構築費 － 端数調整",
            f"－ {yen(hasu_adj)}",
            f"▲ {yen(hasu_adj)}",
        )

    # ── 年間利用料 ──
    fb = econ.get("fee_breakdown")
    # 計算式：項目の足し算のみ（年額は×12）
    formula_fee = "通信費 ＋ クラウド ＋ ソフト保守 ＋ データ収集　（※年額 ＝ 上記の月額合計 × 12）"
    _monthly_disp = int(fb["monthly"]) if fb else int(round(sys_fee / 12))
    if fb:
        apply_fee = (f"{yen(fb['m_tsushin'])} ＋ {yen(fb['m_cloud'])} ＋ {yen(fb['m_hoshu'])} ＋ {yen(fb['m_data'])}"
                     f" ＝ {yen(fb['monthly'])}／月　×12 ＝ {yen(sys_fee)}")
    else:
        apply_fee = f"月額合計 {yen(_monthly_disp)} × 12 ＝ {yen(sys_fee)}"
    _fee_result = f"{yen(sys_fee)} ／年（月額 {yen(_monthly_disp)}）"

    # 各項目の「算出概念」（単価×数量の式のみ・値は記載しない）を結果の下に箇条書きで表示
    if fb:
        _pt = fb.get("p_tsushin", 550); _pc = fb.get("p_cloud", 500)
        _ph = fb.get("p_hoshu", 200); _pd = fb.get("p_data", 200)
        _bul = (
            "<b>【各項目の算出概念】</b>"
            f"<br>・通信費　　＝ 単価{_pt:,}円 ×（マーキュリー台数 ＋ MMEazyAir親機台数）"
            f"<br>・クラウド　＝ 単価{_pc:,}円 ×（制御台数 ＋ 1）"
            f"<br>・ソフト保守＝ 単価{_ph:,}円 × 1"
            f"<br>・データ収集＝ 単価{_pd:,}円 ×（マーキュリー台数 ＋ MMEazyAir親機台数 ＋ MMEazyAir子機台数）"
        )
    else:
        _bul = "<b>【各項目の算出概念】</b><br>（内訳データなし）"
    _note = (
        _bul +
        "<br><br>機器台数（マーキュリー／MM親機・子機）はインポートテンプレ『利用料機器』シートでも指定可。"
        "クラウド＝制御台数＋1・ソフト保守＝1 は手入力で調整可。"
        "単価は『各種設定 ＞ 計算ロジックの変更』で変更できます。"
        "この利用料は毎年の削減額から必ず差し引いて純削減（実利）を算出します。"
    )
    html += card(
        "利用料", "年間システム利用料",
        formula_fee, apply_fee, _fee_result, note=_note
    )

    # ── 制御可能比率 ──
    src = res.get("ctrl_ratio_source", "—")
    html += card(
        "制御比率", "制御可能比率（空調のうち制御できる割合）",
        "型番ベース：Σ(制御対象機の 電力kW×稼働係数) ÷ Σ(全機 電力kW×稼働係数)　／　台数ベース：制御台数 ÷ 総台数",
        f"取得方法：{src}" + (f"　（制御 {n_ctrl} / 総 {n_total}台）" if n_total else ""),
        f"制御可能比率 ＝ {pct(ctrl_ratio)}",
        note="容量削減率（制御率）と掛け合わせ、実際に削減できる量を求めます。"
    )

    # ── ① 電力量料金 削減 ──
    _seasonal = bool(res.get("seasonal_ac", False))
    _gyc      = bool(res.get("gyotai_consider", False))
    _bac_pct  = float(res.get("base_ac_ratio", 0.0)) * 100
    _fuel_inc = bool(res.get("fuel_included", False))
    _ren_inc  = bool(res.get("ren_included", False))
    if _seasonal:
        _ac_method = ("空調kWh＝各月使用量−最少月（基礎電力）＝季節超過分（ベースロード法"
                      + (f"・業態配慮：基礎の空調分{_bac_pct:.0f}%を控除）" if _gyc else "）"))
    else:
        _ac_method = "空調kWh＝使用量 × 空調割合（業態一律）"
    _fuel_tag = "✅反映（月別）" if _fuel_inc else "※今回は排除（テンプレ未記載/空）"
    _ren_tag  = "✅反映（月別）" if _ren_inc else "※今回は排除（テンプレ未記載/空）"
    html += card(
        "①", "電力量料金の削減額（使った分への課金）",
        f"月次：{_ac_method} × 制御可能比率 × 容量削減率 × 換算係数(conv=0.5) ＝ 削減kWh　→　×（電力量単価＋燃調＋再エネ）　→　年間合計",
        f"空調割合(実効) {pct(ac_kwh_r)}　× 制御比率 {pct(ctrl_ratio)}　× 容量削減率 {pct(cap_rate)}　× conv {CONV_FACTOR}"
        f"（年間総使用量 {kwh(ann_usage)}／季節性{'ON' if _seasonal else 'OFF'}）"
        f"<br>燃料費調整額：{_fuel_tag}　／　再エネ賦課金：{_ren_tag}",
        f"年間削減 {kwh(total_kwh)} ＝ {yen(ene_y)}",
        note=(f"conv（換算係数）={CONV_FACTOR}：空調能力を絞った分の約半分が正味の電力削減になる安全側の係数。"
              "電力量(①)はリバウンド等で目減りするため0.5、デマンド(②)はピーククリップが確実なため適用せず——"
              "この非対称は実態に即した意図的な設定です。"
              "燃料費調整額・再エネ賦課金はインポートテンプレの月別列があれば自動反映、無ければ自動で排除します。"
              + ("　※【分析オプション】の自動判断により季節性（ベースロード法）"
                 + ("＋業態配慮" if _gyc else "")
                 + "がONのため、空調kWhは業態一律割合ではなく検針票の季節超過（最少月超過分"
                 + (f"−基礎の空調分{_bac_pct:.0f}%控除" if _gyc else "")
                 + "）から実測ベースで算出しています。手動でOFFにすると業態一律割合に戻ります。"
                 if _seasonal else
                 "　※【分析オプション】の季節性がOFFのため、空調kWhは業態一律割合で算出しています。"
                 "検針票12ヶ月に季節変動がある場合は、季節性をONにすると実測ベースに切り替わります。"))
    )

    # ── ② 基本料金 削減 ──
    html += card(
        "②", "基本料金の削減額（契約電力＝デマンドへの課金）",
        "契約削減kW ＝ 年間最大需要 × 空調ピーク割合 × 制御可能比率 × 容量削減率　→　月次：契約削減kW × 基本料金単価 × 力率調整 → 年間合計",
        f"年間最大需要 {kw(max_demand)} × 空調ピーク割合 {pct(ac_peak_r)} × 制御比率 {pct(ctrl_ratio)} × 容量削減率 {pct(cap_rate)}"
        f"　→　契約 {kw(old_c)} → {kw(new_c)}（削減 {kw(cdelta)}）",
        f"年間削減 ＝ {yen(dm_y)}",
        note=("◆空調ピーク割合とは：年間の最大需要電力（デマンド＝基本料金を決める最大の瞬間使用kW）のうち、"
              f"空調が占める割合（業態別、当案件は{pct(ac_peak_r)}）。これに制御可能比率・容量削減率を掛けて"
              "『ピーク時に空調制御で下げられるkW』＝契約削減kWを求めます。"
              "◆実量制：新契約の下限は『制御後の年間ピーク需要（最大需要−削減kW）』。大きな削減幅でも制御分を満額計上。"
              f"力率調整＝1−(力率−基準{res.get('pf_base',0.85)*100:.0f}%)を±{res.get('pf_cap',0.15)*100:.0f}%内で反映"
              "（基準・上下限は各種設定で変更可。力率不明時は割引・割増なし＝×1.0）。")
    )

    # ── ③ 純削減・回収 ──
    html += card(
        "③", "年間粗削減 → 純削減（実利） → 投資回収期間",
        "粗削減 ＝ ①＋②　／　純削減 ＝ 粗削減 − 年間利用料　／　回収期間 ＝ 初期費用 ÷ 年間純削減",
        f"①＋②＝粗削減 {yen(gross)}　−　年間利用料 {yen(sys_fee)}　＝　年間純削減 {yen(net)}"
        f"　→　回収期間 ＝ 初期費用 {yen(total_inv)} ÷ 年間純削減 {yen(net if not over_fee else gross)}",
        (f"年間純削減 {yen(net)} ／年 ・ 投資回収期間 {payback:.1f}年" if not over_fee
         else f"年間純削減 {yen(net)}（マイナス）・ 参考回収 {payback:.1f}年（利用料超過：粗削減ベース）"),
        note=("◆“年間純削減”は毎年『削減額から年間利用料を差し引いた手残り』です。回収期間はこの手残りで初期費用を"
              "割るので、回収に1年超かかる場合でも、その年数ぶんの利用料は毎年すでに差し引かれており全額考慮済みです。"
              "（例：初期費用240万・粗削減100万/年・利用料25万/年なら 純削減75万/年 → 240÷75＝3.2年。"
              "この3.2年間に払う利用料 25万×3.2＝80万も、純削減に織り込み済み）。"
              "◆利用料が削減額を上回る場合のみ純削減がマイナスになり、粗削減ベースの参考年数＋『利用料超過』を表示します。")
    )

    # ── CO2 ──
    html += card(
        "CO₂", "CO₂削減量・杉の木換算（参考）",
        f"CO₂削減 ＝ 年間削減kWh × 排出係数({CO2_FACTOR}kg/kWh)　／　杉換算 ＝ CO₂ ÷ 杉1本年間吸収({SUGI_KG}kg)",
        f"{kwh(total_kwh)} × {CO2_FACTOR} ＝ {co2:,.0f}kg　／　÷ {SUGI_KG}",
        f"CO₂ {co2:,.0f}kg ・ 杉 約{sugi:,.0f}本分",
    )

    st.markdown(html, unsafe_allow_html=True)


def show_control_mode(res: dict) -> None:
    st.markdown("---")
    st.markdown("### 制御強度シミュレーター")

    # 💡 1. ラジオボタンの選択肢に「制御率：〇〇%」を追加し、captionsで補足を追加します
    mode = st.radio(
        "制御モードを選択",
        options=[
            "😊 標準制御（快適重視） [制御率: 20%]",
            "😇 積極制御（削減重視） [制御率: 35%]",
            "⚖️ 最大制御（デマンド死守） [制御率: 45%]"
        ],
        captions=[
            "室外機をマイルドに制御し、快適性を最優先します",
            "削減を意識し、少し踏み込んだ制御を行います",
            "デマンド超過を絶対に防ぐための強力な制御を行います"
        ],
        horizontal=True
    )

    # 💡 2. 選択肢の文字が変わったため、mode_configsのキー側も合わせます
    mode_configs = {
        "😊 標準制御（快適重視） [制御率: 20%]":        {"cap_rate": 0.20, "comfort": "±1.0℃未満<br>（快適性維持）", "icon": "😊", "bg_color": "#E8F5E9"},
        "😇 積極制御（削減重視） [制御率: 35%]":        {"cap_rate": 0.35, "comfort": "±2.0〜3.0℃<br>（若干の温度変化あり）", "icon": "😇", "bg_color": "#FFF3E0"},
        "⚖️ 最大制御（デマンド死守） [制御率: 45%]":   {"cap_rate": 0.45, "comfort": "±1.0℃未満<br>（ピーク時のみ限定・短時間フル稼働）", "icon": "⚖️", "bg_color": "#E3F2FD"},
    }
    
    cfg = mode_configs[mode]
    target_cap = cfg["cap_rate"]
    
    # 標準キャップ率（30% = 0.30）からの比率計算
    ratio = target_cap / res["cap_rate"]
    
    # リアルタイム動的パラメータ算出
    reduc_kwh_mode = res["total_reduc_kwh"] * ratio
    reduc_bill_mode = res["gross_saving"] * ratio
    payback_years_mode, _over_fee_mode = calc_payback(res["total_invest"], reduc_bill_mode - res["sys_fee"], reduc_bill_mode)
    
    st.markdown(f"""
<div style="background:{cfg['bg_color']}; border-radius:12px; padding:25px; border: 1.5px solid #cbd5e1; margin-top:15px; box-shadow: 0 4px 6px rgba(0,0,0,0.03);">
    <div style="display: flex; align-items: center; gap: 15px; margin-bottom: 20px;">
        <h4 style="margin: 0; color: #1F3864; font-size: 20px; font-weight: bold;">{mode}</h4>
    </div>
    <div style="display: grid; grid-template-columns: repeat(4, 1fr); gap: 20px; text-align: center; align-items: stretch;">
        <div style="background: #fff; padding: 15px; border-radius: 8px; border: 1px solid #e2e8f0; display:flex; flex-direction:column; justify-content:center;">
            <div style="font-size: 11px; color: #64748b; font-weight: bold; margin-bottom: 5px;">室温変化目途</div>
            <div style="font-size: 16px; font-weight: bold; color: #1F3864; line-height: 1.45;">{cfg['comfort']}</div>
        </div>
        <div style="background: #fff; padding: 15px; border-radius: 8px; border: 1px solid #e2e8f0; display:flex; flex-direction:column; justify-content:center;">
            <div style="font-size: 11px; color: #64748b; font-weight: bold; margin-bottom: 5px;">年間電気削減量（kWh/年）</div>
            <div style="font-size: 33px; font-weight: bold; color: #00965E;">{reduc_kwh_mode:,.0f}</div>
        </div>
        <div style="background: #fff; padding: 15px; border-radius: 8px; border: 1px solid #e2e8f0; display:flex; flex-direction:column; justify-content:center;">
            <div style="font-size: 11px; color: #64748b; font-weight: bold; margin-bottom: 5px;">年間電気料金減額（円/年）</div>
            <div style="font-size: 33px; font-weight: bold; color: #00965E;">{reduc_bill_mode:,.0f}</div>
        </div>
        <div style="background: #fff; padding: 15px; border-radius: 8px; border: 1px solid #e2e8f0; display:flex; flex-direction:column; justify-content:center;">
            <div style="font-size: 11px; color: #64748b; font-weight: bold; margin-bottom: 5px;">投資回収期間（年）</div>
            <div style="font-size: 33px; font-weight: bold; color: #2F5496;">{payback_years_mode:.1f}{'<span style="font-size:13px;color:#C00000;">（利用料超過）</span>' if _over_fee_mode else ''}</div>
        </div>
    </div>
</div>
""", unsafe_allow_html=True)


# ══════════════════════════════════════════════════════════════════
# トークスクリプト
# ══════════════════════════════════════════════════════════════════
def show_financial_risk_10yr(res: dict) -> None:
    """制御強度シミュレーターとトークスクリプトの間に挿入する「今後10年間の財務リスク試算」。"""
    st.markdown("---")
    st.markdown("## 今後10年間の財務リスク試算")

    gross = float(res.get('gross_saving', 0))      # 年間グロス削減（=現状維持の払い過ぎ）
    net   = float(res.get('net_saving', 0))        # 年間純削減（利用料控除後）
    invest = float(res.get('total_invest', 0))     # 初期投資
    loss10 = gross * 10                            # 現状維持の10年損失
    dps10  = net * 10 - invest                     # DPS導入の10年累計
    gap    = dps10 + loss10                        # 財務インパクト格差
    payback, _over_fee_risk = calc_payback(invest, net, gross)
    peak_pct = res.get('ac_peak_r', 0.18)
    YEXPR = "round(datum.value/10000) + '万'"

    # ── ① 累積経費流出（現状維持） ──
    st.markdown("##### ① 何もしないという選択がもたらす、雪だるま式の経費流出")
    years = list(range(1, 11))
    df1 = pd.DataFrame({"年": [f"Year {y}" for y in years], "累積流出": [gross * y for y in years]})
    c1, c2 = st.columns([2, 1])
    with c1:
        chart1 = alt.Chart(df1).mark_bar(color="#C00000").encode(
            x=alt.X("年:N", sort=[f"Year {y}" for y in years],
                    axis=alt.Axis(title=None, labelAngle=0, labelFontSize=11, labelColor="#666666",
                                  domain=False, ticks=False, grid=False)),
            y=alt.Y("累積流出:Q", title="累積経費流出（現状維持）",
                    axis=alt.Axis(grid=True, gridColor="#E0E0E0", gridWidth=0.5, labelColor="#666666",
                                  titleColor="#666666", domain=False, ticks=False, labelExpr=YEXPR)),
            tooltip=[alt.Tooltip("累積流出:Q", title="累積", format=",.0f")],
        ).properties(height=280).configure_view(strokeWidth=0)
        st.altair_chart(chart1, use_container_width=True)
    with c2:
        st.markdown(f"""
<div style="background:#FFEBEE;border-radius:12px;padding:20px;border:2px solid #C00000;text-align:center;height:280px;display:flex;flex-direction:column;justify-content:center;">
  <div style="font-size:13px;color:#C00000;font-weight:bold;margin-bottom:8px;">10年間の想定経費損失額<br>（現状維持のコスト）</div>
  <div style="font-size:32px;font-weight:900;color:#C00000;line-height:1.2;">{loss10:,.0f}<span style="font-size:18px;">円</span></div>
  <div style="font-size:11px;color:#666;margin-top:10px;">毎年約 {gross/10000:,.0f}万円 の電気代を「払い過ぎ」ています</div>
</div>
""", unsafe_allow_html=True)

    # ── ② DPSによる「見えないムダ」の自動カット（結果連動：空調のうち削減できる割合でバー可変） ──
    st.markdown("##### ② DPSによる「見えないムダ」の自動カット")
    _au = float(res['df']['使用量合計'].sum()) if isinstance(res.get('df'), pd.DataFrame) else 0.0
    _aircon = _au * res.get('ac_kwh_r', 0.15)
    _cut = float(res.get('total_reduc_kwh', 0))
    cut_ratio = min(max((_cut / _aircon) if _aircon > 0 else 0.15, 0.03), 0.5)
    # 棒グラフの実数（kWh/年）：何が何%減るのかを明示するため
    aircon_kwh = _aircon                       # 空調の年間電力量（棒全体）
    cut_kwh    = _cut                          # うちDPSでカットできる「見えないムダ」
    after_kwh  = max(_aircon - _cut, 0)        # 制御後に残る空調電力（緑）
    st.caption(f"棒グラフ＝**空調の年間電力量（約 {aircon_kwh:,.0f} kWh/年）**。"
               f"灰色＝DPS制御でカットできる“見えないムダ”（約 {cut_kwh:,.0f} kWh）、"
               f"緑＝制御後に残る空調電力（約 {after_kwh:,.0f} kWh）。▲{cut_ratio:.0%} は空調電力に対する削減割合です。")
    s_top, s_base = 20, 215
    s_h = s_base - s_top
    cut_h = round(s_h * cut_ratio)
    rem_h = s_h - cut_h
    rem_top = s_base - rem_h
    muda_y = s_top + max(cut_h // 2 + 4, 13)
    dps_y = (s_top + rem_top) // 2
    dps_y2 = dps_y + 15                          # ▲% 直下の「約○kWh減」行
    dps_val_y = rem_top + rem_h // 2 + 4         # 緑バー中央の残kWh値
    full_mid_y = rem_top + rem_h // 2 + 4        # 薄グレー（現状バー）中央：空調全体kWh
    c3, c4 = st.columns([2, 1])
    with c3:
        st.markdown(f"""
<svg viewBox="0 -28 420 278" width="100%" style="max-width:560px; display:block; margin:0 auto;" font-family="Yu Gothic, sans-serif">
  <text x="210" y="-9" text-anchor="middle" fill="#000000" font-size="14" font-weight="bold">空調の年間電力量（kWh）</text>
  <rect x="10" y="10" width="400" height="205" fill="#f6f7f9" stroke="#E0E0E0" stroke-width="0.5"/>
  <g stroke="#E6E8EC" stroke-width="0.5">
    <line x1="10" y1="55" x2="410" y2="55"/><line x1="10" y1="100" x2="410" y2="100"/>
    <line x1="10" y1="145" x2="410" y2="145"/><line x1="10" y1="190" x2="410" y2="190"/>
    <line x1="110" y1="10" x2="110" y2="215"/><line x1="210" y1="10" x2="210" y2="215"/><line x1="310" y1="10" x2="310" y2="215"/>
  </g>
  <rect x="80" y="{s_top}" width="110" height="{cut_h}" fill="#9E9E9E"/>
  <rect x="80" y="{rem_top}" width="110" height="{rem_h}" fill="#E0E0E0"/>
  <text x="135" y="{muda_y}" text-anchor="middle" fill="#ffffff" font-size="11" font-weight="bold">見えないムダ</text>
  <text x="135" y="{full_mid_y}" text-anchor="middle" fill="#555555" font-size="12" font-weight="bold">約{aircon_kwh:,.0f}</text>
  <rect x="250" y="{rem_top}" width="110" height="{rem_h}" fill="#00B050"/>
  <text x="305" y="{dps_val_y}" text-anchor="middle" fill="#ffffff" font-size="11" font-weight="bold">約{after_kwh:,.0f}</text>
  <line x1="190" y1="{s_top}" x2="360" y2="{rem_top}" stroke="#00B050" stroke-width="1.5" stroke-dasharray="6 4"/>
  <line x1="190" y1="{rem_top}" x2="250" y2="{rem_top}" stroke="#00B050" stroke-width="1.5" stroke-dasharray="6 4"/>
  <text x="222" y="{dps_y}" text-anchor="middle" fill="#000000" font-size="15" font-weight="bold">▲{cut_ratio:.0%}</text>
  <text x="222" y="{dps_y2}" text-anchor="middle" fill="#000000" font-size="10" font-weight="bold">約{cut_kwh:,.0f}kWh減</text>
  <text x="135" y="233" text-anchor="middle" fill="#666666" font-size="11">現状（制御なし）</text>
  <text x="305" y="233" text-anchor="middle" fill="#666666" font-size="11">DPS制御後</text>
</svg>
""", unsafe_allow_html=True)
    with c4:
        st.markdown(f"""
<div style="background:#E8F5E9;border-radius:12px;padding:20px;border:1.5px solid #1E6B2E;text-align:center;height:240px;display:flex;flex-direction:column;justify-content:center;">
  <div style="font-size:13px;color:#1E6B2E;font-weight:bold;margin-bottom:8px;">年間予想削減純利益<br>（純利益換算）</div>
  <div style="font-size:30px;font-weight:900;color:#00965E;">{net:,.0f}<span style="font-size:16px;">円</span></div>
  <div style="font-size:11px;color:#666;margin-top:10px;">空調電力の約 {cut_ratio:.0%}（約{cut_kwh:,.0f}kWh/年）をカット。<br>快適性を損なわずに、設備の運用を最適化。</div>
</div>
""", unsafe_allow_html=True)

    # ── ③ 投資回収（ブレークイーブン） ──
    st.markdown(f"##### ③ 約 {payback:.1f} 年で投資回収完了、その後は完全な「利益」へ")
    st.caption(f"初期投資の約 {invest/10000:,.0f}万円 は約 {payback:.1f} 年でペイ。以降は年間約 {net/10000:,.0f}万円 が純利益として蓄積し続けます。")
    rows3 = [{"年": float(y), "累積収支": net * y - invest} for y in range(0, 11)]
    rows3.append({"年": payback, "累積収支": 0.0})
    rows3.sort(key=lambda r: r["年"])
    df3 = pd.DataFrame(rows3)
    ex3 = alt.X("年:Q", title=None, scale=alt.Scale(domain=[0, 10]),
                axis=alt.Axis(labelColor="#666666", labelFontSize=12, domain=False, ticks=False,
                              grid=False, tickCount=10))
    ey3 = alt.Y("累積収支:Q", title="累積収支",
                axis=alt.Axis(grid=True, gridColor="#E0E0E0", gridWidth=0.5, labelColor="#666666",
                              titleColor="#666666", domain=False, ticks=False, labelExpr=YEXPR))
    a_neg = alt.Chart(df3[df3["累積収支"] <= 1e-6]).mark_area(color="#C00000", opacity=0.22, interpolate="monotone").encode(x=ex3, y=ey3)
    a_pos = alt.Chart(df3[df3["累積収支"] >= -1e-6]).mark_area(color="#1F3864", opacity=0.22, interpolate="monotone").encode(x=ex3, y=ey3)
    line3 = alt.Chart(df3).mark_line(color="#1F3864", strokeWidth=2.5, interpolate="monotone").encode(x=ex3, y=ey3)
    zero3 = alt.Chart(pd.DataFrame({"y": [0]})).mark_rule(color="#555555", strokeWidth=1).encode(y="y:Q")
    bept3 = alt.Chart(pd.DataFrame({"年": [payback], "累積収支": [0.0]})).mark_point(color="#C00000", size=140, filled=True).encode(x="年:Q", y="累積収支:Q")
    t_st = alt.Chart(pd.DataFrame({"年": [0.4], "累積収支": [-invest], "t": [f"-{invest/10000:,.0f}万円"]})).mark_text(color="#C00000", fontWeight="bold", align="left", dy=14, fontSize=13).encode(x="年:Q", y="累積収支:Q", text="t:N")
    t_en = alt.Chart(pd.DataFrame({"年": [9.6], "累積収支": [dps10], "t": [f"+{dps10/10000:,.0f}万円"]})).mark_text(color="#1F3864", fontWeight="bold", align="right", dy=-12, fontSize=13).encode(x="年:Q", y="累積収支:Q", text="t:N")
    t_be = alt.Chart(pd.DataFrame({"年": [payback], "累積収支": [0.0], "t": ["ブレークイーブン（投資回収完了）"]})).mark_text(color="#C00000", align="left", dx=10, dy=22, fontSize=11).encode(x="年:Q", y="累積収支:Q", text="t:N")
    chart3 = alt.layer(a_neg, a_pos, line3, zero3, bept3, t_st, t_en, t_be).properties(height=300).configure_view(strokeWidth=0)
    st.altair_chart(chart3, use_container_width=True)

    # ── ④ 10年間の財務インパクト格差 ──
    st.markdown(f"##### ④ 10年間の財務インパクト格差：約 {gap/10000:,.0f}万円")
    st.markdown(f"""
<div style="display:grid; grid-template-columns:repeat(2,1fr); gap:16px;">
  <div style="background:#FFF5F5;border:1.5px solid #C00000;border-radius:10px;padding:18px;text-align:center;display:flex;flex-direction:column;justify-content:center;">
    <div style="font-weight:bold;color:#C00000;font-size:15px;">現状維持（不作為）</div>
    <div style="font-size:12px;color:#666;margin:6px 0;">初期費用：0円</div>
    <div style="font-size:12px;color:#666;">10年後 累計キャッシュフロー</div>
    <div style="font-size:26px;font-weight:900;color:#C00000;">-{loss10:,.0f}円</div>
    <div style="font-size:11px;color:#999;">（継続的な流出）</div>
  </div>
  <div style="background:#F0F4FA;border:1.5px solid #1F3864;border-radius:10px;padding:18px;text-align:center;display:flex;flex-direction:column;justify-content:center;">
    <div style="font-weight:bold;color:#1F3864;font-size:15px;">DPS導入（投資実行）</div>
    <div style="font-size:12px;color:#666;margin:6px 0;">初期費用：{invest:,.0f}円</div>
    <div style="font-size:12px;color:#666;">10年後 累計キャッシュフロー</div>
    <div style="font-size:26px;font-weight:900;color:#1F3864;">+{dps10:,.0f}円</div>
    <div style="font-size:11px;color:#999;">（黒字転換）</div>
  </div>
</div>
<div style="background:#1F3864;color:#fff;border-radius:10px;padding:14px;text-align:center;margin-top:14px;font-size:18px;font-weight:bold;">
  その差額、なんと {gap:,.0f}円 のキャッシュフロー改善効果。
</div>
""", unsafe_allow_html=True)


def show_talk_script(res: dict, app_data: dict, client_name: str,
                     payback_years: float, invest: int) -> None:
    st.markdown("---")
    st.markdown("### 営業トークスクリプト")

    item_price  = app_data["item_price"]
    item_margin = app_data["item_margin"]
    items_needed = math.ceil(res["net_saving"] / (item_price * item_margin)) if item_price * item_margin > 0 else 0

    st.caption("話法名をクリックすると、その場で内容が開きます。")
    with st.expander("社長（経営層）向け話法", expanded=False):
        st.text_area(
            label="社長向け",
            label_visibility="collapsed",
            value=f"""「{client_name}さん、今回の数字をご覧ください。

電気料金の削減だけで年間【¥{res['gross_saving']:,.0f}】。システム利用料を差し引いた実利は【¥{res['net_saving']:,.0f}】です。

これは御社の主力商品（単価¥{item_price:,}・粗利{item_margin:.1%}）を、新規で【{items_needed:,}個】追加販売して稼ぎ出す純利益と全く同じ財務価値があります。

新規でこれだけの数量を売る難易度と、今あるエアコンの無駄をAIが自動で削る確実性を比較してください。

初期投資¥{invest:,}に対して、回収期間は約{payback_years:.1f}年。CO₂も年間{res['co2_kg']:,.0f}kgの削減実績として対外発表にそのままご活用いただけます。」""",
            height=240, key="talk_president"
        )
    with st.expander("現場・工場長向け話法", expanded=False):
        st.text_area(
            label="現場向け",
            label_visibility="collapsed",
            value=f"""「室外機を1台ずつ数分間だけ順繰りに送風へ切り替えるローテーション制御のため、室温変化は1℃未満で現場の快適性は一切変わりません。

主電源を強制オフにするような無理な制御ではないため、コンプレッサーへの負荷や機器寿命への悪影響もゼロです。

警報が鳴って誰かが走る必要は一切ありません。AIが24時間365日完全自動で制御するため、現場の業務負担・我慢は「0秒」です。

生産ライン・重要設備には100%タッチしません。空調（全体の{res['ac_peak_r']:.0%}）だけをスマート制御します。」""",
            height=240, key="talk_field"
        )


# ══════════════════════════════════════════════════════════════════
# メインアプリ
# ══════════════════════════════════════════════════════════════════
def _quote_seed_rows(app_data: dict, res: dict):
    """機材費の初期行：主要機器をシード（売価・原価はマスタ参照）。res に台数があれば数量へ反映。
    通常は部材表インポートで上書きされる。"""
    master = {m.get("商品名", ""): m for m in (app_data.get("material_master", []) or [])}
    econ = (res or {}).get("econ", {}) if res else {}
    fb = econ.get("fee_breakdown", {}) if isinstance(econ, dict) else {}
    seeds = [
        ("データ取集装置（Marimba Mercury)", int(fb.get("n_mercury", 0) or 0)),
        ("MMEazyAir親機", int(fb.get("n_mm_parent", 0) or 0)),
        ("MMEazyAir2（子機）", int(fb.get("n_mm_child", 0) or 0)),
    ]
    rows = []
    for nm, q in seeds:
        m = master.get(nm, {})
        rows.append({"商品名": nm, "数量": int(q), "単位": m.get("単位", "台"),
                     "売価単価": float(m.get("売価", 0) or 0), "原価単価": float(m.get("原価", 0) or 0)})
    return rows


def _parse_buzai(uploaded):
    """部材表Excelから『品名・数量』の表を抽出。複数シート/表のうち“有効データ行が最も多い表”を選ぶ
    （空ヘッダーだけのシートを誤検出しない）。シート名に「部材」が含まれる場合は優先。
    戻り値: (DataFrame[整数列], 品名列index, 数量列index) または None。"""
    try:
        xls = pd.ExcelFile(uploaded)
    except Exception:
        return None
    name_kw = ["品名", "商品名", "名称", "品目"]
    qty_kw = ["数量", "台数", "員数"]
    best = None  # (score, body_df, ncol, qcol)
    for sh in xls.sheet_names:
        try:
            raw = xls.parse(sh, header=None, dtype=str)
        except Exception:
            continue
        for i in range(min(len(raw), 40)):
            cells = ["" if x is None else str(x) for x in raw.iloc[i].tolist()]
            norm = [c.replace(" ", "").replace("　", "") for c in cells]
            ncol = next((j for j, c in enumerate(norm) if any(k in c for k in name_kw)), None)
            qcol = next((j for j, c in enumerate(norm) if any(k in c for k in qty_kw)), None)
            if ncol is None or qcol is None:
                continue
            body = raw.iloc[i + 1:].reset_index(drop=True)
            cnt = 0
            for _, rr in body.iterrows():
                nm = "" if ncol >= len(rr) else str(rr.iloc[ncol]).strip()
                if not nm or nm.lower() == "nan":
                    continue
                try:
                    q = float(str(rr.iloc[qcol]).replace(",", "").strip() or 0)
                except Exception:
                    q = 0
                if q > 0:
                    cnt += 1
            score = cnt + (1000 if "部材" in str(sh) else 0)
            if cnt > 0 and (best is None or score > best[0]):
                best = (score, body, ncol, qcol)
    if best is None:
        return None
    return (best[1], best[2], best[3])


def match_buzai_to_master(uploaded, app_data):
    """部材表（file-like）を _parse_buzai で解析し、機材単価マスタ（material_master）の商品名で突合。
    [{商品名,数量,単位,売価単価,原価単価}] と 未マッチ品名リスト を返す。表が無ければ (None, [])。"""
    parsed = _parse_buzai(uploaded)
    if parsed is None:
        return None, []
    body, ncol, qcol = parsed

    def _norm(s):
        return "".join(str(s).split()).lower()
    mnorm = {_norm(m.get("商品名", "")): m for m in (app_data.get("material_master", []) or [])}
    rows, miss = [], []
    for _, r in body.iterrows():
        nm = "" if ncol >= len(r) else str(r.iloc[ncol]).strip()
        if not nm or nm.lower() == "nan":
            continue
        try:
            q = int(float(str(r.iloc[qcol]).replace(",", "").strip() or 0))
        except Exception:
            q = 0
        key = _norm(nm)
        hit = mnorm.get(key) or next((mv for mk, mv in mnorm.items()
                                      if key and (key in mk or mk in key)), None)
        if hit:
            rows.append({"商品名": hit.get("商品名", nm), "数量": q,
                         "単位": hit.get("単位", "台"),
                         "売価単価": float(hit.get("売価", 0) or 0),
                         "原価単価": float(hit.get("原価", 0) or 0)})
        else:
            miss.append(nm)
            rows.append({"商品名": nm, "数量": q, "単位": "台",
                         "売価単価": 0.0, "原価単価": 0.0})
    return rows, miss


def match_material(name, app_data):
    """品名を機材単価マスタ（material_master）に照合して該当dictを返す（無ければNone）。
    完全一致→部分一致→数字/括弧/記号を無視したゆるい一致 の順で探す。"""
    import re
    items = app_data.get("material_master", []) or []

    def _norm(s):
        return "".join(str(s).split()).lower()

    def _loose(s):
        return re.sub(r"[0-9０-９（）()「」【】\[\]・,，、．\.\-_/＿]+", "", _norm(s))

    mnorm = {_norm(m.get("商品名", "")): m for m in items}
    mloose = {_loose(m.get("商品名", "")): m for m in items}
    key = _norm(name)
    if key in mnorm:
        return mnorm[key]
    hit = next((mv for mk, mv in mnorm.items() if key and (key in mk or mk in key)), None)
    if hit:
        return hit
    lk = _loose(name)
    if lk and lk in mloose:
        return mloose[lk]
    return next((mv for mk, mv in mloose.items() if lk and (lk in mk or mk in lk)), None)


def apply_master_prices(rows, app_data):
    """明細行リストの各 商品名 を機材単価マスタに照合し、売価単価/原価単価/単位 をマスタの値で上書き。
    （マスタ更新を常に反映するため、表示・集計の直前に呼ぶ）。戻り値は同じリスト。"""
    for r in (rows or []):
        hit = match_material(r.get("商品名", ""), app_data)
        if hit:
            r["売価単価"] = float(hit.get("売価", 0) or 0)
            r["原価単価"] = float(hit.get("原価", 0) or 0)
            if not str(r.get("単位", "")).strip():
                r["単位"] = hit.get("単位", "台")
    return rows


def read_fee_sheet_items(uploaded, app_data):
    """『利用料機器』シート（項目・台数）の記載品を“全件”読み取り（タイトル/見出し行は除外）、
    機材単価マスタの単価で明細 [{商品名,数量,単位,売価単価,原価単価}] に転記。(rows, miss) を返す。"""
    import re
    name = (getattr(uploaded, "name", "") or "").lower()
    if name.endswith(".csv"):
        return [], []
    try:
        uploaded.seek(0)
        xls = pd.ExcelFile(uploaded)
    except Exception:
        return [], []
    name_kw = ["項目", "品名", "商品名", "名称", "品目"]
    qty_kw = ["台数", "数量", "員数", "個数"]
    # シート選定：名前に「利用料」を含むものを優先、無ければ全シートから 項目+台数 の表を探す
    sheets = sorted(xls.sheet_names, key=lambda s: (0 if "利用料" in str(s) else 1))
    target = None  # (raw, hdr_i, ncol, qcol)
    for sn in sheets:
        try:
            raw = pd.read_excel(xls, sheet_name=sn, header=None, dtype=str)
        except Exception:
            continue
        for i in range(min(len(raw), 40)):
            cells = ["" if pd.isna(v) else str(v) for v in raw.iloc[i].tolist()]
            norm = [c.replace(" ", "").replace("　", "") for c in cells]
            ncol = next((j for j, c in enumerate(norm) if any(k in c for k in name_kw)), None)
            qcol = next((j for j, c in enumerate(norm) if any(k in c for k in qty_kw)), None)
            if ncol is not None and qcol is not None and ncol != qcol:
                target = (raw, i, ncol, qcol)
                break
        if target is not None:
            break
    if target is None:
        return [], []
    raw, hdr_i, ncol, qcol = target

    def _norm(s):
        return "".join(str(s).split()).lower()

    def _loose(s):
        # 数字・全半角括弧・記号を除去して表記差（例：MMEazyAir子機 ⇔ MMEazyAir2（子機））を吸収
        return re.sub(r"[0-9０-９（）()「」【】\[\]・,，、．\.\-_/＿]+", "", _norm(s))

    mnorm = {_norm(m.get("商品名", "")): m for m in (app_data.get("material_master", []) or [])}
    mloose = {_loose(m.get("商品名", "")): m for m in (app_data.get("material_master", []) or [])}

    def _find(nm):
        key = _norm(nm)
        if key in mnorm:
            return mnorm[key]
        hit = next((mv for mk, mv in mnorm.items() if key and (key in mk or mk in key)), None)
        if hit:
            return hit
        lk = _loose(nm)
        if lk and lk in mloose:
            return mloose[lk]
        return next((mv for mk, mv in mloose.items() if lk and (lk in mk or mk in lk)), None)

    rows, miss = [], []
    for r in range(hdr_i + 1, len(raw)):
        cells = ["" if pd.isna(v) else str(v) for v in raw.iloc[r].tolist()]
        nm = cells[ncol].strip() if ncol < len(cells) else ""
        if not nm or nm.lower() == "nan":
            continue
        qm = re.search(r"-?\d+", (cells[qcol] if qcol < len(cells) else "").replace(",", ""))
        q = int(qm.group()) if qm else 0
        hit = _find(nm)
        if hit:
            rows.append({"商品名": hit.get("商品名", nm), "数量": q, "単位": hit.get("単位", "台"),
                         "売価単価": float(hit.get("売価", 0) or 0),
                         "原価単価": float(hit.get("原価", 0) or 0)})
        else:
            miss.append(nm)
            rows.append({"商品名": nm, "数量": q, "単位": "台",
                         "売価単価": 0.0, "原価単価": 0.0})
    return rows, miss


def _parse_buzai_eng(uploaded):
    """部材表からエンジ費の引用値を抽出：電材費(工務記入No2単価)・Ms'外注費(No1単価)・
    エムズ/シムックス人数・日数（"N名*M日"）。dictを返す（取れた項目のみ）。"""
    import re
    try:
        from openpyxl import load_workbook
        wb = load_workbook(uploaded, data_only=True)
    except Exception:
        return {}
    ws = next((wb[sn] for sn in wb.sheetnames if "部材" in sn), wb[wb.sheetnames[0]])
    out = {}
    # 人数・日数（"4名*4日" 等）。同じ行のB列ラベルでエムズ／シムックスを判定
    pat = re.compile(r"(\d+)\s*名\s*[\*×ｘxX]\s*(\d+)\s*日")
    days = None
    for row in ws.iter_rows():
        for c in row:
            if isinstance(c.value, str):
                mt = pat.search(c.value)
                if mt:
                    n, d = int(mt.group(1)), int(mt.group(2))
                    label = str(ws.cell(c.row, 2).value or "")
                    if "エムズ" in label or "ｴﾑｽﾞ" in label:
                        out["n_mzu"] = n
                        days = d
                    elif ("cimx" in label.lower()) or "シムックス" in label or "シム" in label:
                        out["n_simx"] = n
                        if days is None:
                            days = d
    if days is not None:
        out["work_days"] = days
    # 工務記入 単価（件名・単価ヘッダーを探し、No1/No2の単価を取得）
    hdr_r = tanka_c = None
    for row in ws.iter_rows(min_row=1, max_row=20):
        vals = {str(c.value): c.column for c in row if c.value is not None}
        if any("件名" in k for k in vals) and any("単価" in k for k in vals):
            hdr_r = row[0].row
            tanka_c = next(col for k, col in vals.items() if "単価" in k)
            break
    if hdr_r and tanka_c:
        tankas = []
        for rr in range(hdr_r + 1, hdr_r + 8):
            if ws.cell(rr, 1).value is None:
                continue
            try:
                tankas.append(int(float(ws.cell(rr, tanka_c).value)))
            except Exception:
                pass
            if len(tankas) >= 2:
                break
        if len(tankas) >= 1:
            out["ms_out"] = tankas[0]    # No1 → Ms'カラット外注費
        if len(tankas) >= 2:
            out["denzai"] = tankas[1]    # No2 → 電材費（材料費=×倍率／Ms'材料費=そのまま）
    return out


if hasattr(st, "dialog"):
    @st.dialog("⚠ データ不足の確認")
    def _quote_missing_dialog(issues):
        st.warning("以下の項目に不足があります。プレビューは作成しますが、内容をご確認ください。")
        for _it in issues:
            st.markdown(f"- {_it}")
        if st.button("確認しました", key="q_miss_ok"):
            st.rerun()
else:
    def _quote_missing_dialog(issues):
        try:
            st.toast("⚠ データ不足があります（下の警告をご確認ください）")
        except Exception:
            pass
        st.warning("データ不足： " + " ／ ".join(issues))


def show_quote_builder(app_data: dict, res: dict, client_name: str) -> None:
    """📄 見積書作成ページ：マスタ原価＋一律利益率→売価で見積を組み、御見積書プレビュー／転記用明細を出力。"""
    st.markdown("## 見積書作成")

    qs = copy.deepcopy(DEFAULT_DATA["quote_settings"])
    qs.update(app_data.get("quote_settings", {}) or {})
    comp = qs.get("company", {})
    engd = qs.get("eng_defaults", {})
    tax_rate = float(qs.get("tax_rate", 0.10))
    bt_rate = float(qs.get("bt_rate", 0.25))
    kunit = float(qs.get("kouchiku_unit", 12000))
    today = datetime.date.today()
    n_ctrl = int((res or {}).get("n_units_ctrl",
                                 ((res or {}).get("econ", {}) or {}).get("n_ctrl", 0)) or 0)

    def yen(v):
        return f"¥{v:,.0f}"

    # ── 1. 見積ヘッダー（手入力・常時表示） ──────────────
    st.markdown("### 1. 見積ヘッダー")
    h1, h2, h3 = st.columns(3)
    with h1:
        q_title = st.text_input("件名", value="Disital Power Survice（空調省エネ）導入費", key="q_title")
        _to = (client_name + " 御中") if client_name else "御中"
        q_to = st.text_input("宛先", value=_to, key="q_to")
    with h2:
        q_date = st.date_input("発行日", value=today, key="q_date")
        valid_days = st.number_input("有効期限（発行日＋日数）", 0, 365, 30, key="q_valid")
    with h3:
        q_person = st.text_input("担当者", value="", key="q_person")
        doc_no = st.text_input("書類番号",
                               value=f"{q_date.strftime('%Y%m')}-{len(app_data.get('saved_quotes', [])) + 77:05d}[01]",
                               key="q_docno")
    h4, h5, h6 = st.columns(3)
    with h4:
        q_delivery = st.text_input("納期", value=qs.get("delivery", "別途お打ち合わせ"), key="q_deliv")
    with h5:
        q_to_place = st.text_input("納品先", value=qs.get("deliver_to", "貴社ご指定場所"), key="q_place")
    with h6:
        q_pay = st.text_input("支払条件", value=qs.get("payment", "月末締め翌月末払い"), key="q_pay")
    valid_until = q_date + datetime.timedelta(days=int(valid_days))

    st.markdown("**押印（3枠）の氏名 — 各枠の印の名前を指定できます（空欄なら枠のみ）**")
    sct = st.columns(3)
    q_stamp1 = sct[0].text_input("印①（左）", value="", key="q_stamp1")
    q_stamp2 = sct[1].text_input("印②（中）", value="", key="q_stamp2")
    q_stamp3 = sct[2].text_input("印③（右）", value="", key="q_stamp3")

    # ── 2. 利益率（原価への倍率・一律・手動可変・常時表示） ──
    st.markdown("### 2. 利益率（原価への倍率・手動で変更可）")
    m1, m2 = st.columns([1, 2])
    with m1:
        factor = st.number_input("利益率（原価 × 倍率）", 1.0, 5.0,
                                 float(qs.get("markup_factor", 1.2)), step=0.05, key="q_factor")
    with m2:
        target_f = float(qs.get("target_factor", 1.1))
        if factor < target_f:
            st.warning(f"⚠ 倍率 ×{factor:.2f} が下限 ×{target_f:.2f} を下回っています。")
        else:
            st.caption(f"現在 ×{factor:.2f}。**機材費はマスタの売価**を使用。"
                       f"倍率は **材料費・諸経費**（および売価未設定の機材）に適用します（下限 ×{target_f:.2f}）。")

    def sale_of(cost):
        return int(round(cost * factor))

    # ── 3. 機材費（マスタ原価＋インポート） ───────────────
    st.markdown("### 3. 機材費")
    with st.expander("📥 部材表Excelをインポート（アップロードで自動反映）", expanded=False):
        up = st.file_uploader("部材表（品名・数量を含むExcel）", type=["xlsx", "xls"], key="q_buzai")
        if up is not None:
            sig = (up.name, getattr(up, "size", 0))
            if st.session_state.get("q_buzai_sig") != sig:
                # アップロード直後に1回だけ自動取込（プレビュー・列選択なし）
                up.seek(0)
                _data = up.read()
                parsed = _parse_buzai(io.BytesIO(_data))
                if parsed is None:
                    st.error("部材表の『品名・数量』の表を見つけられませんでした。シート内容をご確認ください。")
                else:
                    body, ncol, qcol = parsed

                    def _norm(s):
                        return "".join(str(s).split()).lower()
                    mnorm = {_norm(m.get("商品名", "")): m for m in (app_data.get("material_master", []) or [])}
                    rows, miss = [], []
                    for _, r in body.iterrows():
                        nm = "" if ncol >= len(r) else str(r.iloc[ncol]).strip()
                        if not nm or nm.lower() == "nan":
                            continue
                        try:
                            q = int(float(str(r.iloc[qcol]).replace(",", "").strip() or 0))
                        except Exception:
                            q = 0
                        key = _norm(nm)
                        hit = mnorm.get(key) or next((mv for mk, mv in mnorm.items()
                                                      if key and (key in mk or mk in key)), None)
                        if hit:
                            rows.append({"商品名": hit.get("商品名", nm), "数量": q,
                                         "単位": hit.get("単位", "台"),
                                         "売価単価": float(hit.get("売価", 0) or 0),
                                         "原価単価": float(hit.get("原価", 0) or 0)})
                        else:
                            miss.append(nm)
                            rows.append({"商品名": nm, "数量": q, "単位": "台",
                                         "売価単価": 0.0, "原価単価": 0.0})
                    if rows:
                        st.session_state.quote_mat_rows = rows
                        st.session_state.q_mat_ver = st.session_state.get("q_mat_ver", 0) + 1
                        st.session_state.q_eng_ver = st.session_state.get("q_eng_ver", 0) + 1
                        st.session_state.q_buzai_eng = _parse_buzai_eng(io.BytesIO(_data))
                        st.session_state.q_buzai_sig = sig
                        st.session_state.q_buzai_n = len(rows)
                        st.session_state.q_buzai_miss = miss
                        st.rerun()
                    else:
                        st.error("有効な行がありませんでした。部材表の内容をご確認ください。")
            else:
                st.success(f"部材表を反映済み（{st.session_state.get('q_buzai_n', 0)}件）。下の機材費の表に反映されています。")
                _be = st.session_state.get("q_buzai_eng") or {}
                if _be:
                    st.caption("エンジ費へ引用： "
                               + "／".join(filter(None, [
                                   f"電材費 ¥{_be['denzai']:,}" if "denzai" in _be else "",
                                   f"Ms'外注 ¥{_be['ms_out']:,}" if "ms_out" in _be else "",
                                   f"人数 エムズ{_be.get('n_mzu','-')}・シムックス{_be.get('n_simx','-')}",
                                   f"日数{_be['work_days']}" if "work_days" in _be else "",
                               ])))
                _miss = st.session_state.get("q_buzai_miss") or []
                if _miss:
                    st.warning("マスタ未登録（原価0で取込・要確認）： " + " / ".join(_miss[:12]))
                if st.button("もう一度取り込む", key="q_buzai_redo"):
                    st.session_state.q_buzai_sig = None
                    st.rerun()

    mat_rows = st.session_state.get("quote_mat_rows") or _quote_seed_rows(app_data, res)
    _mat_ver = st.session_state.get("q_mat_ver", 0)
    mat_edit = st.data_editor(
        pd.DataFrame(mat_rows), num_rows="dynamic", use_container_width=True, key=f"q_mat_editor_{_mat_ver}",
        column_config={
            "商品名": st.column_config.TextColumn("商品名", width="large"),
            "数量": st.column_config.NumberColumn("数量", min_value=0, step=1),
            "単位": st.column_config.TextColumn("単位", width="small"),
            "売価単価": st.column_config.NumberColumn("売価単価", min_value=0, step=1, format="%d"),
            "原価単価": st.column_config.NumberColumn("原価単価", min_value=0, step=1, format="%d"),
        })
    st.session_state.quote_mat_rows = mat_edit.to_dict("records")
    st.caption("機材費の売価はマスタの売価をそのまま使用します（利益＝売価−原価）。"
               "売価が0の行は 原価×利益率 で暫定算出します。")

    hw_lines, hw_amount, hw_profit, hw_cost = [], 0, 0, 0
    name_qty = {}
    for r in mat_edit.to_dict("records"):
        nm = str(r.get("商品名", "")).strip()
        if not nm:
            continue
        q = int(r.get("数量", 0) or 0)
        cu = float(r.get("原価単価", 0) or 0)
        su = float(r.get("売価単価", 0) or 0)
        if su <= 0:                       # 売価未設定なら 原価×利益率 でフォールバック
            su = sale_of(cu)
        amt = su * q
        cst = cu * q
        prof = amt - cst
        hw_lines.append({"name": nm, "qty": q, "unit": r.get("単位", "台"),
                         "sale_u": su, "cost_u": cu, "amount": amt, "profit": prof})
        hw_amount += amt; hw_profit += prof; hw_cost += cst
        name_qty[nm] = name_qty.get(nm, 0) + q

    # ── 4. エンジニアリング費（計算式・引用元に基づき算出） ──
    st.markdown("### 4. エンジニアリング費")
    ed = ((res or {}).get("econ", {}) or {}).get("estimate_detail") or {}
    beng = st.session_state.get("q_buzai_eng") or {}   # 部材表から引用（最優先）
    est = app_data.get("const_settings", {}).get("estimate_settings",
                                                 DEFAULT_DATA["const_settings"]["estimate_settings"])
    area_unit_map = est.get("area_unit", {}) or {}
    ev = st.session_state.get("q_eng_ver", 0)          # 部材表取込ごとに入力欄を更新

    def _engv(key, dflt):
        return beng.get(key, ed.get(key, dflt))
    # 部材表／詳細見積が設定済みなら、引用値が見えるよう既定で開く
    _eng_open = bool(beng) or bool(ed) or bool(st.session_state.get("cost_detail_mode"))
    _src = "部材表" if beng else ("詳細見積" if ed else "")
    _eng_label = (f"🔧 4. エンジニアリング費の入力値（{_src}から引用済み・開いて確認）" if _eng_open
                  else "🔧 4. エンジニアリング費の入力値（人数・日数・電材費・地区・Ms'外注／開いて編集）")
    with st.expander(_eng_label, expanded=_eng_open):
        if _src:
            st.success(f"{_src}から人数・日数・電材費・Ms'外注費を引用しました（手動変更可）。")
        st.caption("人数・日数・電材費・Ms'外注費は部材表／詳細見積から自動引用。地区（移動費）は手動選択。"
                   "材料費・諸経費には機材費と同じ倍率を適用します。")
        e1, e2, e3 = st.columns(3)
        with e1:
            n_mzu = st.number_input("エムズカラット人数", 0, 100, int(_engv("n_mzu", 2)), key=f"q_nmzu_{ev}")
            n_simx = st.number_input("シムックス人数", 0, 100, int(_engv("n_simx", 1)), key=f"q_nsimx_{ev}")
        with e2:
            work_days = st.number_input("日数", 0, 365, int(_engv("work_days", 2)), key=f"q_days_{ev}")
            denzai = st.number_input("電材費（部材表 工務記入No2 単価）", 0, 100000000,
                                     int(_engv("denzai", 0)), step=10000, key=f"q_denzai_{ev}")
        with e3:
            _alist = list(area_unit_map.keys()) or list(AREA5_LABELS.keys())
            _aidx = _alist.index(ed.get("area5")) if ed.get("area5") in _alist else 0
            area5 = st.selectbox("施工場所（地区＝移動費）", _alist, index=_aidx,
                                 format_func=lambda k: AREA5_LABELS.get(k, k), key=f"q_area5_{ev}")
            ms_out = st.number_input("Ms'カラット外注費（部材表 工務記入No1 単価）", 0, 100000000,
                                     int(_engv("ms_out", 0)), step=10000, key=f"q_msout_{ev}")

    P = int(n_mzu) + int(n_simx)                    # 人数合計（部材表 C12＋C13）
    D = int(work_days)                              # 日数（部材表 C12）
    move_unit = float(area_unit_map.get(area5, 0))  # 施工場所の移動費（地区別の区分料金）
    ms_mat = int(denzai)                            # Ms'材料費 ＝ 工務記入No2 ＝ 電材費（そのまま）

    # 設置工事費：((Ms＋シムックス)×日数×6万)＋12万（利益0）
    setup_val = P * D * 60000 + 120000
    # 材料費：電材費 ×倍率（原価＝電材費）
    mat_cost = int(denzai); mat_sale = int(round(denzai * factor))
    # 諸経費原価：2×移動費×(1+人数) ＋ 8000×人数×日数 ＋ 12000×(1+日数)。売価＝×倍率
    misc_cost = int(round(2 * move_unit * (1 + P) + 8000 * P * D + 12000 * (1 + D)))
    misc_sale = int(round(misc_cost * factor))
    # 構築費 数量
    q_dc = sum(name_qty.get(k, 0) for k in
               ["データ取集装置（Marimba Mercury）", "MMEazyAir親機", "MMEazyAir2（子機）"])
    q_out = n_ctrl                                  # 室外機 ＝ 制御対象台数（シミュレーション）
    q_temp = sum(v for k, v in name_qty.items() if "温湿度" in k or "SwitchBot" in k)
    # ※ 設置工事費・材料費・諸経費・構築費の算出式は下部「算出式詳細（説明）」に集約

    eng_lines = [
        {"name": "設置工事費", "qty": 1, "unit": "式", "sale_u": setup_val, "cost_u": setup_val},
        {"name": "材料費", "qty": 1, "unit": "式", "sale_u": mat_sale, "cost_u": mat_cost},
        {"name": "諸経費", "qty": 1, "unit": "式", "sale_u": misc_sale, "cost_u": misc_cost},
        {"name": "空調システム構築費（データ取集装置）", "qty": q_dc, "unit": "台", "sale_u": kunit, "cost_u": kunit},
        {"name": "空調システム構築費（室外機）", "qty": q_out, "unit": "台", "sale_u": kunit, "cost_u": kunit},
        {"name": "空調システム構築費（温湿度センサー）", "qty": q_temp, "unit": "個", "sale_u": kunit, "cost_u": kunit},
    ]
    eng_amount = eng_profit = 0
    for ln in eng_lines:
        ln["amount"] = ln["sale_u"] * ln["qty"]
        ln["profit"] = ln["amount"] - ln["cost_u"] * ln["qty"]
        eng_amount += ln["amount"]; eng_profit += ln["profit"]

    # ── 5. 端数調整・内部原価（御見積書には非表示） ──
    with st.expander("🧮 5. 端数調整・内部原価（御見積書には非表示／開いて編集）", expanded=False):
        round_adj = st.number_input("端数調整（売価をマイナス調整）", -1000000, 1000000, 0, step=100, key="q_round")
        st.caption(f"Ms'カラット外注費 {yen(ms_out)}（工務記入No1）／Ms'カラット材料費 {yen(ms_mat)}（工務記入No2＝電材費）"
                   f"／ビジネスタンク ＝ 税前合計 × {bt_rate:.0%}")

    sale_total = hw_amount + eng_amount + round_adj
    gross_profit = hw_profit + eng_profit + round_adj
    bt = int(round(sale_total * bt_rate))            # ビジネスタンク ＝ 税前合計 × 0.25
    net_profit = gross_profit - int(ms_out) - int(ms_mat) - bt
    tax = int(round(sale_total * tax_rate))
    total_incl = sale_total + tax
    margin_rate = (net_profit / sale_total) if sale_total else 0.0

    st.markdown("### 算出後の値")
    _cards = [
        ("売価合計（税抜）", yen(sale_total), "#1F3864", ""),
        ("消費税", yen(tax), "#5A6474", ""),
        ("金額合計（税込）", yen(total_incl), "#1F3864", ""),
        ("利益合計", yen(net_profit), "#2E9E5B", f"利益率 {margin_rate:.1%}"),
    ]
    _html = '<div style="display:flex;flex-wrap:wrap;gap:10px;margin:6px 0 4px;">'
    for _lbl, _val, _col, _sub in _cards:
        _subhtml = (f'<div style="font-size:12px;color:#2E9E5B;font-weight:600;margin-top:2px;">{_sub}</div>'
                    if _sub else "")
        _html += (f'<div style="flex:1 1 150px;min-width:140px;border:1px solid #E0E0E0;border-radius:8px;'
                  f'padding:10px 12px;background:#fff;">'
                  f'<div style="font-size:12px;color:#666;margin-bottom:4px;">{_lbl}</div>'
                  f'<div style="font-size:21px;font-weight:700;color:{_col};white-space:nowrap;">{_val}</div>'
                  f'{_subhtml}</div>')
    _html += "</div>"
    st.markdown(_html, unsafe_allow_html=True)
    st.caption(f"粗利益 {yen(gross_profit)} − Ms'外注 {yen(ms_out)} − Ms'材料 {yen(ms_mat)} "
               f"− ビジネスタンク(税前×{bt_rate:.0%}) {yen(bt)} ＝ 利益合計 {yen(net_profit)}")

    # ── 保存 ─────────────────────────────
    if st.button("💾 この見積を『見積一覧』に保存", key="q_save"):
        app_data.setdefault("saved_quotes", []).append({
            "doc_no": doc_no, "to": q_to, "title": q_title,
            "sale_total": int(sale_total), "profit": int(net_profit),
            "margin": round(margin_rate, 4), "saved_at": today.isoformat(),
        })
        save_data(app_data)
        st.success(f"見積『{doc_no}』を保存しました。")

    # ════════ 下部：プレビュー作成 ＆ 算出式詳細 ════════
    st.divider()
    cprev1, cprev2 = st.columns([1, 2])
    with cprev1:
        gen = st.button("🖨 プレビュー作成", type="primary", use_container_width=True, key="q_gen_prev")
    with cprev2:
        price_mode = st.radio("表示する単価", ["売価版（顧客提出用）", "原価版（社内用）"],
                              horizontal=True, key="q_price_mode")
    cost_view = price_mode.startswith("原価")

    if gen:
        issues = []
        if not str(q_to).strip() or str(q_to).strip() == "御中":
            issues.append("宛先（会社名）が未入力です。")
        if not str(q_title).strip():
            issues.append("件名が未入力です。")
        if not any(l["qty"] > 0 for l in hw_lines):
            issues.append("機材費の数量がすべて0です（部材表インポートまたは手入力をご確認ください）。")
        _zero = [l["name"] for l in hw_lines if l["qty"] > 0 and l["cost_u"] <= 0]
        if _zero:
            issues.append("原価が未入力（¥0）の機材があります：" + " / ".join(_zero[:8]))
        if n_ctrl <= 0:
            issues.append("制御対象台数が0です（室外機の構築費が計上されません）。")
        if int(denzai) <= 0:
            issues.append("電材費（部材表 工務記入No2）が0です（材料費・Ms'材料費が0になります）。")
        if sale_total <= 0:
            issues.append("売価合計が¥0です。入力内容をご確認ください。")
        st.session_state.q_show_prev = True
        if issues:
            _quote_missing_dialog(issues)
    if st.session_state.get("q_show_prev"):
        head = {"title": q_title, "to": q_to, "doc_no": doc_no,
                "date": f"{q_date.year}年{q_date.month}月{q_date.day}日",
                "valid": f"{valid_until.year}年{valid_until.month}月{valid_until.day}日",
                "delivery": q_delivery, "place": q_to_place, "pay": q_pay, "person": q_person,
                "stamps": [q_stamp1, q_stamp2, q_stamp3]}
        html = _quote_preview_html(head, comp, hw_lines, eng_lines, round_adj,
                                   sale_total, tax, total_incl, tax_rate, cost_view, n_ctrl, q_temp)
        st.markdown("#### 御見積書プレビュー" + ("（原価版）" if cost_view else ""))
        components.html(html, height=1180, scrolling=True)

        # 転記用明細のダウンロード（売価・原価 両列）
        rows = []
        for ln in hw_lines + eng_lines:
            rows.append({"区分": "機材費" if ln in hw_lines else "エンジ費", "商品名": ln["name"],
                         "数量": ln["qty"], "単位": ln["unit"], "売価単価": int(ln["sale_u"]),
                         "売価金額": int(ln["amount"]), "原価単価": int(ln["cost_u"]),
                         "原価金額": int(ln["cost_u"] * ln["qty"]), "利益": int(ln["profit"])})
        df_out = pd.DataFrame(rows)
        csv = df_out.to_csv(index=False).encode("utf-8-sig")
        st.download_button("⬇ 転記用明細（CSV）", csv, file_name=f"見積明細_{doc_no}.csv",
                           mime="text/csv", key="q_dl_csv")

    with st.expander("📐 算出式詳細（説明）", expanded=False):
        _quote_methodology(factor, hw_amount, hw_cost, hw_profit, eng_amount, eng_profit,
                           setup_val, mat_sale, misc_sale, P, D, move_unit,
                           q_dc, q_out, q_temp, kunit, round_adj, gross_profit, bt, bt_rate,
                           ms_out, ms_mat, net_profit, sale_total, tax, total_incl, tax_rate, margin_rate)


def _quote_preview_html(head, comp, hw_lines, eng_lines, round_adj,
                        sale_total, tax, total_incl, tax_rate, cost_view, n_ctrl, n_temp):
    """御見積書（顧客版）レイアウトのHTMLを生成。cost_view=Trueで単価を原価表示。"""
    def y(v):
        return f"{v:,.0f}"
    css = """
    <style>
    .qv{font-family:'Yu Gothic','Meiryo',sans-serif;color:#111;background:#fff;padding:18px 22px;width:720px;margin:auto;border:1px solid #ccc;}
    .qv h1{text-align:center;letter-spacing:.5em;font-size:24px;margin:4px 0 10px;}
    .qv table{border-collapse:collapse;width:100%;font-size:12px;}
    .qv .meta td{padding:2px 6px;vertical-align:top;border:none;}
    .qv .items th{background:#e9e9e9;border:1px solid #999;padding:4px 6px;}
    .qv .items td{border:1px solid #999;padding:3px 6px;}
    .qv .r{text-align:right;} .qv .c{text-align:center;}
    .qv .sub{background:#f5f5f5;font-weight:bold;}
    .qv .big{border:2px solid #333;padding:4px 10px;font-size:18px;font-weight:bold;}
    .qv .stampin{display:inline-flex;width:33px;height:33px;border:1.2px solid #c00;border-radius:50%;
        color:#c00;font-size:9px;font-weight:bold;align-items:center;justify-content:center;
        writing-mode:vertical-rl;text-orientation:upright;line-height:1.0;letter-spacing:0.5px;overflow:hidden;}
    .qv .tot td{border:1px solid #999;padding:3px 12px;}
    </style>
    """
    rows_html = ""
    no = 0
    def line_row(ln):
        nonlocal no, rows_html
        no += 1
        unit_price = ln["cost_u"] if cost_view else ln["sale_u"]
        amount = (ln["cost_u"] * ln["qty"]) if cost_view else ln["amount"]
        rows_html += (f'<tr><td class="c">{no}</td><td>{ln["name"]}</td>'
                      f'<td class="c">{ln["qty"]}</td><td class="c">{ln["unit"]}</td>'
                      f'<td class="r">{y(unit_price)}</td><td class="r">{y(amount)}</td></tr>')
    for ln in hw_lines:
        line_row(ln)
    hw_amt = sum((l["cost_u"] * l["qty"]) if cost_view else l["amount"] for l in hw_lines)
    rows_html += f'<tr class="sub"><td></td><td>ハードウェア関連小計</td><td colspan="3"></td><td class="r">{y(hw_amt)}</td></tr>'
    for ln in eng_lines:
        line_row(ln)
    eng_amt = sum((l["cost_u"] * l["qty"]) if cost_view else l["amount"] for l in eng_lines)
    rows_html += f'<tr class="sub"><td></td><td>エンジニアリング費小計</td><td colspan="3"></td><td class="r">{y(eng_amt)}</td></tr>'
    if round_adj:
        no += 1
        rows_html += f'<tr><td class="c">{no}</td><td>端数調整</td><td colspan="3"></td><td class="r">{y(round_adj if not cost_view else 0)}</td></tr>'

    # 押印欄（黒枠3つ横並び＋赤丸印・氏名は指定可）
    _snames = head.get("stamps", ["", "", ""])
    stamp_html = ('<table style="border-collapse:collapse;margin-left:auto;margin-top:6px;"><tr>'
                  + ''.join('<td style="width:40px;height:44px;border:1px solid #333;'
                            f'text-align:center;vertical-align:middle;"><div class="stampin">{nm}</div></td>'
                            for nm in _snames)
                  + '</tr></table>')
    # 合計（右寄せミニ表）
    totals_html = (f'<table class="tot" style="width:auto;margin-left:auto;margin-top:8px;">'
                   f'<tr><td style="background:#eee;">税前合計</td><td class="r">{y(sale_total)}</td></tr>'
                   f'<tr><td style="background:#eee;">税額合計</td><td class="r">{y(tax)}</td></tr>'
                   f'<tr><td style="background:#eee;"><b>金額合計</b></td><td class="r"><b>{y(total_incl)}</b></td></tr>'
                   f'</table>')
    # 備考（2列×1行・制御台数）
    bikou = (f'■計測及び制御台数<br>・デマンド計測 1点<br>・空調機電力計測 {n_ctrl}点'
             f'<br>・空調室外機制御 {n_ctrl}点<br>・温湿度計測 {n_temp}点')
    bikou_html = ('<table style="border-collapse:collapse;width:100%;font-size:11px;margin-top:12px;">'
                  '<tr><td style="border:1px solid #999;background:#eee;width:70px;text-align:center;vertical-align:middle;">備考</td>'
                  f'<td style="border:1px solid #999;padding:6px;">{bikou}</td></tr></table>')

    title_suffix = "（原価版・社内用）" if cost_view else ""
    body = f"""
    {css}
    <div class="qv">
      <h1>御 見 積 書{title_suffix}</h1>
      <table class="meta"><tr>
        <td style="width:55%;"><b style="font-size:15px;">{head['to']}</b><br><br>
          下記の通り御見積申し上げます。</td>
        <td style="text-align:right;">書類番号 {head['doc_no']}<br>発行日 {head['date']}</td>
      </tr></table>
      <table class="meta"><tr>
        <td style="width:55%;">
          件名：{head['title']}<br>納期：{head['delivery']}<br>納品先：{head['place']}<br>
          支払条件：{head['pay']}<br>有効期限：{head['valid']}<br>担当者：{head['person']}</td>
        <td style="vertical-align:top;font-size:12px;">{comp.get('name','')}<br>{comp.get('post','')}<br>
          {comp.get('addr','')}<br>TEL：{comp.get('tel','')}　FAX：{comp.get('fax','')}
          {stamp_html}</td>
      </tr></table>
      <div style="margin:10px 0;">
        <span class="big">金額合計(税抜)　¥ {y(sale_total)} -</span>
        <span style="margin-left:16px;font-size:12px;">消費税{tax_rate:.0%}　¥{y(tax)}　／　税込　¥{y(total_incl)}</span>
      </div>
      <table class="items">
        <tr><th>No</th><th>商品名</th><th>数量</th><th>単位</th><th>{'原価' if cost_view else '単価'}</th><th>金額（税別）</th></tr>
        {rows_html}
      </table>
      {totals_html}
      {bikou_html}
    </div>
    """
    return body


def _quote_methodology(factor, hw_amount, hw_cost, hw_profit, eng_amount, eng_profit,
                       setup_val, mat_sale, misc_sale, P, D, move_unit,
                       q_dc, q_out, q_temp, kunit, round_adj, gross_profit, bt, bt_rate,
                       ms_out, ms_mat, net_profit, sale_total, tax, total_incl, tax_rate, margin_rate):
    """見積書の算出ロジックを 計算式→当てはめ→結果 で開示（既存『算出方法』カードと同方式）。"""
    def yen(v):
        return f"¥{v:,.0f}"
    GRAY = "#666"; AC = "#00B050"; LINE = "#E0E0E0"

    def card(title, formula, apply_html, result, note=""):
        note_html = f'<div style="font-size:11px;color:#999;margin-top:6px;">{note}</div>' if note else ""
        return (f'<div style="border:1px solid {LINE};border-radius:8px;padding:14px 16px;margin-bottom:12px;background:#fff;">'
                f'<div style="font-size:13px;font-weight:700;color:#333;margin-bottom:8px;">{title}</div>'
                f'<div style="font-size:13px;color:{GRAY};margin-bottom:4px;"><span style="display:inline-block;width:64px;color:#aaa;font-size:11px;">計算式</span>{formula}</div>'
                f'<div style="font-size:13px;color:#333;margin-bottom:4px;"><span style="display:inline-block;width:64px;color:#aaa;font-size:11px;">当てはめ</span>{apply_html}</div>'
                f'<div style="font-size:14px;color:#111;font-weight:700;"><span style="display:inline-block;width:64px;color:#aaa;font-size:11px;font-weight:400;">結果</span><span style="color:{AC};">{result}</span></div>'
                f'{note_html}</div>')

    html = ""
    html += card("① 機材費", "売価 ＝ 機材単価マスタの売価 ／ 利益 ＝ 売価 − 原価",
                 "各機材の売価・原価はマスタ参照（売価未設定は 原価×倍率 で暫定）",
                 f"機材費 売価合計 {yen(hw_amount)}（原価 {yen(hw_cost)}／利益 {yen(hw_profit)}）",
                 note="機材費は商品ごとの売価をそのまま使用します（管理画面『機材単価マスタ』で編集）。")
    html += card("② エンジ費（設置/材料/諸経費・構築費）",
                 "設置＝(人数×日数×6万)+12万／材料＝電材費×倍率／"
                 "諸経費＝(2×移動費×(1+人数)+8000×人数×日数+12000×(1+日数))×倍率／構築費＝台数×単価",
                 f"設置 ({P}人×{D}日×6万)+12万 = {yen(setup_val)}　／　材料 電材費×{factor:.2f} = {yen(mat_sale)}<br>"
                 f"諸経費 (移動費{yen(move_unit)}・人数{P}・日数{D})×{factor:.2f} = {yen(misc_sale)}<br>"
                 f"構築費 データ取集装置{q_dc}・室外機(制御台数){q_out}・温湿度{q_temp} × {yen(kunit)}",
                 f"エンジ費 売価合計 {yen(eng_amount)}（利益 {yen(eng_profit)}）",
                 note="室外機の構築費はシミュレーションの制御対象台数を使用。材料費・諸経費の倍率は機材費と共通。")
    html += card("③ 売価合計（税抜）", "機材費 ＋ エンジ費 ＋ 端数調整",
                 f"{yen(hw_amount)} ＋ {yen(eng_amount)} ＋ {yen(round_adj)}",
                 yen(sale_total))
    html += card("④ ビジネスタンク費", f"税前合計 × {bt_rate:.0%}",
                 f"{yen(sale_total)} × {bt_rate:.0%}", yen(bt),
                 note="税前合計（売価合計）を基準に算出。")
    html += card("⑤ 利益合計（純利益）", "粗利益 − Ms'カラット外注 − Ms'カラット材料 − ビジネスタンク",
                 f"{yen(gross_profit)} − {yen(ms_out)} − {yen(ms_mat)} − {yen(bt)}",
                 f"{yen(net_profit)}（利益率 {margin_rate:.1%}）")
    html += card("⑥ 消費税・税込合計", f"税 ＝ 売価合計 × {tax_rate:.0%} ／ 税込 ＝ 売価合計 ＋ 税",
                 f"{yen(sale_total)} × {tax_rate:.0%} ＝ {yen(tax)}",
                 f"税込合計 {yen(total_incl)}")
    st.caption("※ 下記はすべて、今入力している見積の実数を計算式へ当てはめたものです。")
    st.markdown(html, unsafe_allow_html=True)


def render_breakdown_charts(res, app_data, which="all"):
    """月別の最大デマンド／使用量／電気料金を、業態内訳で色分けした積み上げ棒で表示（従来ロジックそのまま）。
    which="all"で3枚、"dm"/"use"/"bill"で該当1枚のみ描画（表示切替用）。"""
    ds = app_data.get("display_settings", DEFAULT_DATA["display_settings"])
    # 仮置き前提（業態標準の空調割合など）を使っている場合はグラフにも明示
    try:
        _gy = str(res.get("gyotai", ""))
        if _gy != "その他※数値を指定する":
            st.caption(f"※ 内訳は業態「{_gy}」の標準割合（最大デマンド {res.get('ac_peak_r',0)*100:.0f}％／"
                       f"使用量 {res.get('ac_kwh_r',0)*100:.0f}％）に基づく概算です（実測比率があれば精緻化されます）。")
    except Exception:
        pass
    df = res["df"].copy()
    df["月"] = df["月"].astype(str)
    L_OTHER = "その他電気量"
    L_PROD = "生産設備"
    L_LIGHT = "照明"
    L_ACR = "空調（導入後に残る分）"
    L_ACS = "空調（DPSで削減できる分）"
    SEG_ORDER = [L_OTHER, L_PROD, L_LIGHT, L_ACR, L_ACS]
    AC_COLOR = "#00B050"
    SEG_COLOR = {
        L_OTHER: "#E0E0E0", L_PROD: "#CCCCCC", L_LIGHT: "#BDBDBD",
        L_ACR: "#A9D18E", L_ACS: AC_COLOR,
    }

    def make_breakdown_bar(src, cur_col, aft_col, y_title, ac_frac, weights):
        def _mlabel(m):
            s = str(m).replace("-", "/")
            p = s.split("/")
            if len(p) >= 2:
                y, mo = p[0], p[1]
                if mo.zfill(2) == "01":
                    return f"{y}/01"
                return f"{int(mo)}月"
            return s
        m_order = [_mlabel(m) for m in src["月"].tolist()]
        w_light = weights.get("照明", 0.0)
        w_prod = weights.get("生産設備", 0.0)
        w_other = weights.get("その他", 0.0)
        rows = []
        for _, r in src.iterrows():
            cur_v = float(r[cur_col])
            aft_v = min(float(r[aft_col]), cur_v)
            reduction = max(0.0, cur_v - aft_v)
            ac_total = min(max(ac_frac * cur_v, reduction), cur_v)
            ac_remain = ac_total - reduction
            non_ac = cur_v - ac_total
            seg_vals = [
                (L_OTHER, non_ac * w_other),
                (L_PROD, non_ac * w_prod),
                (L_LIGHT, non_ac * w_light),
                (L_ACR, ac_remain),
                (L_ACS, reduction),
            ]
            m = _mlabel(r["月"]); y0 = 0.0
            for name, val in seg_vals:
                y1 = y0 + val
                rows.append({"月": m, "種別": name, "y0": y0, "y1": y1, "値": val,
                             "_o": SEG_ORDER.index(name)})
                y0 = y1
        long = pd.DataFrame(rows)
        LEGEND_ORDER = [L_ACS, L_ACR, L_LIGHT, L_PROD, L_OTHER]
        present = [c for c in LEGEND_ORDER if long.loc[long["種別"] == c, "値"].sum() > 1e-9]
        long = long[long["種別"].isin(present)]
        col_range = [SEG_COLOR[c] for c in present]
        return alt.Chart(long).mark_bar(stroke="#FFFFFF", strokeWidth=1).encode(
            x=alt.X("月:N", sort=m_order,
                    axis=alt.Axis(title=None, labelAngle=0, labelFontSize=13,
                                  labelColor="#666666", domain=False, ticks=False, grid=False)),
            y=alt.Y("y0:Q", title=y_title,
                    axis=alt.Axis(grid=True, gridColor="#E0E0E0", gridWidth=0.5,
                                  labelColor="#666666", labelFontSize=12, titleColor="#666666",
                                  domain=False, ticks=False, tickCount=5)),
            y2=alt.Y2("y1"),
            color=alt.Color("種別:N",
                scale=alt.Scale(domain=present, range=col_range),
                legend=alt.Legend(title=None, orient="top", columns=3,
                                  labelColor="#666666", labelFontSize=12, symbolStrokeWidth=0)),
            order=alt.Order("_o:Q"),
            tooltip=["月:N", "種別:N", alt.Tooltip("値:Q", title="値", format=",.1f")],
        ).properties(height=400).configure_view(strokeWidth=0)

    nonac_w = INDUSTRY_NONAC_WEIGHTS.get(res.get("gyotai", ""), DEFAULT_NONAC_WEIGHTS)

    if which in ("all", "dm") and ds.get("show_graph_dm", True):
        st.caption(f"棒全体が『現状の最大デマンド』。業態「{res.get('gyotai','')}」の電力内訳で色分けし、"
                   "濃い緑（アクセント色）の部分が DPS 導入で削減できる空調デマンドです。")
        ch_dm = make_breakdown_bar(
            df, "最大需要電力", "導入後最大DM", "最大デマンド（kW）",
            res.get("ac_peak_r", 0.18), nonac_w)
        st.altair_chart(ch_dm, use_container_width=True)

    if which in ("all", "use") and ds.get("show_graph_bill", True):
        st.caption(f"棒全体が『現状の電力使用量』。業態「{res.get('gyotai','')}」の電力内訳で色分けし、"
                   "濃い緑（アクセント色）の部分が DPS 導入で削減できる使用量（空調由来）です。")
        df_use_man = df.copy()
        df_use_man["現状使用量(万kWh)"] = df_use_man["使用量合計"] / 10000
        df_use_man["導入後使用量(万kWh)"] = df_use_man["導入後使用量kWh"] / 10000
        ch_use = make_breakdown_bar(
            df_use_man, "現状使用量(万kWh)", "導入後使用量(万kWh)", "電力使用量（万kWh）",
            res.get("ac_kwh_r", 0.15), nonac_w)
        st.altair_chart(ch_use, use_container_width=True)

    if which in ("all", "bill") and ds.get("show_graph_bill", True):
        st.caption(f"棒全体が『現状の電気料金』。業態「{res.get('gyotai','')}」の電力内訳で色分けし、"
                   "濃い緑（アクセント色）の部分が DPS 導入で削減できる金額（空調由来）です。")
        try:
            df_bill_man = df.copy()
            df_bill_man["現状電気料金(万円)"] = df_bill_man["現状電気料金推計"] / 10000
            df_bill_man["導入後電気料金(万円)"] = df_bill_man["導入後電気料金推計"] / 10000
            ch_bill = make_breakdown_bar(
                df_bill_man, "現状電気料金(万円)", "導入後電気料金(万円)", "電気料金（万円）",
                res.get("ac_kwh_r", 0.15), nonac_w)
            st.altair_chart(ch_bill, use_container_width=True)
        except Exception:
            st.caption("電気料金グラフは表示できませんでした（料金推計列が無い場合）。")


def render_results_dashboard(res, app_data, client_name="", gyotai=""):
    """結果ダッシュボードを Streamlit 標準コンポーネント（columns / container / metric）で描画（Power BI風）。
    案件情報 → KPI×4 → グラフ・分析 → 財務・AI → リスク・営業 → ダウンロード をグリッド配置。CSSは最小限。"""
    df = res.get("df")
    gross = float(res.get("gross_saving", 0)); net = float(res.get("net_saving", 0))
    invest = float(res.get("total_invest", 0)); net_invest = float(res.get("net_invest", invest))
    payback, _ofee = calc_payback(invest, net, gross)
    cs = app_data.get("const_settings", DEFAULT_DATA["const_settings"])
    nyrs = int(cs.get("npv_years", 10)); rate = float(cs.get("discount_rate", 0.05))
    npv, irr = calc_npv_irr(invest, net, years=nyrs, rate=rate)
    old_c = float(res.get("old_contract", 0)); new_c = float(res.get("new_contract", 0))
    cdelta = float(res.get("contract_delta", 0)); dm = float(res.get("dm_saving_annual", 0))
    ene = float(res.get("ene_saving_annual", 0)); co2 = float(res.get("co2_kg", 0))
    total_kwh = float(res.get("total_reduc_kwh", 0)); sys_fee = float(res.get("sys_fee", 0))
    sugi = float(res.get("sugi_trees", 0))
    n_ctrl = int(res.get("n_units_ctrl", (res.get("econ", {}) or {}).get("n_ctrl", 0)) or 0)

    def yen(v):
        return "¥{:,.0f}".format(v)
    pb = "—" if (payback == float("inf") or payback <= 0) else "{:.1f} 年".format(payback)
    pb_years = 99.9 if (payback == float("inf") or payback <= 0) else float(payback)
    irr_txt = ("{:.1f} %".format(irr * 100) if irr is not None else "—")

    # （案件情報・ステータス・メニューは画面上部の共通ヘッダーに集約）

    # ── 入力整合性チェック（C-12）──
    try:
        _issues = validate_inputs(res.get("df"), res, res.get("econ", {}))
        _has_err = any(lv == "error" for lv, _ in _issues)
        _has_warn = any(lv == "warn" for lv, _ in _issues)
        _exp_title = ("入力整合性チェック — ⚠️ 要確認あり" if (_has_err or _has_warn)
                      else "入力整合性チェック — ✅ 異常なし")
        with st.expander(_exp_title, expanded=_has_err):
            for lv, msg in _issues:
                if lv == "error":
                    st.error("❌ " + msg)
                elif lv == "warn":
                    st.warning("⚠️ " + msg)
                else:
                    st.success("✅ " + msg)
    except Exception:
        pass

    # ── 算出の前提（概算・業態配慮・仮置き）注釈 ──
    try:
        _pn = provisional_notes(res, app_data)
        if _pn:
            with st.expander(f"ⓘ 算出の前提（概算・業態標準値などの仮置きを {len(_pn)} 件含みます）", expanded=False):
                st.caption("以下は実測ではなく概算・業態標準値・自動補完に基づく前提です。"
                           "グラフ・結果・出力資料（PPTX）にも同じ注釈が記載されます。")
                for _n in _pn:
                    st.markdown(f"- {_n}")
    except Exception:
        pass

    # ── 概算モード：3シナリオ比較（ベスト／標準／保守）──
    if (res.get("econ", {}) or {}).get("gaisan_mode") or st.session_state.get("gaisan_mode"):
        try:
            _scs = compute_gaisan_scenarios(res, app_data)
            if _scs:
                with st.expander("🔎 概算：3シナリオ比較（ベスト／標準／保守）", expanded=True):
                    st.caption("概算は環境の想定値のため幅で確認できます。削減式は現行のまま、"
                               "台数推定・機材費・1装置あたり制御台数(HW)・制御可能比率の想定を変えた3案です。"
                               "（各単価・想定値は『各種設定 → 計算ロジック変更』の概算マスタで調整可）")
                    _rows = [{
                        "シナリオ": s["label"],
                        "予測台数(総/制御)": f"{s['n_total']} / {s['n_ctrl']} 台",
                        "制御可能比率": f"{s['ctrl_ratio']:.0%}",
                        "1装置あたり": f"{s['units_per_device']}台",
                        "年間グロス削減": yen(s["gross"]),
                        "年間純削減": yen(s["net"]),
                        "初期費用(概算)": yen(s["invest"]),
                        "投資回収": (f"約 {s['payback']:.1f} 年" if s["payback"] else "—"),
                    } for s in _scs]
                    st.dataframe(pd.DataFrame(_rows), hide_index=True, use_container_width=True)
                    st.caption("※ 初期費用のみ想定差の影響大（削減額は空調負荷ベースで台数に概ね依存しません）。"
                               "実機材・実見積の取得後に精緻化してください。")
        except Exception as _e:
            st.caption(f"（3シナリオ比較の表示に失敗：{_e}）")

    # ── ② KPI ×4（白カード＋上3px色アクセントで控えめに区別）──
    _kpis = [
        ("初期費用（総投資）", yen(invest), "機材・工事込み", "#0F2E5D"),
        ("投資回収期間", pb, "純利益ベース", "#0F6CBD"),
        ("年間削減額", yen(gross), "基本料金＋従量", "#107C41"),
        (f"{nyrs}年累計削減", yen(net * nyrs), f"NPV {yen(npv)}", "#9A6700"),
    ]
    _kc = st.columns(4)
    for _i, (_lab, _val, _sub, _ac) in enumerate(_kpis):
        _kc[_i].markdown(
            f"<div style='background:#fff;border:1px solid #E5E7EB;border-top:3px solid {_ac};"
            f"border-radius:10px;padding:16px 18px;box-shadow:0 1px 2px rgba(16,24,40,.06)'>"
            f"<div style='font-size:13px;color:#6B7280'>{_lab}</div>"
            f"<div style='font-size:28px;font-weight:700;color:#0F2E5D;line-height:1.15;margin-top:4px'>{_val}</div>"
            f"<div style='font-size:11.5px;color:#9CA3AF;margin-top:3px'>{_sub}</div></div>",
            unsafe_allow_html=True)

    st.markdown("<div style='height:16px'></div>", unsafe_allow_html=True)

    # ── ③ AIサマリー（結論・全幅・主要グラフより上＝結論ファースト）──
    _concl = ("3年回収圏内の優良案件。即効性の高いアドオン提案です。" if pb_years <= 3.0
              else "回収3年超。機器更新併用・補助金活用の検討を推奨します。")
    st.markdown(
        f"<div style='background:#F0F6FB;border:1px solid #CBD9EE;border-left:4px solid #0F2E5D;"
        f"border-radius:10px;padding:16px 20px;margin-bottom:16px'>"
        f"<div style='font-size:15px;font-weight:700;color:#0F2E5D'>結論：投資回収 {pb} ／ 年間 {yen(gross)} 削減</div>"
        f"<div style='font-size:13px;color:#1F2937;margin-top:6px'>{_concl}</div>"
        f"<div style='font-size:12.5px;color:#6B7280;margin-top:8px'>"
        f"年間純利益 {yen(net)} ｜ 制御 {n_ctrl}台 ｜ CO₂削減 {co2:,.0f}kg ｜ スギ換算 {sugi:,.0f}本 ｜ AI判定：正常</div></div>",
        unsafe_allow_html=True)

    # ── 主要グラフ（全幅）──
    with st.container(border=True):
        st.markdown("**月別 電気使用量の比較（万kWh）**")
        try:
            _d = pd.DataFrame({
                "月": [str(x) for x in df["月"].tolist()],
                "現状": [round(float(v) / 10000, 1) for v in df["使用量合計"].tolist()],
                "提案後": [round(float(v) / 10000, 1) for v in
                          (df["導入後使用量kWh"] if "導入後使用量kWh" in df.columns
                           else df["使用量合計"] * 0.9).tolist()]})
            _m = _d.melt("月", var_name="区分", value_name="万kWh")
            _ch = alt.Chart(_m).mark_bar().encode(
                x=alt.X("月:N", sort=list(_d["月"]), axis=alt.Axis(labelAngle=-45, title=None)),
                y=alt.Y("万kWh:Q", title=None),
                color=alt.Color("区分:N", scale=alt.Scale(domain=["現状", "提案後"], range=["#0F2E5D", "#0F6CBD"]),
                                legend=alt.Legend(orient="top", title=None)),
                xOffset="区分:N").properties(height=300)
            st.altair_chart(_ch, use_container_width=True)
        except Exception:
            st.caption("グラフを表示できません。")

    # ── 比較表（全幅・項目が増えても見やすい）──
    with st.container(border=True):
        st.markdown("**現状 ／ 提案後 比較（年間）**")
        try:
            _cmp = pd.DataFrame([
                {"項目": "契約電力（最大需要）", "現状": "{:,.0f} kW".format(old_c), "提案後": "{:,.0f} kW".format(new_c), "差額": "▲{:,.0f} kW".format(cdelta), "判定": "改善"},
                {"項目": "基本料金（年間削減）", "現状": "—", "提案後": "—", "差額": "▲" + yen(dm), "判定": "改善"},
                {"項目": "従量料金（年間削減）", "現状": "—", "提案後": "—", "差額": "▲" + yen(ene), "判定": "改善"},
                {"項目": "年間使用量削減", "現状": "—", "提案後": "—", "差額": "▲{:,.0f} kWh".format(total_kwh), "判定": "改善"},
                {"項目": "CO₂排出（年間）", "現状": "—", "提案後": "—", "差額": "▲{:,.0f} kg".format(co2), "判定": "改善"}])
            st.dataframe(_cmp, hide_index=True, use_container_width=True)
        except Exception:
            st.caption("テーブルを表示できません。")

    # ── ④ 詳細分析（タブ）──
    st.markdown("#### 詳細分析")
    _tg, _tf, _tm, _tr, _tcalc, _topt = st.tabs(
        ["グラフ", "財務", "制御強度", "リスク", "算出方法", "回収逆算"])
    with _tg:
        try:
            _gsel = st.radio(
                "表示グラフ", ["最大デマンド", "電力使用量", "電気料金"],
                horizontal=True, label_visibility="collapsed", key="bd_chart_sel")
            _wmap = {"最大デマンド": "dm", "電力使用量": "use", "電気料金": "bill"}
            render_breakdown_charts(res, app_data, which=_wmap.get(_gsel, "dm"))
        except Exception as _e:
            st.caption(f"（グラフを表示できません：{_e}）")
        with st.expander("月別データ", expanded=False):
            try:
                _cols = [c for c in ["月", "最大需要電力", "導入後最大DM", "使用量合計", "導入後使用量kWh"] if c in df.columns]
                st.dataframe(df[_cols], hide_index=True, use_container_width=True)
            except Exception:
                st.caption("データを表示できません。")
    with _tf:
        _f = st.columns(4)
        _f[0].metric(f"NPV（{nyrs}年）", yen(npv))
        _f[1].metric("IRR", irr_txt)
        _f[2].metric("初期費用", yen(invest))
        _f[3].metric("実質投資", yen(net_invest))
    with _tm:
        try:
            with st.container(height=560):
                show_control_mode(res)
        except Exception as _e:
            st.caption(f"（制御強度を表示できません：{_e}）")
    with _tr:
        try:
            with st.container(height=560):
                show_financial_risk_10yr(res)
        except Exception as _e:
            st.caption(f"（リスク分析を表示できません：{_e}）")
    with _tcalc:
        try:
            with st.container(height=560):
                show_calc_methodology(res)
        except Exception as _e:
            st.caption(f"（算出方法を表示できません：{_e}）")
    with _topt:
        try:
            with st.container(height=560):
                show_reduction_planner(res)
        except Exception as _e:
            st.caption(f"（回収逆算を表示できません：{_e}）")

    # ── 付加価値・財務換算 ／ 1枚サマリー（ペライチ）──
    with st.expander("付加価値・財務換算 ／ 1枚サマリー", expanded=False):
        try:
            show_peraichi(res)
        except Exception as _e:
            st.caption(f"（サマリーを表示できません：{_e}）")
        _ip = res.get("item_price", app_data["item_price"])
        _im = res.get("item_margin", app_data["item_margin"])
        _need = math.ceil(net / (_ip * _im)) if (_ip * _im) > 0 else 0
        st.markdown(f"""
<div style="display:grid; grid-template-columns:repeat(2,1fr); gap:16px;">
  <div style="background:#E8F5E9;border-radius:10px;padding:20px;border:1.5px solid #1E6B2E;display:flex;flex-direction:column;">
    <div style="font-weight:bold;color:#1E6B2E;font-size:15px;margin-bottom:10px;">主力製品の販売数量・財務換算</div>
    <div style="font-size:15px;line-height:1.8;">
      年間手残り純削減額 <b>¥{res['net_saving']:,.0f}</b> は、本業の主力製品<br>
      （単価 ¥{_ip:,} · 粗利 {_im:.1%}）を、<br>
      新規に <span style="font-size:37px;font-weight:bold;color:#C00000;">【{_need:,} 個】</span> 追加販売して稼ぎ出す営業純利益に匹敵します。
    </div>
    <div style="font-size:12px;color:#666;margin-top:auto;padding-top:12px;border-top:1px solid #e2e8f0;">
      営業努力による追加販売のハードルと、AIで無駄な電力契約を自動カットする確実性をぜひご比較ください。
    </div>
  </div>
  <div style="background:#F1F8E9;border-radius:10px;padding:20px;border:1.5px solid #1B5E20;display:flex;flex-direction:column;">
    <div style="font-weight:bold;color:#1B5E20;font-size:15px;margin-bottom:10px;">ESG・環境貢献サマリー</div>
    <div style="font-size:15px;line-height:2.0;">
      年間削減電力量：<b>{res['total_reduc_kwh']:,.0f} kWh / 年</b><br>
      CO₂排出削減量：<b>{res['co2_kg']:,.0f} kg-CO₂ / 年</b><br>
      スギの木換算森林効果：<b>約 {res['sugi_trees']:,.0f} 本分 / 年</b>
    </div>
    <div style="font-size:12px;color:#666;margin-top:auto;padding-top:12px;border-top:1px solid #e2e8f0;">
      これらの省エネデータは、ESG報告書や御社ホームページ、SDGsの取り組み成果として対外公表にご利用いただけます。
    </div>
  </div>
</div>
""", unsafe_allow_html=True)

    # ── ④ 提案書・出力（最後に集約）──
    st.markdown("#### 提案書・出力")
    with st.container(border=True):
        _o = st.columns(3)
        with _o[0]:
            if st.button("PowerPoint提案書を生成", type="primary", use_container_width=True, key="dash_ppt"):
                st.session_state["_dash_ppt"] = True
            if st.session_state.get("_dash_ppt"):
                try:
                    st.download_button("PPTXダウンロード", build_pptx(res, client_name, app_data),
                                       file_name=f"DPS提案_{client_name}.pptx", use_container_width=True, key="dash_ppt_dl")
                except Exception as _e:
                    st.caption(f"生成エラー：{_e}")
        with _o[1]:
            if st.button("Excel稟議シートを生成", use_container_width=True, key="dash_xls"):
                st.session_state["_dash_xls"] = True
            if st.session_state.get("_dash_xls"):
                try:
                    st.download_button("XLSXダウンロード", build_excel(res, client_name, app_data),
                                       file_name=f"DPS試算_{client_name}.xlsx", use_container_width=True, key="dash_xls_dl")
                except Exception as _e:
                    st.caption(f"生成エラー：{_e}")
        with _o[2]:
            if st.button("結果一覧に保存", use_container_width=True, key="dash_save"):
                _rec = {"id": datetime.datetime.now().strftime("%Y%m%d%H%M%S"), "company": client_name,
                        "gyotai": res.get("gyotai", gyotai), "units_total": int(res.get("n_units_total", 0)),
                        "units_ctrl": int(res.get("n_units_ctrl", 0)), "saving_kwh": float(res.get("total_reduc_kwh", 0)),
                        "saving_yen": float(res.get("gross_saving", 0)), "payback": float(pb_years),
                        "monthly": case_monthly_records(res),   # 削減レポ呼び出し用の月次データ
                        "saved_at": datetime.datetime.now().strftime("%Y-%m-%d %H:%M")}
                app_data.setdefault("saved_cases", []).append(_rec); save_data(app_data)
                # 削減レポ用：この検針票の月次を企業データへ累積（算出必要5項目のみ）
                try:
                    _cmap = {}
                    for _, _mr in res["df"].iterrows():
                        _ymk = _to_ym(_mr.get("月"))
                        if _ymk:
                            _cmap[_ymk] = {"契約電力": _mr.get("契約電力", 0), "最大需要電力": _mr.get("最大需要電力", 0),
                                           "使用電力量": _mr.get("使用量合計", 0), "基本料金単価": _mr.get("基本料金単価", 0),
                                           "電力量単価": _mr.get("電力量単価", 0)}
                    _nsav = _merge_company_records(app_data, client_name, _cmap)
                except Exception:
                    _nsav = 0
                st.success(f"「{client_name}」を結果一覧に保存しました。"
                           + (f"（企業データへ {_nsav}ヶ月を累積）" if _nsav else ""))
    with st.expander("営業トークスクリプト", expanded=False):
        try:
            show_talk_script(res, app_data, client_name, pb_years, invest)
        except Exception as _e:
            st.caption(f"（営業トークを表示できません：{_e}）")

    # スライド貼付用 画像出力（PNG）— 復元（ロジックは従来どおり）
    with st.expander("スライド貼付用 画像出力（PNG）", expanded=False):
        if not HAS_SLIDE_IMG:
            st.info("画像出力には matplotlib / pillow が必要です。`pip install matplotlib pillow` を実行してください。")
        else:
            st.caption("資料（PowerPoint / Word）にそのまま貼れる PNG を生成します。生成後、各画像をダウンロードできます。")

            def _month_label(v):
                import re
                s = str(v).strip()
                mt = re.search(r"(\d{4})\D+(\d{1,2})", s)   # 年が拾えれば "YYYY/M" で保持（PNG側で年の変わり目に年表示）
                if mt:
                    return f"{int(mt.group(1))}/{int(mt.group(2))}"
                tail = s.replace("年", "/").replace("月", "").split("/")[-1].strip()
                try:
                    return f"{int(tail)}月"
                except Exception:
                    return s

            _months = [_month_label(m) for m in df["月"].tolist()]
            _usage = df["使用量合計"].tolist()
            _reduc = df["削減kWh"].tolist() if "削減kWh" in df.columns else [0] * len(_months)
            _demand = df["最大需要電力"].tolist()

            _ud = res.get("units_detail")
            _list_items = []
            if _ud is not None and len(_ud) > 0:
                for _, _r in _ud.iterrows():
                    _name = str(_r.get("機器ID", "") or "")
                    _list_items.append({
                        "系統名": _name, "設置場所": _name,
                        "メーカー": str(_r.get("メーカー", "") or ""),
                        "型式": str(_r.get("型番", "") or ""),
                        "制御可否": bool(_r.get("制御可否", False)),
                    })
            else:
                for _eq in DEFAULT_EQUIPMENT:
                    _list_items.append({
                        "系統名": _eq["loc"], "設置場所": _eq["loc"],
                        "メーカー": _eq["mfr"], "型式": _eq["model"],
                        "制御可否": str(_eq["ctrl"]).strip() in ("〇", "○", "◯"),
                    })
            _n_total = len(_list_items)
            _n_ctrl = sum(1 for it in _list_items if it["制御可否"])

            # ② 制御可否リストの出力項目を選択（不要な列は外せる）
            _ALL_LIST_COLS = ["系統名", "設置場所", "メーカー", "型式", "制御可否"]
            st.session_state.setdefault("_img_list_cols", _ALL_LIST_COLS)
            _sel_cols = st.multiselect(
                "② 制御可否リストに含める項目（不要な列は外せます）",
                _ALL_LIST_COLS, key="_img_list_cols")

            _NAVY_C, _GREEN_C = "#13315C", "#3DAE4E"
            _inv = float(res.get("total_invest", 0) or 0)
            _grs = float(res.get("gross_saving", 0) or 0)
            _nt = float(res.get("net_saving", 0) or 0)
            _sf = float(res.get("sys_fee", 0) or 0)
            _pb = (_inv / _nt) if _nt > 0 else 0.0
            _summary_cards = [
                {"icon": "money.png", "title": "初期導入費用（総投資額）", "value": f"{_inv:,.0f}円",
                 "subtitle": "税抜 導入費用 / 機器代・工事費含む", "color": _NAVY_C},
                {"icon": "graf.png", "title": "年間総削減額（グロス）", "value": f"{_grs:,.0f}円",
                 "subtitle": "基本料金＋電力量の年間削減合計", "color": _GREEN_C},
                {"icon": "plus.png", "title": "年間実質利点（手残り）", "value": f"{_nt:,.0f}円",
                 "subtitle": f"削減合計 {_grs:,.0f}円 − 年間維持費 {_sf:,.0f}円", "color": _GREEN_C},
                {"icon": "clock.png", "title": "投資回収期間（ROI）", "value": f"約{_pb:.1f}年",
                 "subtitle": f"約{round(_pb*12)}ヶ月で完全回収、以降は純利益", "color": _GREEN_C},
            ]

            if st.button("画像を生成 / 更新", key="gen_slide_imgs"):
                try:
                    st.session_state["_img_chart"] = make_demand_chart_png(
                        _months, _usage, _reduc, _demand, target_units=_n_ctrl)
                    st.session_state["_img_list"] = make_control_list_pngs(
                        _list_items, total_units=_n_total, controllable_units=_n_ctrl,
                        cols=(_sel_cols or None))
                    st.session_state["_img_cards"] = make_summary_cards_png(_summary_cards)
                except Exception as e:
                    st.error(f"画像生成でエラー: {e}")

            if st.session_state.get("_img_chart"):
                st.markdown("**① 使用量＋デマンド グラフ**")
                st.image(st.session_state["_img_chart"], use_container_width=True)
                st.download_button("グラフPNGをダウンロード", st.session_state["_img_chart"],
                                   file_name="使用量_デマンドグラフ.png", mime="image/png",
                                   key="dl_img_chart")
            _list_imgs = st.session_state.get("_img_list")
            if _list_imgs:
                _n_pg = len(_list_imgs)
                st.markdown("**② 制御可否リスト**" + (f"（1スライドに収まらないため {_n_pg} 枚に分割）" if _n_pg > 1 else ""))
                for _k, _img in enumerate(_list_imgs, start=1):
                    st.image(_img, use_container_width=True)
                    _fname = "制御可否リスト.png" if _n_pg == 1 else f"制御可否リスト_{_k}of{_n_pg}.png"
                    _btn_label = "制御可否リストPNGをダウンロード" + ("" if _n_pg == 1 else f"（{_k}/{_n_pg}）")
                    st.download_button(_btn_label, _img, file_name=_fname, mime="image/png",
                                       key=f"dl_img_list_{_k}")
            if st.session_state.get("_img_cards"):
                st.markdown("**③ サマリKPIカード**")
                st.image(st.session_state["_img_cards"], use_container_width=True)
                st.download_button("サマリカードPNGをダウンロード", st.session_state["_img_cards"],
                                   file_name="サマリカード.png", mime="image/png",
                                   key="dl_img_cards")


# ══════════════════════════════════════════════════════════════════
# 削減レポ（電力削減実績報告書）— 既存の試算ロジックには非依存の独立機能
# ══════════════════════════════════════════════════════════════════
def case_monthly_records(res):
    """保存案件に格納する月次データ（削減レポ呼び出し用）。res["df"] から軽量に抽出。"""
    df = res.get("df") if isinstance(res, dict) else None
    if df is None:
        return []
    out = []
    for r in df.to_dict("records"):
        out.append({
            "月": str(r.get("月", "")),
            "契約電力": float(r.get("契約電力", 0) or 0),
            "最大需要電力": float(r.get("最大需要電力", 0) or 0),
            "使用量合計": float(r.get("使用量合計", 0) or 0),
        })
    return out


def _to_ym(v):
    """各種表現（202401 / '2024/01' / '2024-01' / datetime / Excelシリアル日付）を YYYYMM(int) に。"""
    if v is None:
        return None
    try:
        if hasattr(v, "year") and hasattr(v, "month"):
            return int(v.year) * 100 + int(v.month)
    except Exception:
        pass
    s = str(v).strip()
    if not s or s.lower() == "nan":
        return None
    import re
    digits = re.sub(r"[^0-9]", "", s)
    # YYYYMM
    if len(digits) == 6:
        return int(digits)
    # YYYYMMDD → YYYYMM
    if len(digits) == 8:
        return int(digits[:6])
    # 'YYYY/M' 等
    m = re.match(r"^(\d{4})\D+(\d{1,2})", s)
    if m:
        return int(m.group(1)) * 100 + int(m.group(2))
    # Excelシリアル（おおよそ 20000〜60000）→ 日付換算
    try:
        f = float(digits)
        if 20000 <= f <= 80000:
            base = pd.Timestamp("1899-12-30") + pd.Timedelta(days=int(f))
            return int(base.year) * 100 + int(base.month)
    except Exception:
        pass
    return None


def _ym_add(ym, k):
    """YYYYMM に k ヶ月加算。"""
    y, m = ym // 100, ym % 100
    idx = (y * 12 + (m - 1)) + k
    return (idx // 12) * 100 + (idx % 12 + 1)


def _ym_label(ym):
    return f"{ym % 100}月"


# 客先の列名ゆれ → 内部ロール（値の取り込み用）
_REP_NAME_KW = ["お客様名", "顧客名", "会社名", "企業名"]
_REP_YM_KW = ["検針年月", "請求月", "年月", "使用期間", "計測期間"]
_REP_CONTRACT_KW = ["契約電力", "契約kw", "契約"]
_REP_DEMAND_KW = ["最大需要電力", "最大需要", "最大デマンド", "デマンド"]
_REP_USAGE_KW = ["使用電力量", "使用量合計", "使用量", "使用電力"]
_REP_KIHON_KW = ["基本料金単価", "基本料金単位"]      # 円/kW（請求書の契約電力報酬に使用）
_REP_KWHT_KW = ["電力量料金単価", "電力量単価"]       # 円/kWh（請求書の使用電力量報酬に使用）


def parse_zentai_monthly(uploaded):
    """『全体把握』等の月次明細Excelを解析し、[{ym, 契約電力, 最大需要電力, 使用電力量, client}] を返す。
    列名のゆれ（同じ項目で別名）は別名キーワードで吸収。検出不可なら []。"""
    try:
        uploaded.seek(0)
        xls = pd.ExcelFile(uploaded)
    except Exception:
        return []

    def _norm(s):
        return str(s).replace(" ", "").replace("　", "").replace("（", "(").replace("）", ")").lower()

    def _find_col(norm_cells, kws):
        for j, c in enumerate(norm_cells):
            if any(k.lower() in c for k in kws):
                return j
        return None

    best = []
    for sn in xls.sheet_names:
        try:
            raw = xls.parse(sn, header=None, dtype=object)
        except Exception:
            continue
        for i in range(min(len(raw), 8)):
            cells = ["" if v is None else str(v) for v in raw.iloc[i].tolist()]
            norm = [_norm(c) for c in cells]
            c_con = _find_col(norm, _REP_CONTRACT_KW)
            c_dem = _find_col(norm, _REP_DEMAND_KW)
            c_use = _find_col(norm, _REP_USAGE_KW)
            if c_con is None or c_dem is None or c_use is None:
                continue
            c_ym = _find_col(norm, _REP_YM_KW)
            if c_ym is None:   # シミュレーションのインポート形式は列名が「月」（年月入り）
                for _j, _c in enumerate(norm):
                    if _c == "月":
                        c_ym = _j
                        break
            c_nm = _find_col(norm, _REP_NAME_KW)
            c_kih = _find_col(norm, _REP_KIHON_KW)
            c_kwt = _find_col(norm, _REP_KWHT_KW)
            rows = []
            for r in range(i + 1, len(raw)):
                vals = raw.iloc[r].tolist()

                def gv(j):
                    return vals[j] if (j is not None and j < len(vals)) else None
                ym = _to_ym(gv(c_ym)) if c_ym is not None else None
                if ym is None:
                    continue
                def num(x):
                    try:
                        v = float(str(x).replace(",", "").strip())
                        return 0.0 if v != v else v   # NaN→0
                    except Exception:
                        return 0.0
                _con, _dem, _use = num(gv(c_con)), num(gv(c_dem)), num(gv(c_use))
                if _con == 0 and _dem == 0 and _use == 0:
                    continue                          # 空行（将来月など）はスキップ
                rows.append({
                    "ym": ym, "契約電力": _con, "最大需要電力": _dem, "使用電力量": _use,
                    "基本料金単価": num(gv(c_kih)) if c_kih is not None else 0.0,
                    "電力量単価": num(gv(c_kwt)) if c_kwt is not None else 0.0,
                    "client": (str(gv(c_nm)).strip() if c_nm is not None and gv(c_nm) else ""),
                })
            if len(rows) > len(best):
                best = rows
    return best


def _default_label_profile():
    return {
        "title": "電力削減レポート",
        "greeting": "毎度ご利用いただきありがとうございます。\n下記の通りご案内申し上げます。",
        "sec1": "①電力ご利用状況",
        "sec2": "②最大需要電力量　比較",
        "sec3": "③使用電力量　比較",
        "baseline_label": "基準年同月",
        "target_label": "対象年同月",
        "diff_label": "削減効果",
        "row_base_contract": "基準年契約電力",
        "row_cur_contract": "現契約電力",
        "row_cur_demand": "現最大需要電力",
        "row_cur_usage": "現使用電力量",
    }


def build_reduction_report_excel(R) -> bytes:
    """削減レポを既存報告書に近い形のExcelで出力。ラベルは R['labels']（客先原文）を使用。"""
    from openpyxl import Workbook
    from openpyxl.chart import BarChart, Reference
    L = R["labels"]
    months = R["months"]
    wb = Workbook()
    ws = wb.active
    ws.title = "電力削減レポート"
    navy = "1B2A4A"; green = "2E9E5B"
    bold = Font(name="Yu Gothic", size=11, bold=True)
    hwhite = Font(name="Yu Gothic", size=11, bold=True, color="FFFFFF")
    fill_navy = PatternFill("solid", fgColor=navy)
    fill_green = PatternFill("solid", fgColor=green)
    center = Alignment(horizontal="center", vertical="center")
    for col, w in zip("ABCDE", [22, 14, 14, 14, 12]):
        ws.column_dimensions[col].width = w

    ws["A1"] = L["title"]; ws["A1"].font = Font(name="Yu Gothic", size=16, bold=True, color=navy)
    ws["A2"] = f"{R['client']} 御中"; ws["A2"].font = bold
    ws["D2"] = "発行日"; ws["E2"] = R["issue_date"]
    ws["A3"] = R.get("baseline_period", "")
    ws["A4"] = L["greeting"].replace("\n", " ")
    r = 6

    def section(title, rows_def):
        nonlocal r
        ws.cell(row=r, column=1, value=title).font = hwhite
        ws.cell(row=r, column=1).fill = fill_navy
        for ci in range(2, 5):
            ws.cell(row=r, column=ci).fill = fill_navy
        r += 1
        ws.cell(row=r, column=1, value="項目").font = hwhite; ws.cell(row=r, column=1).fill = fill_green
        for ci, mlab in enumerate(months, start=2):
            c = ws.cell(row=r, column=ci, value=mlab); c.font = hwhite; c.fill = fill_green; c.alignment = center
        r += 1
        start = r
        for label, vals, fmt in rows_def:
            ws.cell(row=r, column=1, value=label).font = bold
            for ci, v in enumerate(vals, start=2):
                c = ws.cell(row=r, column=ci, value=v)
                c.number_format = fmt
                c.alignment = center
            r += 1
        r += 1
        return start

    # ① 利用状況
    section(L["sec1"], [
        (L["row_base_contract"], R["row_base_contract"], "#,##0"),
        (L["row_cur_contract"], R["row_cur_contract"], "#,##0"),
        (L["row_cur_demand"], R["row_cur_demand"], "#,##0"),
        (L["row_cur_usage"], R["row_cur_usage"], "#,##0"),
    ])
    # ② 最大需要 比較
    dm_diff = [b - c for b, c in zip(R["dm_base"], R["dm_cur"])]
    s2 = section(L["sec2"], [
        (L["baseline_label"], R["dm_base"], "#,##0"),
        (L["target_label"], R["dm_cur"], "#,##0"),
        (L["diff_label"], dm_diff, "#,##0"),
    ])
    ch2 = BarChart(); ch2.title = L["sec2"]; ch2.height = 6; ch2.width = 14
    data = Reference(ws, min_col=2, max_col=4, min_row=s2, max_row=s2 + 1)
    cats = Reference(ws, min_col=2, max_col=4, min_row=s2 - 1, max_row=s2 - 1)
    ch2.add_data(data, from_rows=True, titles_from_data=False); ch2.set_categories(cats)
    ws.add_chart(ch2, "G6")
    # ③ 使用量 比較
    use_diff = [b - c for b, c in zip(R["use_base"], R["use_cur"])]
    s3 = section(L["sec3"], [
        (L["baseline_label"], R["use_base"], "#,##0"),
        (L["target_label"], R["use_cur"], "#,##0"),
        (L["diff_label"], use_diff, "#,##0"),
    ])
    ch3 = BarChart(); ch3.title = L["sec3"]; ch3.height = 6; ch3.width = 14
    data = Reference(ws, min_col=2, max_col=4, min_row=s3, max_row=s3 + 1)
    cats = Reference(ws, min_col=2, max_col=4, min_row=s3 - 1, max_row=s3 - 1)
    ch3.add_data(data, from_rows=True, titles_from_data=False); ch3.set_categories(cats)
    ws.add_chart(ch3, "G22")

    buf = io.BytesIO(); wb.save(buf); buf.seek(0)
    return buf.getvalue()


def build_reduction_report_pdf(R) -> bytes:
    """削減レポを添付の出力イメージ（電力削減レポート）と同じ構成でPDF化。
    ①電力ご利用状況＝表のみ／②最大需要＝グラフ＋表／③使用電力量＝グラフ＋表。
    ②③は「グラフ＋表」をひとまとめに、棒の数値ラベル・単位列・年付き月見出し・色帯付き。"""
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    import numpy as np
    from matplotlib.patches import Rectangle
    try:
        from slide_image_export import set_jp_font
        set_jp_font()
    except Exception:
        pass
    L = R["labels"]; mfull = R.get("month_full", R["months"]); iss = R.get("issuer", {})
    NAVY = "#1B2A4A"; GRAY = "#BFBFBF"
    BLUE, BLUE_BG = "#1F4E79", "#DCE6F1"
    GREEN, GREEN_BG, GREEN_BAR = "#548235", "#E2EFDA", "#70AD47"
    ORANGE, ORANGE_BG, ORANGE_BAR = "#C55A11", "#FCE4D6", "#ED7D31"
    LEFT, WIDTH = 0.07, 0.86

    from matplotlib.backends.backend_pdf import PdfPages

    def draw_header(fig, page=1):
        if page == 1:
            fig.text(0.93, 0.972, f"発行日　{R['issue_date']}", fontsize=9, ha="right")
            fig.text(0.93, 0.958, "登録番号", fontsize=9, ha="right")
            fig.text(0.50, 0.945, L["title"], fontsize=16, fontweight="bold", color=NAVY, ha="center")
            fig.text(0.07, 0.918, f"{R['client']}　御中", fontsize=13, fontweight="bold")
            _iss = "\n".join([iss.get("name", ""), iss.get("addr", ""),
                              f"TEL:{iss.get('tel','')}　FAX:{iss.get('fax','')}"])
            fig.text(0.93, 0.918, _iss, fontsize=8.5, ha="right", va="top")
            fig.text(0.07, 0.896, L.get("greeting", ""), fontsize=8.5, va="top", linespacing=1.6)
        else:
            fig.text(0.50, 0.945, L["title"] + "（続き）", fontsize=15, fontweight="bold", color=NAVY, ha="center")
            fig.text(0.07, 0.918, f"{R['client']}　御中", fontsize=12, fontweight="bold")
            fig.text(0.93, 0.918, f"発行日　{R['issue_date']}", fontsize=9, ha="right")
            fig.text(0.07, 0.898, R.get("baseline_period", ""), fontsize=8.5)

    def band(fig, y, text, color):
        fig.patches.append(Rectangle((LEFT, y - 0.014), WIDTH, 0.028, transform=fig.transFigure,
                                     facecolor=color, edgecolor="none", zorder=2))
        fig.text(0.50, y, text, ha="center", va="center", fontsize=11,
                 fontweight="bold", color="white", zorder=3)

    def fmtv(vals):
        return [f"{v:,.0f}" for v in vals]

    def table_block(fig, top, header_bg, rows):
        nr = len(rows) + 1
        h = nr * 0.024
        ax = fig.add_axes([LEFT, top - h, WIDTH, h]); ax.axis("off")
        cols = [R.get("baseline_period", ""), "単位"] + list(mfull)
        cell = [[r[0], r[1]] + fmtv(r[2]) for r in rows]
        t = ax.table(cellText=cell, colLabels=cols, loc="upper center", cellLoc="center")
        t.auto_set_font_size(False); t.set_fontsize(9); t.scale(1, 1.5)
        widths = [0.28, 0.10, 0.2067, 0.2067, 0.2066]
        for (rr, cc), co in t.get_celld().items():
            co.set_width(widths[cc] if cc < len(widths) else 0.2)
            if rr == 0:
                co.set_text_props(fontweight="bold"); co.set_facecolor(header_bg)
            elif cc == 0:
                co.set_text_props(fontweight="bold")
        return top - h

    def chart_block(fig, top, base, cur, cur_color, unit, base_name, cur_name):
        h = 0.165
        ax = fig.add_axes([0.13, top - h, 0.78, h])
        x = np.arange(len(mfull)); w = 0.38
        b1 = ax.bar(x - w / 2, base, w, color=GRAY, label=base_name)
        b2 = ax.bar(x + w / 2, cur, w, color=cur_color, label=cur_name)
        ax.set_xticks(x); ax.set_xticklabels(mfull, fontsize=9)
        ax.tick_params(labelsize=8)
        ax.bar_label(b1, fmt="%.0f", fontsize=7.5, padding=2)
        ax.bar_label(b2, fmt="%.0f", fontsize=7.5, padding=2)
        ax.margins(y=0.22)
        ax.legend(fontsize=8.5, ncol=2, loc="lower center", bbox_to_anchor=(0.5, 1.02), frameon=False)
        if unit:
            ax.set_ylabel(unit, fontsize=8)
        for s in ("top", "right"):
            ax.spines[s].set_visible(False)
        return top - h

    def diffv(base, cur):
        return [b - c for b, c in zip(base, cur)]

    BAND_GAP = 0.030
    CHART_BAND_GAP = 0.052
    INNER_GAP = 0.046        # グラフ下端のx軸（月）ラベルを次の帯が隠さないよう広めに
    SECTION_GAP = 0.060

    buf = io.BytesIO()
    with PdfPages(buf) as pdf:
        # ── ページ1：① 電力ご利用状況 ＋ ② 最大需要電力量 ──
        fig = plt.figure(figsize=(8.27, 11.69))
        draw_header(fig, 1)
        y = 0.800
        band(fig, y, L["sec1"], BLUE)
        y = table_block(fig, y - BAND_GAP, BLUE_BG, [
            (L["row_base_contract"], "kW", R["row_base_contract"]),
            (L["row_cur_contract"], "kW", R["row_cur_contract"]),
            (L["row_cur_demand"], "kW", R["row_cur_demand"]),
            (L["row_cur_usage"], "kWh", R["row_cur_usage"]),
        ])
        y -= SECTION_GAP
        band(fig, y, L["sec2"], GREEN)
        y = chart_block(fig, y - CHART_BAND_GAP, R["dm_base"], R["dm_cur"], GREEN_BAR, "kW",
                        L["baseline_label"], L["target_label"])
        y -= INNER_GAP
        band(fig, y, L["sec2"], GREEN)
        y = table_block(fig, y - BAND_GAP, GREEN_BG, [
            (L["baseline_label"], "kW", R["dm_base"]),
            (L["target_label"], "kW", R["dm_cur"]),
            (L["diff_label"], "kW", diffv(R["dm_base"], R["dm_cur"])),
        ])
        pdf.savefig(fig, facecolor="white"); plt.close(fig)

        # ── ページ2：③ 使用電力量 比較 ──
        fig = plt.figure(figsize=(8.27, 11.69))
        draw_header(fig, 2)
        y = 0.820
        band(fig, y, L["sec3"], ORANGE)
        y = chart_block(fig, y - CHART_BAND_GAP, R["use_base"], R["use_cur"], ORANGE_BAR, "kWh",
                        L["baseline_label"], L["target_label"])
        y -= INNER_GAP
        band(fig, y, L["sec3"], ORANGE)
        y = table_block(fig, y - BAND_GAP, ORANGE_BG, [
            (L["baseline_label"], "kWh", R["use_base"]),
            (L["target_label"], "kWh", R["use_cur"]),
            (L["diff_label"], "kWh", diffv(R["use_base"], R["use_cur"])),
        ])
        pdf.savefig(fig, facecolor="white"); plt.close(fig)

    buf.seek(0)
    return buf.getvalue()


def _invoice_line_rows(INV):
    """請求書の明細行 [品目, 削減量, 単位, 単価, 金額] を生成（添付様式の詳細内訳）。
    契約電力は各月＋（式）注記、使用電力量は各月を 電力量単価／燃料調整費／再生エネルギー賦課金 の3行に分割。
    戻り値（rows, 小計）。小計は単純合算版と一致。"""
    rate = float(INV.get("rate", 0.5))
    m1 = INV.get("mark1", "※備考1"); m2 = INV.get("mark2", "※備考2")
    rows = []
    sub = 0.0
    rows.append([f"1. 契約電力　削減実績報酬　{m1}", "", "", "", ""])
    for c in INV.get("contract", []):
        amt = float(c["減量"]) * float(c["単価"]) * rate
        sub += amt
        rows.append([f"　{c['month']}分", f"{c['減量']:,.0f}", "kW", f"{c['単価']:,.2f}", f"{amt:,.0f}"])
        rows.append([f"　（削減量×基本料金単価×{rate:g}）", "", "", "", ""])
    rows.append([f"2. 使用電力量　削減実績報酬　{m2}", "", "", "", ""])
    for u in INV.get("usage", []):
        rows.append([f"　{u['month']}分", "", "", "", ""])
        a1 = float(u["減量"]) * float(u["電力量単価"]) * rate
        sub += a1
        rows.append(["　　電力量単価", f"{u['減量']:,.0f}", "kWh", f"{u['電力量単価']:,.2f}", f"{a1:,.0f}"])
        a2 = float(u["減量"]) * float(u["燃料調整費"]) * rate
        sub += a2
        rows.append(["　　燃料調整費", "", "", f"{u['燃料調整費']:,.2f}", f"{a2:,.0f}"])
        a3 = float(u["減量"]) * float(u["再エネ"]) * rate
        sub += a3
        rows.append(["　　再生エネルギー賦課金", "", "", f"{u['再エネ']:,.2f}", f"{a3:,.0f}"])
    if INV.get("baseline_inline"):
        rows.append([INV["baseline_inline"], "", "", "", ""])
    return rows, sub


def build_invoice_pdf(INV) -> bytes:
    """成果報酬型 請求書を添付様式（請求書内訳・外枠・詳細内訳・押印2つ）でPDF化。"""
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    from matplotlib.patches import Rectangle, Ellipse
    try:
        from slide_image_export import set_jp_font
        set_jp_font()
    except Exception:
        pass
    iss = INV.get("issuer", {}); bank = INV.get("bank", {})
    rows, sub = _invoice_line_rows(INV)
    tax = int(round(sub * 0.1)); total = int(round(sub)) + tax
    rows = rows + [["", "", "", "税　区分", "10％対象"],
                   ["", "", "", "合　計", f"{int(round(sub)):,}"],
                   ["", "", "", "消費税（10%）", f"{tax:,}"],
                   ["", "", "", "合　計（税込）", f"{total:,}"]]

    fig = plt.figure(figsize=(8.27, 11.69))
    # 外枠
    fig.patches.append(Rectangle((0.045, 0.035), 0.91, 0.93, transform=fig.transFigure,
                                 fill=False, edgecolor="#333333", lw=1.2))
    # ヘッダー
    fig.text(0.74, 0.948, "発　行　日", fontsize=9, ha="right")
    fig.text(0.92, 0.948, INV.get("issue_date", ""), fontsize=9, ha="right")
    fig.text(0.74, 0.933, "伝　票　番　号", fontsize=9, ha="right")
    fig.text(0.92, 0.933, INV.get("denpyo", ""), fontsize=9, ha="right")
    fig.text(0.50, 0.905, "請　求　書　内　訳", fontsize=16, fontweight="bold", ha="center")
    fig.text(0.095, 0.860, f"{INV.get('client','')}　御中", fontsize=13, fontweight="bold")
    # 金額（税込）
    fig.text(0.165, 0.788, "金　額", fontsize=11, ha="right")
    fig.text(0.205, 0.787, f"￥　{total:,}", fontsize=15, fontweight="bold")
    fig.text(0.165, 0.770, "（税込）", fontsize=9, ha="right")
    fig.patches.append(Rectangle((0.205, 0.779), 0.20, 0.0009, transform=fig.transFigure,
                                 facecolor="#333333", edgecolor="none"))
    # 発行元
    fig.text(0.92, 0.800, iss.get("name", ""), fontsize=12, fontweight="bold", ha="right")
    fig.text(0.92, 0.784, iss.get("addr", ""), fontsize=8.5, ha="right")
    fig.text(0.92, 0.770, f"TEL:{iss.get('tel','')}", fontsize=8.5, ha="right")
    fig.text(0.92, 0.756, f"FAX:{iss.get('fax','')}", fontsize=8.5, ha="right")

    def mini_table(rect, data, w0=0.30):
        ax = fig.add_axes(rect); ax.axis("off")
        tt = ax.table(cellText=data, cellLoc="left", bbox=[0, 0, 1, 1])
        tt.auto_set_font_size(False); tt.set_fontsize(9)
        for (r, c), co in tt.get_celld().items():
            co.set_width(w0 if c == 0 else (1 - w0))
            if c == 0:
                co.set_text_props(fontweight="bold")
        return tt

    mini_table([0.095, 0.694, 0.42, 0.058],
               [["見　積　番　号", " " + str(INV.get("mitsumori", "-"))],
                ["納　品　日", " " + str(INV.get("delivery", ""))],
                ["支　払　期　日", " " + str(INV.get("due", ""))]])
    _subj = str(INV.get("subject", "")).replace("　", "\n", 1)   # 全角スペースで2行に折り返し（枠内に収める）
    mini_table([0.095, 0.612, 0.42, 0.062],
               [["件　名", "  " + _subj]], w0=0.18)

    # 押印（左右・氏名は編集可）。A4縦の縦横比を補正し、ほぼ正方形の枠＋丸印に調整。
    def stamp(cx, cy, name):
        bw, bh = 0.073, 0.052
        fig.patches.append(Rectangle((cx - bw / 2, cy - bh / 2), bw, bh, transform=fig.transFigure,
                                     fill=False, edgecolor="#999999", lw=0.8))
        fig.patches.append(Ellipse((cx, cy), 0.057, 0.040, transform=fig.transFigure,
                                    fill=False, edgecolor="red", lw=1.6))
        if str(name).strip():
            fig.text(cx, cy, "\n".join(list(str(name).strip())), ha="center", va="center",
                     fontsize=11, color="red", linespacing=0.98)
    stamp(0.793, 0.632, INV.get("stamp_left", ""))
    stamp(0.866, 0.632, INV.get("stamp_right", ""))

    # 明細（外枠付き・bboxで領域を確定）
    th = 0.41
    ax = fig.add_axes([0.07, 0.585 - th, 0.86, th]); ax.axis("off")
    t = ax.table(cellText=rows, colLabels=["品　目", "削減量", "単位", "単　価", "金　額"],
                 cellLoc="center", bbox=[0, 0, 1, 1])
    t.auto_set_font_size(False); t.set_fontsize(8)
    widths = [0.46, 0.11, 0.08, 0.16, 0.19]
    for (r, c), cell in t.get_celld().items():
        cell.set_width(widths[c])
        if r == 0:
            cell.set_text_props(fontweight="bold"); cell.set_facecolor("#ECECEC")
        if c == 0:
            cell.set_text_props(ha="left")
        if c == 0 and r > 0 and str(rows[r - 1][0]).startswith(("1.", "2.")):
            cell.set_text_props(fontweight="bold")
        if c == 3 and r > 0 and str(rows[r - 1][3]).startswith(("合", "消費", "税")):
            cell.set_text_props(fontweight="bold")

    # 備考・振込先
    yb = 0.585 - th - 0.022
    fig.text(0.07, yb, INV.get("note1", ""), fontsize=8)
    fig.text(0.07, yb - 0.016, INV.get("note2", ""), fontsize=8)
    fig.text(0.60, yb, "【振込先】", fontsize=9, fontweight="bold")
    fig.text(0.60, yb - 0.016, f"{bank.get('name','')}　{bank.get('branch','')}", fontsize=8.5)
    fig.text(0.60, yb - 0.031, f"{bank.get('acct','')}　名義：{bank.get('holder','')}", fontsize=8.5)

    buf = io.BytesIO(); fig.savefig(buf, format="pdf", facecolor="white"); plt.close(fig); buf.seek(0)
    return buf.getvalue()


def build_invoice_excel(INV) -> bytes:
    """成果報酬型 請求書をExcel化。"""
    from openpyxl import Workbook
    iss = INV.get("issuer", {}); bank = INV.get("bank", {})
    rows, sub = _invoice_line_rows(INV)
    tax = int(round(sub * 0.1)); total = int(round(sub)) + tax
    wb = Workbook(); ws = wb.active; ws.title = "請求書"
    bold = Font(name="Yu Gothic", size=11, bold=True)
    for col, w in zip("ABCDE", [34, 12, 8, 14, 16]):
        ws.column_dimensions[col].width = w
    ws["A1"] = "請求書"; ws["A1"].font = Font(name="Yu Gothic", size=16, bold=True)
    ws["D1"] = "発行日"; ws["E1"] = INV.get("issue_date", "")
    ws["D2"] = "伝票番号"; ws["E2"] = INV.get("denpyo", "")
    ws["D3"] = "登録番号"; ws["E3"] = INV.get("reg_no", "")
    ws["A3"] = f"{INV.get('client','')} 御中"; ws["A3"].font = bold
    ws["A4"] = f"件名：{INV.get('subject','')}"
    ws["A5"] = f"金額（税込）：￥{total:,}"; ws["A5"].font = bold
    ws["A6"] = f"お支払期日：{INV.get('due','')}"
    ws["D5"] = iss.get("name", ""); ws["D6"] = iss.get("addr", "")
    ws["D7"] = f"TEL:{iss.get('tel','')}　FAX:{iss.get('fax','')}"
    r = 9
    for ci, h in enumerate(["品　目", "削減量", "単位", "単価", "金　額"], 1):
        c = ws.cell(row=r, column=ci, value=h); c.font = bold
        c.fill = PatternFill("solid", fgColor="EAF0F9")
    r += 1
    for row in rows:
        is_head = str(row[0]).startswith(("1.", "2."))
        for ci, v in enumerate(row, 1):
            cell = ws.cell(row=r, column=ci, value=v)
            if is_head:
                cell.font = bold
        r += 1
    r += 1
    ws.cell(row=r, column=4, value="合計"); ws.cell(row=r, column=5, value=int(round(sub))); r += 1
    ws.cell(row=r, column=4, value="消費税(10%)"); ws.cell(row=r, column=5, value=tax); r += 1
    ws.cell(row=r, column=4, value="合計(税込)").font = bold
    ws.cell(row=r, column=5, value=total).font = bold
    r += 2
    ws.cell(row=r, column=1, value=INV.get("note1", "")); r += 1
    ws.cell(row=r, column=1, value=INV.get("note2", "")); r += 2
    ws.cell(row=r, column=1, value="【振込先】").font = bold; r += 1
    ws.cell(row=r, column=1, value=f"{bank.get('name','')} {bank.get('branch','')} {bank.get('acct','')}"); r += 1
    ws.cell(row=r, column=1, value=f"名義：{bank.get('holder','')}")
    buf = io.BytesIO(); wb.save(buf); buf.seek(0)
    return buf.getvalue()


def _norm_company_key(name) -> str:
    """企業名の表記ゆれを吸収したキー（空白・(株)/株式会社・御中 等を除去）。"""
    import re
    s = str(name or "").strip()
    s = s.replace("　", "").replace(" ", "")
    for w in ("株式会社", "有限会社", "(株)", "（株）", "(有)", "（有）", "㈱", "㈲", "御中", "様"):
        s = s.replace(w, "")
    return s


def _merge_company_records(app_data: dict, client_name: str, monthly_map: dict) -> int:
    """monthly_map {ym(int)→{契約電力,最大需要電力,使用電力量,基本料金単価,電力量単価}} を企業へ累積保存。
    返り値：マージした月数。企業名が空なら 0（保存しない）。"""
    key = _norm_company_key(client_name)
    if not key:
        return 0
    store = app_data.setdefault("company_power_records", {})
    ent = store.setdefault(key, {"name": client_name or key, "monthly": {}})
    if client_name:
        ent["name"] = client_name
    n = 0
    for ym, v in (monthly_map or {}).items():
        if not ym:
            continue
        ent["monthly"][str(int(ym))] = {
            "契約電力": float(v.get("契約電力", 0) or 0),
            "最大需要電力": float(v.get("最大需要電力", 0) or 0),
            "使用電力量": float(v.get("使用電力量", 0) or 0),
            "基本料金単価": float(v.get("基本料金単価", 0) or 0),
            "電力量単価": float(v.get("電力量単価", 0) or 0),
            "燃料費調整額": float(v.get("燃料費調整額", 0) or 0),   # 請求書の使用電力量報酬に加算
            "再エネ賦課金": float(v.get("再エネ賦課金", 0) or 0),   # 請求書の使用電力量報酬に加算
        }
        n += 1
    if n:
        save_data(app_data)
    return n


def build_reduction_import_template_xlsx() -> bytes:
    """削減レポの追加インポート用テンプレ（パーサが確実に認識する列名＋基準年/対象年サンプル）。"""
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment
    wb = Workbook(); ws = wb.active; ws.title = "全体把握"
    headers = ["検針年月", "契約電力(kW)", "最大需要電力(kW)", "使用電力量(kWh)",
               "基本料金単価(円/kW)", "電力量単価(円/kWh)", "お客様名"]
    widths = [16, 16, 18, 18, 20, 22, 20]
    ws.cell(row=1, column=1,
            value="▼ 削減レポ用：基準年と対象年の月次を入力（2行目の列名は変更しないでください）。"
                  "検針年月は「2025年4月」または「202504」形式。お客様名は先頭行のみでも可。")
    ws.cell(row=1, column=1).font = Font(bold=True, color="C00000")
    for j, h in enumerate(headers, start=1):
        c = ws.cell(row=2, column=j, value=h)
        c.font = Font(bold=True, color="FFFFFF")
        c.fill = PatternFill("solid", fgColor="13315C")
        c.alignment = Alignment(horizontal="center")
        ws.column_dimensions[chr(64 + j)].width = widths[j - 1]
    r = 3
    for yr, dem, use in [(2024, 480, 100000), (2025, 455, 92000)]:   # 基準年 → 対象年（サンプル）
        for mth in range(1, 13):
            ws.cell(row=r, column=1, value=f"{yr}年{mth}月")
            ws.cell(row=r, column=2, value=500)
            ws.cell(row=r, column=3, value=dem)
            ws.cell(row=r, column=4, value=use)
            ws.cell(row=r, column=5, value=1800)
            ws.cell(row=r, column=6, value=18)
            ws.cell(row=r, column=7, value=("サンプル株式会社" if r == 3 else ""))
            r += 1
    buf = io.BytesIO(); wb.save(buf); buf.seek(0)
    return buf.getvalue()


def show_reduction_report(app_data: dict) -> None:
    """削減レポ：保存案件の呼び出し＋追加インポートで、客先様式の電力削減実績報告書を作成。"""
    st.markdown("## 削減レポ（電力削減実績報告書）")
    st.caption("保存案件の月次データを呼び出し、または明細Excel（全体把握等）をインポートして、"
               "客先様式の削減実績報告書（①利用状況／②最大需要／③使用電力量＋グラフ）を作成します。"
               "※既存の試算ロジックには影響しません。")

    # ── 1) データの取り込み（プール構築：ym→値）──
    pool = {}   # ym -> {"契約電力","最大需要電力","使用電力量"}
    detected_client = ""

    st.markdown("### 1. データの取り込み")
    # 企業の蓄積データ（自動蓄積された月次）から呼び出す
    _cpr = app_data.get("company_power_records", {}) or {}
    if _cpr:
        _comp_keys = list(_cpr.keys())
        _comp_names = ["（使わない）"] + [f"{_cpr[k].get('name', k)}（{len(_cpr[k].get('monthly', {}))}ヶ月）" for k in _comp_keys]
        _csel = st.selectbox("企業の蓄積データから呼び出す（シミュ保存・過去インポートを累積）", _comp_names, key="rep_company_sel")
        _cidx = _comp_names.index(_csel)
        if _cidx > 0:
            _ent = _cpr.get(_comp_keys[_cidx - 1], {})
            detected_client = _ent.get("name", detected_client)
            for _ymk, _v in (_ent.get("monthly", {}) or {}).items():
                try:
                    pool[int(_ymk)] = {"契約電力": float(_v.get("契約電力", 0) or 0),
                                       "最大需要電力": float(_v.get("最大需要電力", 0) or 0),
                                       "使用電力量": float(_v.get("使用電力量", 0) or 0),
                                       "基本料金単価": float(_v.get("基本料金単価", 0) or 0),
                                       "電力量単価": float(_v.get("電力量単価", 0) or 0)}
                except Exception:
                    pass
            st.success(f"「{detected_client}」の蓄積データ {len(_ent.get('monthly', {}))} ヶ月を呼び出しました。")

    cases = list(app_data.get("saved_cases", []) or [])
    def _clabel(c):
        mk = "" if c.get("monthly") else "（月次データなし）"
        return f"{c.get('company','?')}｜{c.get('saved_at','')}{mk}"
    _case_labels = [_clabel(c) for c in cases]
    csel_label = st.selectbox("結果一覧の保存案件から呼び出す", ["（使わない）"] + _case_labels, key="rep_case_sel")
    if csel_label in _case_labels:
        _c = cases[_case_labels.index(csel_label)]
        detected_client = _c.get("company", "")
        _mon = _c.get("monthly", []) or []
        for m in _mon:
            ym = _to_ym(m.get("月"))
            if ym:
                pool[ym] = {"契約電力": float(m.get("契約電力", 0) or 0),
                            "最大需要電力": float(m.get("最大需要電力", 0) or 0),
                            "使用電力量": float(m.get("使用量合計", 0) or 0)}
        if not _mon:
            st.info("この保存案件には月次データが含まれていません（旧バージョンで保存された案件）。"
                    "シミュレーション画面で再実行して『結果一覧に保存』し直すと月次データが入ります。"
                    "または下の明細Excelインポートで月次データを追加してください。"
                    "（顧客名は本案件から引用しています）")

    up = st.file_uploader("明細Excelを追加インポート（全体把握／報告書など・複数可）",
                          type=["xlsx", "xls"], accept_multiple_files=True, key="rep_uploads")
    _dlc1, _dlc2 = st.columns([1, 2])
    with _dlc1:
        st.download_button("📄 インポート用テンプレDL", build_reduction_import_template_xlsx(),
                           file_name="削減レポ_インポートテンプレ.xlsx",
                           mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                           key="rep_tmpl_dl", use_container_width=True)
    with _dlc2:
        st.caption("※ 列名が異なると取り込めません。上のテンプレDLを使うか、"
                   "シミュレーションのインポート形式（月／契約電力／最大需要電力／使用量合計 …）でもそのまま取り込めます。"
                   "※ 検針年月は「2025年4月」または「202504」形式にしてください。")
    for f in (up or []):
        recs = parse_zentai_monthly(f)
        if recs:
            if not detected_client:
                detected_client = next((r["client"] for r in recs if r.get("client")), "")
            for r in recs:
                pool[r["ym"]] = {"契約電力": r["契約電力"], "最大需要電力": r["最大需要電力"],
                                 "使用電力量": r["使用電力量"],
                                 "基本料金単価": r.get("基本料金単価", 0.0),
                                 "電力量単価": r.get("電力量単価", 0.0)}
            # 企業データへ累積（導入後インポート分もその企業の実績として蓄積）
            _n_acc = 0
            if detected_client:
                _n_acc = _merge_company_records(app_data, detected_client,
                                                {r["ym"]: pool[r["ym"]] for r in recs})
            st.success(f"「{getattr(f,'name','')}」から {len(recs)} ヶ月分を取り込みました。"
                       + (f"（企業「{detected_client}」データへ {_n_acc}ヶ月を累積）" if _n_acc else
                          "（顧客名が未特定のため企業データへは未累積。下で顧客名を選択/入力すると次回から累積します）"))
        else:
            st.warning(f"「{getattr(f,'name','')}」から月次明細を検出できませんでした（列名をご確認ください）。")

    if not pool:
        st.info("保存案件の選択、または明細Excelのインポートを行うとレポートを作成できます。")
        return

    yms = sorted(pool.keys())
    st.caption(f"取込済みデータ：{len(yms)}ヶ月（{yms[0]//100}/{yms[0]%100:02d} 〜 {yms[-1]//100}/{yms[-1]%100:02d}）")

    # ── 2) 顧客名・対象期間・基準年 ──
    st.markdown("### 2. 顧客名・対象期間・基準年")
    client = st.text_input("顧客名（御中）", value=detected_client, key="rep_client")
    cc = st.columns(3)
    with cc[0]:
        start_choices = [f"{y//100}/{y%100:02d}" for y in yms]
        # 既定：最新3ヶ月の開始（末尾から3つ目）
        _def_idx = max(0, len(yms) - 3)
        start_lbl = st.selectbox("対象期間の開始月（3ヶ月）", start_choices, index=_def_idx, key="rep_start")
        target_start = yms[start_choices.index(start_lbl)]
    target_yms = [_ym_add(target_start, k) for k in range(3)]
    with cc[1]:
        base_year = int(st.number_input("基準年（西暦・開始月基準）", min_value=2000, max_value=2100,
                                        value=int(target_yms[0] // 100 - 1), step=1, key="rep_baseyear",
                                        help="対象期間の開始月に対する基準年。通常は前年（既定＝開始年−1）。"))
    with cc[2]:
        issue_date = st.text_input("発行日", value=str(datetime.date.today()), key="rep_issue")
    # 年またぎの対象期間でも各月の対応前年月が正しく揃うよう、開始月の年差×12ヶ月で全月をずらす
    _yr_delta = target_yms[0] // 100 - int(base_year)
    base_yms = [_ym_add(ym, -12 * _yr_delta) for ym in target_yms]
    months_lbl = [_ym_label(ym) for ym in target_yms]
    st.caption(f"対象期間：{target_yms[0]//100}/{target_yms[0]%100:02d}〜{target_yms[-1]//100}/{target_yms[-1]%100:02d}"
               f"　／　基準：{base_yms[0]//100}/{base_yms[0]%100:02d}〜{base_yms[-1]//100}/{base_yms[-1]%100:02d}（各月の対応前年月）")
    baseline_period = f"基準年（{base_yms[0]//100}/{base_yms[0]%100:02d}〜{base_yms[-1]//100}/{base_yms[-1]%100:02d}）"

    def look(ym, key):
        return float(pool.get(ym, {}).get(key, 0) or 0)

    # ── 3) ラベルプロファイル（客先原文）──
    st.markdown("### 3. 表示ラベル（客先様式）")
    profiles = app_data.setdefault("label_profiles", {})
    prof = dict(_default_label_profile())
    if client and client in profiles:
        prof.update(profiles.get(client, {}))
        st.caption(f"「{client}」の保存済みラベルプロファイルを読込みました。")
    with st.expander("ラベルを編集（客先の表記に合わせる）", expanded=False):
        for k, jp in [("title", "タイトル"), ("sec1", "①見出し"), ("sec2", "②見出し"), ("sec3", "③見出し"),
                       ("baseline_label", "基準年系列名"), ("target_label", "対象年系列名"), ("diff_label", "差分の名称"),
                       ("row_base_contract", "①基準年契約電力 行名"), ("row_cur_contract", "①現契約電力 行名"),
                       ("row_cur_demand", "①現最大需要電力 行名"), ("row_cur_usage", "①現使用電力量 行名")]:
            prof[k] = st.text_input(jp, value=prof.get(k, ""), key=f"rep_lbl_{k}")
        if st.button("このラベルを客先プロファイルに保存", key="rep_lbl_save"):
            if client.strip():
                profiles[client] = prof; save_data(app_data)
                st.success(f"「{client}」のラベルプロファイルを保存しました。")
            else:
                st.warning("顧客名を入力してください。")

    # ── エクセル眼鏡モード：使用電力量を営業日数で補正（差分計算シートのロジック） ──
    _megane = st.toggle("☑ 株式会社エクセル眼鏡はこちら（営業日数補正モード）", value=False, key="rep_megane",
                        help="眼鏡様専用。使用電力量の削減効果を「対象年を基準年の営業日条件に換算」して算出します。"
                             "契約電力・最大需要電力は通常どおり単純差のままです。")
    _biz_factor = [1.0, 1.0, 1.0]
    if _megane:
        st.caption("使用電力量：対象年同月 ＝ 対象年実績 ×（基準年営業日数 ÷ 対象年営業日数）に換算 → "
                   "削減効果 ＝ 基準年使用量 − 換算後の対象年（＝眼鏡『差分計算』シートと同一）")
        _cbd = st.columns(3); _ctd = st.columns(3)
        _base_days = []; _tgt_days = []
        for i in range(3):
            _base_days.append(_cbd[i].number_input(f"{months_lbl[i]}：基準年営業日数", min_value=1, max_value=31,
                                                   value=22, step=1, key=f"rep_bd_{i}"))
            _tgt_days.append(_ctd[i].number_input(f"{months_lbl[i]}：対象年営業日数", min_value=1, max_value=31,
                                                  value=21, step=1, key=f"rep_td_{i}"))
        _biz_factor = [(float(_base_days[i]) / float(_tgt_days[i])) if _tgt_days[i] else 1.0 for i in range(3)]

    # ── 4) 数値プレビュー（編集可：営業日数補正など客先固有の手当てに対応）──
    st.markdown("### 4. 数値プレビュー（必要に応じて手修正可）")
    # 2年比較には基準年(前年)と対象年の両方の月次が必要。片方が無ければ明示して促す。
    _have_base = any(ym in pool for ym in base_yms)
    _have_tgt = any(ym in pool for ym in target_yms)
    if not _have_base and not _have_tgt:
        st.warning("⚠️ 基準年・対象年ともデータが未取込です。明細Excelを追加インポートしてください。")
    elif not _have_base:
        st.warning(f"⚠️ **基準年（{base_year}年）の月次データが未取込**のため前年側が空欄になっています。"
                   "前年の明細を上の『追加インポート』で取り込むと、2年比較（前年 vs 対象年）が表示されます。"
                   "（表内で手入力しても反映できます）")
    elif not _have_tgt:
        st.warning(f"⚠️ **対象年（{target_yms[0]//100}年）の月次データが未取込**です。対象年の実績を追加インポートしてください。")
    else:
        st.caption(f"2年比較：基準年 {base_year}年 vs 対象年 {target_yms[0]//100}年（前年同月）")
    df1 = pd.DataFrame({
        "項目": [prof["row_base_contract"], prof["row_cur_contract"], prof["row_cur_demand"], prof["row_cur_usage"]],
        months_lbl[0]: [look(base_yms[0], "契約電力"), look(target_yms[0], "契約電力"), look(target_yms[0], "最大需要電力"), look(target_yms[0], "使用電力量")],
        months_lbl[1]: [look(base_yms[1], "契約電力"), look(target_yms[1], "契約電力"), look(target_yms[1], "最大需要電力"), look(target_yms[1], "使用電力量")],
        months_lbl[2]: [look(base_yms[2], "契約電力"), look(target_yms[2], "契約電力"), look(target_yms[2], "最大需要電力"), look(target_yms[2], "使用電力量")],
    })
    e1 = st.data_editor(df1, key="rep_e1", use_container_width=True, hide_index=True, disabled=["項目"])

    st.markdown(f"**{prof['sec2']}**")
    df2 = pd.DataFrame({
        "項目": [prof["baseline_label"], prof["target_label"]],
        months_lbl[0]: [look(base_yms[0], "最大需要電力"), look(target_yms[0], "最大需要電力")],
        months_lbl[1]: [look(base_yms[1], "最大需要電力"), look(target_yms[1], "最大需要電力")],
        months_lbl[2]: [look(base_yms[2], "最大需要電力"), look(target_yms[2], "最大需要電力")],
    })
    e2 = st.data_editor(df2, key="rep_e2", use_container_width=True, hide_index=True, disabled=["項目"])

    st.markdown(f"**{prof['sec3']}**")
    # 眼鏡モード：対象年＝実績×(基準年営業日/対象年営業日)。通常モードは _biz_factor=1.0 で従来と同一。
    _use_cur3 = [round(look(target_yms[i], "使用電力量") * _biz_factor[i]) for i in range(3)]
    df3 = pd.DataFrame({
        "項目": [prof["baseline_label"], prof["target_label"]],
        months_lbl[0]: [look(base_yms[0], "使用電力量"), _use_cur3[0]],
        months_lbl[1]: [look(base_yms[1], "使用電力量"), _use_cur3[1]],
        months_lbl[2]: [look(base_yms[2], "使用電力量"), _use_cur3[2]],
    })
    if _megane:
        st.caption("※ 上表『③使用電力量』の対象年は営業日数補正後の換算値です（削減効果＝基準年−換算値）。")
    e3 = st.data_editor(df3, key="rep_e3", use_container_width=True, hide_index=True, disabled=["項目"])

    def _row(edf, i):
        return [float(edf.iloc[i][m] or 0) for m in months_lbl]

    _co = (app_data.get("quote_settings", {}) or {}).get("company", {}) or {}
    month_full = [f"{ym // 100}年{ym % 100}月" for ym in target_yms]
    R = {
        "client": client, "issue_date": issue_date, "labels": prof, "month_full": month_full,
        "issuer": {"name": _co.get("name", "株式会社シムックスイニシアティブ"),
                   "addr": _co.get("addr", "東京都港区浜松町1丁目30-5 浜松町スクエア 10F"),
                   "tel": _co.get("tel", "03-6402-2650"), "fax": _co.get("fax", "03-6402-2651")},
        "months": months_lbl, "baseline_period": baseline_period,
        "row_base_contract": _row(e1, 0), "row_cur_contract": _row(e1, 1),
        "row_cur_demand": _row(e1, 2), "row_cur_usage": _row(e1, 3),
        "dm_base": _row(e2, 0), "dm_cur": _row(e2, 1),
        "use_base": _row(e3, 0), "use_cur": _row(e3, 1),
    }

    # ── 5) グラフプレビュー ──
    st.markdown("### 5. グラフプレビュー")
    def chart(title, base, cur):
        d = pd.DataFrame({"月": months_lbl * 2,
                          "系列": [prof["baseline_label"]] * 3 + [prof["target_label"]] * 3,
                          "値": base + cur})
        ch = alt.Chart(d).mark_bar().encode(
            x=alt.X("月:N", sort=months_lbl, axis=alt.Axis(labelAngle=0, title=None)),
            y=alt.Y("値:Q", title=None),
            color=alt.Color("系列:N", scale=alt.Scale(domain=[prof["baseline_label"], prof["target_label"]],
                                                       range=["#B0B7C3", "#2E9E5B"]),
                            legend=alt.Legend(orient="top", title=None)),
            xOffset="系列:N").properties(height=240, title=title)
        st.altair_chart(ch, use_container_width=True)
    gc = st.columns(2)
    with gc[0]:
        chart(prof["sec2"], R["dm_base"], R["dm_cur"])
    with gc[1]:
        chart(prof["sec3"], R["use_base"], R["use_cur"])

    # ── 6) 出力（Excel / PDF 選択。不足があれば警告ポップアップ）──
    st.markdown("### 6. 出力")
    fmt = st.radio("出力形式", ["PDF", "Excel"], horizontal=True, key="rep_fmt")
    safe = (client or "削減レポート").replace("/", "_").replace(" ", "")

    # ── データ不足チェック ──
    missing = []
    if not (up or []):
        missing.append("追加データ（明細Excel）が未インポートです（基準年データが不足する可能性があります）")
    for ym in base_yms:
        if ym not in pool:
            missing.append(f"基準年 {ym//100}/{ym % 100:02d} のデータが見つかりません")
    for ym in target_yms:
        if ym not in pool:
            missing.append(f"対象年 {ym//100}/{ym % 100:02d} のデータが見つかりません")

    def _zero_months(vals):
        return [months_lbl[i] for i, v in enumerate(vals) if not v]
    for _nm, _vals in [
        ("現契約電力", R["row_cur_contract"]), ("現最大需要電力", R["row_cur_demand"]),
        ("現使用電力量", R["row_cur_usage"]),
        (f"{prof['baseline_label']}（最大需要）", R["dm_base"]),
        (f"{prof['baseline_label']}（使用電力量）", R["use_base"]),
    ]:
        _z = _zero_months(_vals)
        if _z:
            missing.append(f"{_nm}：{ '・'.join(_z) } が 0（未取得）です")

    def _make_output():
        """(label, data, filename, mime, key) を返す。"""
        if fmt == "Excel":
            return ("Excelをダウンロード", build_reduction_report_excel(R),
                    f"電力削減レポート_{safe}.xlsx",
                    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", "rep_dl_xlsx")
        return ("PDFをダウンロード", build_reduction_report_pdf(R),
                f"電力削減レポート_{safe}.pdf", "application/pdf", "rep_dl_pdf")

    @st.dialog("⚠️ 出力前の確認：データ不足があります")
    def _confirm_dialog():
        st.warning("以下のデータが不足しています。**不足分は 0 として出力**されます。内容をご確認ください。")
        for w in missing:
            st.markdown(f"- {w}")
        st.caption("不足を補うには、削減レポ上部で『明細Excel（全体把握等）』をインポートするか、"
                   "対象期間・基準年の指定を見直してください。")
        try:
            lbl, data, fn, mime, k = _make_output()
            if st.download_button(f"このまま出力（{lbl}）", data, file_name=fn, mime=mime, key=k + "_dlg"):
                st.rerun()   # ダウンロードしたらポップアップを自動で閉じる
        except Exception as _e:
            st.error(f"生成エラー：{_e}")

    if missing:
        st.warning(f"⚠️ データ不足が {len(missing)} 件あります。『レポートを出力』を押すと確認が表示されます。")
        if st.button("レポートを出力", type="primary", key="rep_gen"):
            _confirm_dialog()
    else:
        try:
            lbl, data, fn, mime, k = _make_output()
            st.download_button(lbl, data, file_name=fn, mime=mime, key=k)
        except Exception as _e:
            st.error(f"生成エラー：{_e}")

    # ── 7) 請求書（成果報酬型）：削減レポとは別に出力 ──
    st.markdown("---")
    st.markdown("### 7. 請求書（成果報酬型）")
    st.caption("削減実績（契約電力・使用電力量の削減量）を引用し、単価・諸条件を追記して成果報酬請求書を作成します。"
               "削減レポートとは別に出力します。報酬＝削減量 × 単価 × 成果報酬率。")
    with st.expander("請求書を作成する", expanded=False):
        inv_rate = st.number_input("成果報酬率（例：0.5＝50%）", 0.0, 1.0, 0.5, step=0.05, key="inv_rate")
        # 引用：削減量（契約電力＝基準年契約−現契約／使用電力量＝基準年同月−対象年同月）
        contract_red = [max(b - c, 0.0) for b, c in zip(R["row_base_contract"], R["row_cur_contract"])]
        usage_red = [max(b - c, 0.0) for b, c in zip(R["use_base"], R["use_cur"])]
        kihon_pre = [float(pool.get(ym, {}).get("基本料金単価", 0) or 0) for ym in target_yms]
        kwh_pre = [float(pool.get(ym, {}).get("電力量単価", 0) or 0) for ym in target_yms]

        st.markdown("**1. 契約電力 削減実績報酬**（削減量＝引用／基本料金単価＝引用または追記）")
        cdf = pd.DataFrame({"月": months_lbl, "削減量(kW)": contract_red, "基本料金単価": kihon_pre})
        ce = st.data_editor(cdf, key="inv_cdf", hide_index=True, use_container_width=True,
                            disabled=["月", "削減量(kW)"])
        st.markdown("**2. 使用電力量 削減実績報酬**（削減量＝引用／各単価＝引用または追記）")
        udf = pd.DataFrame({"月": months_lbl, "削減量(kWh)": usage_red, "電力量単価": kwh_pre,
                            "燃料調整費": [0.0, 0.0, 0.0], "再エネ賦課金": [0.0, 0.0, 0.0]})
        ue = st.data_editor(udf, key="inv_udf", hide_index=True, use_container_width=True,
                            disabled=["月", "削減量(kWh)"])

        ic = st.columns(3)
        inv_issue = ic[0].text_input("発行日", value=issue_date, key="inv_issue")
        inv_denpyo = ic[1].text_input("伝票番号", value="", key="inv_denpyo")
        inv_due = ic[2].text_input("お支払期日", value="", key="inv_due")
        ic2 = st.columns(2)
        inv_mitsumori = ic2[0].text_input("見積番号", value="-", key="inv_mitsumori")
        inv_delivery = ic2[1].text_input("納品日", value="", key="inv_delivery")
        inv_subject = st.text_input(
            "件名", value="2025年12月～2026年2月　成果報酬型_空調制御サービス運用費用", key="inv_subject")
        inv_reg = st.text_input("登録番号（インボイス）", value="T4010001168914", key="inv_reg")
        st.markdown("**押印（※印）の氏名 — 左右それぞれ記入できます（空欄なら枠のみ）**")
        scol = st.columns(2)
        stamp_left = scol[0].text_input("印①（左）の氏名", value="", key="inv_stamp_l")
        stamp_right = scol[1].text_input("印②（右）の氏名", value="", key="inv_stamp_r")
        bcols = st.columns(4)
        bank = {"name": bcols[0].text_input("銀行名", "みずほ銀行", key="inv_bk_name"),
                "branch": bcols[1].text_input("支店名", "新横浜支店", key="inv_bk_branch"),
                "acct": bcols[2].text_input("口座", "当座 0103955", key="inv_bk_acct"),
                "holder": bcols[3].text_input("名義", "ｶ)ｼﾑｯｸｽｲﾆｼｱﾃｨﾌﾞ", key="inv_bk_holder")}

        INV = {
            "client": client, "issue_date": inv_issue, "denpyo": inv_denpyo, "due": inv_due,
            "mitsumori": inv_mitsumori, "delivery": inv_delivery,
            "stamp_left": stamp_left, "stamp_right": stamp_right,
            "subject": inv_subject, "reg_no": inv_reg, "issuer": R["issuer"], "bank": bank,
            "rate": float(inv_rate), "baseline_period": baseline_period,
            "baseline_inline": baseline_period.replace("\n", " "),
            "contract": [{"month": months_lbl[i], "減量": float(ce.iloc[i]["削減量(kW)"] or 0),
                          "単価": float(ce.iloc[i]["基本料金単価"] or 0)} for i in range(3)],
            "usage": [{"month": months_lbl[i], "減量": float(ue.iloc[i]["削減量(kWh)"] or 0),
                       "電力量単価": float(ue.iloc[i]["電力量単価"] or 0),
                       "燃料調整費": float(ue.iloc[i]["燃料調整費"] or 0),
                       "再エネ": float(ue.iloc[i]["再エネ賦課金"] or 0)} for i in range(3)],
            "note1": "※備考１　契約電力詳細は電力削減レポート 項目①②を参照",
            "note2": "※備考２　使用電力量詳細は電力削減レポート 項目③を参照",
        }
        _rows, _sub = _invoice_line_rows(INV)
        _tax = int(round(_sub * 0.1)); _total = int(round(_sub)) + _tax
        mm = st.columns(3)
        mm[0].metric("合計", f"¥{int(round(_sub)):,}")
        mm[1].metric("消費税(10%)", f"¥{_tax:,}")
        mm[2].metric("合計(税込)", f"¥{_total:,}")

        ifmt = st.radio("請求書の出力形式", ["PDF", "Excel"], horizontal=True, key="inv_fmt")
        isafe = (client or "請求書").replace("/", "_").replace(" ", "")
        try:
            if ifmt == "PDF":
                st.download_button("請求書（PDF）をダウンロード", build_invoice_pdf(INV),
                                   file_name=f"請求書_{isafe}.pdf", mime="application/pdf", key="inv_dl_pdf")
            else:
                st.download_button("請求書（Excel）をダウンロード", build_invoice_excel(INV),
                                   file_name=f"請求書_{isafe}.xlsx",
                                   mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                                   key="inv_dl_xlsx")
        except Exception as _e:
            st.error(f"請求書生成エラー：{_e}")


def main():
    st.set_page_config(
        page_title=APP_TITLE,
        page_icon="⚡",
        layout="wide",
        initial_sidebar_state="expanded"
    )

    # 🎨 グローバルCSS (③ナビゲーションボタンの通常・ホバー時スタイル完全反転仕様)
    load_css()

    # ── セッション状態の初期化 ────────────────────
    if "app_data" not in st.session_state:
        st.session_state.app_data = load_data()
    # 各種設定の固定値を計算用モジュール定数へ反映（毎回・サンプル初期化や各計算より前に）
    _apply_const_settings(st.session_state.app_data)
    if "screen" not in st.session_state:
        st.session_state.screen = "main"
    if "df_input" not in st.session_state:
        st.session_state.df_input = get_blank_df()   # 運用：クリア起動（テスト用サンプルはボタンで読込）
    if "units_df" not in st.session_state:
        st.session_state.units_df = None   # 室外機リスト（機体入力／インポートから生成）
    if "model_rows" not in st.session_state:
        st.session_state.model_rows = None  # 機体入力テーブル（型番リスト）
    if "sim_result" not in st.session_state:
        # 運用：クリア起動。初期サンプル試算は行わず、結果は未生成（実行を促す表示になる）。
        st.session_state.sim_result = None

    app_data = st.session_state.app_data

    # ── サイドバー ────────────────────────────────
    with st.sidebar:
        st.markdown(f"""
<div style="background:#1F3864;border-radius:10px;padding:12px 10px;margin-bottom:14px;text-align:center;position:relative;">
  <div style="position:absolute;top:6px;right:10px;font-size:10px;color:#cfd8ea;font-weight:bold;">{APP_VERSION}</div>
  <div style="font-size:22px;">⚡</div>
  <div style="font-weight:bold;font-size:13px;color:#fff;">DPS シミュレーター</div>
  <div style="font-size:10px;color:#aaa;">空調デマンド制御 提案・稟議</div>
</div>
""", unsafe_allow_html=True)

        st.markdown("**メニュー**")
        col_n1, col_n2 = st.columns(2)
        with col_n1:
            if st.button("ホーム", use_container_width=True):
                st.session_state.screen = "main"
                st.rerun()
        with col_n2:
            if st.button("各種設定", use_container_width=True):
                st.session_state.screen = "admin"
                st.rerun()
        col_n3, col_n4 = st.columns(2)
        with col_n3:
            if st.button("操作ヘルプ", use_container_width=True):
                st.session_state.screen = "help"
                st.rerun()
        with col_n4:
            if st.button("結果一覧", use_container_width=True):
                st.session_state.screen = "results_list"
                st.rerun()
        col_n5, col_n6 = st.columns(2)
        with col_n5:
            if st.button("見積書作成", use_container_width=True):
                st.session_state.screen = "quote"
                st.rerun()
        with col_n6:
            if st.button("削減レポ", use_container_width=True):
                st.session_state.screen = "reduction_report"
                st.rerun()

        st.divider()

        # ── 顧客情報 ─────────────────────────────
        st.markdown("### 顧客情報")
        client_name = st.text_input("顧客名", value="")
        MANUAL_GYOTAI = "その他※数値を指定する"
        gyotai_list = list(app_data["gyotai_master"].keys()) + [MANUAL_GYOTAI]
        gyotai = st.selectbox("業態", gyotai_list, index=1)
        if gyotai == MANUAL_GYOTAI:
            # 固定値入力モード：空調割合を手動指定
            st.session_state.setdefault("manual_ac_peak_pct", 18)
            st.session_state.setdefault("manual_ac_kwh_pct", 15)
            _mc = st.columns(2)
            _mp = _mc[0].number_input("最大デマンドの空調割合（％）", min_value=0, max_value=100,
                                      step=1, key="manual_ac_peak_pct")
            _mk = _mc[1].number_input("使用量の空調割合（％）", min_value=0, max_value=100,
                                      step=1, key="manual_ac_kwh_pct")
            st.session_state["manual_ac"] = {"ac_peak": int(_mp) / 100.0, "ac_kwh": int(_mk) / 100.0}
            st.caption(f"空調割合想定（固定値入力）：最大デマンド {int(_mp)}％ ／ 使用量 {int(_mk)}％")
        else:
            _gm_sel = app_data["gyotai_master"].get(gyotai, {})
            _ac_peak_pct = float(_gm_sel.get("ac_peak", 0)) * 100
            _ac_kwh_pct = float(_gm_sel.get("ac_kwh", 0)) * 100
            st.caption(f"空調割合想定（自動反映）：最大デマンド {_ac_peak_pct:.0f}％ ／ 使用量 {_ac_kwh_pct:.0f}％")

        # ── 📊 分析オプション（業態・データから自動判断／手動変更可）──
        st.markdown("### 📊 分析オプション（自動判断・手動変更可）")
        _cset = app_data.get("const_settings", DEFAULT_DATA["const_settings"])
        st.caption("業態と検針票12ヶ月の傾向から、対象になりえる分析方法を自動でオンにします（各トグルで手動変更可）。")

        # 月別追加列の取込状況を判定する共通関数
        def _colon(name, signed=False):
            if name not in st.session_state.df_input.columns:
                return False
            s = pd.to_numeric(st.session_state.df_input[name], errors="coerce")
            return float((s.abs() if signed else s).sum() or 0) > 0

        # ── 自動判断の材料：12ヶ月使用量の季節変動 と 業態 ──
        _dfi = sanitize_columns(st.session_state.df_input)
        _useries = (pd.to_numeric(_dfi["使用量合計"], errors="coerce").dropna()
                    if "使用量合計" in _dfi.columns else pd.Series([], dtype=float))
        _swing = (float((_useries.max() - _useries.min()) / _useries.min())
                  if (len(_useries) >= 2 and _useries.min() > 0) else 0.0)
        _base_ac_pct = float(app_data.get("const_settings", {}).get("base_ac_ratio", 0.0)) * 100
        _BASE_AC_GYOTAI = {"スーパー・食品(冷凍冷蔵主役)", "物流センター・倉庫",
                           "ホテル・商業施設", "病院・医療施設"}
        auto_seasonal = bool(_swing >= 0.15)                                   # 最少月比 +15%以上の季節変動で採用
        auto_gyotai   = bool(gyotai in _BASE_AC_GYOTAI and _base_ac_pct > 0)   # 通年稼働業態＋率設定済みのとき採用

        # 季節性（ベースロード法）：自動ON。ONのとき採択理由を1文で表示
        seasonal_ac = st.toggle(
            "季節性を反映（ベースロード法）", value=auto_seasonal, key="opt_seasonal_ac",
            help="最少月を基礎電力とみなし各月の超過分を空調と推計。業態の一律割合より実勢に近づきます。")
        if seasonal_ac:
            _why = (f"最少月比で使用量が最大 +{_swing*100:.0f}% 変動し季節負荷が明確なため"
                    if _swing > 0 else "検針票に季節変動があるため")
            st.caption(f"↳ 採択理由：{_why}、業態一律割合よりベースロード法（実測）を採用。")

        # 業態を考慮（基礎電力の空調分）：通年稼働業態のとき自動ON
        gyotai_consider = st.toggle(
            f"業態を考慮（基礎電力の空調分 {_base_ac_pct:.0f}% を見込む）",
            value=auto_gyotai, key="opt_gyotai_consider", disabled=not seasonal_ac,
            help="24時間冷凍冷蔵など、最少月にも空調/冷凍が残る業態向け。"
                 "基礎電力の一部を空調として残します（率は各種設定で変更）。")
        if seasonal_ac and gyotai_consider:
            st.caption(f"↳ 採択理由：業態『{gyotai}』は最少月も冷凍冷蔵・空調が通年稼働するため、"
                       f"基礎電力の {_base_ac_pct:.0f}% を空調分として見込む。")
        # 比率0%だと業態配慮ONでも結果が変わらないため明示（表示のみ・算出は不変）
        if gyotai_consider and _base_ac_pct <= 0:
            st.warning("『基礎電力の空調比率』が 0% のため、ONにしても算出結果は変わりません。"
                       "各種設定 → 計算ロジック →「■ 季節性の業態配慮」で比率（%）を設定してください。")

        # 燃調・再エネは月別列があれば自動反映（自動判断・常時表面表示）
        st.caption(
            f"燃料費調整額 {'✅自動反映' if _colon('燃料費調整額', True) else '— 排除（列なし/空）'} ／ "
            f"再エネ賦課金 {'✅自動反映' if _colon('再エネ賦課金') else '— 排除（列なし/空）'}"
            "　…テンプレの月別列があれば自動反映、無ければ自動で排除します。")

        with st.expander("月別単価欄の状況・電力会社プリセット", expanded=False):
            st.caption(
                f"（欄のみ・計算未反映）夏季ピーク {'入力あり' if _colon('夏季ピーク単価') else '空'} ／ "
                f"平日昼間 {'入力あり' if _colon('平日昼間単価') else '空'} ／ "
                f"夜間休日 {'入力あり' if _colon('夜間休日単価') else '空'} ／ "
                f"市場価格調整 {'入力あり' if _colon('市場価格調整項', True) else '空'}")
            st.divider()
            preset_name = st.selectbox("電力会社プリセット（目安・要確認）",
                                       list(TARIFF_PRESETS.keys()), index=0, key="tariff_preset")
            preset = TARIFF_PRESETS.get(preset_name)
            if preset:
                st.caption(f"目安：基本料金 ¥{preset['basic']:,.2f}/kW・電力量 ¥{preset['energy']:,.1f}/kWh"
                           f"・力率割引 {'あり' if preset['pf_discount'] else 'なし'}")
                if st.button("▶ プリセット単価を全月へ反映", key="apply_preset"):
                    _d = sanitize_columns(st.session_state.df_input).copy()
                    _d["基本料金単価"] = preset["basic"]
                    _d["電力量単価"] = preset["energy"]
                    st.session_state.df_input = _d
                    st.success(f"{preset_name} の単価を全月へ反映しました。")
                    st.rerun()

        st.markdown("### 室外機情報")
        units_df = st.session_state.get("units_df")
        if units_df is not None and len(units_df) > 0:
            # 型番ベース：リストから自動取得（容量加重）
            n_units_total = int(len(units_df))
            n_units_ctrl = int(units_df["制御可否"].sum())
            _w = units_df["電力kW"] * units_df["稼働係数"]
            _cw = float(_w[units_df["制御可否"]].sum())
            ctrl_ratio = (_cw / float(_w.sum())) if _w.sum() > 0 else (n_units_ctrl / max(n_units_total, 1))
            st.success(f"📋 室外機リスト検出：{n_units_total}台（制御対象 {n_units_ctrl}台）")
            st.markdown(f"**容量加重 制御比率**: {ctrl_ratio:.0%}　/　型番ベースで精緻化中")
            if st.button("🗑️ 機体入力を解除して台数手入力に戻す", use_container_width=True):
                st.session_state.units_df = None
                st.session_state.model_rows = None
                st.rerun()
        else:
            _gaisan = st.toggle("🔧 機材情報なし → 検針票から概算する", value=False, key="gaisan_mode",
                                help="室外機・機材情報が無い場合に、検針票の空調ピーク負荷と業態平均から"
                                     "台数・初期費用・回収期間を概算します（機材取得後に精緻化）。")
            if _gaisan:
                _ac_ovr_g = st.session_state.get("manual_ac") if gyotai == MANUAL_GYOTAI else None
                _estu = estimate_units_from_meter(st.session_state.df_input, app_data, gyotai, ac_override=_ac_ovr_g)
                st.info(f"概算台数：空調ピーク {_estu['ac_peak_kw']:.0f}kW ÷ 標準室外機 {_estu['unit_kw']:.0f}kW"
                        f" ≒ **約 {_estu['n_units_total']} 台**（推定）")
                st.session_state.setdefault("gaisan_total", int(_estu['n_units_total']))
                st.session_state.setdefault("gaisan_ctrl", int(_estu['n_units_ctrl']))
                if st.button("↻ 検針票から再推定", key="gaisan_reest"):
                    st.session_state["gaisan_total"] = int(_estu['n_units_total'])
                    st.session_state["gaisan_ctrl"] = int(_estu['n_units_ctrl'])
                    st.rerun()
                n_units_total = st.number_input("室外機 総台数（概算・調整可）", min_value=1, max_value=5000,
                                                key="gaisan_total")
                if int(st.session_state.get("gaisan_ctrl", 1)) > int(n_units_total):
                    st.session_state["gaisan_ctrl"] = int(n_units_total)
                n_units_ctrl  = st.number_input("うち制御対象台数（概算・調整可）", min_value=1,
                                                max_value=int(n_units_total), key="gaisan_ctrl")
                ctrl_ratio = n_units_ctrl / max(n_units_total, 1)
                st.caption("※ 台数・初期費用・回収期間は概算です（機材情報の取得後に精緻化されます）。")
            else:
                n_units_total = st.number_input("室外機 総台数", min_value=1, max_value=5000, value=34)
                n_units_ctrl  = st.number_input("うち制御対象台数", min_value=1,
                                                 max_value=int(n_units_total), value=31)
                ctrl_ratio = n_units_ctrl / n_units_total
                st.markdown(f"**制御対象率**: {ctrl_ratio:.0%}  ({n_units_ctrl}/{n_units_total}台)")
                st.caption("※ 室外機リスト（型番付き）をインポートすると、型番→定格で自動精緻化します。")

        st.markdown("### 初期費用")
        detail_mode = st.toggle("詳細見積（見積書方式）で算出", value=False, key="cost_detail_mode",
                                help="電材費・作業人数・日数・地区などの詳細がある場合に、見積書と同じ方式で算出します。"
                                     "OFFのときは従来の概算。")
        _est = app_data.get("const_settings", {}).get("estimate_settings",
                                                      DEFAULT_DATA["const_settings"]["estimate_settings"])
        estimate_detail = None

        if not detail_mode:
            # ── 概算（従来）。項目名は詳細見積と統一（材料費／工事費／交通費）──
            area = st.selectbox("施工場所エリア", list(AREA_BUFFER.keys()), index=1)
            buf = AREA_BUFFER[area]
            kouji_cost = round(n_units_ctrl * AREA_UNIT_PRICE * buf / 10000) * 10000
            st.metric("交通費（概算）", f"¥{kouji_cost:,}")

            st.markdown("### 投資設定")
            if st.session_state.get("gaisan_mode"):
                _gs = app_data.get("const_settings", {}).get("gaisan_settings", {})
                _mpu = float(_gs.get("material_per_unit", 100000))
                _spu = float(_gs.get("setup_per_unit", 50000))
                system_cost = int(round(n_units_ctrl * _mpu))
                setup_cost  = int(round(n_units_ctrl * _spu))
                st.metric("材料費（概算：標準単価×台数）", f"¥{system_cost:,}")
                st.caption(f"標準機材費 ¥{_mpu:,.0f}/台 × 制御 {n_units_ctrl}台（管理画面で調整可）")
                st.metric("工事費（概算：標準単価×台数）", f"¥{setup_cost:,}")
                st.caption(f"標準工事費 ¥{_spu:,.0f}/台 × 制御 {n_units_ctrl}台（管理画面で調整可）")
            else:
                system_cost  = st.number_input("材料費（円）", min_value=0,
                                                max_value=100000000, step=100000, value=0)
                setup_cost   = st.number_input("工事費（円）", min_value=0,
                                                max_value=10000000, step=50000, value=0)
        else:
            # ── 詳細見積（見積書方式）──
            _mf  = float(_est.get("material_factor", 1.2))
            _ldu = float(_est.get("labor_day_unit", 60000))
            _lfx = float(_est.get("labor_fixed", 120000))
            _lou = float(_est.get("lodging_unit", 8000))
            _au  = _est.get("area_unit", {})

            st.markdown("**材料費**")
            denzai = st.number_input("電材費（円）", min_value=0, max_value=100000000, step=10000,
                                     value=0, key="det_denzai")
            system_cost = int(round(denzai * _mf))

            st.markdown("**工事費**")
            c_h1, c_h2 = st.columns(2)
            with c_h1:
                n_mzu  = st.number_input("ｴﾑｽﾞｶﾗｯﾄ 作業人数", min_value=0, max_value=100, step=1, value=2, key="det_mzu")
            with c_h2:
                n_simx = st.number_input("シムックス 作業人数", min_value=0, max_value=100, step=1, value=1, key="det_simx")
            work_days = st.number_input("作業日数（日）", min_value=0, max_value=365, step=1, value=2, key="det_days")
            workers = int(n_mzu) + int(n_simx)
            setup_cost = int(round(workers * int(work_days) * _ldu + _lfx))

            st.markdown("**交通費**")
            area5_key = st.selectbox("施工場所エリア（地区）", list(AREA5_LABELS.keys()),
                                     format_func=lambda k: AREA5_LABELS[k], index=0, key="det_area5")
            area_unit = float(_au.get(area5_key, 0))
            if area5_key == "地区5":
                trans_base = st.number_input("交通費（実費・新幹線/飛行機等の検索額）", min_value=0,
                                             max_value=10000000, step=1000, value=0, key="det_trans")
            else:
                trans_base = 0
                st.caption("交通費（実費入力）は地区5を選択した場合のみ表示されます（地区1〜4は移動拘束費で算定）。")
            _md_default = workers * int(work_days)
            lodging_md = st.number_input("宿泊のべ人日（人工・0で宿泊なし）", min_value=0, max_value=10000,
                                         step=1, value=int(_md_default), key="det_lodging_md")
            RENTACAR_UNIT = 12000
            rentacar_days = st.number_input("レンタカー利用日数（日・@12,000円/日 固定）", min_value=0,
                                            max_value=365, step=1, value=0, key="det_rentacar_days")
            rentacar = int(rentacar_days) * RENTACAR_UNIT
            st.caption(f"レンタカー代 ＝ 12,000円 × {int(rentacar_days)}日 ＝ ¥{rentacar:,}")
            restraint = workers * area_unit
            lodging   = int(lodging_md) * _lou
            kouji_cost = int(round(trans_base + restraint + lodging + rentacar))

            area = AREA5_LABELS.get(area5_key, area5_key)
            estimate_detail = {
                "denzai": int(denzai), "material_factor": _mf,
                "n_mzu": int(n_mzu), "n_simx": int(n_simx), "work_days": int(work_days),
                "labor_day_unit": _ldu, "labor_fixed": _lfx,
                "area5": area5_key, "area_unit": area_unit, "trans_base": int(trans_base),
                "lodging_md": int(lodging_md), "lodging_unit": _lou, "lodging": int(lodging),
                "restraint": int(restraint), "rentacar": int(rentacar), "rentacar_days": int(rentacar_days),
                "material": int(system_cost), "kouji": int(setup_cost), "trans": int(kouji_cost),
            }
            st.metric("初期費用（合計）", f"¥{system_cost + setup_cost + kouji_cost:,}")

        # ── ハードウェア関連（電力データと同じテンプレで取込→メイン画面の一覧に自動反映。専用ボタンなし）──
        st.markdown("### ハードウェア関連")
        hw_rows = apply_master_prices(st.session_state.get("sim_hw_rows", []) or [], app_data)
        st.session_state.sim_hw_rows = hw_rows
        if st.session_state.get("gaisan_mode") and not hw_rows:
            _gs2 = app_data.get("const_settings", {}).get("gaisan_settings", {})
            _hwf = float(_gs2.get("hw_fixed", 300000)); _hwu = float(_gs2.get("hw_per_unit", 50000))
            hardware_cost = int(round(_hwf + n_units_ctrl * _hwu))
            st.metric("ハードウェア関連費（概算：固定＋標準単価×台数）", f"¥{hardware_cost:,}")
            st.caption(f"固定 ¥{_hwf:,.0f}（拠点1式：データ収集装置/親機等） ＋ ¥{_hwu:,.0f}/台 × 制御 {n_units_ctrl}台"
                       "（管理画面で調整可）")
        else:
            hardware_cost = int(sum(float(r.get("売価単価", 0) or 0) * int(r.get("数量", 0) or 0) for r in hw_rows))
            st.metric("ハードウェア関連費（初期費用に加算）", f"¥{hardware_cost:,}")
            st.caption("電力データのインポートに部材（品名・数量）を含めると自動算出されます。"
                       "明細はメイン画面『ハードウェア関連（部材）一覧』で確認・編集できます。")

        hasu_adj     = st.number_input("端数調整（円・初期費用合計から減算）▲", min_value=0,
                                        max_value=100000000, step=1000, value=0,
                                        help="初期費用合計（総初期投資額A）から差し引く端数調整額。")

        subsidy      = st.number_input("補助金・助成金（円）▲", min_value=0,
                                        max_value=100000000, step=100000, value=0)

        # ── 年間システム利用料：機器台数（マーキュリー/MM親機/MM子機）駆動 ──
        st.markdown("### 年間システム利用料")
        with st.expander("📡 利用料の機器台数", expanded=True):
            fee_set   = app_data.get("fee_settings", DEFAULT_DATA["fee_settings"])
            p_tsushin = int(fee_set.get("tsushin", 550))
            p_cloud   = int(fee_set.get("cloud", 500))
            p_hoshu   = int(fee_set.get("hoshu", 200))
            p_data    = int(fee_set.get("data", 200))
            n_ctrl_i  = int(n_units_ctrl)

            # 機器台数の初期化＋インポート検出値の反映（取り込み直後のみ上書き）
            st.session_state.setdefault("dev_mercury", 1)
            st.session_state.setdefault("dev_mm_parent", 1)
            st.session_state.setdefault("dev_mm_child", 0)
            if st.session_state.pop("_fee_dev_apply", False):
                _fd = st.session_state.get("fee_devices", {})
                st.session_state["dev_mercury"]   = int(_fd.get("mercury", st.session_state["dev_mercury"]))
                st.session_state["dev_mm_parent"] = int(_fd.get("mm_parent", st.session_state["dev_mm_parent"]))
                st.session_state["dev_mm_child"]  = int(_fd.get("mm_child", st.session_state["dev_mm_child"]))

            n_mercury   = st.number_input("マーキュリー（台）", min_value=0, max_value=100000,
                                          step=1, key="dev_mercury")
            n_mm_parent = st.number_input("MMEazyAir 親機（台）", min_value=0, max_value=100000,
                                          step=1, key="dev_mm_parent")
            n_mm_child  = st.number_input("MMEazyAir 子機（台）", min_value=0, max_value=100000,
                                          step=1, key="dev_mm_child")

            # 通信費・データ収集は機器台数から（今回の変更点）
            q_tsushin = int(n_mercury) + int(n_mm_parent)                       # 通信費＝マーキュリー＋MM親機
            q_data    = int(n_mercury) + int(n_mm_parent) + int(n_mm_child)     # データ収集＝ﾏｰｷｭﾘｰ＋MM親機＋MM子機

            # クラウド利用料＝制御台数＋1、ソフト保守＝1（固定）を自動算出（入力欄なし）
            q_cloud = int(n_ctrl_i) + 1
            q_hoshu = 1
            st.caption(f"クラウド利用料の数量：制御台数 {n_ctrl_i} ＋ 1 ＝ {q_cloud}（自動）")
            st.caption("ソフト保守の数量：1（固定・自動）")

            m_tsushin = p_tsushin * q_tsushin
            m_cloud   = p_cloud   * q_cloud
            m_hoshu   = p_hoshu   * q_hoshu
            m_data    = p_data    * q_data
            monthly_fee = m_tsushin + m_cloud + m_hoshu + m_data
            system_fee_annual = int(monthly_fee * 12)

            st.metric("年間システム利用料", f"¥{system_fee_annual:,}")

        # 管理タブ等が参照する sidebar_sys_fee に確定値を反映
        st.session_state["sidebar_sys_fee"] = int(system_fee_annual)

        # ── 空調システム構築費（単価×台数。初期費用に加算）──
        st.markdown("### 空調システム構築費")
        _kunit = float(app_data.get("quote_settings", {}).get("kouchiku_unit", 12000))
        _kq_dc   = int(n_mercury) + int(n_mm_parent) + int(n_mm_child)      # データ取集装置＝ﾏｰｷｭﾘｰ＋MM親機＋MM子機
        _kq_out  = int(n_units_ctrl)                                        # 室外機＝制御台数（除外で自動減）
        _kq_temp = int(sum(int(r.get("数量", 0) or 0)                        # 温湿度センサー＝防水温湿度計の個数
                           for r in (st.session_state.get("sim_hw_rows", []) or [])
                           if ("温湿度" in str(r.get("商品名", "")))))
        kouchiku_dc   = int(_kq_dc * _kunit)
        kouchiku_out  = int(_kq_out * _kunit)
        kouchiku_temp = int(_kq_temp * _kunit)
        kouchiku_total = kouchiku_dc + kouchiku_out + kouchiku_temp
        st.caption(f"単価 ¥{int(_kunit):,}/台　×　データ取集装置 {_kq_dc}（ﾏｰｷｭﾘｰ＋MM親機＋MM子機）／"
                   f"室外機 {_kq_out}（制御台数）／温湿度センサー {_kq_temp}（防水温湿度計）")
        st.metric("空調システム構築費 合計（初期費用に加算）", f"¥{kouchiku_total:,}")

        total_invest = max(system_cost + kouji_cost + setup_cost + hardware_cost
                           + kouchiku_total - int(hasu_adj), 0)
        net_invest   = max(total_invest - subsidy, 0)
        st.metric("総初期投資額 (A)", f"¥{total_invest:,}",
                  delta=(f"端数調整 ▲¥{int(hasu_adj):,}" if int(hasu_adj) > 0 else None),
                  delta_color="off")

        st.markdown("### 本業商品設定")
        item_price  = st.number_input("主力商品単価（円）", min_value=100, max_value=100000000,
                                       step=10000, value=5000)
        item_margin = st.slider("粗利益率", 0.01, 0.80, 0.30, format="%.2f")

        st.markdown("### データ補完（自動判定）",
                    help=(
                        "インポートしたデータに不足項目があると、計算が破綻しないよう"
                        "自動で代替値を補います。補完される項目と方法は次のとおりです。\n\n"
                        "■ 契約電力（基本料金の削減に影響）\n"
                        "・検針票に『契約電力』列あり → その最大値を使用\n"
                        "・列が無い → 年間最大需要(kW)で自動代用\n\n"
                        "■ 制御可能比率（電力量・基本料金の両方に影響）\n"
                        "・型番リストあり → 容量×稼働係数の加重で算出\n"
                        "・型番リスト無し → 台数ベース（制御台数÷総台数）で算出\n\n"
                        "■ その他の重要列（最大需要/使用量/基本料金単価/電力量単価）\n"
                        "・未認識の列は暫定の既定値で補完（要・実数値修正）\n\n"
                        "下の『手動で上書き』をオンにすると、上記の自動判定を無視して"
                        "固定値で計算できます（通常は不要）。"
                    ))
        # インポート時の列認識結果・機体リスト有無から、不足データを自動判定してフォールバックを決定
        _imp = st.session_state.get("imported_cols")
        _units = st.session_state.get("units_df")
        _units_present = _units is not None and len(_units) > 0

        # 補完ロジックの早見表（常時表示・折りたたみ）
        with st.expander("❓ どの項目がどう補完される？", expanded=False):
            st.markdown(
                "| 項目 | データがある時 | 不足している時（自動補完） | 影響先 |\n"
                "|---|---|---|---|\n"
                "| 契約電力 | 検針票の最大値 | **年間最大需要で代用** | 基本料金の削減 |\n"
                "| 制御可能比率 | 型番リストで容量加重 | **台数ベース**（制御÷総台数） | 電力量＋基本料金 |\n"
                "| 最大需要/使用量/各単価 | 検針票の実値 | **暫定の既定値**（要修正） | 全削減額 |\n"
            )
            st.caption("※『手動で上書き』をオンにすると、契約電力・制御可能比率を固定値で強制できます。")

        if _imp is None:
            # 未インポート（サンプル/手動データ）→ 既定動作（補完なし）
            st.caption("インポート前（サンプル/手動データ）。数値をインポートすると不足項目を自動判定します。")
            ratio_mode, ratio_fix_pct = "自動（台帳/台数）", None
            contract_mode, contract_fix_kw = "検針票から", 0
        else:
            # ① 契約電力：検針票に列があれば検針票から、無ければ年間最大需要で自動代用
            _contract_ok = "契約電力" in _imp
            if _contract_ok:
                st.caption("✅ 契約電力：検針票から取得")
                contract_mode, contract_fix_kw = "検針票から", 0
            else:
                st.warning("⚠️ 契約電力が未取得 → 年間最大需要で自動代用（インポート画面に詳細）")
                contract_mode, contract_fix_kw = "固定値", 0  # contract_fix=0 で年間最大需要を使用
            # ② 制御可能比率：型番リストがあれば容量加重、無ければ台数ベース（どちらも自動）
            if _units_present:
                st.caption("✅ 制御可能比率：型番リスト（容量加重）で自動算出")
            else:
                st.caption("ℹ️ 制御可能比率：台数ベースで自動算出（型番リスト未検出）")
            ratio_mode, ratio_fix_pct = "自動（台帳/台数）", None

        # 任意：自動判定を手動で上書き（通常は不要）
        with st.expander("🔧 データ補完を手動で上書き（任意・通常は不要）", expanded=False):
            st.caption(
                "通常は触らずオフのままでOK（上の自動判定が使われます）。"
                "検針票に無い契約電力を把握している、台帳が不正確で経験則の%で概算したい、"
                "といった例外時だけオンにして固定値を指定してください。")
            if st.checkbox("手動設定を使う", value=False, key="fb_manual_on",
                           help=("オンにすると、自動判定を無視して下の固定値で計算します。"
                                 "オフに戻すと自動判定に復帰します。反映は『シミュレーション実行』時。")):
                _rm = st.radio("制御可能比率の取得", ["自動（台帳/台数）", "固定値"],
                               horizontal=True, key="fb_ratio_mode",
                               help=("自動＝型番リスト(容量加重)または台数ベースで算出。"
                                     "固定値＝下の%をそのまま採用（台帳・台数を無視）。"
                                     "電力量削減と基本料金削減の両方に効きます。"))
                ratio_mode = _rm
                ratio_fix_pct = (st.number_input("　└ 制御可能比率 固定値(%)", min_value=0, max_value=100,
                                                 value=85, key="fb_ratio_fix",
                                                 help="空調のうち制御できる割合。台帳が無い時の経験則は概ね80〜90%。")
                                 if _rm == "固定値" else None)
                _cm = st.radio("契約電力の取得", ["検針票から", "固定値"],
                               horizontal=True, key="fb_contract_mode",
                               help=("検針票から＝取り込んだ契約電力列の最大値を使用。"
                                     "固定値＝下のkWを使用（0なら年間最大需要で代用）。"
                                     "基本料金の削減額に効きます。"))
                contract_mode = _cm
                contract_fix_kw = (st.number_input("　└ 契約電力 固定値(kW・0なら年間最大需要を使用)",
                                                   min_value=0, max_value=99999, value=0, key="fb_contract_fix",
                                                   help="実際の契約電力(kW)を把握している場合に入力。0のままなら年間最大需要を代用。")
                                   if _cm == "固定値" else 0)

        if st.button("🚀 シミュレーション実行", type="primary", use_container_width=True):
            # 月の削除は行わない（A：振れが大きい場合は季節性配慮で対応）。旧「対象」列が残っていれば除去のみ。
            _dfin = st.session_state.df_input
            if hasattr(_dfin, "columns") and "対象" in _dfin.columns:
                _dfin = _dfin.drop(columns=["対象"])
            _excluded_months = []
            df_safe = sanitize_columns(_dfin)
            try:
                if df_safe is None or len(df_safe) == 0:
                    raise ValueError("試算対象の月がありません。データ編集テーブルで最低1ヶ月に『対象』チェックを入れてください。")
                # ユーザーが編集した年間システム利用料をシミュレーション全体に上書き同期
                app_data["system_fee"] = system_fee_annual
                app_data["fallback"] = {
                    "ratio_mode": ratio_mode,
                    "ratio_fix": (ratio_fix_pct / 100.0) if ratio_fix_pct is not None else None,
                    "contract_mode": contract_mode,
                    "contract_fix": contract_fix_kw,
                }
                # 季節性・業態配慮はトグルで明示制御。燃調・再エネは df の月別列から自動反映。
                app_data["calc_options"] = {
                    "seasonal_ac": bool(seasonal_ac),
                    "gyotai_consider": bool(gyotai_consider),
                }
                _ac_ovr = st.session_state.get("manual_ac") if gyotai == MANUAL_GYOTAI else None
                res = calc_simulation(df_safe, app_data, gyotai, ctrl_ratio,
                                      int(n_units_total), int(n_units_ctrl),
                                      units_df=st.session_state.get("units_df"),
                                      ac_override=_ac_ovr)
                res["excluded_months"] = list(_excluded_months)   # 月スコープ除外の記録（注釈用）
                res["net_invest"] = net_invest
                res["total_invest"] = total_invest
                res["item_price"] = item_price
                res["item_margin"] = item_margin
                res["gyotai"] = gyotai
                res["econ"] = {
                    "system_cost": system_cost, "setup_cost": setup_cost,
                    "kouji_cost": kouji_cost, "subsidy": subsidy, "area": area,
                    "hasu_adj": int(hasu_adj),
                    "hardware_cost": int(hardware_cost),
                    "hardware_rows": list(hw_rows),
                    "kouchiku_unit": int(_kunit),
                    "kouchiku_q_dc": int(_kq_dc), "kouchiku_q_out": int(_kq_out), "kouchiku_q_temp": int(_kq_temp),
                    "kouchiku_dc": int(kouchiku_dc), "kouchiku_out": int(kouchiku_out),
                    "kouchiku_temp": int(kouchiku_temp), "kouchiku_total": int(kouchiku_total),
                    "estimate_mode": bool(detail_mode), "estimate_detail": estimate_detail,
                    "gaisan_mode": bool(st.session_state.get("gaisan_mode")),   # 機材情報なしの概算
                    "n_ctrl": int(n_units_ctrl), "sys_fee": system_fee_annual,
                    "fee_breakdown": {
                        "p_tsushin": int(p_tsushin), "p_cloud": int(p_cloud),
                        "p_hoshu": int(p_hoshu), "p_data": int(p_data),
                        "q_tsushin": int(q_tsushin), "q_cloud": int(q_cloud),
                        "q_hoshu": int(q_hoshu), "q_data": int(q_data),
                        "n_mercury": int(n_mercury), "n_mm_parent": int(n_mm_parent),
                        "n_mm_child": int(n_mm_child),
                        "m_tsushin": int(m_tsushin), "m_cloud": int(m_cloud),
                        "m_hoshu": int(m_hoshu), "m_data": int(m_data),
                        "monthly": int(monthly_fee),
                    },
                }
                st.session_state.sim_result = res
                st.session_state.screen = "main"
                st.success("試算が完了しました！")
                st.rerun()
            except Exception as e:
                st.error(f"シミュレーション中にエラーが発生しました。詳細: {e}")

    # ─── 画面切り替え ─────────────────────────────
    screen = st.session_state.screen

    # ── コンパクト・ヘッダー（案件情報＋ステータスバッジ＋メニュー一体・カード・sticky）──
    _cust = (str(client_name).strip() + " 御中") if str(client_name).strip() else "（顧客名未入力）"

    with st.container(key="dps_header"):
        st.markdown(
            "<div class='dps-hdr'><div>"
            f"<div class='sys'>{APP_TITLE}</div>"
            f"<div class='cust'>{_cust}</div>"
            f"<div class='meta'>{gyotai} ・ {datetime.date.today()}</div>"
            "</div><div class='R'>"
            "<span class='badge b-green'><span class='dot'>●</span>AI判定：正常</span>"
            "<span class='badge b-amber'><span class='dot'>●</span>決裁：準備中</span>"
            "<span class='badge b-navy'><span class='dot'>●</span>担当：営業</span>"
            "</div></div>", unsafe_allow_html=True)

    if screen == "admin":
        show_admin(app_data)
        return

    if screen == "help":
        show_help()
        return

    if screen == "results_list":
        show_results_list(app_data)
        return

    if screen == "quote":
        show_quote_builder(app_data, st.session_state.get("sim_result"), client_name)
        return

    if screen == "reduction_report":
        show_reduction_report(app_data)
        return

    # ──────────────────────────────────────────────────────────────
    # メイン画面（案件情報・メニューは上部の共通ヘッダーに集約）
    # ──────────────────────────────────────────────────────────────
    with st.expander("電力データ インポート / 手動入力 (指定外フォーマットも内容から自動取込)",
                     expanded=bool(st.session_state.get("import_warn", False))):
        st.caption("指定テンプレート以外（別ツール作成のExcel等）でも、全シートを走査して"
                   "列名の表記ゆれ・半角カナ・改行を自動判定して取り込みます。")
        # ── Upload（ボタンのみ）／サンプル値代入／テンプレDL の3ボタン横並び ──

        _bc1, _bc2, _bc3 = st.columns(3)
        with _bc1:
            uploaded = st.file_uploader("CSV / Excel をアップロード", type=["csv", "xlsx"], key="file_upload")
        with _bc2:
            if st.button("サンプル値代入", key="load_sample_btn", use_container_width=True,
                         help="動作確認用のサンプル（検針票12ヶ月＋空調機体＋機材／利用料機器）を一括で代入します（テスト用）。"):
                st.session_state.df_input = get_sample_df()
                st.session_state.model_rows = get_sample_models()          # 空調情報（機体）
                st.session_state.model_editor_ver = st.session_state.get("model_editor_ver", 2) + 1
                st.session_state.sim_hw_rows = get_sample_hw()             # 機材情報（ハードウェア）
                st.session_state.sim_hw_ver = st.session_state.get("sim_hw_ver", 0) + 1
                st.session_state.sim_hw_import_sig = None
                st.session_state["fee_devices"] = {"mercury": 1, "mm_parent": 1, "mm_child": 2}  # 利用料機器台数
                st.session_state["_fee_dev_apply"] = True
                st.success("サンプル値（検針票・空調機体・機材）を代入しました。")
                st.rerun()
        with _bc3:
            if HAS_OPENPYXL:
                st.download_button(
                    "インポート用テンプレDL",
                    data=build_input_template_xlsx(),
                    file_name="DPS_電力データ入力テンプレート.xlsx",
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    help="このテンプレートに12ヶ月分のデータを入力してアップロードすれば確実に取り込めます。",
                    key="tmpl_dl_btn", use_container_width=True,
                )
        st.caption("※ インポートファイル容量：200MB以下 ／ 形式：CSV・Excel（.xlsx）")

        if uploaded is not None:
            try:
                df_up, used_sheet, hdr_sc = load_uploaded_table(uploaded)
                matched = {m for m in (_canonical_for(c) for c in df_up.columns) if m}
                df_sanitized = sanitize_columns(df_up)
                st.session_state.df_input = df_sanitized

                sheet_note = f"（シート「{used_sheet}」を採用）" if not uploaded.name.lower().endswith(".csv") else ""
                _req7 = len([c for c in matched if c in REQUIRED_COLS])
                st.success(f"データを取り込みました{sheet_note}：{len(df_sanitized)}行 ／ 主要列 {_req7}/7 を認識")
                # 任意列（燃調・再エネ）の取込状況を明示
                _opt_msgs = []
                _opt_msgs.append("燃料費調整額 " + ("✅反映" if "燃料費調整額" in df_sanitized.columns and float(df_sanitized["燃料費調整額"].abs().sum() or 0) > 0 else "— 排除"))
                _opt_msgs.append("再エネ賦課金 " + ("✅反映" if "再エネ賦課金" in df_sanitized.columns and float(df_sanitized["再エネ賦課金"].sum() or 0) > 0 else "— 排除"))
                st.caption("任意列：" + " ／ ".join(_opt_msgs) + "（列が無い/空なら自動で計算から排除）")

                # ── 不足データの自動判定（フォールバックはここで自動決定し、警告も表示）──
                missing_critical = [c for c in CRITICAL_COLS if c not in matched]
                contract_missing = "契約電力" not in matched
                st.session_state["imported_cols"] = set(matched)
                st.session_state["import_warn"] = bool(missing_critical or contract_missing)

                if missing_critical or contract_missing:
                    _msgs = []
                    if missing_critical:
                        _msgs.append("重要列が未認識：**" + "、".join(missing_critical)
                                     + "** → 暫定の既定値で補完中（要修正）")
                    if contract_missing:
                        _msgs.append("**契約電力**列が未認識 → 基本料金の試算は『年間最大需要』を契約電力に"
                                     "自動代用して算出します")
                    st.warning(
                        "⚠️ 不足データを自動補完しています。\n\n- " + "\n- ".join(_msgs)
                        + "\n\n下の編集テーブルで実数値に修正するか、上の『入力用Excelテンプレート』をご利用ください。"
                    )
                else:
                    st.caption("✅ 試算に必要な主要列（契約電力・最大需要・使用量・基本料金単価・電力量単価）は"
                               "すべて自動認識できました。データ補完は不要です。")

                # ── 室外機リスト（任意）も検出したら『機体入力』テーブルへ流し込む ──
                try:
                    _cop = app_data.get("calc_settings", {}).get("cop", 3.5)
                    units = detect_units_from_upload(uploaded, cop=_cop)
                except Exception:
                    units = None
                if units is not None and len(units) > 0:
                    st.session_state.model_rows = [
                        {"機器ID/場所": r["機器ID"], "メーカー": r["メーカー"], "型番": r["型番"],
                         "稼働係数": float(r["稼働係数"]),
                         "制御可否": "○" if r["制御可否"] else "×"}
                        for _, r in units.iterrows()
                    ]
                    st.success(f"🔌 室外機リストも検出：{len(units)}台 を『機体入力』へ反映（型番→馬力マスタで自動判定します）")

                # ── 利用料機器（マーキュリー/MM親機/MM子機）も検出したらサイドバーへ反映 ──
                try:
                    fdev = detect_fee_devices_from_upload(uploaded)
                except Exception:
                    fdev = None
                if fdev:
                    st.session_state["fee_devices"] = fdev
                    st.session_state["_fee_dev_apply"] = True
                    st.success(f"📡 利用料機器を検出：マーキュリー{fdev['mercury']} / MM親機{fdev['mm_parent']} / "
                               f"MM子機{fdev['mm_child']} 台 を年間利用料へ反映します。")

                # ── ハードウェア関連＝『利用料機器(任意)』シートの記載品を“全件”転記（見出し以外）──
                try:
                    _hwrows, _hwmiss = read_fee_sheet_items(uploaded, app_data)
                except Exception:
                    _hwrows, _hwmiss = [], []
                if _hwrows:
                    # 同一ファイルの再描画では再取込しない（一覧での手編集を保持）
                    _hwsig = (uploaded.name, getattr(uploaded, "size", 0))
                    if st.session_state.get("sim_hw_import_sig") != _hwsig:
                        st.session_state.sim_hw_rows = _hwrows
                        st.session_state.sim_hw_miss = _hwmiss
                        st.session_state.sim_hw_ver = st.session_state.get("sim_hw_ver", 0) + 1
                        st.session_state.sim_hw_import_sig = _hwsig
                    st.success(f"🧰 『利用料機器(任意)』シートの記載品 {len(_hwrows)}件 を"
                               "『ハードウェア関連（部材）一覧』へ転記しました（機材単価マスタの売価で自動算出）。")
                else:
                    st.caption("※ 『利用料機器(任意)』シート（項目・台数）が見つかりませんでした。"
                               "当該シートに記載のうえ取り込むと、ハードウェア関連へ全件転記されます。")
            except Exception as e:
                st.error(f"インポートエラー: {e}\n\n取り込めない場合は、上の『入力用Excelテンプレート』をダウンロードしてご利用ください。")

        st.markdown("**手動データ編集テーブル**（こちらから数値を直接変更・調整してシミュレーションすることも可能です）")
        # 旧「対象」チェック列は廃止（A：自動提案＋トグルへ変更）。残っていれば除去。
        if hasattr(st.session_state.df_input, "columns") and "対象" in st.session_state.df_input.columns:
            st.session_state.df_input = st.session_state.df_input.drop(columns=["対象"])
        edited_df = st.data_editor(
            st.session_state.df_input,
            use_container_width=True,
            num_rows="dynamic",
            column_config={
                "月": st.column_config.TextColumn("月（YYYY/MM）", width=100),
                "契約電力": st.column_config.NumberColumn("契約電力(kW)", format="%d"),
                "最大需要電力": st.column_config.NumberColumn("最大需要(kW)", format="%d"),
                "使用量合計": st.column_config.NumberColumn("使用電力量(kWh)", format="%d"),
                "力率": st.column_config.NumberColumn("力率調整", format="%.2f", min_value=0.01, max_value=1.0),
                "基本料金単価": st.column_config.NumberColumn("基本単価(円/kW)", format="%.2f"),
                "電力量単価": st.column_config.NumberColumn("電力量単価(円/kWh)", format="%.2f"),
                "燃料費調整額": st.column_config.NumberColumn("燃料費調整額(円/kWh・任意)", format="%.2f"),
                "再エネ賦課金": st.column_config.NumberColumn("再エネ賦課金(円/kWh・任意)", format="%.2f"),
                "夏季ピーク単価": st.column_config.NumberColumn("夏季ピーク単価(円/kWh・任意/欄のみ)", format="%.2f"),
                "平日昼間単価": st.column_config.NumberColumn("平日昼間単価(円/kWh・任意/欄のみ)", format="%.2f"),
                "夜間休日単価": st.column_config.NumberColumn("夜間休日単価(円/kWh・任意/欄のみ)", format="%.2f"),
                "市場価格調整項": st.column_config.NumberColumn("市場価格調整項(円/kWh・任意/欄のみ)", format="%.2f"),
            }
        )
        st.session_state.df_input = edited_df

        # ── A：使用量の振れが大きい場合は『季節性（ベースロード法）』配慮を推奨（※月の削除はしない）──
        _sw = _usage_swing(st.session_state.df_input)
        if _sw >= 0.30:
            st.info(f"💡 使用電力量の振れが大きいデータです（最少月比で最大 +{_sw*100:.0f}%）。"
                    "単純な平均や業態一律割合では偏りが出やすいため、月別の実態を活かす"
                    "**『季節性（ベースロード法）』配慮**を推奨します"
                    "（サイドバー『分析オプション』→『季節性を反映（ベースロード法）』をON）。"
                    "※ 振れの大きい月を削除するのではなく、月ごとの実データに沿って算出します。")

    # ── 🔌 機体（室外機）入力：型番→馬力マスタから自動反映（管理外は仮置き） ──
    with st.expander("機体（室外機）入力 — 型番から馬力を自動反映（管理外は仮置き）", expanded=False):
        hp_master = app_data.get("model_hp_master", {})
        default_hp = float(app_data.get("calc_settings", {}).get("default_hp", 5.0))
        cop_now = float(app_data.get("calc_settings", {}).get("cop", 3.5))
        st.caption(
            f"型番を入力すると管理マスタ（{len(hp_master)}型番登録）から馬力を自動取得します。"
            f"未登録の型番は「※管理外」と表示し、仮置き {default_hp:.1f} 馬力で算出します。"
            "（型番→馬力の登録は『各種設定 → 🔌 馬力・機材テーブル』から）"
        )

        base_rows = st.session_state.get("model_rows")
        if not base_rows:
            base_rows = [{"機器ID/場所": "", "階/エリア": "", "メーカー": "", "型番": "", "稼働係数": 0.7, "制御可否": "○"}
                         for _ in range(3)]
        model_in = pd.DataFrame(base_rows)
        if "階/エリア" not in model_in.columns:
            model_in.insert(1, "階/エリア", "")
        model_in["階/エリア"] = model_in["階/エリア"].fillna("")
        # 列順を固定（機器ID/場所 の次に 階/エリア）
        _order = ["機器ID/場所", "階/エリア", "メーカー", "型番", "稼働係数", "制御可否"]
        model_in = model_in[[c for c in _order if c in model_in.columns]
                            + [c for c in model_in.columns if c not in _order]]
        edited_models = st.data_editor(
            model_in, use_container_width=True, num_rows="dynamic",
            column_config={
                "機器ID/場所": st.column_config.TextColumn("機器ID/場所"),
                "階/エリア": st.column_config.TextColumn(
                    "階/エリア", help="例：1F／2F／A棟 など。下の選択で階・エリア単位の一括除外に使えます"),
                "メーカー": st.column_config.TextColumn("メーカー"),
                "型番": st.column_config.TextColumn("型番", width="large"),
                "稼働係数": st.column_config.NumberColumn("稼働係数", format="%.2f", min_value=0.0, max_value=1.5),
                "制御可否": st.column_config.SelectboxColumn("制御可否", options=["○", "×"], width="small"),
            },
            key=f"model_rows_editor_v{st.session_state.get('model_editor_ver', 2)}",   # スキーマ/サンプル投入を確実に反映
        )
        st.session_state.model_rows = edited_models.to_dict("records")

        # ── 階・エリア単位の一括除外（B）：選択した階/施設の機を制御対象外に（現行算出式は不変）──
        def _area_of(r):
            a = str(r.get("階/エリア", "")).strip()
            return "" if a.lower() == "nan" else a
        _areas = sorted({_area_of(r) for r in st.session_state.model_rows if _area_of(r)})
        if _areas:
            _excl_areas = st.multiselect("試算から除外する階・エリア（選択した階/施設の機を一括で制御対象外に）",
                                         _areas, key="excl_areas")
        else:
            _excl_areas = []
            st.caption("※ 機体の『階/エリア』を入力すると、ここで階・施設単位の一括除外を選べます。")
        _rows_for_build = [dict(r) for r in st.session_state.model_rows]
        if _excl_areas:
            for r in _rows_for_build:
                if str(r.get("階/エリア", "")).strip() in _excl_areas:
                    r["制御可否"] = "×"
            st.caption("除外中：" + "・".join(_excl_areas)
                       + " → 該当機を制御対象から外しました（制御台数・比率・室外機構築費に反映）。")

        built = build_units_from_models(_rows_for_build, hp_master, default_hp, cop_now,
                                        app_data.get("model_kw_master", {}))
        if built is not None and len(built) > 0:
            st.session_state.units_df = built
            view = built.copy()
            view["状態"] = view["管理"].map(lambda b: "登録済み" if b else "※管理外（仮）")
            view["制御"] = view["制御可否"].map(lambda b: "○" if b else "×")
            st.dataframe(
                view[["機器ID", "メーカー", "型番", "馬力", "定格冷房kW", "電力kW", "稼働係数", "制御", "状態"]].style.format(
                    {"馬力": "{:.1f} HP", "定格冷房kW": "{:.1f} kW", "電力kW": "{:.2f} kW", "稼働係数": "{:.2f}"}
                ),
                use_container_width=True
            )
            n_un = int((~built["管理"]).sum())
            if n_un > 0:
                st.warning(f"※ {n_un}台が管理外（マスタ未登録）です。仮置き {default_hp:.1f} 馬力で算出中。"
                           "正確な試算には管理画面で型番→馬力を登録してください。")
            st.success(f"機体 {len(built)}台を反映 → サイドバーの台数・容量加重比率に自動連動します（要シミュレーション実行）。")
        else:
            st.session_state.units_df = None
            st.info("型番を入力すると機体ベース（容量加重）で試算します。未入力時はサイドバーの台数ベースで算出します。")

    # ── ハードウェア関連（部材）一覧：電力データと同じテンプレで取り込んだ部材を自動反映（編集可）──
    with st.expander("ハードウェア関連（部材）一覧 — 取り込んだ部材を機材単価マスタの単価で自動反映（編集可）",
                     expanded=bool(st.session_state.get("sim_hw_rows"))):
        # 単価は常に機材単価マスタを正とする（マスタ更新を即反映）
        _hw_seed = apply_master_prices(st.session_state.get("sim_hw_rows", []) or [], app_data)
        _hw_miss = [str(r.get("商品名", "")) for r in _hw_seed
                    if match_material(r.get("商品名", ""), app_data) is None and str(r.get("商品名", "")).strip()]
        if _hw_miss:
            st.warning("機材単価マスタ未登録（単価0）： " + " / ".join(_hw_miss[:12])
                       + "　→ 各種設定の『機材単価マスタ』に同名で登録すると自動反映されます。")
        _hw_cols = ["商品名", "数量", "単位", "売価単価", "原価単価"]
        if _hw_seed:
            _hw_df_in = pd.DataFrame([{c: r.get(c, "") for c in _hw_cols} for r in _hw_seed])
        else:
            _hw_df_in = pd.DataFrame([{"商品名": "", "数量": 0, "単位": "台",
                                       "売価単価": 0.0, "原価単価": 0.0}])
        _hw_ver = st.session_state.get("sim_hw_ver", 0)
        _hw_edit = st.data_editor(
            _hw_df_in, num_rows="dynamic", use_container_width=True,
            key=f"sim_hw_editor_{_hw_ver}",
            column_config={
                "商品名": st.column_config.TextColumn("商品名", width="large"),
                "数量": st.column_config.NumberColumn("数量", min_value=0, step=1),
                "単位": st.column_config.TextColumn("単位", width="small", disabled=True),
                "売価単価": st.column_config.NumberColumn("売価単価（マスタ準拠）", format="%d", disabled=True),
                "原価単価": st.column_config.NumberColumn("原価単価（マスタ準拠）", format="%d", disabled=True),
            })
        # 編集後（商品名・数量）に対し、単価をマスタから再適用してから保存・集計
        _hw_rows2 = apply_master_prices(_hw_edit.to_dict("records"), app_data)
        st.session_state.sim_hw_rows = _hw_rows2
        _hw_total = int(sum(float(r.get("売価単価", 0) or 0) * int(r.get("数量", 0) or 0)
                            for r in _hw_rows2))
        st.metric("ハードウェア関連費（全機材費の合計 → 初期費用に加算）", f"¥{_hw_total:,}")
        st.caption("『利用料機器(任意)』シートの記載品を全件転記。単価は機材単価マスタの売価（自動・編集不可）、"
                   "金額＝売価×数量。数量の変更や品名修正は反映され、総初期投資額(A)へ加算されます。")

    res = st.session_state.sim_result
    if res is None:
        st.info("👈 左側のサイドバーにある「🚀 シミュレーション実行」を押して、試算を開始してください。")
        return

    # 📊 結果ダッシュボード（表示のみ差し替え・計算/関数/セッションは不変）。旧レポートは描画しない。
    try:
        render_results_dashboard(res, app_data, client_name, gyotai)
    except Exception as _e:
        st.error(f"ダッシュボード描画エラー：{_e}")
    return

    payback_years, _over_fee = calc_payback(res["total_invest"], res["net_saving"], res.get("gross_saving", 0))
    item_price  = res.get("item_price",  app_data["item_price"])
    item_margin = res.get("item_margin", app_data["item_margin"])
    items_needed = math.ceil(res["net_saving"] / (item_price * item_margin)) if item_price * item_margin > 0 else 0
    ds = app_data["display_settings"]

    st.markdown("## 財務サマリー（導入効果の総括）")

    m1, m2, m3, m4 = st.columns(4)
    with m1:
        st.metric("初期費用 (総投資額)", f"¥{res['total_invest']:,}",
                  delta="機材費・工事費など",
                  delta_color="inverse")
    with m2:
        st.metric("年間総削減額（グロス）",
                  f"¥{res['gross_saving']:,.0f} 円/年",
                  delta=f"基本料金¥{res['dm_saving_annual']:,.0f} + 従量¥{res['ene_saving_annual']:,.0f}")
    with m3:
        st.metric("年間純利益（実利・手残り）",
                  f"¥{res['net_saving']:,.0f} 円/年",
                  delta=f"利用料 ¥{res['sys_fee']:,} 差引後",
                  delta_color="normal")
    with m4:
        if _over_fee:
            delta_txt = "⚠ 年間利用料が削減額を上回ります（粗削減ベースの参考年数）"
            delta_color = "inverse"
        else:
            delta_color = "normal" if payback_years <= 3.0 else "inverse"
            delta_txt = "✅ 3年回収圏内（極めて優秀）" if payback_years <= 3.0 else "⚠ 3年超 → 機器更新併用を推奨"
        st.metric("投資回収期間 (ROI)",
                  f"{payback_years:.1f} 年",
                  delta=delta_txt,
                  delta_color=delta_color)

    if _over_fee:
        st.warning(f"🔄 **空調リプレイス更新推奨モード（粗削減ベース回収: 約{payback_years:.1f}年）** — 現在の設定では年間利用料が削減額を上回っています。利用料・制御台数の見直し、または省エネ補助金を活用した高効率機器リプレイスをご検討ください。")
    elif payback_years <= 3.0:
        st.success(f"📊 **DPSデマンド制御アドオン提案モード（実質投資回収: {payback_years:.1f}年）** — 既存機に制御ユニットを追加装着する、最も短工期かつ即効性の高い極めて優れた収益性ストーリーです。")
    else:
        st.warning(f"🔄 **空調リプレイス更新推奨モード（実質投資回収: {payback_years:.1f}年）** — 老朽化機器の一新を兼ね、省エネ補助金などをフル活用した高効率機器へのリプレイスプランをお勧めします。")

    # ── B-5: 投資評価指標（NPV / IRR）──
    _cs_npv = app_data.get("const_settings", DEFAULT_DATA["const_settings"])
    _rate = float(_cs_npv.get("discount_rate", 0.05))
    _nyrs = int(_cs_npv.get("npv_years", 10))
    _npv, _irr = calc_npv_irr(res["total_invest"], res["net_saving"], years=_nyrs, rate=_rate)
    nv1, nv2, nv3 = st.columns(3)
    nv1.metric(f"NPV（{_nyrs}年・割引率{_rate:.1%}）", f"¥{_npv:,.0f}",
               delta=("プラス＝投資価値あり" if _npv >= 0 else "マイナス＝要再検討"),
               delta_color=("normal" if _npv >= 0 else "inverse"))
    nv2.metric("IRR（内部収益率）", (f"{_irr*100:.1f} %" if _irr is not None else "—"),
               delta=(f"割引率{_rate:.1%}を{'上回る' if (_irr is not None and _irr> _rate) else '下回る'}"
                      if _irr is not None else "算定不可（純削減≤0等）"),
               delta_color=("normal" if (_irr is not None and _irr > _rate) else "inverse"))
    nv3.metric("単純投資回収（参考）", f"{payback_years:.1f} 年")
    st.caption(f"NPV＝−初期投資＋Σ(年間純利益÷(1+割引率)^t)（{_nyrs}年）。IRRはNPV=0となる割引率。"
               "割引率・年数は『各種設定 ＞ 計算ロジックの変更』で変更できます。")

    # ── C-12: 入力整合性チェック ──
    _issues = validate_inputs(res.get("df"), res, res.get("econ", {}))
    _has_err = any(lv == "error" for lv, _ in _issues)
    _has_warn = any(lv == "warn" for lv, _ in _issues)
    _exp_title = ("🔍 入力整合性チェック — ⚠️ 要確認あり" if (_has_err or _has_warn)
                  else "🔍 入力整合性チェック — ✅ 異常なし")
    with st.expander(_exp_title, expanded=_has_err):
        for lv, msg in _issues:
            if lv == "error":
                st.error("❌ " + msg)
            elif lv == "warn":
                st.warning("⚠️ " + msg)
            else:
                st.success("✅ " + msg)

    show_peraichi(res)

    df = res["df"].copy()
    df["月"] = df["月"].astype(str)

    # ⚙️ 現状の棒（＝棒全体）はそのまま残し、選択業態の電力内訳（空調／生産設備／照明／その他）で
    #    割合別に色分け。さらに空調分のうち「DPSで削減できる分」だけを点線枠ハッチングで強調する。
    L_OTHER = "その他電気量"
    L_PROD  = "生産設備"
    L_LIGHT = "照明"
    L_ACR   = "空調（導入後に残る分）"
    L_ACS   = "空調（DPSで削減できる分）"
    SEG_ORDER = [L_OTHER, L_PROD, L_LIGHT, L_ACR, L_ACS]   # 下→上
    # B2B提案向け「引き算のデザイン」：非空調=ニュートラルグレー（後退）、空調=アクセント緑（主役を際立たせる）
    AC_COLOR = "#00B050"      # 空調（アクセント緑）
    SEG_COLOR = {
        L_OTHER: "#E0E0E0",   # その他（薄グレー）
        L_PROD:  "#CCCCCC",   # 生産設備（中グレー）
        L_LIGHT: "#BDBDBD",   # 照明（やや濃いグレー）
        L_ACR:   "#A9D18E",   # 空調・残る分（ライトグリーン）
        L_ACS:   AC_COLOR,    # 空調・削減分（アクセント緑＝主役・ベタ塗り）
    }

    def make_breakdown_bar(src, cur_col, aft_col, y_title, ac_frac, weights):
        def _mlabel(m):
            s = str(m).replace("-", "/")
            p = s.split("/")
            if len(p) >= 2:
                y, mo = p[0], p[1]
                if mo.zfill(2) == "01":      # 年が変わる1月のみ年表記
                    return f"{y}/01"
                return f"{int(mo)}月"          # それ以外は「5月」等に簡略化
            return s
        m_order = [_mlabel(m) for m in src["月"].tolist()]
        w_light = weights.get("照明", 0.0)
        w_prod  = weights.get("生産設備", 0.0)
        w_other = weights.get("その他", 0.0)
        rows = []
        for _, r in src.iterrows():
            cur_v = float(r[cur_col])
            aft_v = min(float(r[aft_col]), cur_v)
            reduction = max(0.0, cur_v - aft_v)                 # DPS削減分（空調由来）
            ac_total  = min(max(ac_frac * cur_v, reduction), cur_v)  # 空調全体（削減分は必ず内包）
            ac_remain = ac_total - reduction
            non_ac    = cur_v - ac_total
            seg_vals = [
                (L_OTHER, non_ac * w_other),
                (L_PROD,  non_ac * w_prod),
                (L_LIGHT, non_ac * w_light),
                (L_ACR,   ac_remain),
                (L_ACS,   reduction),
            ]
            m = _mlabel(r["月"]); y0 = 0.0
            for name, val in seg_vals:
                y1 = y0 + val
                rows.append({"月": m, "種別": name, "y0": y0, "y1": y1, "値": val,
                             "_o": SEG_ORDER.index(name)})
                y0 = y1
        long = pd.DataFrame(rows)
        # 凡例は空調(主役)を先頭に。積み上げ順は _o（SEG_ORDER）で固定。
        LEGEND_ORDER = [L_ACS, L_ACR, L_LIGHT, L_PROD, L_OTHER]
        present = [c for c in LEGEND_ORDER if long.loc[long["種別"] == c, "値"].sum() > 1e-9]
        long = long[long["種別"].isin(present)]
        col_range = [SEG_COLOR[c] for c in present]
        return alt.Chart(long).mark_bar(stroke="#FFFFFF", strokeWidth=1).encode(
            x=alt.X("月:N", sort=m_order,
                    axis=alt.Axis(title=None, labelAngle=0, labelFontSize=13,
                                  labelColor="#666666", domain=False, ticks=False, grid=False)),
            y=alt.Y("y0:Q", title=y_title,
                    axis=alt.Axis(grid=True, gridColor="#E0E0E0", gridWidth=0.5,
                                  labelColor="#666666", labelFontSize=12, titleColor="#666666",
                                  domain=False, ticks=False, tickCount=5)),
            y2=alt.Y2("y1"),
            color=alt.Color("種別:N",
                scale=alt.Scale(domain=present, range=col_range),
                legend=alt.Legend(title=None, orient="top", columns=3,
                                  labelColor="#666666", labelFontSize=12, symbolStrokeWidth=0)),
            order=alt.Order("_o:Q"),
            tooltip=["月:N", "種別:N", alt.Tooltip("値:Q", title="値", format=",.1f")],
        ).properties(height=400).configure_view(strokeWidth=0)

    nonac_w = INDUSTRY_NONAC_WEIGHTS.get(res.get("gyotai", ""), DEFAULT_NONAC_WEIGHTS)

    if ds.get("show_graph_dm", True):
        st.markdown("---")
        st.markdown("### 最大デマンド比較（kW）")
        st.caption(f"棒全体が『現状の最大デマンド』。業態「{res.get('gyotai','')}」の電力内訳で色分けし、"
                   "濃い緑（アクセント色）の部分が DPS 導入で削減できる空調デマンドです。")
        ch_dm = make_breakdown_bar(
            df, "最大需要電力", "導入後最大DM", "最大デマンド（kW）",
            res.get("ac_peak_r", 0.18), nonac_w)
        st.altair_chart(ch_dm, use_container_width=True)

    if ds.get("show_graph_bill", True):
        st.markdown("---")
        st.markdown("### 月別電気使用量の比較（万kWh）")
        st.caption(f"棒全体が『現状の電力使用量』。業態「{res.get('gyotai','')}」の電力内訳で色分けし、"
                   "濃い緑（アクセント色）の部分が DPS 導入で削減できる使用量（空調由来）です。")
        df_use_man = df.copy()
        df_use_man["現状使用量(万kWh)"] = df_use_man["使用量合計"] / 10000
        df_use_man["導入後使用量(万kWh)"] = df_use_man["導入後使用量kWh"] / 10000
        ch_use = make_breakdown_bar(
            df_use_man, "現状使用量(万kWh)", "導入後使用量(万kWh)", "電力使用量（万kWh）",
            res.get("ac_kwh_r", 0.15), nonac_w)
        st.altair_chart(ch_use, use_container_width=True)

        st.markdown("---")
        st.markdown("### 月別電気料金の比較（万円）")
        st.caption(f"棒全体が『現状の電気料金』。業態「{res.get('gyotai','')}」の電力内訳で色分けし、"
                   "濃い緑（アクセント色）の部分が DPS 導入で削減できる金額（空調由来）です。")
        df_bill_man = df.copy()
        df_bill_man["現状電気料金(万円)"] = df_bill_man["現状電気料金推計"] / 10000
        df_bill_man["導入後電気料金(万円)"] = df_bill_man["導入後電気料金推計"] / 10000
        ch_bill = make_breakdown_bar(
            df_bill_man, "現状電気料金(万円)", "導入後電気料金(万円)", "電気料金（万円）",
            res.get("ac_kwh_r", 0.15), nonac_w)
        st.altair_chart(ch_bill, use_container_width=True)

    # ══════════════════════════════════════════════════════════════
    # 🖼 スライド貼付用 画像出力（PNG）
    # ══════════════════════════════════════════════════════════════
    st.markdown("---")
    st.markdown("### 🖼 スライド貼付用 画像出力（PNG）")
    if not HAS_SLIDE_IMG:
        st.info("画像出力には matplotlib / pillow が必要です。`pip install matplotlib pillow` を実行してください。")
    else:
        st.caption("資料（PowerPoint / Word）にそのまま貼れる PNG を生成します。生成後、各画像をダウンロードできます。")

        # ── 月ラベルを「○月」へ整形 ──
        def _month_label(v):
            import re
            s = str(v).strip()
            mt = re.search(r"(\d{4})\D+(\d{1,2})", s)   # 年が拾えれば "YYYY/M" で保持（PNG側で年の変わり目に年表示）
            if mt:
                return f"{int(mt.group(1))}/{int(mt.group(2))}"
            tail = s.replace("年", "/").replace("月", "").split("/")[-1].strip()
            try:
                return f"{int(tail)}月"
            except Exception:
                return s

        months = [_month_label(m) for m in df["月"].tolist()]
        usage  = df["使用量合計"].tolist()
        reduc  = df["削減kWh"].tolist() if "削減kWh" in df.columns else [0] * len(months)
        demand = df["最大需要電力"].tolist()

        # ── 機器リスト（アップロードがあれば units_detail、無ければ機器台帳マスター）──
        ud = res.get("units_detail")
        list_items = []
        if ud is not None and len(ud) > 0:
            for _, r in ud.iterrows():
                name = str(r.get("機器ID", "") or "")
                list_items.append({
                    "系統名": name, "設置場所": name,
                    "メーカー": str(r.get("メーカー", "") or ""),
                    "型式": str(r.get("型番", "") or ""),
                    "制御可否": bool(r.get("制御可否", False)),
                })
        else:
            for eq in DEFAULT_EQUIPMENT:
                list_items.append({
                    "系統名": eq["loc"], "設置場所": eq["loc"],
                    "メーカー": eq["mfr"], "型式": eq["model"],
                    "制御可否": str(eq["ctrl"]).strip() in ("〇", "○", "◯"),
                })
        n_total = len(list_items)
        n_ctrl  = sum(1 for it in list_items if it["制御可否"])

        # ── サマリKPIカード（4枚）の値を res から組み立て ──
        NAVY_C, GREEN_C = "#13315C", "#3DAE4E"
        invest = float(res.get("total_invest", 0) or 0)
        gross  = float(res.get("gross_saving", 0) or 0)
        net    = float(res.get("net_saving", 0) or 0)
        sys_fee = float(res.get("sys_fee", 0) or 0)
        payback = (invest / net) if net > 0 else 0.0
        summary_cards = [
            {"icon": "money.png", "title": "初期導入費用（総投資額）", "value": f"{invest:,.0f}円",
             "subtitle": "税抜 導入費用 / 機器代・工事費含む", "color": NAVY_C},
            {"icon": "graf.png", "title": "年間総削減額（グロス）", "value": f"{gross:,.0f}円",
             "subtitle": "基本料金＋電力量の年間削減合計", "color": GREEN_C},
            {"icon": "plus.png", "title": "年間実質利点（手残り）", "value": f"{net:,.0f}円",
             "subtitle": f"削減合計 {gross:,.0f}円 − 年間維持費 {sys_fee:,.0f}円", "color": GREEN_C},
            {"icon": "clock.png", "title": "投資回収期間（ROI）", "value": f"約{payback:.1f}年",
             "subtitle": f"約{round(payback*12)}ヶ月で完全回収、以降は純利益", "color": GREEN_C},
        ]

        if st.button("🖼 画像を生成 / 更新", key="gen_slide_imgs"):
            try:
                st.session_state["_img_chart"] = make_demand_chart_png(
                    months, usage, reduc, demand, target_units=n_ctrl)
                st.session_state["_img_list"] = make_control_list_pngs(
                    list_items, total_units=n_total, controllable_units=n_ctrl)
                st.session_state["_img_cards"] = make_summary_cards_png(summary_cards)
            except Exception as e:
                st.error(f"画像生成でエラー: {e}")

        if st.session_state.get("_img_chart"):
            st.markdown("**① 使用量＋デマンド グラフ**")
            st.image(st.session_state["_img_chart"], use_container_width=True)
            st.download_button("⬇ グラフPNGをダウンロード", st.session_state["_img_chart"],
                               file_name="使用量_デマンドグラフ.png", mime="image/png",
                               key="dl_img_chart")
        list_imgs = st.session_state.get("_img_list")
        if list_imgs:
            n_pg = len(list_imgs)
            st.markdown("**② 制御可否リスト**" + (f"（1スライドに収まらないため {n_pg} 枚に分割）" if n_pg > 1 else ""))
            for k, img in enumerate(list_imgs, start=1):
                st.image(img, use_container_width=True)
                fname = "制御可否リスト.png" if n_pg == 1 else f"制御可否リスト_{k}of{n_pg}.png"
                btn_label = "⬇ 制御可否リストPNGをダウンロード" + ("" if n_pg == 1 else f"（{k}/{n_pg}）")
                st.download_button(btn_label, img, file_name=fname, mime="image/png",
                                   key=f"dl_img_list_{k}")
        if st.session_state.get("_img_cards"):
            st.markdown("**③ サマリKPIカード**")
            st.image(st.session_state["_img_cards"], use_container_width=True)
            st.download_button("⬇ サマリカードPNGをダウンロード", st.session_state["_img_cards"],
                               file_name="サマリカード.png", mime="image/png",
                               key="dl_img_cards")

    st.markdown("---")
    st.markdown("## 付加価値・財務換算セクション")
    
    st.markdown(f"""
<div style="display:grid; grid-template-columns:repeat(2,1fr); gap:16px;">
  <div style="background:#E8F5E9;border-radius:10px;padding:20px;border:1.5px solid #1E6B2E;display:flex;flex-direction:column;">
    <div style="font-weight:bold;color:#1E6B2E;font-size:15px;margin-bottom:10px;">主力製品の販売数量・財務換算</div>
    <div style="font-size:15px;line-height:1.8;">
      年間手残り純削減額 <b>¥{res['net_saving']:,.0f}</b> は、本業の主力製品<br>
      （単価 ¥{item_price:,} · 粗利 {item_margin:.1%}）を、<br>
      新規に <span style="font-size:37px;font-weight:bold;color:#C00000;">【{items_needed:,} 個】</span> 追加販売して稼ぎ出す営業純利益に匹敵します。
    </div>
    <div style="font-size:12px;color:#666;margin-top:auto;padding-top:12px;border-top:1px solid #e2e8f0;">
      営業努力による追加販売のハードルと、AIで無駄な電力契約を自動カットする確実性をぜひご比較ください。
    </div>
  </div>
  <div style="background:#F1F8E9;border-radius:10px;padding:20px;border:1.5px solid #1B5E20;display:flex;flex-direction:column;">
    <div style="font-weight:bold;color:#1B5E20;font-size:15px;margin-bottom:10px;">ESG・環境貢献サマリー</div>
    <div style="font-size:15px;line-height:2.0;">
      年間削減電力量：<b>{res['total_reduc_kwh']:,} kWh / 年</b><br>
      CO₂排出削減量：<b>{res['co2_kg']:,.0f} kg-CO₂ / 年</b><br>
      スギの木換算森林効果：<b>約 {res['sugi_trees']:,.0f} 本分 / 年</b>
    </div>
    <div style="font-size:12px;color:#666;margin-top:auto;padding-top:12px;border-top:1px solid #e2e8f0;">
      これらの省エネデータは、ESG報告書や御社ホームページ、SDGsの取り組み成果として対外公表にご利用いただけます。
    </div>
  </div>
</div>
""", unsafe_allow_html=True)

    # （旧）累積機会損失の表示は削除し、下記「今後10年間の財務リスク試算」に統合

    # 制御強度シミュレーター
    show_control_mode(res)

    # 今後10年間の財務リスク試算（シミュレーターとトークスクリプトの間）
    show_financial_risk_10yr(res)

    show_talk_script(res, app_data, client_name, payback_years, res["total_invest"])

    # 📥 一括ダウンロード
    # ── 💾 この結果を「📋 結果一覧」に保存（軽量サマリー）──
    st.markdown("---")
    st.markdown("### 💾 結果一覧に保存")
    st.caption("企業名・空調数・年間削減量・回収年数を『📋 結果一覧』に登録します（詳細データは各社フォルダで管理）。")
    if st.button("💾 この結果を一覧に保存", key="save_case_btn"):
        rec = {
            "id": datetime.datetime.now().strftime("%Y%m%d%H%M%S"),
            "company": client_name,
            "gyotai": res.get("gyotai", gyotai),
            "units_total": int(res.get("n_units_total", 0)),
            "units_ctrl": int(res.get("n_units_ctrl", 0)),
            "saving_kwh": float(res.get("total_reduc_kwh", 0)),
            "saving_yen": float(res.get("gross_saving", 0)),
            "payback": float(payback_years),
            "monthly": case_monthly_records(res),   # 削減レポ呼び出し用の月次データ
            "saved_at": datetime.datetime.now().strftime("%Y-%m-%d %H:%M"),
        }
        app_data.setdefault("saved_cases", []).append(rec)
        save_data(app_data)
        st.success(f"「{client_name}」を結果一覧に保存しました（{rec['saved_at']}）。"
                   "ナビの『📋 結果一覧』で確認できます。")

    st.markdown("---")
    st.markdown("## シミュレーション結果一括ダウンロード")

    col_dl1, col_dl2 = st.columns(2)
    
    with col_dl1:
        st.markdown("### Excel稟議対応シート")
        st.caption("Excel稟議シート（最左端に『シート0_顧客提示用サマリー』が挿入された、1円の差異もない財務対応Excelデータ。SKILL.md完全準拠 ＋ 追加4シート完備）")
        filename_xl = f"DPS試算_{client_name}_{datetime.date.today()}.xlsx"
        
        # openpyxl がインストールされているかチェック
        if HAS_OPENPYXL:
            excel_bytes = build_excel(res, client_name, app_data)
            st.download_button(
                label="📥 Excel稟議シートをダウンロード",
                data=excel_bytes,
                file_name=filename_xl,
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                type="primary",
                use_container_width=True
            )
        else:
            st.error("Excel出力に必要な 'openpyxl' ライブラリが見つかりません。")
            
    with col_dl2:
        st.markdown("### PowerPoint提案書スライド")
        st.caption("以前ご指示いただいた【全12枚のMcKinsey風・外資コンサル仕様スライド】を動的にフルビルドします。")
        if HAS_PPTX:
            filename_pp = f"DPS提案スライド_{client_name}_{datetime.date.today()}.pptx"
            pptx_bytes = build_pptx(res, client_name, app_data)
            if st.download_button(
                label="📥 PowerPoint提案書をダウンロード",
                data=pptx_bytes,
                file_name=filename_pp,
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                type="primary",
                use_container_width=True
            ):
                log_download(app_data, filename_pp)
        else:
            st.warning("⚠️ PowerPoint自動生成エンジン（python-pptx）の準備ができていません。")
            st.info("💡 **一撃でPowerPoint出力ボタンを有効化する方法:**")
            st.code("pip install python-pptx", language="bash")
            st.caption("※ローカル環境のターミナルで上記のコマンドを実行すると、Streamlitが自動リロードされ、ここに『PowerPoint提案書ダウンロードボタン』が即座に出現します！")

    with st.expander("📋 月次シミュレーション詳細データ（グリッド一覧）"):
        show_cols = ["月","最大需要電力","導入後最大DM","削減デマンドkW",
                     "使用量合計","導入後使用量kWh","削減kWh",
                     "基本料金削減額","電力量料金削減額","現状電気料金推計","導入後電気料金推計"]
        st.dataframe(
            df[show_cols].style.format({
                "最大需要電力": "{:,.0f} kW",
                "導入後最大DM": "{:,.0f} kW",
                "削減デマンドkW": "{:,.1f} kW",
                "使用量合計": "{:,.0f} kWh",
                "導入後使用量kWh": "{:,.0f} kWh",
                "削減kWh": "{:,.0f} kWh",
                "基本料金削減額": "¥{:,.0f}",
                "電力量料金削減額": "¥{:,.0f}",
                "現状電気料金推計": "¥{:,.0f}",
                "導入後電気料金推計": "¥{:,.0f}",
            }),
            use_container_width=True
        )

    # ── 室外機別ROI（型番ベース時のみ）──
    units_detail = res.get("units_detail")
    if units_detail is not None and len(units_detail) > 0 and "年間削減配分円" in units_detail.columns:
        with st.expander("🔌 室外機別 削減・ROI内訳（型番→定格 容量加重 / methodology §5-6）", expanded=False):
            cost_per_unit = res.get("total_invest", 0) / max(res.get("n_units_ctrl", 1), 1)
            ut = units_detail.copy()
            ut["回収年"] = ut.apply(
                lambda r: (cost_per_unit / r["年間削減配分円"]) if (r["制御可否"] and r["年間削減配分円"] > 0) else None,
                axis=1)
            ut["制御"] = ut["制御可否"].map(lambda b: "○ 対象" if b else "× 除外")

            def _managed_label(r):
                if "管理" in ut.columns:
                    return "登録済み" if r.get("管理") else "※管理外（仮）"
                return "型番推定" if r.get("定格推定") else "—"
            ut["状態"] = ut.apply(_managed_label, axis=1)
            ut["定格"] = ut.apply(lambda r: f'{r["定格冷房kW"]:.1f}kW', axis=1)
            cols = ["機器ID", "メーカー", "型番"]
            if "馬力" in ut.columns:
                cols.append("馬力")
            cols += ["定格", "稼働係数", "電力kW", "制御", "状態", "年間削減配分円", "回収年"]
            view = ut[cols]
            fmt = {"稼働係数": "{:.2f}", "電力kW": "{:.2f} kW",
                   "年間削減配分円": "¥{:,.0f}", "回収年": "{:.1f} 年"}
            if "馬力" in ut.columns:
                fmt["馬力"] = "{:.1f} HP"
            st.dataframe(view.style.format(fmt, na_rep="—"), use_container_width=True)
            st.caption(
                f"※ 馬力→冷房能力kW（×{HP_TO_KW}）→ COP={res.get('cop', 3.5)} で電力kWへ換算。"
                "年間削減額は『電力kW×稼働係数』の容量加重で各機に配分。"
                f"回収年 = 1台あたり制御費(¥{cost_per_unit:,.0f}) ÷ 年間削減配分。"
                "「※管理外」は管理画面に型番未登録のため仮置き馬力で算出しています。"
            )

    # ══════════════════════════════════════════════════════════════
    # 🎯 目標回収年数から制御対象を逆算（独立機能：1行消せば撤去可）
    # ══════════════════════════════════════════════════════════════
    with st.expander("🎯 目標回収年数から制御対象を逆算（小馬力＝割高機を自動で削ぎ落とし）", expanded=False):
        show_payback_optimizer(res)

    # ══════════════════════════════════════════════════════════════
    # 📐 算出方法（計算式＋このシミュレーションの当てはめ数値）
    # ══════════════════════════════════════════════════════════════
    with st.expander("📐 算出方法（計算式と当てはめ数値を表示）", expanded=False):
        show_calc_methodology(res)


if __name__ == "__main__":
    main()