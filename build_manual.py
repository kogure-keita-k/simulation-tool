# -*- coding: utf-8 -*-
"""DPS試算アプリ 操作マニュアル(PPTX)を生成する。
各画面スライドは「スクショ貼付枠＋番号バッジ＋番号付き説明リスト」で構成（方式B）。
使い方: py -3 build_manual.py  → DPS試算アプリ_操作マニュアル.pptx を出力。
実スクショは、生成後に各スライドの点線枠へ貼り付け、番号バッジを該当箇所へドラッグしてください。
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
try:
    from pptx.enum.line import MSO_LINE_DASH_STYLE
    _DASH = MSO_LINE_DASH_STYLE.DASH
except Exception:
    _DASH = None

NAVY = RGBColor(0x13, 0x33, 0x1C) if False else RGBColor(0x13, 0x33, 0x5C)
GREEN = RGBColor(0x2E, 0x9E, 0x5B)
GRAY = RGBColor(0x60, 0x66, 0x70)
LGRAY = RGBColor(0xF2, 0xF4, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = 'Yu Gothic'
SW, SH = 13.333, 7.5

prs = Presentation()
prs.slide_width = Inches(SW); prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]


def _txt(slide, x, y, w, h, text, size=14, bold=False, color=NAVY, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.text = text
    p.font.name = FONT; p.font.size = Pt(size); p.font.bold = bold
    p.font.color.rgb = color; p.alignment = align
    return tb


def _header(slide, title, subtitle=None):
    b1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(0.45), Inches(0.13), Inches(0.5))
    b1.fill.solid(); b1.fill.fore_color.rgb = NAVY; b1.line.fill.background()
    b2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.86), Inches(0.45), Inches(0.13), Inches(0.5))
    b2.fill.solid(); b2.fill.fore_color.rgb = GREEN; b2.line.fill.background()
    _txt(slide, 1.12, 0.42, 11.5, 0.6, title, size=24, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        _txt(slide, 1.12, 1.15, 11.5, 0.45, subtitle, size=12.5, color=GRAY)


def _badge(slide, x, y, n, d=0.34, color=GREEN):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    c.fill.solid(); c.fill.fore_color.rgb = color; c.line.color.rgb = WHITE; c.line.width = Pt(1.25)
    p = c.text_frame.paragraphs[0]; p.text = str(n)
    p.font.name = FONT; p.font.size = Pt(13); p.font.bold = True
    p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
    c.text_frame.margin_top = c.text_frame.margin_bottom = 0
    return c


def cover():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(SW), Inches(SH))
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
    _txt(s, 1.0, 2.5, 11.3, 1.4, "DPS 空調デマンド制御\n試算アプリ 操作マニュアル",
         size=40, bold=True, color=WHITE)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.05), Inches(4.15), Inches(2.2), Inches(0.09))
    bar.fill.solid(); bar.fill.fore_color.rgb = GREEN; bar.line.fill.background()
    _txt(s, 1.0, 4.45, 11.3, 0.6, "検針票から試算・提案書・見積・削減レポ・請求書までを1つで。",
         size=15, color=RGBColor(0xCB, 0xDA, 0xF0))
    _txt(s, 1.0, 6.5, 11.3, 0.4, "株式会社シムックスイニシアティブ", size=12, color=RGBColor(0x9F, 0xB2, 0xCE))


def toc():
    s = prs.slides.add_slide(BLANK)
    _header(s, "目次")
    items = [
        "1. 最低限の操作フロー（まずはこれだけ）",
        "2. 画面①：試算入力（サイドバー）",
        "3. 画面②：結果ダッシュボード",
        "4. 概算モードの使い方（機材情報が無いとき）",
        "5. 画面③：各種設定（管理画面）",
        "6. 画面④：見積書",
        "7. 画面⑤：削減レポ（削減実績報告書）",
        "8. 画面⑥：請求書（成果報酬型）",
        "9. 出力ファイルの種類 / 困ったとき",
    ]
    y = 2.0
    for it in items:
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.95), Inches(y + 0.06), Inches(0.14), Inches(0.14))
        dot.fill.solid(); dot.fill.fore_color.rgb = GREEN; dot.line.fill.background()
        _txt(s, 1.25, y, 11.0, 0.4, it, size=15, color=NAVY)
        y += 0.52


def flow():
    s = prs.slides.add_slide(BLANK)
    _header(s, "1. 最低限の操作フロー", "まずはこの5ステップ。詳細は各画面のページを参照してください。")
    steps = [
        ("① 検針票を取込", "サイドバー『検針票データ』でUpload。無ければ『サンプル値代入』／『テンプレDL』。"),
        ("② 顧客名・業態", "『顧客情報』で顧客名（御中）と業態を選択（空調比率に使用）。"),
        ("③ 室外機の台数", "『室外機情報』で総台数・制御台数を入力。機材情報が無ければ『概算する』をON。"),
        ("④ 初期費用", "材料費・工事費・HW・構築費を入力。概算モードなら標準単価で自動。"),
        ("⑤ 試算実行→確認", "『試算実行』→結果ダッシュボードでKPI・グラフ確認、提案書PPTX等を出力。"),
    ]
    n = len(steps); gap = 0.3
    cw = (SW - 1.4 - gap * (n - 1)) / n
    y = 2.2; ch = 3.6
    for i, (t, d) in enumerate(steps):
        x = 0.7 + i * (cw + gap)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(cw), Inches(ch))
        card.fill.solid(); card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RGBColor(0xD5, 0xDD, 0xE8); card.line.width = Pt(1)
        _num = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + cw / 2 - 0.35), Inches(y + 0.3), Inches(0.7), Inches(0.7))
        _num.fill.solid(); _num.fill.fore_color.rgb = GREEN; _num.line.fill.background()
        pp = _num.text_frame.paragraphs[0]; pp.text = str(i + 1)
        pp.font.name = FONT; pp.font.size = Pt(26); pp.font.bold = True; pp.font.color.rgb = WHITE; pp.alignment = PP_ALIGN.CENTER
        _txt(s, x + 0.15, y + 1.15, cw - 0.3, 0.6, t[2:], size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        _txt(s, x + 0.2, y + 1.8, cw - 0.4, 1.6, d, size=10.5, color=GRAY, align=PP_ALIGN.CENTER)
        if i < n - 1:
            ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + cw + 0.02), Inches(y + ch / 2 - 0.13), Inches(gap - 0.04), Inches(0.26))
            ar.fill.solid(); ar.fill.fore_color.rgb = NAVY; ar.line.fill.background()
    _txt(s, 0.7, 6.5, 12.0, 0.4, "※ 結果は『結果一覧に保存』でき、削減レポ・請求書で再利用できます。", size=10, color=GRAY)


def screen_slide(title, subtitle, shot_label, items, must=None):
    """items: [(見出し, 説明), ...]（番号は自動①②…）。左に貼付枠＋番号バッジ、右に説明リスト。"""
    s = prs.slides.add_slide(BLANK)
    _header(s, title, subtitle)
    # 左：スクショ貼付枠（点線）
    fx, fy, fw, fh = 0.7, 2.0, 7.3, 4.7
    frame = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(fx), Inches(fy), Inches(fw), Inches(fh))
    frame.fill.solid(); frame.fill.fore_color.rgb = LGRAY
    frame.line.color.rgb = GRAY; frame.line.width = Pt(1.5)
    if _DASH is not None:
        try:
            frame.line.dash_style = _DASH
        except Exception:
            pass
    _txt(s, fx + 0.3, fy + fh / 2 - 0.5, fw - 0.6, 1.0,
         f"▼ ここに『{shot_label}』の\nスクリーンショットを貼り付け\n（右の番号バッジを該当箇所へドラッグ）",
         size=14, bold=True, color=GRAY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 番号バッジ（枠の左端に縦に並べる＝ドラッグして使う）
    for i in range(len(items)):
        _badge(s, fx + 0.12, fy + 0.15 + i * 0.5, i + 1)
    # 右：番号付き説明リスト
    rx, ry, rw = 8.25, 1.95, 4.4
    y = ry
    for i, (h, d) in enumerate(items):
        _badge(s, rx, y, i + 1, d=0.3)
        _txt(s, rx + 0.42, y - 0.02, rw - 0.42, 0.32, h, size=12, bold=True, color=NAVY)
        _txt(s, rx + 0.42, y + 0.28, rw - 0.42, 0.42, d, size=9.5, color=GRAY)
        y += 0.68
    if must:
        mb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(fx), Inches(6.85), Inches(fw), Inches(0.42))
        mb.fill.solid(); mb.fill.fore_color.rgb = RGBColor(0xE6, 0xF2, 0xE9)
        mb.line.color.rgb = GREEN; mb.line.width = Pt(1); mb.shadow.inherit = False
        p = mb.text_frame.paragraphs[0]; p.text = "★ 最低限：" + must
        p.font.name = FONT; p.font.size = Pt(10.5); p.font.bold = True; p.font.color.rgb = RGBColor(0x1E, 0x6B, 0x3B)
        mb.text_frame.margin_left = Inches(0.2)


# ── 生成 ──
cover()
toc()
flow()

screen_slide(
    "2. 画面①：試算入力（サイドバー）",
    "アプリ左側のサイドバーで、試算に必要な情報を入力します。",
    "試算入力（サイドバー）",
    [
        ("メニュー / ナビ", "見積書・削減レポ・結果一覧・各種設定へ移動するボタン。"),
        ("顧客情報", "顧客名（御中）と業態を選択。業態は空調比率の想定に使います。"),
        ("検針票データ", "Upload（CSV/Excel）／サンプル値代入／テンプレDL。12ヶ月の使用量・最大デマンド・単価を取込。"),
        ("室外機情報", "総台数・制御対象台数を入力。機材情報が無い場合は『概算する』トグルON（次頁）。"),
        ("初期費用", "材料費・工事費・ハードウェア費・構築費を入力（概算モードは標準単価で自動）。"),
        ("算出オプション", "季節性・業態配慮のトグル（必要な場合のみ）。"),
        ("試算実行", "入力後このボタンで計算 → 結果ダッシュボードへ。"),
    ],
    must="② 顧客名・業態 ／ ③ 検針票 ／ ④ 台数 ／ ⑦ 試算実行。この4つだけで試算できます。",
)

screen_slide(
    "3. 画面②：結果ダッシュボード",
    "試算結果の確認と、提案書・試算表などの出力を行います。",
    "結果ダッシュボード",
    [
        ("入力整合性チェック", "入力の警告・注意を表示（要確認があれば展開）。"),
        ("算出前提の注釈", "概算・業態標準値などの仮置き前提を明示。"),
        ("概算3シナリオ比較", "概算モード時のみ。ベスト／標準／保守の3案を並べて確認。"),
        ("KPI ×4", "初期費用・投資回収期間・年間削減額・10年累計。"),
        ("グラフ・分析", "月別削減・デマンド推移などのグラフ。"),
        ("ダウンロード", "提案書PPTX／試算XLSX／『結果一覧に保存』。"),
        ("画像出力（PNG）", "スライド貼付用の画像を個別に出力。"),
    ],
    must="⑥ ダウンロードから提案書PPTX・試算XLSXを出力。『結果一覧に保存』で後日再利用。",
)

screen_slide(
    "4. 概算モードの使い方（機材情報が無いとき）",
    "室外機・機材情報が無くても、検針票と業態標準値から概算できます。",
    "概算モード（室外機情報＋結果）",
    [
        ("概算トグル", "『室外機情報』の『機材情報なし→検針票から概算する』をON。"),
        ("台数の自動推定", "最大デマンド×業態空調比率÷標準室外機kWで台数を推定（調整可）。"),
        ("初期費用も概算", "材料費・工事費・HW費を標準単価×台数で自動算出。"),
        ("3シナリオ比較", "結果画面でベスト／標準／保守の幅を確認。"),
        ("資料の注記", "提案書には概算注記＋概算説明ページが自動で入ります。"),
    ],
    must="トグルONにするだけ。単価・標準室外機kWは『各種設定→計算ロジック変更』で調整可。",
)

screen_slide(
    "5. 画面③：各種設定（管理画面）",
    "計算に使う単価・係数・マスタを管理します（通常は初期値のままでOK）。",
    "各種設定（管理画面）",
    [
        ("計算ロジック変更", "容量削減率・各種固定値・概算マスタ（標準室外機kW／機材費／HW費など）。"),
        ("機材単価マスタ", "見積の機材費に使う単価。"),
        ("馬力・機材テーブル", "型番→定格などの対応。"),
        ("業態マスタ", "業態ごとの空調比率（ピーク／使用量）。"),
        ("表示項目変更", "画面・出力の表示切替。"),
        ("履歴ログ", "変更・ダウンロード履歴。"),
        ("画像出力（手動値）", "任意の数値で画像を生成。"),
    ],
)

screen_slide(
    "6. 画面④：見積書",
    "機材・工事の見積書を作成し、PDF / Excel で出力します。",
    "見積書",
    [
        ("件名・宛先", "件名・宛先（御中）・担当・書類番号など。"),
        ("明細", "機材・工事・ハードウェアの明細（単価はマスタ引用）。"),
        ("押印欄", "押印（※印）の氏名を左右それぞれ編集可。"),
        ("出力", "PDF または Excel でダウンロード。"),
    ],
)

screen_slide(
    "7. 画面⑤：削減レポ（削減実績報告書）",
    "運用後の実績を、客先様式の削減実績報告書として作成します（前年 vs 対象年）。",
    "削減レポ",
    [
        ("データ取込", "企業の蓄積データ／保存案件／追加インポート／テンプレDL。2年分あると比較が揃います。"),
        ("対象期間・基準年", "対象期間の開始月を選択（基準年は前年を自動、変更は折りたたみ）。"),
        ("表示ラベル", "客先の表記に合わせてラベルを調整・保存。"),
        ("エクセル眼鏡モード", "眼鏡様専用。使用電力量を営業日数で補正するトグル。"),
        ("数値プレビュー", "前年 vs 対象年の2年比較。表は手修正も可。"),
        ("出力", "PDF または Excel でダウンロード。"),
    ],
    must="① で2年分のデータを取り込む → ② 対象期間を選ぶ → ⑥ 出力。",
)

screen_slide(
    "8. 画面⑥：請求書（成果報酬型）",
    "削減実績を引用し、成果報酬型の請求書を作成します。",
    "請求書",
    [
        ("成果報酬率", "例：0.5＝50%。"),
        ("削減量の引用", "契約電力・使用電力量の削減量を削減レポから自動引用。"),
        ("単価・調整", "電力量単価＋燃料調整費＋再エネ賦課金（引用＋追記）。"),
        ("押印", "押印（※印）の氏名を左右それぞれ記入。"),
        ("出力", "PDF または Excel でダウンロード。"),
    ],
)

# 最終：出力物・困ったとき
s = prs.slides.add_slide(BLANK)
_header(s, "9. 出力ファイルの種類 / 困ったとき")
_txt(s, 0.7, 1.9, 12.0, 0.4, "■ 出力できるもの", size=15, bold=True, color=NAVY)
outs = [
    "提案書（PowerPoint / .pptx）… 表紙・収益サマリー・分析方法・制御可否・10年収支・環境価値 等。概算時は概算説明ページ付き。",
    "試算表（Excel / .xlsx）… 前提・計算・顧客提示サマリー。",
    "削減実績報告書（PDF / Excel）… 客先様式の①②③＋グラフ。",
    "見積書・請求書（PDF / Excel）。",
    "スライド貼付用 画像（PNG）。",
]
y = 2.35
for t in outs:
    d = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.95), Inches(y + 0.05), Inches(0.12), Inches(0.12))
    d.fill.solid(); d.fill.fore_color.rgb = GREEN; d.line.fill.background()
    _txt(s, 1.2, y, 11.4, 0.4, t, size=11.5, color=GRAY)
    y += 0.46
_txt(s, 0.7, 5.0, 12.0, 0.4, "■ 困ったとき", size=15, bold=True, color=NAVY)
tips = [
    "起動できない … 『はじめにお読みください.txt』の手順で Python を入れ、コマンドで起動。",
    "台数や機材が不明 … 『概算する』トグルで検針票から概算（提案書に概算注記が付きます）。",
    "削減レポで前年が空 … 前年の明細を追加インポート（または企業の蓄積データから呼び出し）。",
]
y = 5.45
for t in tips:
    d = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.95), Inches(y + 0.05), Inches(0.12), Inches(0.12))
    d.fill.solid(); d.fill.fore_color.rgb = NAVY; d.line.fill.background()
    _txt(s, 1.2, y, 11.4, 0.4, t, size=11.5, color=GRAY)
    y += 0.46

# ページ番号
for i, sl in enumerate(prs.slides, start=1):
    if i == 1:
        continue
    tb = sl.shapes.add_textbox(Inches(0.7), Inches(7.06), Inches(1.4), Inches(0.32))
    p = tb.text_frame.paragraphs[0]; p.text = f"P. {i:02d}"
    p.font.name = FONT; p.font.size = Pt(9); p.font.bold = True; p.font.color.rgb = NAVY

out = "DPS試算アプリ_操作マニュアル.pptx"
prs.save(out)
print("saved", out, "slides=", len(prs.slides._sldIdLst))
