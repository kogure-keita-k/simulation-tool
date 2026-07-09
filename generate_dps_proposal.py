# -*- coding: utf-8 -*-
import os
import sys
import math
import calendar
import openpyxl
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData

sys.stdout.reconfigure(encoding='utf-8')

# Colors
PRIMARY = RGBColor(11, 28, 63)       # Main Navy
ACCENT = RGBColor(0, 150, 94)        # Accent Green
ALERT = RGBColor(220, 53, 69)        # Alert Red
GRID_COLOR = RGBColor(241, 245, 249) # Grid Line Color (Light Grey)
TEXT_COLOR = RGBColor(51, 65, 85)    # Slate-Dark for text on white slides
WHITE = RGBColor(255, 255, 255)
SHADOW_COLOR = RGBColor(226, 232, 240) # Shadow Color (Light Grey)
LIGHT_GREY = RGBColor(241, 245, 249)
BROWN = RGBColor(101, 67, 33)

def draw_background_grid(prs, slide):
    # Draw horizontal and vertical lines at 5mm intervals (approx 14.17pt = 180000 EMU)
    step = 180000
    for x in range(0, prs.slide_width, step):
        c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x, 0, x, prs.slide_height)
        c.line.color.rgb = GRID_COLOR
        c.line.width = Pt(0.25)
    for y in range(0, prs.slide_height, step):
        c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, 0, y, prs.slide_width, y)
        c.line.color.rgb = GRID_COLOR
        c.line.width = Pt(0.25)

def add_left_accent_bar(prs, slide):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Pt(12), prs.slide_height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = PRIMARY
    rect.line.color.rgb = PRIMARY

def add_run_with_font(p, text, size=12, bold=False, italic=False, color=TEXT_COLOR):
    run = p.add_run()
    run.text = text
    run.font.name = "Meiryo UI"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return run

def make_slide_skipped_support(slide):
    slide.set_skipped = lambda skipped: slide.element.set('show', '0' if skipped else '1')

def add_slide_header_3layer(slide, action_title, key_takeaway, border_color=ACCENT):
    # Layer 1: [ACTION TITLE] (28pt bold Navy)
    title_box = slide.shapes.add_textbox(Pt(36), Pt(20), Pt(900), Pt(45))
    title_box.fill.solid()
    title_box.fill.fore_color.rgb = WHITE
    title_box.line.fill.background()
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    add_run_with_font(p, action_title, size=28, bold=True, color=PRIMARY)
    
    # Layer 2: [HIGHLIGHT BOX] (14pt bold Green/Red border)
    hb_bg = RGBColor(240, 248, 244) if border_color == ACCENT else RGBColor(254, 242, 242)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(36), Pt(75), Pt(900), Pt(36))
    box.fill.solid()
    box.fill.fore_color.rgb = hb_bg
    box.line.color.rgb = border_color
    box.line.width = Pt(1.5)
    
    tf_box = box.text_frame
    tf_box.word_wrap = True
    p_box = tf_box.paragraphs[0]
    p_box.alignment = PP_ALIGN.LEFT
    add_run_with_font(p_box, "  " + key_takeaway, size=14, bold=True, color=border_color)

def add_shadow_card(slide, left, top, width, height, bg_color=WHITE, border_color=None):
    # 1. Shadow shape
    shadow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + Pt(4), top + Pt(4), width, height)
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = SHADOW_COLOR
    shadow.line.color.rgb = SHADOW_COLOR
    
    # 2. Front card
    front = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    front.fill.solid()
    front.fill.fore_color.rgb = bg_color
    if border_color:
        front.line.color.rgb = border_color
        front.line.width = Pt(1.2)
    else:
        front.line.color.rgb = bg_color
    return front

def format_text_in_card(shape, title_text, value_text, unit_text, value_color=ACCENT, icon_glyph=None, icon_color=None):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(10)
    tf.margin_bottom = Pt(10)
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(10)
    
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.CENTER
    if icon_glyph:
        run_icon = p0.add_run()
        run_icon.text = icon_glyph + "\n"
        run_icon.font.name = "Segoe MDL2 Assets"
        run_icon.font.size = Pt(28)
        run_icon.font.color.rgb = icon_color or value_color
        
    add_run_with_font(p0, title_text, size=12, bold=True, color=PRIMARY)
    
    p1 = tf.add_paragraph()
    p1.alignment = PP_ALIGN.CENTER
    add_run_with_font(p1, "\n" + value_text, size=28, bold=True, color=value_color)
    
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    add_run_with_font(p2, unit_text, size=12, bold=False, color=TEXT_COLOR)

def add_bullet_item(tf, text, icon_glyph="\uE73E", icon_color=ACCENT, size=12, bold_text=False):
    p = tf.add_paragraph()
    p.space_after = Pt(8)
    
    # Add icon run
    run_icon = p.add_run()
    run_icon.text = icon_glyph + "  "
    run_icon.font.name = "Segoe MDL2 Assets"
    run_icon.font.size = Pt(size if size >= 12 else 12)
    run_icon.font.color.rgb = icon_color
    
    # Add text run
    run_text = p.add_run()
    run_text.text = text
    run_text.font.name = "Meiryo UI"
    run_text.font.size = Pt(size)
    run_text.font.bold = bold_text
    run_text.font.color.rgb = TEXT_COLOR
    return p

def draw_geometric_cedar_tree(slide, left, top, width, height):
    # Trunk (stem)
    trunk_w = width * 0.2
    trunk_h = height * 0.3
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + (width - trunk_w)/2, top + height - trunk_h, trunk_w, trunk_h)
    r.fill.solid()
    r.fill.fore_color.rgb = BROWN
    r.line.color.rgb = BROWN
    
    # Overlapping leaf triangles
    # Bottom leaf
    t3 = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, left + Pt(5), top + height * 0.4, width - Pt(10), height * 0.35)
    t3.fill.solid()
    t3.fill.fore_color.rgb = RGBColor(0, 120, 75)
    t3.line.color.rgb = RGBColor(0, 120, 75)
    
    # Middle leaf
    t2 = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, left + Pt(15), top + height * 0.2, width - Pt(30), height * 0.3)
    t2.fill.solid()
    t2.fill.fore_color.rgb = RGBColor(0, 150, 94)
    t2.line.color.rgb = RGBColor(0, 150, 94)
    
    # Top leaf
    t1 = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, left + Pt(25), top, width - Pt(50), height * 0.25)
    t1.fill.solid()
    t1.fill.fore_color.rgb = RGBColor(16, 185, 129)
    t1.line.color.rgb = RGBColor(16, 185, 129)

def draw_geometric_stacked_bar_chart(slide, left, top, width, height, sorted_mon_data, energy_ratio, controllable_ratio, conv, is_mini=False):
    from pptx.enum.dml import MSO_LINE_DASH_STYLE
    baseline_y = top + height - (Pt(15) if is_mini else Pt(30))
    chart_h = height - (Pt(30) if is_mini else Pt(60))
    
    usage_vals = [x["usage"] for x in sorted_mon_data]
    max_val = max(usage_vals) if usage_vals else 1.0
    ymax = math.ceil(max_val / 10000.0) * 10000.0 if max_val > 0 else 100.0
    scale = chart_h / ymax
    
    # Calculate costs: total_cost = usage * rate + basic * contract
    costs = [x["usage"] * x["rate"] + x["basic"] * x["contract"] for x in sorted_mon_data]
    max_cost = max(costs) if costs else 1.0
    ymax_cost = math.ceil(max_cost / 500000.0) * 500000.0 if max_cost > 0 else 100000.0
    scale_cost = chart_h / ymax_cost
    
    # Baseline
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left, baseline_y, left + width, baseline_y)
    line.line.color.rgb = TEXT_COLOR
    line.line.width = Pt(1)
    
    col_w = width / len(sorted_mon_data)
    bar_w = col_w * 0.55
    
    cost_pts = []
    
    for i in range(len(sorted_mon_data)):
        col_left = left + i * col_w
        
        # Monthly labels
        m_str = sorted_mon_data[i]["month"]
        m_val = m_str.split("/")[1] if "/" in m_str else m_str.split("-")[1]
        m_label = m_val.lstrip("0") + "月"
        
        if not is_mini:
            lbl_box = slide.shapes.add_textbox(col_left, baseline_y + Pt(2), col_w, Pt(20))
            lbl_box.fill.solid()
            lbl_box.fill.fore_color.rgb = WHITE
            lbl_box.line.fill.background()
            p = lbl_box.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            add_run_with_font(p, m_label, size=12, bold=True, color=TEXT_COLOR)
        
        # Calculate three layers
        total_usage = usage_vals[i]
        aircon_usage = total_usage * energy_ratio * controllable_ratio
        base_load = total_usage - aircon_usage
        saved_margin = aircon_usage * 0.30
        post_aircon_load = aircon_usage - saved_margin
        
        h_base = base_load * scale
        h_aircon = post_aircon_load * scale
        h_saved = saved_margin * scale
        
        bar_left = col_left + (col_w - bar_w) / 2
        
        # 1. Base Load (Navy)
        if h_base > Pt(1):
            r1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bar_left, baseline_y - h_base, bar_w, h_base)
            r1.fill.solid()
            r1.fill.fore_color.rgb = PRIMARY
            r1.line.color.rgb = PRIMARY
            
            # Add usage text inside the base load bar (Navy) in White text
            if not is_mini:
                v_box = slide.shapes.add_textbox(bar_left - Pt(15), baseline_y - h_base + Pt(5), bar_w + Pt(30), Pt(16))
                v_box.fill.background()
                v_box.line.fill.background()
                p_v = v_box.text_frame.paragraphs[0]
                p_v.alignment = PP_ALIGN.CENTER
                add_run_with_font(p_v, f"{total_usage/1000:.0f}k", size=12, bold=True, color=WHITE)
            
        # 2. Aircon Load (Light Blue)
        if h_aircon > Pt(1):
            r2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bar_left, baseline_y - h_base - h_aircon, bar_w, h_aircon)
            r2.fill.solid()
            r2.fill.fore_color.rgb = RGBColor(59, 130, 246)
            r2.line.color.rgb = RGBColor(59, 130, 246)
            
        # 3. DPS Saved Margin (Accent Green with dashed border)
        if h_saved > Pt(1):
            top_y = baseline_y - h_base - h_aircon - h_saved
            r3 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bar_left, top_y, bar_w, h_saved)
            r3.fill.solid()
            r3.fill.fore_color.rgb = ACCENT
            r3.line.color.rgb = ACCENT
            
            # Dashed top line specifically
            line_top = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, bar_left, top_y, bar_left + bar_w, top_y)
            line_top.line.color.rgb = ACCENT
            line_top.line.width = Pt(1.5)
            line_top.line.dash_style = MSO_LINE_DASH_STYLE.DASH
            
        # Store cost point coordinates for the line chart
        cost_val = costs[i]
        cost_y = baseline_y - cost_val * scale_cost
        cost_x = col_left + col_w / 2
        cost_pts.append((cost_x, cost_y, cost_val))
        
    # Draw line chart on top of bars
    if not is_mini:
        # Connect points with lines
        for i in range(len(cost_pts) - 1):
            x1, y1, _ = cost_pts[i]
            x2, y2, _ = cost_pts[i+1]
            line_conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
            line_conn.line.color.rgb = RGBColor(249, 115, 22)  # Orange
            line_conn.line.width = Pt(2.5)
            
        # Add markers and value labels
        for x, y, val in cost_pts:
            # Marker dot
            marker = slide.shapes.add_shape(MSO_SHAPE.OVAL, x - Pt(4), y - Pt(4), Pt(8), Pt(8))
            marker.fill.solid()
            marker.fill.fore_color.rgb = RGBColor(249, 115, 22)
            marker.line.color.rgb = RGBColor(255, 255, 255)  # White border to pop
            marker.line.width = Pt(1)
            
            # Value label above the dot
            c_box = slide.shapes.add_textbox(x - Pt(20), y - Pt(18), Pt(40), Pt(16))
            c_box.fill.background()
            c_box.line.fill.background()
            p_c = c_box.text_frame.paragraphs[0]
            p_c.alignment = PP_ALIGN.CENTER
            add_run_with_font(p_c, f"{val/10000:.0f}万", size=12, bold=True, color=RGBColor(249, 115, 22))

# Load data from Excel
def load_data():
    excel_path = "西富士事業所_空調デマンド制御シミュレーション.xlsx"
    if not os.path.exists(excel_path):
        print(f"Error: {excel_path} not found.")
        sys.exit(1)
        
    wb = openpyxl.load_workbook(excel_path, data_only=True)
    
    # 1. Parse client name
    ws0 = wb["シート0_顧客提示用サマリー"]
    title_val = ws0["A1"].value
    company_name = "西富士事業所"
    if title_val and "　　" in title_val:
        parts = title_val.split("　　")
        if len(parts) > 1 and "様" in parts[1]:
            company_name = parts[1].split("様")[0]
            
    # 2. Parse monthly billing data
    ws_mon = wb["月次・電気料金"]
    mon_data = []
    for r in range(4, 200):
        val = ws_mon.cell(r, 1).value
        if val in (None, "") or "計" in str(val) or "最大" in str(val):
            break
        mon_data.append({
            "month": str(ws_mon.cell(r, 1).value),
            "days": float(ws_mon.cell(r, 2).value or 30),
            "usage": float(ws_mon.cell(r, 3).value or 0),
            "demand": float(ws_mon.cell(r, 4).value or 0),
            "rate": float(ws_mon.cell(r, 6).value or 0),
            "basic": float(ws_mon.cell(r, 7).value or 0),
            "contract": float(ws_mon.cell(r, 8).value or 0),
        })
        
    # 3. Load parameters
    ws_set = wb["前提・制御条件"]
    peak_ratio = float(ws_set["B5"].value or 0.20)
    energy_ratio = float(ws_set["B6"].value or 0.15)
    conv = float(ws_set["B9"].value or 0.50)
    item_price = float(ws_set["B18"].value or 5000)
    item_margin = float(ws_set["B19"].value or 0.30)
    co2_factor = float(ws_set["B20"].value or 0.451)
    
    # 4. Load equipment list
    ws_eq = wb["機器台帳"]
    eq_list = []
    total_ac_count = 0
    control_ac_count = 0
    total_capacity = 0.0
    controllable_capacity = 0.0
    for r in range(5, 500):
        val = ws_eq.cell(r, 1).value
        if val in (None, ""):
            if ws_eq.cell(r, 2).value == "合計":
                break
            continue
        total_ac_count += 1
        ctrl = ws_eq.cell(r, 5).value or "×"
        rt = ws_eq.cell(r, 6).value
        rt_val = float(rt) if (rt is not None and rt != "") else 0.0
        if ctrl == "〇":
            control_ac_count += 1
            controllable_capacity += rt_val
        if rt_val > 0.0:
            total_capacity += rt_val
        eq_list.append((ws_eq.cell(r, 2).value, ws_eq.cell(r, 3).value, ws_eq.cell(r, 4).value, ctrl, rt_val))
        
    controllable_ratio = controllable_capacity / total_capacity if total_capacity > 0 else 0.0
    
    # 5. Evaluate standard annual savings (30% Cap)
    pk_list = [x["demand"] for x in mon_data]
    kwh_list = [x["usage"] for x in mon_data]
    bu_list = [x["basic"] for x in mon_data]
    rate_list = [x["rate"] for x in mon_data]
    
    max_demand = max(pk_list)
    total_kwh = sum(kwh_list)
    avg_basic_rate = sum(bu_list) / len(bu_list)
    avg_energy_rate = sum(rate_list) / len(rate_list)
    
    ctrl_peak_kw = max_demand * peak_ratio * controllable_ratio
    ctrl_energy_kwh = total_kwh * energy_ratio * controllable_ratio
    
    saving_demand_annual = ctrl_peak_kw * 0.30 * avg_basic_rate * 12
    saving_energy_annual = ctrl_energy_kwh * 0.30 * conv * avg_energy_rate
    total_annual_saving_30 = saving_demand_annual + saving_energy_annual
    
    # 6. Evaluate ROI adopted units
    units_roi = []
    for loc, mk, md, ctrl, rt in eq_list:
        if ctrl == "〇" and rt > 0.0:
            units_roi.append({"loc": loc, "md": md, "kw": rt, "cf": 0.8})
            
    units_roi.sort(key=lambda x: -(x["kw"] * x["cf"]))
    total_weight = sum(x["kw"] * x["cf"] for x in units_roi)
    
    adopted_count = 0
    adopted_savings = 0.0
    cost_per_unit = 80000
    for x in units_roi:
        weight = x["kw"] * x["cf"]
        unit_saving = total_annual_saving_30 * weight / total_weight if total_weight > 0 else 0.0
        unit_payback = cost_per_unit / unit_saving if unit_saving > 0 else 999.0
        if unit_payback <= 7.0:
            adopted_count += 1
            adopted_savings += unit_saving
            
    adopted_investment = adopted_count * cost_per_unit
    adopted_payback = adopted_investment / adopted_savings if adopted_savings > 0 else 0.0
    
    # 7. Sheet 0 actual savings
    sheet0_annual_savings = 0.0
    for x in mon_data:
        c_r = round(x["demand"] * peak_ratio * 0.30 * controllable_ratio, 1)
        f_r = round(x["usage"] * energy_ratio * 0.30 * controllable_ratio * conv, 0)
        j_r = round(c_r * x["basic"] + f_r * x["rate"], 0)
        sheet0_annual_savings += j_r
        
    total_reduc_kwh = sum(round(x["usage"] * energy_ratio * 0.30 * controllable_ratio * conv, 0) for x in mon_data)
    co2_savings_kg = total_reduc_kwh * co2_factor
    cedar_trees = co2_savings_kg / 8.8
    item_sales_count = math.ceil(sheet0_annual_savings / (item_price * item_margin))
    
    return {
        "company_name": company_name,
        "mon_data": mon_data,
        "total_ac_count": total_ac_count,
        "control_ac_count": control_ac_count,
        "controllable_ratio": controllable_ratio,
        "peak_ratio": peak_ratio,
        "energy_ratio": energy_ratio,
        "conv": conv,
        "item_price": item_price,
        "item_margin": item_margin,
        "co2_factor": co2_factor,
        "total_annual_saving_30": total_annual_saving_30,
        "adopted_count": adopted_count,
        "adopted_investment": adopted_investment,
        "adopted_savings": adopted_savings,
        "adopted_payback": adopted_payback,
        "sheet0_annual_savings": sheet0_annual_savings,
        "total_reduc_kwh": total_reduc_kwh,
        "co2_savings_kg": co2_savings_kg,
        "cedar_trees": cedar_trees,
        "item_sales_count": item_sales_count,
        "max_demand": max_demand,
        "total_kwh": total_kwh,
        "avg_basic_rate": avg_basic_rate,
        "avg_energy_rate": avg_energy_rate,
    }

def get_payback_warning_reason(payback_years, avg_energy_rate, controllable_ratio, basic_rate):
    if payback_years <= 3.0:
        return ""
    reasons = []
    if avg_energy_rate < 15.0:
        reasons.append("現状の電力量料金単価が極めて安い（15円/kWh未満）ため、使用量削減による金銭メリットが小さくなっています。")
    if controllable_ratio < 0.50:
        reasons.append("制御対象となる空調室外機の割合が低い（50%未満）ため、削減可能容量が制限されています。")
    if basic_rate < 1200.0:
        reasons.append("基本料金単価が安いため、デマンドカットによる基本料金の削減幅が小さくなっています。")
    if not reasons:
        reasons.append("現状の空調稼働率が極めて低いか、または基本料金が実量制のミニマム制限に達しているためです。")
    return "【回収期間3年超の原因分析】\n" + "\n".join(reasons)

def build_proposal():
    res = load_data()
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    
    # -------------------------------------------------------------
    # Slide 1: Title Slide (Clean White Grid background)
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    draw_background_grid(prs, slide1)
    
    title_box = slide1.shapes.add_textbox(Pt(54), Pt(180), Pt(800), Pt(150))
    title_box.fill.solid()
    title_box.fill.fore_color.rgb = WHITE
    title_box.line.fill.background()
    tf = title_box.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    add_run_with_font(p0, "デジタルパワーサービス（DPS） 提案書", size=32, bold=True, color=PRIMARY)
    
    p1 = tf.add_paragraph()
    p1.text = ""
    add_run_with_font(p1, "\n～過去の電気料金実績に基づく空調デマンド制御シミュレーション提案～", size=18, bold=False, color=PRIMARY)
    
    info_box = slide1.shapes.add_textbox(Pt(54), Pt(450), Pt(600), Pt(150))
    info_box.fill.solid()
    info_box.fill.fore_color.rgb = WHITE
    info_box.line.fill.background()
    tf_info = info_box.text_frame
    p_comp = tf_info.paragraphs[0]
    add_run_with_font(p_comp, f"{res['company_name']} 御中", size=22, bold=True, color=PRIMARY)
    
    p_sub = tf_info.add_paragraph()
    p_sub.text = ""
    add_run_with_font(p_sub, "\n提案主体: 株式会社シムックスイニシアティブ", size=14, bold=False, color=PRIMARY)


    # -------------------------------------------------------------
    # Slide 2: Summary Slide
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    draw_background_grid(prs, slide2)
    add_left_accent_bar(prs, slide2)
    
    annual_savings = res["adopted_savings"]
    system_fee = 241800
    net_profit = annual_savings - system_fee
    payback_years = res["adopted_investment"] / net_profit if net_profit > 0 else 0.0
    
    net_profit_tenk = net_profit / 10000
    
    add_slide_header_3layer(
        slide2, 
        "導入効果の総括", 
        f"約 {payback_years:.1f} 年で投資回収、以降は毎年約 {net_profit_tenk:.0f} 万円の純利益を生む”収益試算”へ"
    )
    
    # 4 metrics cards
    card_w = Pt(200)
    card_h = Pt(160)
    spacing = Pt(24)
    start_left = Pt(44)
    top_pos = Pt(180)
    
    metrics = [
        ("① 初期費用", f"{int(round(res['adopted_investment'])):,}", f"円\n(制御対象 {res['control_ac_count']}台)", PRIMARY, "\uE8A1", PRIMARY),
        ("② 年間削減額", f"{int(round(annual_savings)):,}", "円 / 年", ACCENT, "\uE9D2", ACCENT),
        ("③ 年間実利", f"{int(round(net_profit)):,}", f"円 / 年\n(削減額: {int(round(annual_savings)):,}円 - 利用料: {system_fee:,}円)", ACCENT, "\uE9D2", ACCENT),
        ("④ 投資回収期間", f"{payback_years:.1f}", "年", PRIMARY, "\uE73E", PRIMARY)
    ]
    
    for idx, (title, val, unit, color, icon_glyph, icon_color) in enumerate(metrics):
        left_pos = start_left + idx * (card_w + spacing)
        card = add_shadow_card(slide2, left_pos, top_pos, card_w, card_h)
        format_text_in_card(card, title, val, unit, color, icon_glyph=icon_glyph, icon_color=icon_color)
        
    # Payback period > 3.0 warning check
    warning_text = get_payback_warning_reason(
        payback_years, 
        res["avg_energy_rate"], 
        res["controllable_ratio"], 
        res["avg_basic_rate"]
    )
    if warning_text:
        warn_box = slide2.shapes.add_textbox(Pt(44), Pt(370), Pt(872), Pt(110))
        warn_box.fill.solid()
        warn_box.fill.fore_color.rgb = WHITE
        warn_box.line.fill.background()
        tf_warn = warn_box.text_frame
        tf_warn.word_wrap = True
        p_warn = tf_warn.paragraphs[0]
        add_run_with_font(p_warn, warning_text, size=12, bold=True, color=ALERT)
    else:
        # Additional notes
        note_box = slide2.shapes.add_textbox(Pt(44), Pt(370), Pt(872), Pt(110))
        note_box.fill.solid()
        note_box.fill.fore_color.rgb = WHITE
        tf_note = note_box.text_frame
        tf_note.word_wrap = True
        tf_note.margin_left = Pt(10)
        tf_note.margin_right = Pt(10)
        tf_note.margin_top = Pt(10)
        
        p_note = tf_note.paragraphs[0]
        add_run_with_font(p_note, "■ 試算の妥当性と保証条件", size=12, bold=True, color=PRIMARY)
        
        add_bullet_item(tf_note, "過去12ヶ月の確定検針票データをもとに、確立された空調削減ロジックを適用し精緻に算出", icon_glyph="\uE73E", icon_color=ACCENT)
        add_bullet_item(tf_note, "削減計算式：[各月電力量 × 空調比率(15%) × 制御対象比率(100%) × 制御削減率(30%) × 換算 conv]", icon_glyph="\uE73E", icon_color=PRIMARY)
        add_bullet_item(tf_note, "独自のローテーション制御技術により、工場内環境・品質や生産ラインへの影響はゼロであることを保証", icon_glyph="\uE73E", icon_color=ACCENT)

    # -------------------------------------------------------------
    # Slide 3: Methodology & Assumptions
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    draw_background_grid(prs, slide3)
    add_left_accent_bar(prs, slide3)
    add_slide_header_3layer(
        slide3,
        "試算の分析方法および算出前提の定義",
        "不確実な未来予測を排除し、電力使用量の配分推計と稼働データから削減余地を逆算"
    )
    
    # Left Column: Allocation definition and mix cards
    left_card = add_shadow_card(slide3, Inches(0.8), Pt(160), Pt(400), Pt(280))
    tf_l = left_card.text_frame
    tf_l.vertical_anchor = MSO_ANCHOR.TOP
    tf_l.word_wrap = True
    tf_l.margin_left = Pt(15)
    tf_l.margin_right = Pt(15)
    tf_l.margin_top = Pt(15)
    tf_l.margin_bottom = Pt(15)
    
    p_l0 = tf_l.paragraphs[0]
    p_l0.alignment = PP_ALIGN.CENTER
    add_run_with_font(p_l0, "■ 電力使用量配分の定義と前提", size=13, bold=True, color=PRIMARY)
    
    # Render the mix using geometric block cards inside Slide 3
    # HVAC 15.0%, Production 65.0%, Lighting 15.0%, Office 5.0%, Others 5.0%
    alloc_mix = [
        ("空調設備負荷 (HVAC)", "15.0%", RGBColor(240, 253, 244), ACCENT, True),
        ("生産設備・機械動力", "65.0%", WHITE, PRIMARY, False),
        ("照明・コンセント負荷", "15.0%", WHITE, PRIMARY, False),
        ("事務・OA機器類", "5.0%", WHITE, PRIMARY, False),
        ("その他付随設備", "5.0%", WHITE, PRIMARY, False)
    ]
    
    for idx, (name, pct, bg_col, text_col, is_highlight) in enumerate(alloc_mix):
        card_y = Pt(215) + idx * Pt(38)
        shape_c = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8) + Pt(15), card_y, Pt(370), Pt(32))
        shape_c.fill.solid()
        shape_c.fill.fore_color.rgb = bg_col
        if is_highlight:
            shape_c.line.color.rgb = ACCENT
            shape_c.line.width = Pt(1.2)
        else:
            shape_c.line.color.rgb = LIGHT_GREY
            shape_c.line.width = Pt(1)
            
        tf_c = shape_c.text_frame
        tf_c.word_wrap = True
        tf_c.margin_left = Pt(8)
        tf_c.margin_top = Pt(6)
        
        p = tf_c.paragraphs[0]
        add_run_with_font(p, f"{name}  ", size=12, bold=is_highlight, color=text_col)
        run_pct = p.add_run()
        run_pct.text = f" {pct}"
        run_pct.font.name = "Meiryo UI"
        run_pct.font.size = Pt(12)
        run_pct.font.bold = True
        run_pct.font.color.rgb = text_col
        
    # Right Column: 3 Analysis Methods
    right_card = add_shadow_card(slide3, Pt(500), Pt(160), Pt(400), Pt(280))
    tf_r = right_card.text_frame
    tf_r.vertical_anchor = MSO_ANCHOR.TOP
    tf_r.word_wrap = True
    tf_r.margin_left = Pt(15)
    tf_r.margin_right = Pt(15)
    tf_r.margin_top = Pt(15)
    tf_r.margin_bottom = Pt(15)
    
    p_r0 = tf_r.paragraphs[0]
    add_run_with_font(p_r0, "■ 3つの主要稼働分析アプローチ", size=13, bold=True, color=PRIMARY)
    
    add_bullet_item(tf_r, "A. ベースライン分析 (Baseline Load Analysis)", icon_glyph="\uE9CA", icon_color=ACCENT)
    p_ab = tf_r.add_paragraph()
    add_run_with_font(p_ab, "   ・常時必要な非空調基本負荷を推定\n   ・最低電力量との差分から空調動力量を逆算", size=12, color=TEXT_COLOR)
    
    add_bullet_item(tf_r, "B. 稼働時間分析 (Operating Rate Analysis)", icon_glyph="\uE9D2", icon_color=ACCENT)
    p_op = tf_r.add_paragraph()
    add_run_with_font(p_op, "   ・操業カレンダーとシフト体制から稼働時間を特定\n   ・実質的な制御ポテンシャルを精緻化", size=12, color=TEXT_COLOR)
    
    add_bullet_item(tf_r, "C. ピーク需要分析 (Peak Demand Analysis)", icon_glyph="\uE73E", icon_color=ACCENT)
    p_pk = tf_r.add_paragraph()
    add_run_with_font(p_pk, "   ・契約電力と最大デマンド値からピークを特定\n   ・基本料金の削減効果を最大化", size=12, color=TEXT_COLOR)
    
    p_ind = tf_r.add_paragraph()
    p_ind.space_before = Pt(8)
    add_bullet_item(tf_r, "業態特性に応じたパラメータ調整により推計精度を向上", icon_glyph="\uE73E", icon_color=PRIMARY)
    
    # Proof footnote on Slide 3
    footnote_box = slide3.shapes.add_textbox(Pt(60), Pt(455), Pt(840), Pt(30))
    footnote_box.fill.solid()
    footnote_box.fill.fore_color.rgb = WHITE
    footnote_box.line.color.rgb = LIGHT_GREY
    footnote_box.line.width = Pt(1)
    tf_fn = footnote_box.text_frame
    tf_fn.word_wrap = True
    tf_fn.margin_left = Pt(8)
    tf_fn.margin_top = Pt(6)
    p_fn = tf_fn.paragraphs[0]
    add_run_with_font(p_fn, "※The simulation engine applies a manufacturing-model deduction framework. Based on verified calendar inputs, minimum baseline volumes are extracted to calculate true net aircon loads.", size=12, bold=False, color=TEXT_COLOR)

    # -------------------------------------------------------------
    # Slide 4: Feasibility (Feasibility)
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    draw_background_grid(prs, slide4)
    add_left_accent_bar(prs, slide4)
    add_slide_header_3layer(
        slide4,
        "制御可否：対象拠点における機器仕分け実態",
        "現場の操業を止めない、ロジカルな制御対象台数の割り出し"
    )
    
    # Large number display card on the left
    card = add_shadow_card(slide4, Pt(60), Pt(160), Pt(360), Pt(280))
    format_text_in_card(card, "制御可能空調比率", f"{res['control_ac_count']} / {res['total_ac_count']}台", f"(容量比率: {res['controllable_ratio']:.1%})", value_color=ACCENT, icon_glyph="\uE8A1", icon_color=PRIMARY)
    
    # Right explanation texts
    desc_box2 = slide4.shapes.add_textbox(Pt(500), Pt(160), Pt(400), Pt(280))
    desc_box2.fill.solid()
    desc_box2.fill.fore_color.rgb = WHITE
    desc_box2.line.fill.background()
    tf_desc2 = desc_box2.text_frame
    tf_desc2.word_wrap = True
    
    p_title = tf_desc2.paragraphs[0]
    add_run_with_font(p_title, "【工場長様向け・安心安全な制御スキーム】", size=14, bold=True, color=PRIMARY)
    
    # Spacing
    p_space = tf_desc2.add_paragraph()
    p_space.text = ""
    
    # Dynamic exclusion item
    add_bullet_item(tf_desc2, f"制御除外室外機 (計 {res['total_ac_count'] - res['control_ac_count']}台): サーバー室・重要設備などの室外機は100%除外して温度安全を完全保証。", icon_glyph="\uE8A1", icon_color=ALERT, bold_text=True)
    
    # Rotation control items
    add_bullet_item(tf_desc2, f"段階的送風切替（ローテーション制御）: {res['control_ac_count']}台を同時停止せず、1台ずつ数分間ずつ順次送風へ切り替え。", icon_glyph="\uE73E", icon_color=ACCENT)
    add_bullet_item(tf_desc2, "操業ライン稼働維持：室温変化は1℃以下に抑制。工場の生産活動や作業環境への影響は一切ありません。", icon_glyph="\uE73E", icon_color=ACCENT)

    # -------------------------------------------------------------
    # Slide 5: Simulation Detail (Simulation)
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    draw_background_grid(prs, slide5)
    add_left_accent_bar(prs, slide5)
    add_slide_header_3layer(
        slide5,
        "シミュレーション：企業別・月別料金比較詳細",
        "月別の電力高低差を1円単位で徹底検証"
    )
    
    # Sort months April to March
    sorted_mon_data = sorted(res["mon_data"], key=lambda x: ((int(x["month"].split("/")[1]) - 4) % 12 if "/" in x["month"] else (int(x["month"].split("-")[1]) - 4) % 12))
    
    # Draw geometric stacked bar chart
    draw_geometric_stacked_bar_chart(slide5, Pt(60), Pt(180), Pt(360), Pt(280), sorted_mon_data, res["energy_ratio"], res["controllable_ratio"], res["conv"], is_mini=False)
    
    # Draw Legend for Stacked Chart
    legend_labels = [
        ("ベース負荷", PRIMARY, "bar"),
        ("空調負荷 (制御後)", RGBColor(59, 130, 246), "bar"),
        ("DPS削減分", ACCENT, "bar"),
        ("電気料金 (右軸)", RGBColor(249, 115, 22), "line")
    ]
    legend_lefts = [Pt(60), Pt(155), Pt(300), Pt(405)]
    for idx, (lbl, col, mode) in enumerate(legend_labels):
        leg_left = legend_lefts[idx]
        leg_top = Pt(145)
        if mode == "bar":
            box = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, leg_left, leg_top + Pt(2), Pt(12), Pt(12))
            box.fill.solid()
            box.fill.fore_color.rgb = col
            box.line.color.rgb = col
        else:
            conn = slide5.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, leg_left, leg_top + Pt(8), leg_left + Pt(16), leg_top + Pt(8))
            conn.line.color.rgb = col
            conn.line.width = Pt(2)
            dot = slide5.shapes.add_shape(MSO_SHAPE.OVAL, leg_left + Pt(5), leg_top + Pt(5), Pt(6), Pt(6))
            dot.fill.solid()
            dot.fill.fore_color.rgb = col
            dot.line.color.rgb = col
            
        # text
        tb = slide5.shapes.add_textbox(leg_left + Pt(18), leg_top - Pt(4), Pt(110), Pt(20))
        tb.fill.solid()
        tb.fill.fore_color.rgb = WHITE
        tb.line.fill.background()
        p = tb.text_frame.paragraphs[0]
        add_run_with_font(p, lbl, size=12, bold=True, color=TEXT_COLOR)
        
    # Right column: scenario comparison table
    table_shape = slide5.shapes.add_table(6, 4, Pt(450), Pt(180), Pt(450), Pt(220))
    table = table_shape.table
    table.columns[0].width = Pt(110)
    table.columns[1].width = Pt(40)
    table.columns[2].width = Pt(150)
    table.columns[3].width = Pt(150)
    
    dem_max = res['max_demand']
    dem_factor = res['peak_ratio'] * res['controllable_ratio']
    
    def get_dem_str(pct):
        red_kw = dem_max * dem_factor * pct
        post_kw = dem_max - red_kw
        return f"{int(round(dem_max)):,} kW > {int(round(post_kw)):,} kW (▲{int(round(red_kw)):,} kW)"
        
    rows_data = [
        ("制御モード", "削減率", "デマンド推移（現状値 ＞ 制御後）", "快適性指標"),
        ("保守制御", "20%", get_dem_str(0.20), "影響なし (\\uE899)"),
        ("標準制御", "30%", get_dem_str(0.30), "影響なし (\\uE899)"),
        ("積極制御", "40%", get_dem_str(0.40), "注意 (\\uE7BA)"),
        ("限界制御", "50%", get_dem_str(0.50), "警告 (\\uEB90)"),
        ("参考：現状の確定実績 - 総使用電力量: 2,390,591 kWh / 最大デマンド: 595 kW", "", "", "")
    ]
    
    for row_idx, row_vals in enumerate(rows_data):
        if row_idx == 5:
            cell = table.cell(row_idx, 0)
            cell.merge(table.cell(row_idx, 3))
            cell.text = ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE
            p = cell.text_frame.paragraphs[0]
            add_run_with_font(p, row_vals[0], size=12, bold=True, color=PRIMARY)
        else:
            for col_idx, val in enumerate(row_vals):
                cell = table.cell(row_idx, col_idx)
                cell.text = ""
                cell.fill.solid()
                
                # Highlight Row 2 (Standard) with a soft light green background
                if row_idx == 2:
                    cell.fill.fore_color.rgb = RGBColor(240, 253, 244)
                elif row_idx == 0:
                    cell.fill.fore_color.rgb = LIGHT_GREY
                else:
                    cell.fill.fore_color.rgb = WHITE
                    
                p = cell.text_frame.paragraphs[0]
                
                # Column 3 is the comfort index glyph (MDL2)
                if col_idx == 3 and row_idx > 0:
                    p.alignment = PP_ALIGN.LEFT
                    if "影響なし" in val:
                        prefix = "影響なし ("
                        glyph = "\uE899"
                        color = ACCENT
                    elif "注意" in val:
                        prefix = "注意 ("
                        glyph = "\uE7BA"
                        color = RGBColor(249, 115, 22)
                    else:
                        prefix = "警告 ("
                        glyph = "\uEB90"
                        color = ALERT
                        
                    # prefix
                    r1 = p.add_run()
                    r1.text = prefix
                    r1.font.name = "Meiryo UI"
                    r1.font.size = Pt(12)
                    r1.font.bold = (row_idx == 2)
                    r1.font.color.rgb = PRIMARY if row_idx == 2 else TEXT_COLOR
                    
                    # glyph
                    r2 = p.add_run()
                    r2.text = glyph
                    r2.font.name = "Segoe MDL2 Assets"
                    r2.font.size = Pt(12)
                    r2.font.bold = (row_idx == 2)
                    r2.font.color.rgb = color
                    
                    # suffix
                    r3 = p.add_run()
                    r3.text = ")"
                    r3.font.name = "Meiryo UI"
                    r3.font.size = Pt(12)
                    r3.font.bold = (row_idx == 2)
                    r3.font.color.rgb = PRIMARY if row_idx == 2 else TEXT_COLOR
                else:
                    # Standard text
                    add_run_with_font(p, val, size=12, bold=(row_idx == 0 or row_idx == 2), color=PRIMARY if (row_idx == 0 or row_idx == 2) else TEXT_COLOR)

    # -------------------------------------------------------------
    # Slide 6: Investment and Cost (Investment)
    # -------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    draw_background_grid(prs, slide6)
    add_left_accent_bar(prs, slide6)
    add_slide_header_3layer(
        slide6,
        "導入方法と費用",
        "既存設備をそのまま活かす『配管切断なし工法』と明瞭な費用ブレイクダウン"
    )
    
    # Top-Left: Visual hardware box placeholder card
    hw_card = add_shadow_card(slide6, Pt(60), Pt(160), Pt(400), Pt(140))
    tf_hw = hw_card.text_frame
    tf_hw.word_wrap = True
    p_hw = tf_hw.paragraphs[0]
    
    # Title with MDL2 icon
    run_icon = p_hw.add_run()
    run_icon.text = "\uE99A  "
    run_icon.font.name = "Segoe MDL2 Assets"
    run_icon.font.size = Pt(14)
    run_icon.font.color.rgb = PRIMARY
    add_run_with_font(p_hw, "制御機器：DPSユニット (ハードウェア)", size=13, bold=True, color=PRIMARY)
    
    add_bullet_item(tf_hw, "室外機制御用スマートボード搭載", icon_glyph="\uE73E", icon_color=ACCENT)
    add_bullet_item(tf_hw, "配管切断なし・アドオン装着仕様（工期短縮）", icon_glyph="\uE73E", icon_color=ACCENT)
    add_bullet_item(tf_hw, "セーフティオートシャットダウン自動復旧回路内蔵", icon_glyph="\uE73E", icon_color=ACCENT)
    
    # Top-Right: Standard installation text run card
    inst_card = add_shadow_card(slide6, Pt(500), Pt(160), Pt(400), Pt(140))
    tf_inst = inst_card.text_frame
    tf_inst.word_wrap = True
    p_inst = tf_inst.paragraphs[0]
    add_run_with_font(p_inst, "■標準取付工法適用仕様", size=13, bold=True, color=PRIMARY)
    
    add_bullet_item(tf_inst, "通常の室外機信号線へのアドオン工事を適用", icon_glyph="\uE73E", icon_color=ACCENT)
    add_bullet_item(tf_inst, "完全屋外施工：室内への立ち入りは原則不要", icon_glyph="\uE73E", icon_color=ACCENT)
    add_bullet_item(tf_inst, "稼働停止リスクゼロ：生産ラインを止める必要はありません", icon_glyph="\uE73E", icon_color=ACCENT)
    
    # Bottom Span: Wide cost breakdown matrix table
    table_shape2 = slide6.shapes.add_table(5, 2, Pt(60), Pt(320), Pt(840), Pt(150))
    table2 = table_shape2.table
    table2.columns[0].width = Pt(440)
    table2.columns[1].width = Pt(400)
    
    cost_items = [
        ("初期費用内訳", "金額"),
        (f"制御機器費 ({res['adopted_count']}台分)", "1,400,000 円"),
        ("標準取付工事費", "600,000 円"),
        ("初期設定・検証調整費", "400,000 円"),
        ("自己負担総額 (税別)", f"{int(round(res['adopted_investment'])):,} 円")
    ]
    
    for row_idx, (k, v) in enumerate(cost_items):
        for col_idx, text in enumerate([k, v]):
            cell = table2.cell(row_idx, col_idx)
            cell.text = ""
            cell.fill.solid()
            if row_idx == 0:
                cell.fill.fore_color.rgb = LIGHT_GREY
            elif row_idx == 4:
                cell.fill.fore_color.rgb = RGBColor(226, 245, 232)
            else:
                cell.fill.fore_color.rgb = WHITE
            p = cell.text_frame.paragraphs[0]
            add_run_with_font(p, text, size=12, bold=(row_idx == 0 or row_idx == 4), color=PRIMARY if row_idx == 0 else (ACCENT if row_idx == 4 else TEXT_COLOR))

    # -------------------------------------------------------------
    # Slide 7: Aircon Replacement Addendum (Hidden)
    # -------------------------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    make_slide_skipped_support(slide7)
    slide7.set_skipped(True)
    draw_background_grid(prs, slide7)
    add_left_accent_bar(prs, slide7)
    add_slide_header_3layer(
        slide7,
        "＜アディショナル＞空調交換提案（リプレイス・シナリオ）",
        "効果見込み薄の旧型機を、公的補助金を利用して最新省エネ空調へ置き換えた場合の効果"
    )
    
    # Left Card: Scenario A (DPS Add-on)
    card_left = add_shadow_card(slide7, Pt(60), Pt(160), Pt(400), Pt(280))
    tf_l = card_left.text_frame
    tf_l.word_wrap = True
    tf_l.margin_left = Pt(20)
    tf_l.margin_right = Pt(20)
    tf_l.margin_top = Pt(20)
    
    p_l0 = tf_l.paragraphs[0]
    add_run_with_font(p_l0, "【シナリオA：DPSデマンド制御（現状維持＋アドオン）】", size=14, bold=True, color=PRIMARY)
    
    add_bullet_item(tf_l, "初期投資額：約 240 万円", icon_glyph="\uE73E", icon_color=PRIMARY)
    add_bullet_item(tf_l, "公的補助金：なし（投資規模が対象外）", icon_glyph="\uE73E", icon_color=PRIMARY)
    add_bullet_item(tf_l, f"年間削減額：約 {int(round(res['adopted_savings']/10000))} 万円", icon_glyph="\uE73E", icon_color=ACCENT)
    add_bullet_item(tf_l, f"投資回収期間：約 {res['adopted_payback']:.1f} 年", icon_glyph="\uE73E", icon_color=ACCENT)
    add_bullet_item(tf_l, "特徴：既存機を活用、配管工事なしで即導入可能", icon_glyph="\uE73E", icon_color=PRIMARY)
    
    # Right Card: Scenario B (Replacement + Subsidy)
    card_right = add_shadow_card(slide7, Pt(500), Pt(160), Pt(400), Pt(280))
    tf_r = card_right.text_frame
    tf_r.word_wrap = True
    tf_r.margin_left = Pt(20)
    tf_r.margin_right = Pt(20)
    tf_r.margin_top = Pt(20)
    
    p_r0 = tf_r.paragraphs[0]
    add_run_with_font(p_r0, "【シナリオB：高効率空調更新＋省エネ補助金適用】", size=14, bold=True, color=PRIMARY)
    
    add_bullet_item(tf_r, "初期投資額：約 2,400 万円（30台分）", icon_glyph="\uE73E", icon_color=PRIMARY)
    add_bullet_item(tf_r, "公的補助金：約 800 万円（1/3補助想定）", icon_glyph="\uE73E", icon_color=PRIMARY)
    add_bullet_item(tf_r, "自己負担額：約 1,600 万円", icon_glyph="\uE73E", icon_color=PRIMARY)
    add_bullet_item(tf_r, "投資回収期間：約 8.5 年（省エネ更新効果）", icon_glyph="\uE73E", icon_color=ALERT)
    add_bullet_item(tf_r, "特徴：老朽化した空調の更新用、回収は長期化", icon_glyph="\uE73E", icon_color=PRIMARY)

    # -------------------------------------------------------------
    # Slide 8: Next Steps Timeline (Visible)
    # -------------------------------------------------------------
    slide8 = prs.slides.add_slide(blank_layout)
    draw_background_grid(prs, slide8)
    add_left_accent_bar(prs, slide8)
    add_slide_header_3layer(
        slide8,
        "今後の流れ：最終商談後のプロジェクトスケジュール",
        "合意から本稼働まで約1ヶ月、スムーズな移行ステップのご案内"
    )
    
    steps = [
        ("1：概算試算", "検針票データより削減ポテンシャルを机上試算。", LIGHT_GREY, TEXT_COLOR, "完了"),
        ("2：現場調査（現地確認）", "全室外機の型番・状態を確認し施工計画を策定。", LIGHT_GREY, TEXT_COLOR, "完了"),
        ("3：最終提案・ご契約", "実地調査に基づく最終シミュレーション提示。", LIGHT_GREY, TEXT_COLOR, "現フェーズ"),
        ("4：導入工事・テスト運用", "配管切断なしの取付工事と通信テスト実施。", PRIMARY, WHITE, "次フェーズ"),
        ("5：本運用開始", "本運用を開始し、月次レポートで効果可視化。", PRIMARY, WHITE, "本番稼働")
    ]
    
    card_w = Pt(150)
    card_h = Pt(260)
    spacing_x = Pt(22)
    start_left = Pt(61)
    top_pos = Pt(180)
    
    for idx, (title, desc, bg_col, txt_col, status) in enumerate(steps):
        left_pos = start_left + idx * (card_w + spacing_x)
        card = add_shadow_card(slide8, left_pos, top_pos, card_w, card_h, bg_color=bg_col)
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(12)
        tf.margin_right = Pt(12)
        tf.margin_top = Pt(15)
        
        is_navy = (bg_col == PRIMARY)
        p0 = tf.paragraphs[0]
        add_run_with_font(p0, title, size=13, bold=True, color=WHITE if is_navy else PRIMARY)
        
        p_stat = tf.add_paragraph()
        stat_color = WHITE if is_navy else (ALERT if status in ["現フェーズ", "次フェーズ"] else (ACCENT if status == "本番稼働" else RGBColor(100, 116, 139)))
        add_run_with_font(p_stat, f"\n【{status}】", size=12, bold=True, color=stat_color)
        
        p1 = tf.add_paragraph()
        add_run_with_font(p1, "\n" + desc, size=12, bold=False, color=WHITE if is_navy else TEXT_COLOR)
        
        if idx < 4:
            arrow_left = left_pos + card_w
            arrow_w = spacing_x
            tb_arr = slide8.shapes.add_textbox(arrow_left, top_pos + card_h/2 - Pt(15), arrow_w, Pt(30))
            tb_arr.fill.solid()
            tb_arr.fill.fore_color.rgb = WHITE
            tb_arr.line.fill.background()
            p_arr = tb_arr.text_frame.paragraphs[0]
            p_arr.alignment = PP_ALIGN.CENTER
            add_run_with_font(p_arr, "➔", size=14, bold=True, color=PRIMARY)
            
    # -------------------------------------------------------------
    # Slide 9: 10-Year Opportunity Loss (Hidden)
    # -------------------------------------------------------------
    slide9 = prs.slides.add_slide(blank_layout)
    make_slide_skipped_support(slide9)
    slide9.set_skipped(True)
    draw_background_grid(prs, slide9)
    add_left_accent_bar(prs, slide9)
    add_slide_header_3layer(
        slide9,
        "＜アディショナル＞導入比較：対策有無による10年間の利益差異",
        "いただきたい検針票月から比較するシステム導入有無の影響"
    )
    
    # Top Body Area: Miniature 12-month stacked bar chart
    draw_geometric_stacked_bar_chart(slide9, Pt(60), Pt(140), Pt(460), Pt(130), sorted_mon_data, res["energy_ratio"], res["controllable_ratio"], res["conv"], is_mini=True)
    
    # Calculation basis text box next to the chart
    basis_box = slide9.shapes.add_textbox(Pt(560), Pt(140), Pt(340), Pt(130))
    basis_box.fill.solid()
    basis_box.fill.fore_color.rgb = WHITE
    basis_box.line.color.rgb = LIGHT_GREY
    basis_box.line.width = Pt(1)
    tf_basis = basis_box.text_frame
    tf_basis.word_wrap = True
    tf_basis.margin_left = Pt(8)
    tf_basis.margin_top = Pt(8)
    p_basis = tf_basis.paragraphs[0]
    add_run_with_font(p_basis, "■ 試算計算根拠\n・単年削減額 × 年数\n・過去12ヶ月実績の月次削減値を累計\n・実測ベースの精緻な手残り予測", size=12, bold=False, color=TEXT_COLOR)
    
    # Two horizontal numerical comparison tracks
    saving_annual = res["sheet0_annual_savings"]
    inv = res["adopted_investment"]
    
    # Track 1: 現状維持 (Compounding Cash Deficit)
    track1_y = Pt(285)
    # Label card
    lbl1 = add_shadow_card(slide9, Pt(60), track1_y, Pt(200), Pt(65), bg_color=RGBColor(254, 242, 242), border_color=ALERT)
    tf_lbl1 = lbl1.text_frame
    tf_lbl1.word_wrap = True
    tf_lbl1.margin_left = Pt(10)
    tf_lbl1.margin_top = Pt(12)
    p_lbl1 = tf_lbl1.paragraphs[0]
    add_run_with_font(p_lbl1, "現状維持 (対策なし)\n【累積機会損失額】", size=12, bold=True, color=ALERT)
    
    # Year 1 Card
    y1_1 = add_shadow_card(slide9, Pt(280), track1_y, Pt(190), Pt(65), bg_color=WHITE)
    tf_y1_1 = y1_1.text_frame
    tf_y1_1.word_wrap = True
    tf_y1_1.margin_top = Pt(8)
    p_y1_1 = tf_y1_1.paragraphs[0]
    p_y1_1.alignment = PP_ALIGN.CENTER
    add_run_with_font(p_y1_1, "1年目", size=12, bold=False, color=TEXT_COLOR)
    p_y1_1_val = tf_y1_1.add_paragraph()
    p_y1_1_val.alignment = PP_ALIGN.CENTER
    add_run_with_font(p_y1_1_val, f"▲{int(round(saving_annual)):,} 円", size=14, bold=True, color=TEXT_COLOR)
    
    # Year 5 Card
    y5_1 = add_shadow_card(slide9, Pt(490), track1_y, Pt(190), Pt(65), bg_color=WHITE)
    tf_y5_1 = y5_1.text_frame
    tf_y5_1.word_wrap = True
    tf_y5_1.margin_top = Pt(8)
    p_y5_1 = tf_y5_1.paragraphs[0]
    p_y5_1.alignment = PP_ALIGN.CENTER
    add_run_with_font(p_y5_1, "5年目", size=12, bold=False, color=TEXT_COLOR)
    p_y5_1_val = tf_y5_1.add_paragraph()
    p_y5_1_val.alignment = PP_ALIGN.CENTER
    add_run_with_font(p_y5_1_val, f"▲{int(round(saving_annual*5)):,} 円", size=14, bold=True, color=TEXT_COLOR)
    
    # Year 10 Card (Massive Red Deficit)
    y10_1 = add_shadow_card(slide9, Pt(700), track1_y, Pt(200), Pt(65), bg_color=RGBColor(254, 226, 226), border_color=ALERT)
    tf_y10_1 = y10_1.text_frame
    tf_y10_1.word_wrap = True
    tf_y10_1.margin_top = Pt(8)
    p_y10_1 = tf_y10_1.paragraphs[0]
    p_y10_1.alignment = PP_ALIGN.CENTER
    add_run_with_font(p_y10_1, "10年目 (最大の損失)", size=12, bold=True, color=ALERT)
    p_y10_1_val = tf_y10_1.add_paragraph()
    p_y10_1_val.alignment = PP_ALIGN.CENTER
    add_run_with_font(p_y10_1_val, f"▲{int(round(saving_annual*10)):,} 円", size=16, bold=True, color=ALERT)
    
    # Track 2: DPS導入 (Compounding Cash Success)
    track2_y = Pt(365)
    # Label card
    lbl2 = add_shadow_card(slide9, Pt(60), track2_y, Pt(200), Pt(65), bg_color=RGBColor(240, 253, 244), border_color=ACCENT)
    tf_lbl2 = lbl2.text_frame
    tf_lbl2.word_wrap = True
    tf_lbl2.margin_left = Pt(10)
    tf_lbl2.margin_top = Pt(12)
    p_lbl2 = tf_lbl2.paragraphs[0]
    add_run_with_font(p_lbl2, "DPS導入 (対策あり)\n【累積手残り純利】", size=12, bold=True, color=PRIMARY)
    
    # Year 1 Card
    y1_2 = add_shadow_card(slide9, Pt(280), track2_y, Pt(190), Pt(65), bg_color=WHITE)
    tf_y1_2 = y1_2.text_frame
    tf_y1_2.word_wrap = True
    tf_y1_2.margin_top = Pt(8)
    p_y1_2 = tf_y1_2.paragraphs[0]
    p_y1_2.alignment = PP_ALIGN.CENTER
    add_run_with_font(p_y1_2, "1年目", size=12, bold=False, color=TEXT_COLOR)
    p_y1_2_val = tf_y1_2.add_paragraph()
    p_y1_2_val.alignment = PP_ALIGN.CENTER
    add_run_with_font(p_y1_2_val, f"▲{int(round(inv - saving_annual)):,} 円", size=14, bold=True, color=TEXT_COLOR)
    
    # Year 5 Card
    y5_2 = add_shadow_card(slide9, Pt(490), track2_y, Pt(190), Pt(65), bg_color=WHITE)
    tf_y5_2 = y5_2.text_frame
    tf_y5_2.word_wrap = True
    tf_y5_2.margin_top = Pt(8)
    p_y5_2 = tf_y5_2.paragraphs[0]
    p_y5_2.alignment = PP_ALIGN.CENTER
    add_run_with_font(p_y5_2, "5年目", size=12, bold=False, color=TEXT_COLOR)
    p_y5_2_val = tf_y5_2.add_paragraph()
    p_y5_2_val.alignment = PP_ALIGN.CENTER
    add_run_with_font(p_y5_2_val, f"¥{int(round(saving_annual*5 - inv)):,} 円", size=14, bold=True, color=ACCENT)
    
    # Year 10 Card (Massive Green Net Profit)
    y10_2 = add_shadow_card(slide9, Pt(700), track2_y, Pt(200), Pt(65), bg_color=RGBColor(220, 252, 231), border_color=ACCENT)
    tf_y10_2 = y10_2.text_frame
    tf_y10_2.word_wrap = True
    tf_y10_2.margin_top = Pt(8)
    p_y10_2 = tf_y10_2.paragraphs[0]
    p_y10_2.alignment = PP_ALIGN.CENTER
    add_run_with_font(p_y10_2, "10年目 (手残り利益)", size=12, bold=True, color=ACCENT)
    p_y10_2_val = tf_y10_2.add_paragraph()
    p_y10_2_val.alignment = PP_ALIGN.CENTER
    add_run_with_font(p_y10_2_val, f"¥{int(round(saving_annual*10 - inv)):,} 円", size=16, bold=True, color=ACCENT)
    
    # Footnote note box
    fn_box = slide9.shapes.add_textbox(Pt(60), Pt(445), Pt(840), Pt(30))
    fn_box.fill.solid()
    fn_box.fill.fore_color.rgb = WHITE
    fn_box.line.color.rgb = LIGHT_GREY
    fn_box.line.width = Pt(1)
    tf_fn = fn_box.text_frame
    tf_fn.word_wrap = True
    tf_fn.margin_left = Pt(8)
    tf_fn.margin_top = Pt(6)
    p_fn = tf_fn.paragraphs[0]
    add_run_with_font(p_fn, "※Formula Basis: Single-year certified savings delta multiplied by cumulative operational lifespan years. Deducts upfront元本 of 2,400,000 JPY to isolate pure cash-in-hand.", size=12, bold=False, color=TEXT_COLOR)

    # -------------------------------------------------------------
    # Slide 10: Product Sales Equivalent (Hidden)
    # -------------------------------------------------------------
    slide10 = prs.slides.add_slide(blank_layout)
    make_slide_skipped_support(slide10)
    slide10.set_skipped(True)
    draw_background_grid(prs, slide10)
    add_left_accent_bar(prs, slide10)
    add_slide_header_3layer(
        slide10,
        "＜アディショナル＞純利益換算：削減額の本当の財務価値",
        "年間実利を、御社の本業の製品販売数量に逆換算する"
    )
    
    # Left Column: Visual Badge (Numeric Callout Card)
    badge_card = add_shadow_card(slide10, Pt(60), Pt(180), Pt(260), Pt(250))
    format_text_in_card(badge_card, "主力製品 追加販売換算", f"{res['item_sales_count']} 個", "本業の営業活動による利益価値", value_color=ACCENT, icon_glyph="\uE9D2", icon_color=ACCENT)
    
    # Right Column: Explanation text
    card_sales = add_shadow_card(slide10, Pt(360), Pt(180), Pt(540), Pt(250))
    tf_sales = card_sales.text_frame
    tf_sales.word_wrap = True
    tf_sales.margin_left = Pt(25)
    tf_sales.margin_right = Pt(25)
    tf_sales.margin_top = Pt(25)
    tf_sales.margin_bottom = Pt(25)
    
    p = tf_sales.paragraphs[0]
    add_run_with_font(p, f"本システムによる年間純削減額（約 {res['sheet0_annual_savings']:,.0f} 円）の価値分析：", size=13, bold=True, color=PRIMARY)
    
    add_bullet_item(tf_sales, f"年間純削減額は、本業の製品利益に換算すると 新規 {res['item_sales_count']} 個 の販売に相当", icon_glyph="\uE73E", icon_color=ACCENT)
    add_bullet_item(tf_sales, f"利益計算基礎：製品単価 {res['item_price']:,.0f} 円 × 粗利益率 {res['item_margin']:.0%} = {int(res['item_price'] * res['item_margin']):,} 円 / 個", icon_glyph="\uE73E", icon_color=PRIMARY)
    add_bullet_item(tf_sales, "新規顧客を開拓し、追加販売することの難易度との比較材料に", icon_glyph="\uE73E", icon_color=PRIMARY)
    add_bullet_item(tf_sales, "エアコンの無駄をAIで削る経費削減は、全自動で確実な利益となります", icon_glyph="\uE73E", icon_color=ACCENT)

    # -------------------------------------------------------------
    # Slide 11: ESG / CO2 reduction (Hidden)
    # -------------------------------------------------------------
    slide11 = prs.slides.add_slide(blank_layout)
    make_slide_skipped_support(slide11)
    slide11.set_skipped(True)
    draw_background_grid(prs, slide11)
    add_left_accent_bar(prs, slide11)
    add_slide_header_3layer(
        slide11,
        "＜アディショナル＞CO2削減：省エネと環境価値（ESG投資）",
        "財務改善の副産物として得られる、企業の社会的責任（CSR）への貢献"
    )
    
    card_co2 = add_shadow_card(slide11, Pt(60), Pt(180), Pt(280), Pt(250), border_color=ACCENT)
    format_text_in_card(card_co2, "年間 CO2排出削減量", f"{res['co2_savings_kg']:,.0f} kg-CO2 / 年", "(省エネ法やサステナビリティ実績へ公表可能)", value_color=ACCENT, icon_glyph="\uEC05", icon_color=ACCENT)
    
    card_tree = add_shadow_card(slide11, Pt(370), Pt(180), Pt(280), Pt(250), border_color=ACCENT)
    format_text_in_card(card_tree, "スギの木換算森林効果", f"約 {res['cedar_trees']:,.0f} 本分 / 年", "(杉1本あたり年間 8.8kg 吸収として算出)", value_color=ACCENT, icon_glyph="\uEC05", icon_color=ACCENT)
    
    # Corporate value caption at the bottom
    pr_box = slide11.shapes.add_textbox(Pt(60), Pt(460), Pt(840), Pt(45))
    pr_box.fill.solid()
    pr_box.fill.fore_color.rgb = WHITE
    pr_box.line.color.rgb = LIGHT_GREY
    pr_box.line.width = Pt(1)
    tf_pr = pr_box.text_frame
    tf_pr.word_wrap = True
    tf_pr.margin_left = Pt(10)
    tf_pr.margin_top = Pt(10)
    p_pr = tf_pr.paragraphs[0]
    add_run_with_font(p_pr, "・本データは省エネ法定期報告書の作成に活用可能\n・御社HP等でのESG・サステナビリティ実績として公表可能", size=12, bold=True, color=PRIMARY)
    
    # Render three cedar tree icons geometrically on the right
    draw_geometric_cedar_tree(slide11, Pt(680), Pt(200), Pt(60), Pt(110))
    draw_geometric_cedar_tree(slide11, Pt(760), Pt(230), Pt(50), Pt(80))
    draw_geometric_cedar_tree(slide11, Pt(830), Pt(200), Pt(60), Pt(110))

    # -------------------------------------------------------------
    # Slide 12: FAQ Slide (Hidden)
    # -------------------------------------------------------------
    slide12 = prs.slides.add_slide(blank_layout)
    make_slide_skipped_support(slide12)
    slide12.set_skipped(True)
    draw_background_grid(prs, slide12)
    add_left_accent_bar(prs, slide12)
    add_slide_header_3layer(
        slide12,
        "＜アディショナル＞QA・よくある質問とトラブルシューティング",
        "現場・工場長からの懸念を100%先回りして解消する安心担保"
    )
    
    row_h = Pt(90)
    spacing_y = Pt(15)
    start_top = Pt(160)
    
    for idx, (q, a_bullets) in enumerate([
        ("Q1. 温度管理：室温や製品品質への影響はありませんか？", [
            "A1. 送風切替時間は数分間のローテーション制御。室温変化は1℃以下に抑制。",
            "工場全体の操業環境や、精密部品の品質管理にも影響はありません。"
        ]),
        ("Q2. 機器寿命：室外機のコンプレッサー等の摩耗への懸念について", [
            "A2. 主電源の強制オンオフではなく、電子信号による送風・冷却の切り替え。",
            "コンプレッサーへの負荷や機器寿命への悪影響はゼロであることを保証。"
        ]),
        ("Q3. 緊急対応：システム障害時のフェイルセーフ機能について", [
            "A3. セーフティオートシャットダウン自動復旧回路を標準内蔵。",
            "通信エラー等の発生時は、即座に自動で通常のエアコン運転へ復旧。"
        ])
    ]):
        row_y = start_top + idx * (row_h + spacing_y)
        
        # 3D Shadow Card for the row
        row_card = add_shadow_card(slide12, Pt(60), row_y, Pt(840), row_h, border_color=PRIMARY)
        tf_row = row_card.text_frame
        tf_row.word_wrap = True
        tf_row.margin_left = Pt(15)
        tf_row.margin_right = Pt(15)
        tf_row.margin_top = Pt(10)
        tf_row.margin_bottom = Pt(10)
        
        p_q = tf_row.paragraphs[0]
        add_run_with_font(p_q, q, size=13, bold=True, color=PRIMARY)
        
        # Add answer bullets
        for b_idx, bullet in enumerate(a_bullets):
            p_b = tf_row.add_paragraph()
            p_b.space_before = Pt(4)
            # MDL2 icon checkmark for bullets
            run_ic = p_b.add_run()
            run_ic.text = "\uE73E  " if b_idx > 0 else "➔  "
            run_ic.font.name = "Segoe MDL2 Assets"
            run_ic.font.size = Pt(12)
            run_ic.font.color.rgb = ACCENT if b_idx > 0 else PRIMARY
            
            run_tx = p_b.add_run()
            run_tx.text = bullet
            run_tx.font.name = "Meiryo UI"
            run_tx.font.size = Pt(12)
            run_tx.font.bold = (b_idx == 0)
            run_tx.font.color.rgb = TEXT_COLOR
        
    # Style override for textboxes to prevent double borders and white backgrounds
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == 17:  # 17 is TEXT_BOX
                # Clear background fill and line border to let text float over grid
                shape.fill.background()
                shape.line.fill.background()
                # Set text frame margins to Inches(0.1)
                shape.text_frame.margin_left = Inches(0.1)
                shape.text_frame.margin_right = Inches(0.1)
                shape.text_frame.margin_top = Inches(0.1)
                shape.text_frame.margin_bottom = Inches(0.1)
                
    prs.save("西富士事業所_提案書.pptx")
    print("OK: PowerPoint file generated successfully as '西富士事業所_提案書.pptx'")

if __name__ == "__main__":
    build_proposal()
