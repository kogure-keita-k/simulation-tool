# -*- coding: utf-8 -*-
"""DPS試算アプリ 導入マニュアル（PowerPoint）生成スクリプト — デザイン強化版＋リンク付き"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

NAVY    = RGBColor(0x1F, 0x38, 0x64)
NAVY_DK = RGBColor(0x0D, 0x18, 0x30)
NAVY_MD = RGBColor(0x2A, 0x46, 0x7A)
GREEN   = RGBColor(0x2E, 0x9E, 0x5B)
GREEN_BR= RGBColor(0x36, 0xC0, 0x6B)
ICE     = RGBColor(0xC9, 0xDA, 0xF5)
LIGHT   = RGBColor(0xF4, 0xF7, 0xFB)
CARD    = RGBColor(0xEC, 0xF1, 0xF9)
GTINT   = RGBColor(0xE4, 0xF4, 0xEA)
GRAY    = RGBColor(0x5A, 0x64, 0x74)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
DARKTX  = RGBColor(0x24, 0x2C, 0x3A)
CODEBG  = RGBColor(0x10, 0x1B, 0x2E)
WMARK   = RGBColor(0xEC, 0xF1, 0xF9)

HEAD = "Cambria"     # 見出し（セーフなセリフ）
BODY = "Calibri"
MONO = "Consolas"
LOGO = "company_logo_192x192.png"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
blank = prs.slide_layouts[6]


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def shadow(sp, blur=110000, dist=40000, alpha=22000, col='1F3864'):
    el = sp._element.spPr
    for e in el.findall(qn('a:effectLst')):
        el.remove(e)
    eff = el.makeelement(qn('a:effectLst'), {})
    sh = eff.makeelement(qn('a:outerShdw'),
                         {'blurRad': str(blur), 'dist': str(dist), 'dir': '5400000', 'rotWithShape': '0'})
    c = sh.makeelement(qn('a:srgbClr'), {'val': col})
    a = c.makeelement(qn('a:alpha'), {'val': str(alpha)})
    c.append(a); sh.append(c); eff.append(sh); el.append(eff)


def box(slide, x, y, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE, sh=False):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if sh:
        shadow(sp)
    return sp


def txt(slide, x, y, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp_after=4, line_sp=1.06):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.04); tf.margin_top = tf.margin_bottom = 0
    if isinstance(paras[0], tuple):
        paras = [paras]
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sp_after); p.line_spacing = line_sp
        for run in para:
            s, sz, col, bold = run[0], run[1], run[2], run[3]
            font = run[4] if len(run) > 4 and run[4] else BODY
            link = run[5] if len(run) > 5 else None
            r = p.add_run(); r.text = s
            r.font.size = Pt(sz); r.font.bold = bold; r.font.name = font
            if link:
                r.hyperlink.address = link
            r.font.color.rgb = col   # リンク後に色指定（緑/白を維持）
    return tb


def badge(slide, x, y, d, n, emoji=None, fill=GREEN, num_sz=40, shp=MSO_SHAPE.ROUNDED_RECTANGLE):
    b = box(slide, x, y, d, d, fill=fill, shape=shp, sh=True)
    tf = b.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = tf.margin_bottom = tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    has_num = str(n) != ""
    if emoji:
        emoji_sz = max(14, num_sz - 6) if has_num else round(d * 30)
        r0 = p.add_run(); r0.text = (emoji + "  ") if has_num else emoji
        r0.font.size = Pt(emoji_sz); r0.font.name = BODY
    if has_num:
        r = p.add_run(); r.text = str(n); r.font.size = Pt(num_sz); r.font.bold = True
        r.font.color.rgb = WHITE; r.font.name = HEAD
    return b


def dot(slide, x, y, d=0.14, c=GREEN):
    box(slide, x, y, d, d, fill=c, shape=MSO_SHAPE.OVAL)


def head(slide, t, sub=None, emoji=None):
    dot(slide, 0.72, 0.66, 0.16, GREEN)
    txt(slide, 1.0, 0.45, 11.6, 0.7, [[((emoji + "  " if emoji else "") + t, 30, NAVY, True, HEAD)]])
    if sub:
        txt(slide, 1.0, 1.18, 11.6, 0.4, [[(sub, 14, GRAY, False)]])


def code_box(slide, x, y, w, lines):
    h = 0.46 + 0.36 * len(lines)
    box(slide, x, y, w, h, fill=CODEBG, shape=MSO_SHAPE.ROUNDED_RECTANGLE, sh=True)
    # mac風ドット
    for k, cc in enumerate([RGBColor(0xFF,0x5F,0x57), RGBColor(0xFE,0xBC,0x2E), RGBColor(0x28,0xC8,0x40)]):
        box(slide, x + 0.22 + k*0.22, y + 0.16, 0.12, 0.12, fill=cc, shape=MSO_SHAPE.OVAL)
    tb = slide.shapes.add_textbox(Inches(x + 0.22), Inches(y + 0.42), Inches(w - 0.44), Inches(h - 0.5))
    tf = tb.text_frame; tf.word_wrap = True; tf.margin_left = tf.margin_top = tf.margin_bottom = 0
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.15
        r = p.add_run(); r.text = ln; r.font.size = Pt(13.5); r.font.name = MONO
        r.font.color.rgb = GREEN_BR if ln.lstrip().startswith(("pip", "streamlit", "python", "cd", ">")) else ICE
    return y + h


def note(slide, x, y, w, text, icon="💡", fill=GTINT, tc=RGBColor(0x1E,0x6B,0x2E)):
    h = 0.74
    box(slide, x, y, w, h, fill=fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(slide, x + 0.22, y, w - 0.4, h, [[(icon + " ", 14, tc, True), (text, 13.5, DARKTX, False)]],
        anchor=MSO_ANCHOR.MIDDLE)


def logo(slide, x=12.4, y=6.92, h=0.4):
    if os.path.exists(LOGO):
        slide.shapes.add_picture(LOGO, Inches(x), Inches(y), height=Inches(h))


def wmark_num(slide, n):
    # 右側に薄い特大数字（装飾）
    tb = slide.shapes.add_textbox(Inches(9.6), Inches(2.2), Inches(3.6), Inches(4.4))
    tf = tb.text_frame; tf.margin_left = tf.margin_top = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    r = p.add_run(); r.text = str(n); r.font.size = Pt(260); r.font.bold = True
    r.font.color.rgb = WMARK; r.font.name = HEAD


# ════════════════ Slide 1: 表紙 ════════════════
s = prs.slides.add_slide(blank); bg(s, NAVY_DK)
box(s, 0, 0, 13.333, 7.5, fill=NAVY_DK)
# 装飾リング
ring = box(s, 8.4, -1.6, 6.2, 6.2, line=NAVY_MD, line_w=2.0, shape=MSO_SHAPE.OVAL)
ring2 = box(s, 9.8, 3.0, 5.4, 5.4, line=GREEN, line_w=1.5, shape=MSO_SHAPE.OVAL)
for i in range(4):
    dot(s, 0.95 + i*0.32, 2.05, 0.16, GREEN if i == 0 else NAVY_MD)
txt(s, 0.9, 2.45, 11.2, 2.0, [
    [("DPS 空調デマンド制御", 46, WHITE, True, HEAD)],
    [("試算アプリ　導入マニュアル", 46, WHITE, True, HEAD)],
], sp_after=2, line_sp=1.04)
txt(s, 0.92, 4.55, 10.5, 0.6, [[("検針票から削減効果・投資回収を自動試算するブラウザアプリ", 16, ICE, False)]])
box(s, 0.95, 5.25, 0.5, 0.06, fill=GREEN)
txt(s, 0.92, 5.42, 10.5, 0.5, [[("初めての方向け  セットアップ ＆ 基本操作ガイド", 14, GREEN_BR, True)]])
if os.path.exists(LOGO):
    s.shapes.add_picture(LOGO, Inches(11.55), Inches(0.55), height=Inches(0.66))

# ════════════════ Slide 2: できること（2x2グリッド） ════════════════
s = prs.slides.add_slide(blank); bg(s, WHITE)
head(s, "このアプリでできること", "検針票12ヶ月分のデータから、提案・稟議に使う数字を自動で作成", emoji="✨")
items = [
    ("📉", "削減効果を自動試算", "空調デマンド制御による電気代の削減額・投資回収期間を算出"),
    ("💰", "財務試算まで一気通貫", "機材費・工事費・利用料・補助金 → 純削減／NPV・IRR／10年・15年収支"),
    ("📑", "提案資料を自動生成", "Excel稟議シート・PowerPoint提案書をワンクリックで出力"),
    ("🗂️", "案件管理・各種設定", "結果一覧に保存。固定値（単価・係数）は管理画面で調整可能"),
]
gx, gy, cw, ch, gp = 0.9, 1.95, 5.65, 2.1, 0.35
for i, (em, h, d) in enumerate(items):
    x = gx + (i % 2) * (cw + gp); y = gy + (i // 2) * (ch + 0.3)
    box(s, x, y, cw, ch, fill=LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE, sh=True)
    badge(s, x + 0.3, y + 0.32, 0.95, "", emoji=em, fill=GREEN, num_sz=2)
    txt(s, x + 1.5, y + 0.36, cw - 1.7, 0.5, [[(h, 19, NAVY, True, HEAD)]])
    txt(s, x + 1.5, y + 0.92, cw - 1.7, 1.0, [[(d, 13.5, GRAY, False)]])
logo(s)

# ════════════════ Slide 3: 動作環境 ════════════════
s = prs.slides.add_slide(blank); bg(s, WHITE)
head(s, "動作環境・必要なもの", "用意するのは次の4つだけ。アプリ一式フォルダは別途共有されます", emoji="🧩")
cards = [
    ("💻", "Windows PC", "インターネット接続が\nできるパソコン"),
    ("🐍", "Python（無料）", "次ページの手順で\nインストールします"),
    ("📁", "アプリ一式フォルダ", "app.py・画像・設定ファイル\n※別途共有予定"),
    ("🌐", "Webブラウザ", "Chrome / Edge など\n（標準のものでOK）"),
]
cx, cw, gap, cy, ch = 0.9, 2.85, 0.27, 1.95, 2.7
for i, (em, h, d) in enumerate(cards):
    x = cx + i * (cw + gap)
    box(s, x, cy, cw, ch, fill=LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE, sh=True)
    badge(s, x + cw/2 - 0.45, cy + 0.3, 0.9, "", emoji=em, fill=NAVY, num_sz=2)
    txt(s, x + 0.12, cy + 1.34, cw - 0.24, 0.5, [[(h, 15.5, NAVY, True, HEAD)]], align=PP_ALIGN.CENTER)
    txt(s, x + 0.12, cy + 1.82, cw - 0.24, 0.9,
        [[(l, 12.5, GRAY, False)] for l in d.split("\n")], align=PP_ALIGN.CENTER, sp_after=1)
box(s, 0.9, 5.4, 11.55, 0.85, fill=GTINT, shape=MSO_SHAPE.ROUNDED_RECTANGLE, sh=True)
txt(s, 1.2, 5.4, 11.0, 0.85,
    [[("📁  ", 16, GREEN, True), ("「アプリ一式フォルダ」は管理者から別途共有されます", 15, RGBColor(0x1E,0x6B,0x2E), True),
      ("（app.py・画像png・dps_app_data.json などを同梱）", 13, GRAY, False)]],
    anchor=MSO_ANCHOR.MIDDLE)
logo(s)

# ════════════════ Slide 4: 導入の流れ ════════════════
s = prs.slides.add_slide(blank); bg(s, WHITE)
head(s, "導入の流れ（4ステップ）", "最初の1回だけ設定すれば、次回からは起動コマンドだけで使えます", emoji="🚀")
steps = [("🐍", "Pythonを\nインストール"), ("📦", "ライブラリを\n導入"),
         ("📁", "フォルダを\n配置"), ("▶", "アプリを\n起動")]
sx, sw, sgap, sy, dd = 0.95, 2.7, 0.42, 2.55, 1.3
for i, (em, h) in enumerate(steps):
    x = sx + i * (sw + sgap)
    badge(s, x + sw/2 - dd/2, sy, dd, i + 1, emoji=em, fill=GREEN if i % 2 == 0 else NAVY, num_sz=40)
    txt(s, x, sy + dd + 0.2, sw, 0.4, [[(f"STEP {i+1}", 13, GREEN, True, HEAD)]], align=PP_ALIGN.CENTER)
    txt(s, x, sy + dd + 0.58, sw, 0.8,
        [[(ln, 16, NAVY, True, HEAD)] for ln in h.split("\n")], align=PP_ALIGN.CENTER, sp_after=0)
    if i < 3:
        box(s, x + sw + 0.03, sy + dd/2 - 0.13, sgap - 0.06, 0.26, fill=ICE, shape=MSO_SHAPE.CHEVRON)
logo(s)

# ════════════════ STEP slides 共通 ════════════════
def step_slide(num, emoji, title, sub, bullets, codelines=None, code_lbl="コマンド（コピーして実行）", note_txt=None):
    s = prs.slides.add_slide(blank); bg(s, WHITE)
    wmark_num(s, num)
    badge(s, 0.7, 0.5, 1.0, num, emoji=emoji, fill=GREEN, num_sz=34)
    txt(s, 1.92, 0.5, 10.5, 0.6, [[(title, 28, NAVY, True, HEAD)]])
    txt(s, 1.94, 1.18, 10.5, 0.4, [[(f"STEP {num}　", 13, GREEN, True, HEAD), (sub, 14, GRAY, False)]])
    yy = 2.05
    for b in bullets:
        dot(s, 0.95, yy + 0.11, 0.15, GREEN)
        txt(s, 1.28, yy - 0.02, 11.1, 0.6, [b], sp_after=0)
        yy += 0.64
    if codelines:
        txt(s, 0.95, yy + 0.04, 8, 0.32, [[("▼ " + code_lbl, 12, GRAY, True)]])
        yy = code_box(s, 0.95, yy + 0.4, 11.45, codelines) + 0.22
    if note_txt:
        note(s, 0.95, max(yy, 5.95), 11.45, note_txt, icon="⚠", fill=RGBColor(0xFF,0xF5,0xDC),
             tc=RGBColor(0xB8,0x86,0x00))
    logo(s)
    return s

PY_URL = "https://www.python.org/downloads/"
LOCAL  = "http://localhost:8501"

step_slide(1, "🐍", "Python をインストール", "プログラムを動かす土台（無料）。最初の1回だけ。",
    [[("公式サイト ", 15, DARKTX, False),
      ("python.org/downloads", 15, GREEN, True, HEAD, PY_URL),
      (" から Windows 用インストーラーをダウンロード", 15, DARKTX, False)],
     [("インストール画面で ", 15, DARKTX, False), ("「Add Python to PATH」に必ずチェック", 15, NAVY, True), ("（重要）", 15, DARKTX, False)],
     [("「Install Now」で進めれば完了", 15, DARKTX, False)]],
    codelines=["python --version", "  → Python 3.xx と表示されればOK"],
    note_txt="「Add Python to PATH」のチェックを忘れると、後の手順でコマンドが認識されません。")

step_slide(2, "📦", "必要なライブラリを導入", "アプリが使う部品をまとめて入れます。最初の1回だけ。",
    [[("スタートメニューで ", 15, DARKTX, False), ("「コマンドプロンプト」", 15, NAVY, True), (" を開く", 15, DARKTX, False)],
     [("下のコマンドを貼り付けて Enter（数分かかります）", 15, DARKTX, False)],
     [("ネット接続が必要。完了すれば次回からは不要", 15, DARKTX, False)]],
    codelines=["pip install streamlit pandas altair openpyxl python-pptx pillow"],
    note_txt="エラーが出る場合は pip を python -m pip に置き換えて再実行してください（例：python -m pip install …）。")

step_slide(3, "📁", "アプリ一式フォルダを配置", "管理者から共有されるフォルダを、PCの好きな場所に置きます。",
    [[("共有された ", 15, DARKTX, False), ("「アプリ一式フォルダ」", 15, NAVY, True), (" を任意の場所（例：デスクトップ）に置く", 15, DARKTX, False)],
     [("フォルダ内には app.py / 画像(png) / dps_app_data.json などが入っています", 15, DARKTX, False)],
     [("フォルダの場所（パス）を控えておく（次の起動で使います）", 15, DARKTX, False)]],
    note_txt="フォルダ内のファイル構成は変更しないでください（app.py と同じ場所に画像・設定が必要です）。")

step_slide(4, "▶", "アプリを起動する", "次回以降はこの手順だけ。コマンド2つで起動します。",
    [[("コマンドプロンプトでフォルダへ移動（cd のあとに半角スペース＋パス）", 15, DARKTX, False)],
     [("起動コマンドを実行 → ", 15, DARKTX, False), ("自動でブラウザが開きます", 15, NAVY, True)],
     [("開かないときは ", 15, DARKTX, False), (LOCAL, 15, GREEN, True, HEAD, LOCAL), (" をブラウザに入力／終了は Ctrl + C", 15, DARKTX, False)]],
    codelines=["cd  （アプリ一式フォルダのパス）", "streamlit run app.py"],
    note_txt="起動後はブラウザのタブで操作します。コマンドプロンプトは開いたままにしてください。")

# ════════════════ Slide 9: 基本操作（タイムライン） ════════════════
s = prs.slides.add_slide(blank); bg(s, WHITE)
head(s, "基本操作：試算の流れ", "サイドバー（左側）に入力 →「シミュレーション実行」で結果が出ます", emoji="🧭")
flow = [
    ("👤", "顧客情報を入力", "顧客名・業態を選択"),
    ("📥", "電力データを取込", "検針票をインポート（または手入力）"),
    ("🔌", "機体を入力（任意）", "室外機の型番・台数で精緻化"),
    ("⚡", "シミュレーション実行", "ボタンを押すと一括計算"),
    ("📊", "結果を確認", "削減額・回収期間・グラフを表示"),
]
yy = 1.95
for i, (em, h, d) in enumerate(flow):
    badge(s, 0.85, yy, 0.62, "", emoji=em, fill=NAVY if i % 2 else GREEN, num_sz=2)
    txt(s, 1.7, yy + 0.02, 4.6, 0.55, [[(h, 17.5, NAVY, True, HEAD)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 6.4, yy + 0.02, 6.1, 0.55, [[(d, 14, GRAY, False)]], anchor=MSO_ANCHOR.MIDDLE)
    if i < 4:
        box(s, 1.13, yy + 0.62, 0.06, 0.36, fill=ICE)
    yy += 0.98
logo(s)

# ════════════════ Slide 10: 出力物 ════════════════
s = prs.slides.add_slide(blank); bg(s, WHITE)
head(s, "出力できる資料", "試算結果はそのまま提案・稟議の資料として書き出せます", emoji="📤")
outs = [
    ("📗", "Excel 稟議シート", "1円差異のない財務対応シート。\n稟議・社内承認にそのまま使用可"),
    ("📙", "PowerPoint 提案書", "自動生成スライド。\n顧客提示用にすぐ展開可能"),
    ("🗂️", "結果一覧に保存", "企業名・空調数・年間削減量・\n回収年数を一覧で蓄積"),
]
cx, cw, gap, cy, ch = 0.95, 3.75, 0.3, 2.05, 3.05
for i, (em, h, d) in enumerate(outs):
    x = cx + i * (cw + gap)
    box(s, x, cy, cw, ch, fill=LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE, sh=True)
    badge(s, x + 0.3, cy + 0.32, 0.9, "", emoji=em, fill=GREEN, num_sz=2)
    txt(s, x + 0.3, cy + 1.4, cw - 0.6, 0.5, [[(h, 17, NAVY, True, HEAD)]])
    txt(s, x + 0.3, cy + 1.95, cw - 0.55, 1.0,
        [[(l, 13.5, GRAY, False)] for l in d.split("\n")], sp_after=2)
logo(s)

# ════════════════ Slide 11: 各種設定 ════════════════
s = prs.slides.add_slide(blank); bg(s, WHITE)
head(s, "各種設定（管理者向け・任意）", "固定値の調整はサイドバー「各種設定」から。通常は変更不要です", emoji="⚙️")
rows = [
    ("🧮", "計算ロジックの変更", "容量削減率・利用料単価・換算係数・力率割引・詳細見積の係数 など"),
    ("🏷️", "型番→馬力マスタ", "室外機の型番ごとの馬力・COP・仮置き値の登録"),
    ("🏭", "業態マスタ", "業態別の空調割合（電力量・ピーク）の調整"),
    ("📈", "表示設定 / 履歴", "グラフ表示のON/OFF・変更履歴の確認"),
]
yy = 1.95
for em, h, d in rows:
    box(s, 0.9, yy, 11.5, 0.98, fill=LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE, sh=True)
    badge(s, 1.12, yy + 0.22, 0.55, "", emoji=em, fill=GREEN, num_sz=2)
    txt(s, 1.95, yy, 3.7, 0.98, [[(h, 16, NAVY, True, HEAD)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 5.7, yy, 6.5, 0.98, [[(d, 13.5, GRAY, False)]], anchor=MSO_ANCHOR.MIDDLE)
    yy += 1.12
logo(s)

# ════════════════ Slide 12: 困ったとき ════════════════
s = prs.slides.add_slide(blank); bg(s, WHITE)
head(s, "困ったとき（トラブルシュート）", "よくあるつまずきと対処法", emoji="🛟")
qa = [
    ("「streamlit は認識されません」", "STEP2 の pip install を再実行（python -m pip install …）", None),
    ("ブラウザが開かない", LOCAL + " をブラウザに直接入力", LOCAL),
    ("コマンドが動かない", "アプリ一式フォルダ内でコマンドを実行しているか確認（cd で移動）", None),
    ("数値に違和感がある", "検針票の単価・力率・使用量の入力値を確認（空欄は既定値で補完）", None),
]
yy = 1.95
for q, a, link in qa:
    box(s, 0.9, yy, 11.5, 0.95, fill=LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE, sh=True)
    txt(s, 1.15, yy + 0.1, 11.0, 0.4, [[("Q. ", 15, GREEN, True, HEAD), (q, 15, NAVY, True, HEAD)]])
    a_runs = [("A. ", 13.5, GRAY, True)]
    if link:
        a_runs.append((a, 13.5, GREEN, True, HEAD, link))
    else:
        a_runs.append((a, 13.5, DARKTX, False))
    txt(s, 1.4, yy + 0.5, 10.9, 0.4, [a_runs])
    yy += 1.1
logo(s)

# ════════════════ Slide 13: まとめ ════════════════
s = prs.slides.add_slide(blank); bg(s, NAVY_DK)
ring = box(s, 8.6, 2.4, 6.2, 6.2, line=NAVY_MD, line_w=2.0, shape=MSO_SHAPE.OVAL)
for i in range(4):
    dot(s, 0.95 + i*0.32, 1.7, 0.16, GREEN if i == 0 else NAVY_MD)
txt(s, 0.9, 2.05, 11.5, 1.2, [[("4ステップで、誰でも使えます", 38, WHITE, True, HEAD)]])
txt(s, 0.92, 3.35, 11.2, 0.6, [[("インストール → ライブラリ導入 → フォルダ配置 → ", 17, ICE, False),
                                ("streamlit run app.py", 17, GREEN_BR, True, MONO)]])
for i, t in enumerate([
    "2回目以降は「フォルダへ移動 → 起動コマンド」だけ",
    "アプリ一式フォルダは管理者から別途共有されます",
    "困ったときは本マニュアルのトラブルシュートを参照",
]):
    dot(s, 0.95, 4.35 + i*0.55, 0.14, GREEN_BR)
    txt(s, 1.25, 4.22 + i*0.55, 11.0, 0.45, [[(t, 15, ICE, False)]])
if os.path.exists(LOGO):
    s.shapes.add_picture(LOGO, Inches(11.55), Inches(6.5), height=Inches(0.55))

prs.save("DPS試算アプリ_導入マニュアル.pptx")
print("saved slides:", len(prs.slides._sldIdLst))
